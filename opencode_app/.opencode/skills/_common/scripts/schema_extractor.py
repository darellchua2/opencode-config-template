"""
schema_extractor.py
===================
US-1.1 extraction engine — reads any ``.pptx`` and emits a structured JSON
representation conforming to ``schemas/template_schema.json`` (the normalized
"proposed schema" from ``docs/user-stories/chenyu-user-stories.md``).

Unlike the existing :mod:`template_introspector` (which emits a lightweight
"fingerprint contract" for the renderer), this module emits the *full* component
model: every slide element (placeholders AND freeform shapes, images, charts,
tables, groups, connectors) becomes a ``component`` with a normalized ``polygon``
(0.0-1.0), a ``type`` enum value, ``z_order``, and (on text-bearing components)
a ``font`` object. The slide master is parsed explicitly, and every layout under
it is enumerated.

This module is a **parallel, non-invasive** path: it does NOT import or modify
``template_introspector`` / ``ppt_builder``. Per PLAN-GIT-48 Decision #4, it
duplicates slide-size / theme-parse logic during the coexistence period;
consolidation is deferred to the migration epic.

Public API
----------
    extract_schema(pptx_path) -> dict               # the proposed-schema JSON
    validate_template_schema(schema_dict) -> ValidationResult
    ValidationResult.is_valid / .errors / .warnings

CLI
---
    python schema_extractor.py --input X.pptx --output schema.json
"""
from __future__ import annotations

import argparse
import json
import logging
import os
import re
import tempfile
import zipfile
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, NamedTuple, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

# US-4.6 (m2): polygon/EMU primitives relocated to the shared pure ``geometry``
# module. Re-exported here so ``schema_extractor.normalize_polygon`` /
# ``_EMU_PER_INCH`` / ``_compute_ratio`` keep resolving (parity-unchanged).
from geometry import _EMU_PER_INCH, _compute_ratio, normalize_polygon  # noqa: F401

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SCHEMA_VERSION = "1.1.0"
GENERATED_BY = "opencode-pptx-subagent/schema_extractor"

# US-1.5: the canonical path of the embedded schema inside the PPTX zip, and the
# OPC [Content_Types].xml namespace for the json-Default injection.
_EMBEDDED_SCHEMA_PATH = "ppt/template_schema.json"
_CONTENT_TYPES_XML = "[Content_Types].xml"
_CONTENT_TYPES_NS = "{http://schemas.openxmlformats.org/package/2006/content-types}"
_JSON_CONTENT_TYPE = "application/json"


class EmbeddedSchemaResult(NamedTuple):
    """Outcome of :func:`embed_schema` (US-1.5); doubles as the testable AC4 signal."""

    output_path: str
    original_bytes: int
    new_bytes: int
    delta_bytes: int


# US-3.1: title-inference provenance. ``TITLE_SOURCES`` is the single source of
# truth consumed by :func:`_infer_title`, :func:`validate_template_schema`'s enum
# check, and ``template_schema.json``'s enum — keeping all three in sync
# (architecture review MINOR-5 / MAJOR-2). ``"user"`` is only ever set by the
# generate-template-skill layer after a user override, never by inference itself.
TITLE_SOURCES: frozenset = frozenset({"core_xml", "slide1", "filename", "user"})


class TitleInference(NamedTuple):
    """Result of :func:`_infer_title`: the inferred title plus its provenance."""

    title: str
    source: str

# Self-documenting enums (mirrors chenyu-user-stories.md Reference schema).
COMPONENT_TYPE_ENUM: List[str] = [
    "textbox", "image", "table", "video", "shape",
    "chart", "group", "smartart", "placeholder", "audio",
]

PLACEHOLDER_TYPE_ENUM: List[Optional[str]] = [
    "title", "subtitle", "body", "picture", "chart", "table",
    "media", "date", "slide_number", "footer", "header", None,
]

# Component types that may carry a ``font`` object (C1 cardinality rule).
_TEXT_TYPES = {"textbox", "placeholder"}

# Confidence values emitted on every component's ``type_confidence`` (US-1.3).
# Centralized so the classifier, the validator, and template_schema.json's enum
# stay in sync (single source of truth for the {"high","low"} domain).
_CONFIDENCE_HIGH = "high"
_CONFIDENCE_LOW = "low"

# US-1.4: fonts considered universally available (MS-Office / cross-platform).
# Membership drives is_available + missing_fonts. Curated, extensible.
_BUILTIN_FONTS: frozenset = frozenset({
    "Arial", "Calibri", "Calibri Light", "Cambria", "Candara", "Century Gothic",
    "Comic Sans MS", "Consolas", "Constantia", "Corbel", "Courier New", "Georgia",
    "Impact", "Segoe UI", "Tahoma", "Times New Roman", "Trebuchet MS", "Verdana",
})

# US-1.4: common non-built-in -> built-in fallback (values must be in
# _BUILTIN_FONTS). Unmapped non-built-in families fall back to the theme body
# font (if built-in) else "Arial" — see _font_fallback().
_FONT_FALLBACK_MAP: Dict[str, str] = {
    # sans-serif -> Arial
    "Helvetica": "Arial", "Helvetica Neue": "Arial", "SF Pro": "Arial",
    "SF Pro Display": "Arial", "SF Pro Text": "Arial", "Roboto": "Arial",
    "Open Sans": "Arial", "Inter": "Arial", "Lato": "Arial", "Montserrat": "Arial",
    "Proxima Nova": "Arial", "Avenir": "Arial", "Avenir Next": "Arial",
    "Ubuntu": "Arial", "Noto Sans": "Arial", "Source Sans Pro": "Arial",
    # serif -> Times New Roman / Georgia
    "Garamond": "Times New Roman", "Garamond Premier Pro": "Times New Roman",
    "Baskerville": "Times New Roman", "Didot": "Georgia",
    # monospace -> Consolas
    "Source Code Pro": "Consolas", "Fira Code": "Consolas",
    "Monaco": "Consolas", "Menlo": "Consolas",
}

# PP_ALIGN enum-name -> schema alignment string.
_ALIGNMENT_MAP: Dict[str, str] = {
    "LEFT": "left", "CENTER": "center", "RIGHT": "right",
    "JUSTIFY": "justify", "JUSTIFY_LOW": "justify",
    "DISTRIBUTE": "distribute", "START": "start", "END": "end",
}

# Default fallback used when the theme body font is not itself a built-in.
_DEFAULT_FALLBACK_FONT = "Arial"

# Initial theme semantic-role mapping (refined in US-3.4). Maps the OOXML
# clrScheme role produced by :func:`_raw_theme_colors` to a semantic role.
_THEME_ROLE_MAP = {
    "primary": "dk2",
    "secondary": "lt2",
    "accent": "accent1",
    "background": "lt1",
    "text_color": "dk1",
}

_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


class TemplateExtractionError(Exception):
    """Domain error raised when a .pptx cannot be extracted (bad zip, no master)."""


class SchemaVersionError(Exception):
    """Raised when an embedded schema's major version is incompatible (BT-142 Phase 1.4).

    ``read_embedded_schema`` compares the embedded ``schema_version`` against the
    current ``SCHEMA_VERSION`` constant. A **major** mismatch (e.g. embedded 2.x
    vs current 1.x) raises this error — the consumer should not attempt to use a
    schema from an incompatible generation. Minor and patch mismatches do not
    raise: patch is forward-compatible (silent auto-upgrade on next embed);
    minor produces a ``logger.warning`` but returns the data (additive fields
    are expected).
    """


# ---------------------------------------------------------------------------
# Task 5 (defined early so extract_schema can reuse the issue types)
# ---------------------------------------------------------------------------
class ValidationIssue:
    """A single validation finding (error or warning)."""

    def __init__(self, reason: str, *, field_path: str = "", severity: str = "error"):
        self.field_path = field_path
        self.reason = reason
        self.severity = severity  # "error" | "warning"

    @property
    def is_error(self) -> bool:
        return self.severity == "error"

    def format(self) -> str:
        path = f".{self.field_path}" if self.field_path else ""
        return f"schema{path}: {self.reason}"

    def __repr__(self) -> str:  # pragma: no cover - debug helper
        return f"<{self.severity.upper()} {self.format()}>"


class ValidationResult:
    """Aggregated result of validating a template schema."""

    def __init__(self) -> None:
        self.issues: List[ValidationIssue] = []

    @property
    def errors(self) -> List[ValidationIssue]:
        return [i for i in self.issues if i.is_error]

    @property
    def warnings(self) -> List[ValidationIssue]:
        return [i for i in self.issues if not i.is_error]

    @property
    def is_valid(self) -> bool:
        return not self.errors

    def add(self, issue: ValidationIssue) -> None:
        self.issues.append(issue)

    def error_messages(self) -> List[str]:
        return [e.format() for e in self.errors]

    def warning_messages(self) -> List[str]:
        return [w.format() for w in self.warnings]

    def __repr__(self) -> str:  # pragma: no cover - debug helper
        return (
            f"<ValidationResult valid={self.is_valid} "
            f"errors={len(self.errors)} warnings={len(self.warnings)}>"
        )


# ---------------------------------------------------------------------------
# Task 2: element type mapper
# ---------------------------------------------------------------------------
# Map python-pptx placeholder type -> placeholder_type enum value.
_PLACEHOLDER_TYPE_MAP: Dict[Any, str] = {
    PP_PLACEHOLDER.TITLE: "title",
    PP_PLACEHOLDER.CENTER_TITLE: "title",
    PP_PLACEHOLDER.VERTICAL_TITLE: "title",
    PP_PLACEHOLDER.SUBTITLE: "subtitle",
    PP_PLACEHOLDER.BODY: "body",
    PP_PLACEHOLDER.OBJECT: "body",
    PP_PLACEHOLDER.VERTICAL_BODY: "body",
    PP_PLACEHOLDER.VERTICAL_OBJECT: "body",
    PP_PLACEHOLDER.PICTURE: "picture",
    PP_PLACEHOLDER.BITMAP: "picture",
    PP_PLACEHOLDER.CHART: "chart",
    PP_PLACEHOLDER.ORG_CHART: "chart",
    PP_PLACEHOLDER.TABLE: "table",
    PP_PLACEHOLDER.MEDIA_CLIP: "media",
    PP_PLACEHOLDER.DATE: "date",
    PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
    PP_PLACEHOLDER.FOOTER: "footer",
    PP_PLACEHOLDER.HEADER: "header",
}


def _media_classification(shape: Any) -> Tuple[str, bool]:
    """Classify a MEDIA shape's subtype from its OOXML (US-1.3).

    python-pptx collapses audio and video into ``MSO_SHAPE_TYPE.MEDIA`` and
    exposes no audio/video subtype. The subtype is recoverable from OOXML: a
    media ``<p:pic>`` carries ``<a:audioFile>`` or ``<a:videoFile>`` under
    ``<p:nvPicPr>/<p:nvPr>``.

    Returns ``(subtype, had_marker)`` where ``subtype`` is ``"audio"`` or
    ``"video"`` and ``had_marker`` is True when an ``<a:audioFile>``/
    ``<a:videoFile>`` element was found. Returns ``("video", False)`` when the
    marker is absent or the shape has no ``_element`` (defensive; reachable only
    via synthetic test shapes — real python-pptx shapes always carry ``_element``).
    """
    el = getattr(shape, "_element", None)
    if el is not None:
        if el.find(".//" + qn("a:audioFile")) is not None:
            return "audio", True
        if el.find(".//" + qn("a:videoFile")) is not None:
            return "video", True
    return "video", False


def _classify_shape(shape: Any) -> Tuple[str, str]:
    """Classify a shape into ``(component_type, type_confidence)``.

    ``type_confidence`` is ``"high"`` whenever the element was positively
    recognized (any non-``None`` ``shape_type``, or a table/chart/placeholder),
    and ``"low"`` only for genuinely unrecognized shapes (``shape_type`` is
    ``None`` or unreadable) and for MEDIA whose audio/video subtype could not be
    determined (the US-1.3 "flagged for review" signal).

    Per architecture review MAJOR-1 (PLAN-GIT-52 v2): there is NO whitelist —
    every non-``None`` ``shape_type`` is treated as recognized, so recognized
    members that are not explicitly mapped (``AUTO_SHAPE``, ``FREEFORM``,
    ``LINE``, ``CALLOUT``, ``LINKED_PICTURE``, ``TEXT_EFFECT``, …) fall through
    to ``("shape", "high")`` rather than being mislabeled ``"low"``.
    """
    if getattr(shape, "is_placeholder", False) and shape.placeholder_format is not None:
        return "placeholder", _CONFIDENCE_HIGH
    if getattr(shape, "has_table", False):
        return "table", _CONFIDENCE_HIGH
    if getattr(shape, "has_chart", False):
        return "chart", _CONFIDENCE_HIGH
    try:
        st = shape.shape_type
    except Exception:
        return "shape", _CONFIDENCE_LOW  # cannot even read shape_type -> unrecognized
    if st is None:
        return "shape", _CONFIDENCE_LOW  # genuinely unrecognized
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image", _CONFIDENCE_HIGH
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group", _CONFIDENCE_HIGH
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "textbox", _CONFIDENCE_HIGH
    if st == MSO_SHAPE_TYPE.IGX_GRAPHIC:
        return "smartart", _CONFIDENCE_HIGH
    if st == MSO_SHAPE_TYPE.WEB_VIDEO:
        return "video", _CONFIDENCE_HIGH
    if st == MSO_SHAPE_TYPE.MEDIA:
        sub, had_marker = _media_classification(shape)
        return (sub, _CONFIDENCE_HIGH) if had_marker else ("video", _CONFIDENCE_LOW)
    # Any other non-None shape_type (AUTO_SHAPE, FREEFORM, LINE, CALLOUT,
    # LINKED_PICTURE, TEXT_EFFECT, INK, DIAGRAM, …) is a recognized shape that
    # the enum does not sub-type further.
    return "shape", _CONFIDENCE_HIGH


def map_shape_type(shape: Any) -> str:
    """Map a python-pptx shape to a component ``type`` enum value.

    Backward-compatible wrapper over :func:`_classify_shape` (returns only the
    type); see that function for the full precedence chain and confidence rules.
    """
    return _classify_shape(shape)[0]


def map_placeholder_type(ph: Any) -> Optional[str]:
    """Map a placeholder to its ``placeholder_type`` enum value.

    An untyped placeholder (``type is None``) maps to ``"body"`` (generic
    content), mirroring the introspector's None -> OBJECT convention.
    """
    pf = ph.placeholder_format
    if pf is None:
        return None
    pt = pf.type
    if pt is None:
        return "body"
    return _PLACEHOLDER_TYPE_MAP.get(pt)


# ---------------------------------------------------------------------------
# Task 3: polygon normalizer — relocated to ``geometry`` (US-4.6 m2).
# normalize_polygon / _clamp_unit / _compute_ratio are imported at the top.
# ---------------------------------------------------------------------------


def _build_slide_dimensions(prs: Presentation) -> Dict[str, Any]:
    w_emu = int(prs.slide_width)
    h_emu = int(prs.slide_height)
    return {
        "width_emu": w_emu,
        "height_emu": h_emu,
        "width_inches": round(w_emu / _EMU_PER_INCH, 4),
        "height_inches": round(h_emu / _EMU_PER_INCH, 4),
        "aspect_ratio": _compute_ratio(w_emu, h_emu),
    }


def parse_theme_xml(theme_xml_bytes: bytes) -> Tuple[Dict[str, str], Dict[str, str]]:
    """Parse clrScheme + fontScheme from raw theme1.xml bytes (US-4.8 / CRIT-2).

    Pure function — does NOT require a ``Presentation`` or ``prs.slide_masters``
    access. Used by ``_raw_theme_colors_and_fonts`` (the legacy prs-based
    delegate) and by ``master_repairer._salvage_theme_part`` (zip-level salvage
    that has no prs at all).

    Returns (colors_by_role, fonts_by_role). Best-effort: returns empty dicts
    on parse failure.
    """
    colors: Dict[str, str] = {}
    fonts: Dict[str, str] = {}
    try:
        theme_xml = etree.parse(BytesIO(theme_xml_bytes)).getroot()
        elements = theme_xml.find(f"{{{_NS_A}}}themeElements")
        if elements is not None:
            clr_scheme = elements.find(f"{{{_NS_A}}}clrScheme")
            if clr_scheme is not None:
                for child in clr_scheme:
                    role = etree.QName(child).localname
                    for color_elem in child:
                        color_tag = etree.QName(color_elem).localname
                        if color_tag == "srgbClr":
                            colors[role] = "#" + color_elem.get("val", "")
                        elif color_tag == "sysClr":
                            val = color_elem.get("lastClr") or color_elem.get("val", "")
                            if val:
                                colors[role] = "#" + val
            font_scheme = elements.find(f"{{{_NS_A}}}fontScheme")
            if font_scheme is not None:
                for label, key in (("majorFont", "major_latin"), ("minorFont", "minor_latin")):
                    font_elem = font_scheme.find(f"{{{_NS_A}}}{label}")
                    if font_elem is not None:
                        latin = font_elem.find(f"{{{_NS_A}}}latin")
                        if latin is not None and latin.get("typeface"):
                            fonts[key] = latin.get("typeface")
    except Exception:  # pragma: no cover - defensive
        pass
    return colors, fonts


def _raw_theme_colors_and_fonts(prs: Presentation) -> Tuple[Dict[str, str], Dict[str, str]]:
    """Extract raw clrScheme roles + major/minor Latin fonts from theme1.xml.

    Returns (colors_by_role, fonts_by_role). Best-effort: returns empty dicts on
    failure (logged at warning level). Delegates to the pure
    :func:`parse_theme_xml` for the actual XML walking (US-4.8 / CRIT-2).
    """
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(_THEME_REL)
        return parse_theme_xml(theme_part.blob)
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Theme extraction failed (%s); emitting empty theme", exc)
        return {}, {}


def _build_theme(prs: Presentation) -> Dict[str, Any]:
    """Map raw theme colors/fonts to semantic roles (initial heuristic, US-3.4 refines)."""
    colors, fonts = _raw_theme_colors_and_fonts(prs)
    theme: Dict[str, Any] = {
        "primary_color": colors.get(_THEME_ROLE_MAP["primary"], ""),
        "secondary_color": colors.get(_THEME_ROLE_MAP["secondary"], ""),
        "accent_color": colors.get(_THEME_ROLE_MAP["accent"], ""),
        "background_color": colors.get(_THEME_ROLE_MAP["background"], ""),
        "text_color": colors.get(_THEME_ROLE_MAP["text_color"], ""),
    }
    font_palette: Dict[str, str] = {
        "heading": fonts.get("major_latin", ""),
        "body": fonts.get("minor_latin", ""),
        "accent": fonts.get("major_latin", ""),
    }
    theme["font_palette"] = font_palette
    return theme


def _infer_title(prs: Presentation, path: str) -> TitleInference:
    """Infer the template title plus its provenance.

    Order (US-3.2 AC2): ``docProps/core.xml`` title -> first slide title text ->
    filename stem. The returned :class:`TitleInference.source` is one of
    ``TITLE_SOURCES`` (minus ``"user"``, which is only set by the skill layer
    after a user override). The filename fallback guarantees a non-empty title.
    """
    try:
        title = prs.core_properties.title
        if title and title.strip():
            return TitleInference(title.strip(), "core_xml")
    except Exception:
        pass
    try:
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        return TitleInference(txt[:100], "slide1")
    except Exception:
        pass
    return TitleInference(Path(path).stem, "filename")


def _detect_header_footer(prs: Presentation) -> Dict[str, bool]:
    """Scan the slide master for header/footer placeholders (US-2.1 AC1).

    Returns ``{has_header, has_footer}`` booleans. **Master-only** (per chenyu's
    wording; layout-only chrome is out of scope). Best-effort (never blocks
    extraction).
    """
    has_header = False
    has_footer = False
    try:
        for ph in prs.slide_masters[0].placeholders:
            ptype = ph.placeholder_format.type
            if ptype == PP_PLACEHOLDER.HEADER:
                has_header = True
            elif ptype == PP_PLACEHOLDER.FOOTER:
                has_footer = True
    except Exception:
        pass
    return {"has_header": has_header, "has_footer": has_footer}


def _build_metadata(prs: Presentation, path: str) -> Dict[str, Any]:
    inferred = _infer_title(prs, path)
    return {
        "title": inferred.title,
        "title_source": inferred.source,
        "schema_version": SCHEMA_VERSION,
        "generated_by": GENERATED_BY,
        "generated_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "slide_dimensions": _build_slide_dimensions(prs),
        "missing_fonts": [],  # initially empty; populated below in extract_schema
        "header_footer": _detect_header_footer(prs),  # US-2.1 AC1
    }


# ---------------------------------------------------------------------------
# Component extraction (Task 4)
# ---------------------------------------------------------------------------
class _IdCounter:
    """Global comp_NNN counter spanning master + all layouts."""

    def __init__(self) -> None:
        self._n = 0

    def next_id(self) -> str:
        self._n += 1
        return f"comp_{self._n:03d}"


def _alignment_str(align: Any) -> Optional[str]:
    """Map a ``PP_ALIGN`` value to the schema alignment string (US-1.4).

    python-pptx enum ``str`` is like ``"left (1)"``; take the name token.
    """
    if align is None:
        return None
    name = str(align).split(" (")[0].strip().upper()
    return _ALIGNMENT_MAP.get(name, name.lower())


def _weight_from_bold(bold: Any) -> Optional[str]:
    """Map ``run.font.bold`` (tri-state bool|None) to a weight string.

    True -> "bold"; False/None -> None (explicit-only; "regular" is implied).
    """
    return "bold" if bold is True else None


def _run_color_hex(run: Any) -> Optional[str]:
    """Hex color of a run's font, only when it is an explicit RGB (US-1.4).

    ``run.font.color.rgb`` raises on theme/None colors, so guard on
    ``color.type == MSO_COLOR_TYPE.RGB``. Returns "#RRGGBB" or None.
    """
    try:
        color = run.font.color
        if color.type == MSO_COLOR_TYPE.RGB:
            return "#" + str(color.rgb)
    except Exception:  # pragma: no cover - defensive
        pass
    return None


def _font_fallback(family: Optional[str], default_body: str) -> Optional[str]:
    """Built-in substitute for a non-built-in family (US-1.4, AC4).

    Returns None when ``family`` is None or built-in; the mapped built-in name
    for known non-built-ins; else ``default_body`` (the theme body font when it
    is itself built-in, else "Arial"). The result is always a built-in font name
    or None.
    """
    if family is None or family in _BUILTIN_FONTS:
        return None
    mapped = _FONT_FALLBACK_MAP.get(family)
    if mapped is not None:
        return mapped
    # default_body should already be built-in (resolved by the caller); guard so
    # AC4 holds regardless of caller.
    return default_body if default_body in _BUILTIN_FONTS else _DEFAULT_FALLBACK_FONT


# ---------------------------------------------------------------------------
# US-4.6 (Phase 1) — paragraph bullet/spacing capture + image-asset capture.
# These feed the coordinate-placement path's styling re-application (AC4) and
# image-byte sourcing (C1). Bullets read the OOXML ``<a:pPr>`` children because
# python-pptx exposes no high-level bullet API.
# ---------------------------------------------------------------------------
def _line_spacing_to_json(ls: Any) -> Optional[float]:
    """Serialize ``paragraph.line_spacing`` (float multiple | Length | None).

    A bare float means a line-spacing multiple (the common case); an exact/
    AtLeast Length is serialized as its pt value (number). ``None`` -> ``None``.
    """
    if ls is None:
        return None
    if isinstance(ls, (int, float)):
        return float(ls)
    try:
        return round(float(ls.pt), 2)
    except Exception:
        return None


def _spacing_to_json(sp: Any) -> Optional[float]:
    """Serialize ``space_before``/``space_after`` (Length | None) to pt (float)."""
    if sp is None:
        return None
    try:
        return round(float(sp.pt), 2)
    except Exception:
        return None


def _extract_paragraph_format(tf: Any) -> List[Dict[str, Any]]:
    """Per-paragraph bullet/spacing info for US-4.6 AC4 (bullets re-application).

    Returns a list of ``{level, type, char, font, line_spacing, space_before,
    space_after}`` for paragraphs that explicitly set any of these; empty list
    when everything is default (the renderer then inherits). ``type`` is one of
    ``char`` (buChar), ``autonum`` (buAutoNum), ``none`` (buNone), or ``null``.
    """
    out: List[Dict[str, Any]] = []
    for p in tf.paragraphs:
        info: Dict[str, Any] = {
            "level": None, "type": None, "char": None, "font": None,
            "line_spacing": None, "space_before": None, "space_after": None,
        }
        explicit = False

        lvl = p.level or 0
        if lvl:
            info["level"] = lvl
            explicit = True
        ls = _line_spacing_to_json(p.line_spacing)
        if ls is not None:
            info["line_spacing"] = ls
            explicit = True
        sb = _spacing_to_json(p.space_before)
        if sb is not None:
            info["space_before"] = sb
            explicit = True
        sa = _spacing_to_json(p.space_after)
        if sa is not None:
            info["space_after"] = sa
            explicit = True

        # Bullet type/char/font live under <a:pPr> (python-pptx has no API).
        pPr = p._p.find(qn("a:pPr"))
        if pPr is not None:
            if pPr.find(qn("a:buNone")) is not None:
                info["type"] = "none"
                explicit = True
            else:
                bu_char = pPr.find(qn("a:buChar"))
                bu_autonum = pPr.find(qn("a:buAutoNum"))
                if bu_char is not None:
                    info["type"] = "char"
                    info["char"] = bu_char.get("char")
                    explicit = True
                elif bu_autonum is not None:
                    info["type"] = "autonum"
                    explicit = True
            bu_font = pPr.find(qn("a:buFont"))
            if bu_font is not None and bu_font.get("typeface"):
                info["font"] = bu_font.get("typeface")

        if explicit:
            out.append(info)
    return out


def _extract_image_properties(shape: Any) -> Dict[str, Any]:
    """Capture enough to re-source an image's bytes (US-4.6 C1).

    Returns ``{partname, content_type, [width_px, height_px]}`` — ``partname``
    is the stable ``/ppt/media/imageN.ext`` path the executor resolves against
    the template package. ``{}`` when the shape is not a resolvable picture.

    python-pptx's ``shape.image`` exposes ``content_type``/``size`` but NOT a
    ``partname``; the part path is resolved via the ``<a:blip r:embed>`` rId ->
    ``part.related_part(rId)``.
    """
    info: Dict[str, Any] = {}
    # content_type + pixel dims from the high-level image accessor.
    try:
        img = shape.image
        info["content_type"] = img.content_type
        try:
            w, h = img.size
            info["width_px"] = int(w)
            info["height_px"] = int(h)
        except Exception:
            pass
    except Exception:
        pass
    # partname via blip r:embed -> related ImagePart (python-pptx Image has no
    # partname attribute, so resolve through the owning part's relationship).
    try:
        blip = shape._element.find(".//" + qn("a:blip"))
        if blip is not None:
            rId = blip.get(qn("r:embed"))
            if rId:
                part = shape.part.related_part(rId)
                info["partname"] = str(part.partname)
    except Exception:
        pass
    return info


def _extract_text_fonts(
    shape: Any, default_body: str
) -> Tuple[Dict[str, Any], List[Dict[str, Any]]]:
    """Extract per-run font metadata from a text-bearing shape (US-1.4).

    Returns ``(font_summary, runs)``:
      - ``font_summary``: dict keyed for all target fields, summarizing the
        **first text run**'s explicit props (+ the first **explicit** paragraph
        alignment); null where not explicitly set. is_available/fallback are
        derived. The summary is the first run (possibly empty) — ``runs[]`` is
        authoritative for visible text.
      - ``runs``: list of ``{text, font: {family, size_pt, weight, color}}`` for
        runs that carry text.

    Explicit-only (inherited/None -> null); Latin only (``<a:latin>``).
    """
    summary: Dict[str, Any] = {
        "family": None, "size_pt": None, "weight": None, "color": None,
        "alignment": None, "is_available": True, "fallback": None,
    }
    runs: List[Dict[str, Any]] = []

    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return summary, runs

    paragraphs = list(tf.paragraphs)
    # alignment: first paragraph that explicitly sets one
    for p in paragraphs:
        if p.alignment is not None:
            summary["alignment"] = _alignment_str(p.alignment)
            break

    first_seen = False
    for p in paragraphs:
        for run in p.runs:
            family = run.font.name
            size = run.font.size
            bold = run.font.bold
            color = _run_color_hex(run)
            size_pt = float(size.pt) if size is not None else None
            if not first_seen:
                summary["family"] = family
                summary["size_pt"] = size_pt
                summary["weight"] = _weight_from_bold(bold)
                summary["color"] = color
                first_seen = True
            text = run.text or ""
            if text:
                runs.append({
                    "text": text,
                    "font": {
                        "family": family,
                        "size_pt": size_pt,
                        "weight": _weight_from_bold(bold),
                        "color": color,
                    },
                })

    fam = summary["family"]
    summary["is_available"] = (fam is None) or (fam in _BUILTIN_FONTS)
    summary["fallback"] = _font_fallback(fam, default_body)
    return summary, runs


def _build_component(
    shape: Any,
    slide_w_emu: int,
    slide_h_emu: int,
    z_order: int,
    counter: _IdCounter,
    default_body: str = _DEFAULT_FALLBACK_FONT,
) -> Optional[Dict[str, Any]]:
    """Build a component record from a single shape, or ``None`` to skip.

    Groups recurse: the group becomes a ``group`` component whose children are
    appended to the same layout's component list (flattened with z_order), so the
    polygon/z_order of nested shapes are captured individually.
    """
    comp_id = counter.next_id()
    comp_type, confidence = _classify_shape(shape)

    component: Dict[str, Any] = {
        "id": comp_id,
        "type": comp_type,
        "type_confidence": confidence,
        "name": shape.name or "",
        "polygon": normalize_polygon(shape, slide_w_emu, slide_h_emu),
        "z_order": z_order,
    }

    # placeholder_type
    if comp_type == "placeholder":
        component["placeholder_type"] = map_placeholder_type(shape)
    else:
        component["placeholder_type"] = None

    # font + runs + text_properties.bullets: present ONLY on text-bearing
    # components (C1). Populated per-run (US-1.4); bullets/spacing (US-4.6 AC4).
    if comp_type in _TEXT_TYPES:
        font_summary, runs = _extract_text_fonts(shape, default_body)
        component["font"] = font_summary
        component["runs"] = runs
        tf = getattr(shape, "text_frame", None)
        component["text_properties"] = {
            "bullets": _extract_paragraph_format(tf) if tf is not None else []
        }

    # image_properties: present ONLY on image components (US-4.6 C1). Captures
    # the asset's partname/content_type so the coordinate path can re-source
    # bytes from the template package.
    if comp_type == "image":
        component["image_properties"] = _extract_image_properties(shape)

    # content_template: text-bearing components get a simple placeholder marker.
    if comp_type in _TEXT_TYPES:
        component["content_template"] = "{{content}}"
    else:
        component["content_template"] = None

    return component


def _extract_components(
    shapes: Any,
    slide_w_emu: int,
    slide_h_emu: int,
    counter: _IdCounter,
    default_body: str = _DEFAULT_FALLBACK_FONT,
) -> List[Dict[str, Any]]:
    """Extract components from an iterable of shapes (master or a layout).

    Groups are recursed: the group itself becomes a ``group`` component, then its
    children are appended (flattened) so each nested shape is captured.
    ``z_order`` is assigned by the final flatten order, so it is monotonic and
    unique within each master/layout (the source-shape enumerate index is not
    used, to avoid collisions when a sibling follows a flattened group).
    """
    flat: List[Dict[str, Any]] = []
    for shape in shapes:
        comp = _build_component(shape, slide_w_emu, slide_h_emu, 0, counter, default_body)
        if comp is None:
            continue
        flat.append(comp)
        # Recurse into groups: append nested children after the group.
        if comp["type"] == "group":
            try:
                children = shape.shapes
            except Exception:
                children = []
            flat.extend(_extract_components(
                children, slide_w_emu, slide_h_emu, counter, default_body
            ))
    # Assign z_order by final flatten index (monotonic + unique).
    for i, comp in enumerate(flat):
        comp["z_order"] = i
    return flat


def _slugify(name: str) -> str:
    """Turn a layout/master name into a stable layout_id slug."""
    slug = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_").lower()
    return slug or "layout"


# ---------------------------------------------------------------------------
# Public API — extract_schema (Task 4)
# ---------------------------------------------------------------------------
def extract_schema(pptx_path: str) -> Dict[str, Any]:
    """Read ``pptx_path`` and return the proposed-schema JSON dict.

    Parses the slide master and every layout, emitting a structure conforming to
    ``schemas/template_schema.json``. Raises :class:`TemplateExtractionError` on
    unreadable / non-PPTX input.
    """
    path = Path(pptx_path)
    if not path.exists():
        raise TemplateExtractionError(f"file not found: {pptx_path}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise TemplateExtractionError(
            f"could not open as PPTX ({exc.__class__.__name__}: {exc})"
        ) from exc

    dims = _build_slide_dimensions(prs)
    slide_w_emu = dims["width_emu"]
    slide_h_emu = dims["height_emu"]
    counter = _IdCounter()

    # Theme is built FIRST (it has no component dependency) so the theme body
    # font can drive the per-component font fallback default (US-1.4 MINOR-3).
    theme = _build_theme(prs)
    body_font = (theme.get("font_palette") or {}).get("body") or ""
    default_body = body_font if body_font in _BUILTIN_FONTS else _DEFAULT_FALLBACK_FONT

    # Slide master (AC#2): parse explicitly. A master may legally have zero
    # shapes (e.g., a synthetic minimal deck).
    # US-4.8 (MAJOR-4): tolerate missing master — emit a placeholder entry
    # instead of raising. This makes extract_schema safe to call on a
    # masterless file (Scenario A). The normal pipeline repairs the file
    # BEFORE extraction, but defensive tolerance prevents crashes if the
    # call order ever changes.
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    if not masters:
        slide_master = {
            "name": "(no master)",
            "components": [],
        }
    else:
        master = masters[0]
        master_components = _extract_components(
            master.shapes, slide_w_emu, slide_h_emu, counter, default_body
        )
        slide_master = {
            "name": getattr(master, "name", "Slide Master") or "Slide Master",
            "components": master_components,
        }

    # Layouts (AC#2): enumerate every layout under the master.
    slide_layouts: List[Dict[str, Any]] = []
    seen_ids: Dict[str, int] = {}
    for index, layout in enumerate(prs.slide_layouts):
        components = _extract_components(
            layout.shapes, slide_w_emu, slide_h_emu, counter, default_body
        )
        base_id = _slugify(getattr(layout, "name", "") or f"layout_{index}")
        # Ensure uniqueness across layouts with duplicate names.
        if base_id in seen_ids:
            seen_ids[base_id] += 1
            layout_id = f"{base_id}_{seen_ids[base_id]}"
        else:
            seen_ids[base_id] = 0
            layout_id = base_id
        slide_layouts.append({
            "layout_id": layout_id,
            "layout_name": getattr(layout, "name", "") or f"layout_{index}",
            "layout_index": index,
            "components": components,
        })

    schema: Dict[str, Any] = {
        "template_metadata": _build_metadata(prs, str(path)),
        "slide_master": slide_master,
        "slide_layouts": slide_layouts,
        "component_type_enum": list(COMPONENT_TYPE_ENUM),
        "placeholder_type_enum": list(PLACEHOLDER_TYPE_ENUM),
    }
    schema["theme"] = theme

    # US-1.4: aggregate deduped non-built-in fonts across all text components.
    missing_fonts: List[Dict[str, Any]] = []
    seen_fonts: set = set()
    all_components = list(master_components)
    all_components.extend(c for L in slide_layouts for c in L["components"])
    for comp in all_components:
        font = comp.get("font")
        if not isinstance(font, dict):
            continue
        candidates = [font.get("family")]
        candidates.extend(r.get("font", {}).get("family") for r in comp.get("runs", []))
        for family in candidates:
            if family is None or family in _BUILTIN_FONTS or family in seen_fonts:
                continue
            seen_fonts.add(family)
            missing_fonts.append({
                "family": family,
                "is_available": False,
                "fallback": _font_fallback(family, default_body),
                "download_url": None,
            })
    schema["template_metadata"]["missing_fonts"] = missing_fonts
    if missing_fonts:
        names = ", ".join(m["family"] for m in missing_fonts)
        logger.warning(
            "Template depends on %d non-built-in font(s): %s", len(missing_fonts), names
        )

    logger.info(
        "Extracted schema for %s: master=%d components, %d layouts, %d missing fonts",
        path.name, len(master_components), len(slide_layouts), len(missing_fonts),
    )
    return schema


# ---------------------------------------------------------------------------
# Task 5: validate_template_schema (lightweight, no jsonschema dependency)
# ---------------------------------------------------------------------------
_TOP_LEVEL_REQUIRED = {
    "template_metadata",
    "slide_master",
    "slide_layouts",
    "component_type_enum",
    "placeholder_type_enum",
}
_METADATA_REQUIRED = {
    "title", "schema_version", "generated_by", "generated_at", "slide_dimensions",
}
_DIMENSIONS_REQUIRED = {
    "width_emu", "height_emu", "width_inches", "height_inches", "aspect_ratio",
}
_COMPONENT_REQUIRED = {"id", "type", "name", "polygon", "z_order"}


def _is_number(v: Any) -> bool:
    return isinstance(v, (int, float)) and not isinstance(v, bool)


# Winding check threshold (US-1.2). Below any physically meaningful shape area
# in normalized [0,1] coords (a 1px-tall divider on a 7.5" slide normalizes to
# ~1.4e-3 height; even sub-pixel shapes yield area well above this); above
# float noise (~1e-15 for [0,1] coords over a few terms).
_WINDING_EPSILON = 1e-9


def _signed_area(polygon: List[Dict[str, float]]) -> float:
    """Shoelace signed area of a polygon (normalized coords).

    Positive => algebraically counter-clockwise (the canonical winding emitted by
    :func:`normalize_polygon`, order TL->TR->BR->BL). Negative => reversed.
    ~0 => degenerate/collinear. Works for any n-point simple polygon.
    """
    n = len(polygon)
    if n < 3:
        return 0.0
    s = 0.0
    for i in range(n):
        x1, y1 = polygon[i]["x"], polygon[i]["y"]
        x2, y2 = polygon[(i + 1) % n]["x"], polygon[(i + 1) % n]["y"]
        s += x1 * y2 - x2 * y1
    return s / 2.0


def _validate_component(comp: Any, path: str, result: ValidationResult) -> None:
    if not isinstance(comp, dict):
        result.add(ValidationIssue(f"component must be an object, got {type(comp).__name__}", field_path=path))
        return
    for field in _COMPONENT_REQUIRED:
        if field not in comp:
            result.add(ValidationIssue(f"missing required field '{field}'", field_path=path))
    ctype = comp.get("type")
    if ctype is not None and ctype not in COMPONENT_TYPE_ENUM:
        result.add(ValidationIssue(
            f"type '{ctype}' not in component_type_enum", field_path=f"{path}.type"
        ))
    # type_confidence enum legality + shape/low "flagged for review" (US-1.3).
    tc = comp.get("type_confidence")
    if tc is not None and tc not in (_CONFIDENCE_HIGH, _CONFIDENCE_LOW):
        result.add(ValidationIssue(
            f"type_confidence '{tc}' must be 'high' or 'low'",
            field_path=f"{path}.type_confidence",
        ))
    if ctype == "shape" and tc == _CONFIDENCE_LOW:
        # Scope is intentionally shape-only (locked decision 6, PLAN-GIT-52 v2):
        # an unrecognized element emitted as shape/low is surfaced for review
        # (mirrors the US-1.2 degenerate-polygon warning; is_valid stays True).
        # Indeterminate MEDIA (video/low) is NOT flagged here by design.
        result.add(ValidationIssue(
            "shape with type_confidence 'low' (unrecognized element) -- flagged for review",
            field_path=f"{path}.type_confidence", severity="warning",
        ))
    # placeholder_type enum legality
    ptype = comp.get("placeholder_type")
    if ptype is not None and ptype not in PLACEHOLDER_TYPE_ENUM:
        result.add(ValidationIssue(
            f"placeholder_type '{ptype}' not in placeholder_type_enum",
            field_path=f"{path}.placeholder_type",
        ))
    # polygon: exactly 4 points in [0,1]
    polygon = comp.get("polygon")
    if polygon is not None:
        if not isinstance(polygon, list) or len(polygon) != 4:
            result.add(ValidationIssue(
                f"polygon must have exactly 4 points, got {len(polygon) if isinstance(polygon, list) else 'non-array'}",
                field_path=f"{path}.polygon",
            ))
        else:
            for i, pt in enumerate(polygon):
                if not isinstance(pt, dict) or "x" not in pt or "y" not in pt:
                    result.add(ValidationIssue(
                        f"polygon[{i}] must be {{x,y}}", field_path=f"{path}.polygon[{i}]"
                    ))
                elif not _is_number(pt["x"]) or not _is_number(pt["y"]):
                    result.add(ValidationIssue(
                        f"polygon[{i}] x/y must be numbers", field_path=f"{path}.polygon[{i}]"
                    ))
                else:
                    for axis in ("x", "y"):
                        if not (0.0 <= pt[axis] <= 1.0):
                            result.add(ValidationIssue(
                                f"polygon[{i}].{axis}={pt[axis]} out of [0,1]",
                                field_path=f"{path}.polygon[{i}].{axis}",
                            ))
            # Winding check (US-1.2): only when every point is a numeric {x,y}.
            # Positive signed area => canonical TL->TR->BR->BL (algebraic CCW).
            if all(isinstance(p, dict) and _is_number(p.get("x")) and _is_number(p.get("y"))
                   for p in polygon):
                area = _signed_area(polygon)
                if area < -_WINDING_EPSILON:
                    result.add(ValidationIssue(
                        f"polygon has reversed winding (signed area {area:.4g} < 0); "
                        f"expected canonical TL->TR->BR->BL",
                        field_path=f"{path}.polygon",
                    ))
                elif abs(area) <= _WINDING_EPSILON:
                    result.add(ValidationIssue(
                        "polygon is degenerate/zero-area (collinear/coincident points)",
                        field_path=f"{path}.polygon",
                        severity="warning",
                    ))
    # z_order
    zo = comp.get("z_order")
    if zo is not None and not (isinstance(zo, int) and not isinstance(zo, bool) and zo >= 0):
        result.add(ValidationIssue(f"z_order must be a non-negative int", field_path=f"{path}.z_order"))
    # C1 font-cardinality rule: non-text components must NOT carry font.
    if ctype is not None and ctype not in _TEXT_TYPES and "font" in comp:
        result.add(ValidationIssue(
            f"non-text component type '{ctype}' must not carry a 'font' field",
            field_path=f"{path}.font",
        ))
    # US-4.6 cross-field rules (mirror C1): text_properties only on text-bearing
    # components; image_properties only on image components.
    if ctype is not None and ctype not in _TEXT_TYPES and "text_properties" in comp:
        result.add(ValidationIssue(
            f"non-text component type '{ctype}' must not carry a 'text_properties' field",
            field_path=f"{path}.text_properties",
        ))
    if ctype is not None and ctype != "image" and "image_properties" in comp:
        result.add(ValidationIssue(
            f"non-image component type '{ctype}' must not carry an 'image_properties' field",
            field_path=f"{path}.image_properties",
        ))
    # US-4.6 bullets enum check (when text_properties.bullets is present).
    tp = comp.get("text_properties")
    if isinstance(tp, dict):
        bullets = tp.get("bullets")
        if bullets is not None:
            if not isinstance(bullets, list):
                result.add(ValidationIssue(
                    "text_properties.bullets must be a list",
                    field_path=f"{path}.text_properties.bullets",
                ))
            else:
                _VALID_BULLET_TYPES = {"char", "autonum", "none"}
                for i, b in enumerate(bullets):
                    if not isinstance(b, dict):
                        continue
                    bt = b.get("type")
                    if bt is not None and bt not in _VALID_BULLET_TYPES:
                        result.add(ValidationIssue(
                            f"bullets[{i}].type '{bt}' must be one of "
                            f"{sorted(_VALID_BULLET_TYPES)} or null",
                            field_path=f"{path}.text_properties.bullets[{i}].type",
                            severity="warning",
                        ))
    # US-1.4 font checks (only when a font object is present).
    font = comp.get("font")
    if isinstance(font, dict):
        ia = font.get("is_available")
        if ia is not None and not isinstance(ia, bool):
            result.add(ValidationIssue(
                "font.is_available must be a boolean", field_path=f"{path}.font.is_available"
            ))
        sp = font.get("size_pt")
        if sp is not None and not (_is_number(sp) and not isinstance(sp, bool)):
            result.add(ValidationIssue(
                "font.size_pt must be a number", field_path=f"{path}.font.size_pt"
            ))
        # String-or-null fields (m4: symmetric type checks).
        for key in ("family", "weight", "color", "alignment"):
            v = font.get(key)
            if v is not None and not isinstance(v, str):
                result.add(ValidationIssue(
                    f"font.{key} must be a string or null", field_path=f"{path}.font.{key}"
                ))
        fb = font.get("fallback")
        # AC4: a non-null fallback must always be a built-in font name.
        if fb is not None and fb not in _BUILTIN_FONTS:
            result.add(ValidationIssue(
                f"font.fallback '{fb}' is not a built-in font name (AC4)",
                field_path=f"{path}.font.fallback",
            ))
        # Invariant: is_available == (fallback is None). Covers all four quadrants
        # (incl. the is_available=False / fallback=None violation).
        if ia is not None and ia != (fb is None):
            result.add(ValidationIssue(
                "font.is_available must equal (fallback is None)",
                field_path=f"{path}.font", severity="warning",
            ))


def validate_template_schema(schema_dict: Any) -> ValidationResult:
    """Validate a schema dict against the structural rules of template_schema.json.

    Lightweight (no ``jsonschema`` dependency). Checks top-level keys, metadata,
    slide_master, every slide_layout, and each component (incl. the font-
    cardinality rule from architecture review C1 and polygon value ranges).
    """
    result = ValidationResult()
    if not isinstance(schema_dict, dict):
        result.add(ValidationIssue(f"schema must be a JSON object, got {type(schema_dict).__name__}"))
        return result

    for key in _TOP_LEVEL_REQUIRED:
        if key not in schema_dict:
            result.add(ValidationIssue(f"missing top-level key '{key}'", field_path=key))

    # template_metadata
    meta = schema_dict.get("template_metadata")
    if isinstance(meta, dict):
        for key in _METADATA_REQUIRED:
            if key not in meta:
                result.add(ValidationIssue(f"missing metadata field '{key}'", field_path=f"template_metadata.{key}"))
        dims = meta.get("slide_dimensions")
        if isinstance(dims, dict):
            for key in _DIMENSIONS_REQUIRED:
                if key not in dims:
                    result.add(ValidationIssue(
                        f"missing slide_dimensions field '{key}'",
                        field_path=f"template_metadata.slide_dimensions.{key}",
                    ))
        if "title" in meta and (not isinstance(meta["title"], str) or not meta["title"].strip()):
            result.add(ValidationIssue("title must be a non-empty string", field_path="template_metadata.title"))
        # US-3.1 (MAJOR-2): title_source is runtime-enforced (unlike the general
        # additionalProperties case) so the spec's enum and the validator stay in
        # sync for this field. Keyed off the shared TITLE_SOURCES constant.
        ts = meta.get("title_source")
        if ts is not None and ts not in TITLE_SOURCES:
            result.add(ValidationIssue(
                f"title_source '{ts}' not in {sorted(TITLE_SOURCES)}",
                field_path="template_metadata.title_source",
            ))

    # slide_master
    master = schema_dict.get("slide_master")
    if isinstance(master, dict):
        if not master.get("name"):
            result.add(ValidationIssue("slide_master missing 'name'", field_path="slide_master.name"))
        comps = master.get("components")
        if not isinstance(comps, list):
            result.add(ValidationIssue("slide_master.components must be an array", field_path="slide_master.components"))
        else:
            for i, c in enumerate(comps):
                _validate_component(c, f"slide_master.components[{i}]", result)
    else:
        result.add(ValidationIssue("slide_master must be an object", field_path="slide_master"))

    # slide_layouts
    layouts = schema_dict.get("slide_layouts")
    if isinstance(layouts, list):
        for li, layout in enumerate(layouts):
            if not isinstance(layout, dict):
                result.add(ValidationIssue(f"layout must be an object", field_path=f"slide_layouts[{li}]"))
                continue
            for key in ("layout_id", "layout_name", "layout_index", "components"):
                if key not in layout:
                    result.add(ValidationIssue(f"missing layout field '{key}'", field_path=f"slide_layouts[{li}].{key}"))
            comps = layout.get("components")
            if not isinstance(comps, list):
                result.add(ValidationIssue(
                    "layout.components must be an array", field_path=f"slide_layouts[{li}].components"
                ))
            else:
                for ci, c in enumerate(comps):
                    _validate_component(c, f"slide_layouts[{li}].components[{ci}]", result)
    else:
        result.add(ValidationIssue("slide_layouts must be an array", field_path="slide_layouts"))

    # enums self-document
    for key, valid in (("component_type_enum", COMPONENT_TYPE_ENUM),
                       ("placeholder_type_enum", PLACEHOLDER_TYPE_ENUM)):
        val = schema_dict.get(key)
        if not isinstance(val, list):
            result.add(ValidationIssue(f"{key} must be an array", field_path=key))

    # US-1.4 (AC3): surface each non-built-in font as a non-fatal warning so the
    # subagent/consumer is "flagged for review" (mirrors US-1.2/1.3 warnings).
    meta = schema_dict.get("template_metadata")
    if isinstance(meta, dict):
        missing = meta.get("missing_fonts")
        if isinstance(missing, list):
            for entry in missing:
                if isinstance(entry, dict) and entry.get("family"):
                    result.add(ValidationIssue(
                        f"non-built-in font '{entry['family']}' is required by the template "
                        f"(fallback: {entry.get('fallback')}) -- install it to avoid substitution",
                        field_path="template_metadata.missing_fonts", severity="warning",
                    ))

    return result


# ---------------------------------------------------------------------------
# US-1.5: embed / read the proposed schema inside the PPTX zip
# ---------------------------------------------------------------------------
def _inject_json_default(content_types_xml: bytes) -> bytes:
    """Return ``[Content_Types].xml`` bytes with a ``<Default Extension="json"/>`` added.

    Idempotent: if a ``json`` Default is already declared, the input is returned
    unchanged. Otherwise a ``<Default Extension="json" ContentType="application/json"/>``
    is inserted as the first child of ``<Types>`` (OPC: Defaults precede Overrides).
    Deliberately modifies ``[Content_Types].xml`` (architecture review MAJOR-2): the
    bundled template declares no ``json`` Default and strict PowerPoint may otherwise
    offer to repair an undeclared ``ppt/template_schema.json`` part.
    """
    root = etree.fromstring(content_types_xml)
    for default in root.findall(_CONTENT_TYPES_NS + "Default"):
        if default.get("Extension") == "json":
            return content_types_xml  # already declared
    new_default = etree.Element(_CONTENT_TYPES_NS + "Default")
    new_default.set("Extension", "json")
    new_default.set("ContentType", _JSON_CONTENT_TYPE)
    root.insert(0, new_default)  # Defaults precede Overrides
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def embed_schema(
    pptx_path: str, schema: Dict[str, Any], output_pptx_path: str
) -> EmbeddedSchemaResult:
    """Embed ``schema`` into a COPY of the PPTX at ``pptx_path`` (US-1.5).

    Order-preserving full zip rewrite: ``[Content_Types].xml`` first (with the
    ``json`` Default injected), then every other original entry decompressed-content-
    identical in original order (an existing ``ppt/template_schema.json`` is skipped
    so embedding is idempotent — MAJOR-3), then the minified schema appended last.
    Written atomically (temp file + ``os.replace`` — MINOR-6). Never modifies the
    input in place. Returns an :class:`EmbeddedSchemaResult` (AC4 size delta).
    """
    in_path = Path(pptx_path)
    out_path = Path(output_pptx_path)
    original_bytes = in_path.stat().st_size
    schema_payload = json.dumps(
        schema, separators=(",", ":"), ensure_ascii=False
    ).encode("utf-8")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp_name = tempfile.mkstemp(suffix=".pptx", dir=str(out_path.parent))
    os.close(fd)
    try:
        with zipfile.ZipFile(str(in_path), "r") as zin, \
                zipfile.ZipFile(tmp_name, "w", zipfile.ZIP_DEFLATED) as zout:
            names = zin.namelist()
            # [Content_Types].xml first (with the json Default injected).
            if _CONTENT_TYPES_XML in names:
                ct_info = zin.getinfo(_CONTENT_TYPES_XML)
                ct_out = zipfile.ZipInfo(filename=_CONTENT_TYPES_XML, date_time=ct_info.date_time)
                ct_out.compress_type = ct_info.compress_type
                ct_out.external_attr = ct_info.external_attr
                zout.writestr(ct_out, _inject_json_default(zin.read(_CONTENT_TYPES_XML)))
            # Every other original entry, decompressed-content-identical, original
            # order. Skip the schema path (idempotency) and [Content_Types].xml.
            for name in names:
                if name == _CONTENT_TYPES_XML or name == _EMBEDDED_SCHEMA_PATH:
                    continue
                info = zin.getinfo(name)
                out_info = zipfile.ZipInfo(filename=name, date_time=info.date_time)
                out_info.compress_type = info.compress_type
                out_info.external_attr = info.external_attr
                zout.writestr(out_info, zin.read(name))
            # The new schema part, last.
            schema_info = zipfile.ZipInfo(filename=_EMBEDDED_SCHEMA_PATH)
            schema_info.compress_type = zipfile.ZIP_DEFLATED
            zout.writestr(schema_info, schema_payload)
        os.replace(tmp_name, str(out_path))
    except Exception:
        if os.path.exists(tmp_name):
            os.remove(tmp_name)
        raise

    new_bytes = out_path.stat().st_size
    logger.info(
        "embedded schema into %s: %d -> %d bytes (%+d)",
        out_path, original_bytes, new_bytes, new_bytes - original_bytes,
    )
    return EmbeddedSchemaResult(str(out_path), original_bytes, new_bytes, new_bytes - original_bytes)


def _parse_version(v: str) -> Optional[tuple]:
    """Parse a ``"major.minor.patch"`` string into a 3-tuple of ints. None on failure."""
    try:
        parts = v.split(".")
        if len(parts) != 3:
            return None
        return (int(parts[0]), int(parts[1]), int(parts[2]))
    except (ValueError, AttributeError):
        return None


def _check_schema_version(data: Dict[str, Any], pptx_path: str) -> None:
    """BT-142 Phase 1.4 — negotiate the embedded schema version against ``SCHEMA_VERSION``.

    Behavior:
      - patch mismatch (e.g. embedded 1.1.1 vs current 1.1.0) → silent (forward-compatible)
      - minor mismatch (e.g. embedded 1.0.x vs current 1.1.x) → ``logger.warning`` + return data
      - major mismatch (e.g. embedded 2.x.x vs current 1.x.x) → raise :class:`SchemaVersionError`
      - missing / malformed version field → ``logger.warning`` + return data (defensive)

    Callers should let :class:`SchemaVersionError` propagate — it indicates the
    schema is from an incompatible generation and consumers will misinterpret fields.
    """
    metadata = data.get("template_metadata") if isinstance(data, dict) else None
    if not isinstance(metadata, dict):
        logger.warning(
            "embedded schema at %s has no template_metadata; skipping version check",
            pptx_path,
        )
        return

    embedded_raw = metadata.get("schema_version")
    embedded = _parse_version(embedded_raw) if isinstance(embedded_raw, str) else None
    current = _parse_version(SCHEMA_VERSION)

    if embedded is None:
        logger.warning(
            "embedded schema at %s has missing/malformed schema_version (%r); skipping version check",
            pptx_path, embedded_raw,
        )
        return
    if current is None:  # defensive — SCHEMA_VERSION should always parse
        return

    major_e, minor_e, patch_e = embedded
    major_c, minor_c, patch_c = current

    if major_e != major_c:
        raise SchemaVersionError(
            f"embedded schema major version {embedded_raw} is incompatible with "
            f"current SCHEMA_VERSION {SCHEMA_VERSION} (major mismatch). "
            f"Re-extract the template with the current engine, or pin to a compatible version."
        )
    if minor_e != minor_c:
        logger.warning(
            "embedded schema minor version %s differs from current %s (additive fields may be missing/extra); proceeding",
            embedded_raw, SCHEMA_VERSION,
        )
        return
    if patch_e != patch_c:
        # Patch is forward-compatible — silent. The next embed re-writes the current version.
        return


def read_embedded_schema(pptx_path: str) -> Optional[Dict[str, Any]]:
    """Read the embedded ``ppt/template_schema.json`` from a PPTX (US-1.5).

    Error contract (architecture review MINOR-3/4 + BT-142 Phase 1.4):
      - valid zip + present        -> the schema dict
      - valid zip + absent         -> ``None``
      - malformed / non-object JSON -> ``logger.warning`` + ``None`` (corrupt-as-absent)
      - corrupt zip / unreadable    -> raise :class:`TemplateExtractionError`
      - major version mismatch      -> raise :class:`SchemaVersionError` (BT-142 Phase 1.4)
      - minor version mismatch      -> ``logger.warning`` + return data
      - patch version mismatch      -> silent (forward-compatible)
    """
    try:
        with zipfile.ZipFile(str(pptx_path), "r") as z:
            if _EMBEDDED_SCHEMA_PATH not in z.namelist():
                return None
            raw = z.read(_EMBEDDED_SCHEMA_PATH)
    except zipfile.BadZipFile as exc:
        raise TemplateExtractionError(
            f"not a valid zip/pptx: {pptx_path} ({exc.__class__.__name__}: {exc})"
        ) from exc
    except OSError as exc:  # FileNotFoundError, permission errors, ...
        raise TemplateExtractionError(
            f"could not read pptx: {pptx_path} ({exc.__class__.__name__}: {exc})"
        ) from exc

    try:
        data = json.loads(raw.decode("utf-8"))
    except (ValueError, UnicodeDecodeError) as exc:
        logger.warning(
            "embedded schema at %s is malformed (%s); treating as absent",
            _EMBEDDED_SCHEMA_PATH, exc,
        )
        return None
    if not isinstance(data, dict):
        logger.warning(
            "embedded schema at %s is not a JSON object; treating as absent",
            _EMBEDDED_SCHEMA_PATH,
        )
        return None
    _check_schema_version(data, str(pptx_path))  # BT-142 Phase 1.4
    return data


# ---------------------------------------------------------------------------
# US-3.1: human-readable extraction summary (US-3.3 AC2)
# ---------------------------------------------------------------------------
def build_extraction_summary(schema: Dict[str, Any]) -> str:
    """Render a multi-line, human-readable summary of an extracted schema.

    Pure (no I/O, no logging) so it stays unit-testable. Consumed by the
    generate-template-skill (Stage 4) and the CLI ``--summary`` flag. The exact
    wording is **not** part of any contract; callers may reformat.
    """
    meta = schema.get("template_metadata") or {}
    dims = meta.get("slide_dimensions") or {}
    master = schema.get("slide_master") or {}
    master_comps = master.get("components") or []
    layouts = schema.get("slide_layouts") or []
    theme = schema.get("theme") or {}
    palette = theme.get("font_palette") or {}
    missing = meta.get("missing_fonts") or []

    lines: List[str] = []
    title = meta.get("title") or "(untitled)"
    source = meta.get("title_source")
    lines.append(f"Title: {title}" + (f"  [source: {source}]" if source else ""))
    w = dims.get("width_inches")
    h = dims.get("height_inches")
    ar = dims.get("aspect_ratio")
    size_str = ""
    if w is not None and h is not None:
        size_str = f"{w:g} x {h:g} in"
    if ar:
        size_str = f"{size_str}  ({ar})" if size_str else str(ar)
    if size_str:
        lines.append(f"Slide size: {size_str}")
    lines.append(f"Slide master: {master.get('name') or '(unnamed)'} "
                 f"({len(master_comps)} component{'s' if len(master_comps) != 1 else ''})")
    lines.append(f"Layouts: {len(layouts)}")
    for layout in layouts:
        if isinstance(layout, dict):
            name = layout.get("layout_name") or layout.get("layout_id") or "(unnamed)"
            ncomp = len(layout.get("components") or [])
            lines.append(f"  - {name}: {ncomp} component{'s' if ncomp != 1 else ''}")

    color_keys = ("primary_color", "secondary_color", "accent_color", "background_color", "text_color")
    colors = [f"{k.replace('_color', '')}={theme[k]}" for k in color_keys if theme.get(k)]
    if colors:
        lines.append("Theme colors: " + ", ".join(colors))
    font_keys = ("heading", "body", "accent")
    fonts = [f"{k}={palette[k]}" for k in font_keys if palette.get(k)]
    if fonts:
        lines.append("Font palette: " + ", ".join(fonts))

    if missing:
        names = ", ".join(str(m.get("family")) for m in missing if isinstance(m, dict) and m.get("family"))
        lines.append(f"Missing fonts: {len(missing)}" + (f" ({names})" if names else ""))
    else:
        lines.append("Missing fonts: 0")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# US-2.1: header/footer consumer-facing helpers
# ---------------------------------------------------------------------------

def needs_header_footer_prompt(schema: Dict[str, Any]) -> bool:
    """True when the template has **neither** header nor footer (US-2.1 AC2).

    Pure function on the extracted schema dict; consumed by
    ``generate-template-skill`` (full prompt + inject) and ``pptx-subagent``
    (light informational note). Headless/subagent callers skip the prompt.
    """
    hf = (schema.get("template_metadata") or {}).get("header_footer") or {}
    return not hf.get("has_header") and not hf.get("has_footer")


def inject_default_header_zone(schema: Dict[str, Any]) -> None:
    """Inject a default header zone into the schema metadata (US-2.1 AC3).

    **Schema-only**; never touches the PPTX. The polygon is a 4-point
    normalised top-strip per the US-1.2 model (TL→TR→BR→BL, ``[0,1]`` range).
    Called by ``generate-template-skill`` when the user opts to add a header.
    """
    meta = schema.setdefault("template_metadata", {})
    hf = meta.setdefault("header_footer", {})
    hf.setdefault("header", {
        "source": "user_default",
        "polygon": [
            {"x": 0, "y": 0}, {"x": 1, "y": 0},
            {"x": 1, "y": 0.05}, {"x": 0, "y": 0.05},
        ],
        "note": "Default header zone (top strip); metadata only — not rendered "
                "into the PPTX until a real header placeholder is added",
    })


# ---------------------------------------------------------------------------
# Task 7: CLI entry
# ---------------------------------------------------------------------------
def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        description="Extract a PPTX into the proposed template schema JSON (US-1.1)."
    )
    parser.add_argument("--input", "-i", required=True, help="path to the input .pptx")
    parser.add_argument("--output", "-o", required=True, help="path to write the schema JSON")
    parser.add_argument("--log-level", default="info", help="log level (debug/info/warn/error)")
    parser.add_argument(
        "--embed", action="store_true",
        help="also embed the schema into a PPTX copy (US-1.5); see --output-pptx",
    )
    parser.add_argument(
        "--output-pptx", default=None,
        help="destination for the embedded PPTX when --embed is set (default: <input>.templated.pptx)",
    )
    parser.add_argument(
        "--summary", action="store_true",
        help="print a human-readable extraction summary to stdout (US-3.1 / US-3.3 AC2)",
    )
    args = parser.parse_args(argv)

    if args.output_pptx and not args.embed:
        logger.error("--output-pptx requires --embed")
        return 2

    level_name = str(args.log_level).upper()
    if level_name not in {"DEBUG", "INFO", "WARN", "WARNING", "ERROR"}:
        logger.error("invalid --log-level '%s'", args.log_level)
        return 2
    logging.basicConfig(level=level_name)

    try:
        schema = extract_schema(args.input)
    except TemplateExtractionError as exc:
        logger.error("extraction failed: %s", exc)
        return 2  # runtime error
    validation = validate_template_schema(schema)
    if not validation.is_valid:
        logger.error("extracted schema failed validation (%d errors):", len(validation.errors))
        for msg in validation.error_messages():
            logger.error("  %s", msg)
        return 1  # validation error

    try:
        Path(args.output).write_text(
            json.dumps(schema, indent=2, ensure_ascii=False), encoding="utf-8"
        )
    except OSError as exc:
        logger.error("could not write output '%s': %s", args.output, exc)
        return 2
    logger.info("wrote schema to %s", args.output)

    if args.summary:
        print(build_extraction_summary(schema))

    if args.embed:
        out_pptx = args.output_pptx
        if out_pptx is None:
            in_pptx = Path(args.input)
            out_pptx = str(in_pptx.with_name(in_pptx.stem + ".templated.pptx"))
        try:
            result = embed_schema(args.input, schema, out_pptx)
        except (OSError, TemplateExtractionError) as exc:
            logger.error("embedding failed: %s", exc)
            return 2
        logger.info(
            "wrote embedded PPTX to %s (%d -> %d bytes%+d)",
            result.output_path, result.original_bytes, result.new_bytes, result.delta_bytes,
        )
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
