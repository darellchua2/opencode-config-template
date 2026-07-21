"""US-3.5 tests — slide-master text-defaults capture.

Covers ``_extract_master_text_styles`` and its integration into
``extract_schema`` (``slide_master.text_defaults``). These defaults are the
"default font" a user inherits when typing into a placeholder or plain text
box with no run-level override — read from the master's ``<p:txStyles>``.
"""
import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

_HERE = Path(__file__).resolve().parent
_COMMON_SCRIPTS = _HERE.parent
if str(_COMMON_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_COMMON_SCRIPTS))

from schema_extractor import (  # noqa: E402
    SCHEMA_VERSION,
    _extract_master_text_styles,
    extract_schema,
)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
FONT = "Century Gothic"


def _set_master_txstyles(prs, title_tf, body_tf, other_tf):
    """Replace the slide master's <p:txStyles> with explicit typefaces."""
    master = prs.slide_masters[0]

    def lvl(tf, color="18181B", sz=1600, bold=False):
        return (
            '<a:lvl1pPr marL="0" indent="0">'
            '<a:defRPr sz="%d" b="%d">'
            '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
            '<a:latin typeface="%s"/>'
            '</a:defRPr></a:lvl1pPr>'
        ) % (sz, 1 if bold else 0, color, tf)

    txs_xml = (
        '<p:txStyles xmlns:a="%s" xmlns:p="%s">'
        '<p:titleStyle>%s</p:titleStyle>'
        '<p:bodyStyle>%s</p:bodyStyle>'
        '<p:otherStyle>%s</p:otherStyle>'
        '</p:txStyles>'
    ) % (_NS_A, _NS_P,
         lvl(title_tf, color="18181B", sz=3200, bold=True),
         lvl(body_tf, color="18181B", sz=1600, bold=False),
         lvl(other_tf, color="18181B", sz=1400, bold=False))
    new = etree.fromstring(txs_xml)
    existing = master.element.find("{%s}txStyles" % _NS_P)
    if existing is not None:
        master.element.replace(existing, new)
    else:
        master.element.append(new)
    return master


def test_explicit_typefaces_captured():
    prs = Presentation()
    master = _set_master_txstyles(prs, FONT, FONT, FONT)
    td = _extract_master_text_styles(master, {"major": "Calibri Light", "minor": "Calibri"})
    assert set(td) == {"title", "body", "other"}
    assert td["title"] == {"font": FONT, "size_pt": 32.0, "color": "#18181B", "bold": True}
    assert td["body"] == {"font": FONT, "size_pt": 16.0, "color": "#18181B", "bold": False}
    assert td["other"]["font"] == FONT
    assert td["other"]["size_pt"] == 14.0


def test_theme_reference_resolution():
    prs = Presentation()
    master = _set_master_txstyles(prs, "+mj-lt", "+mn-lt", "+mn-lt")
    theme_fonts = {"major": "Calibri Light", "minor": "Calibri"}
    td = _extract_master_text_styles(master, theme_fonts)
    assert td["title"]["font"] == "Calibri Light"
    assert td["body"]["font"] == "Calibri"
    assert td["other"]["font"] == "Calibri"


def test_missing_txstyles_returns_empty():
    prs = Presentation()
    master = prs.slide_masters[0]
    txs = master.element.find("{%s}txStyles" % _NS_P)
    if txs is not None:
        master.element.remove(txs)
    td = _extract_master_text_styles(master, {"major": "X", "minor": "Y"})
    assert td == {}


def test_extract_schema_emits_text_defaults(tmp_path):
    prs = Presentation()
    _set_master_txstyles(prs, FONT, FONT, FONT)
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    schema = extract_schema(str(out))
    assert schema["template_metadata"]["schema_version"] == SCHEMA_VERSION
    td = schema["slide_master"].get("text_defaults")
    assert isinstance(td, dict)
    assert td["title"]["font"] == FONT
    assert td["body"]["size_pt"] == 16.0


def test_masterless_schema_has_empty_text_defaults():
    schema = {"slide_master": {"name": "(no master)", "components": [], "text_defaults": {}}}
    assert schema["slide_master"]["text_defaults"] == {}
