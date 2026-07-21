"""US-4.8 Phase 1 tests — master repair cascade (Scenario A).

Tests the three-level cascade (Chain of Responsibility):
  Level 1 — salvage ``ppt/theme/theme1.xml`` from the zip (exact fidelity)
  Level 2 — scavenge explicit styles from slide XML (best-effort)
  Level 3 — fallback to in-code minimal theme (BT-142 Phase 1.5: was default.pptx)

Fixture strategy (MAJOR-6 / R-FIXTURE-1): masterless ``.pptx`` files cannot be
created by python-pptx — they are built by zip surgery that strips the slide
master part + rels from a normal ``.pptx``.

BT-142 Phase 1.5: tests no longer depend on a bundled ``template/default.pptx``.
The ``default_pptx`` fixture synthesizes a minimal valid PPTX in ``tmp_path``
via ``_build_minimal_pptx_bytes(None)`` — same in-code fallback the production
code uses when ``default_template_path=None``.
"""
import shutil
import sys
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from lxml import etree
from pptx import Presentation

_HERE = Path(__file__).resolve().parent
_COMMON_SCRIPTS = _HERE.parent  # .../_common/scripts — where master_repairer.py lives
_SKILLS = _COMMON_SCRIPTS.parent.parent  # .../.opencode/skills
_REPO_ROOT = _SKILLS.parent.parent
_FILLER_SCRIPTS = _SKILLS / "generate-slide-skill" / "scripts"
for _p in (str(_COMMON_SCRIPTS), str(_FILLER_SCRIPTS)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from master_repairer import (  # noqa: E402
    RepairResult,
    _build_minimal_pptx_bytes,
    _salvage_theme_part,
    _scavenge_slide_styles,
    repair_if_needed,
)
from schema_extractor import parse_theme_xml  # noqa: E402

# BT-142 Phase 1.5: no longer reads ``template/default.pptx``. The ``default_pptx``
# fixture (session-scoped) synthesizes a minimal valid PPTX via the same in-code
# fallback the production code uses.

# Namespaces for zip-level XML surgery.
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Fixture builder (R-FIXTURE-1 / T5.1.1)
# ---------------------------------------------------------------------------

def _make_masterless_fixture(
    src_pptx: str,
    dst_pptx: str,
    keep_theme: bool = True,
) -> str:
    """Strip the slide master from ``src_pptx`` to create a masterless fixture.

    Zip-level surgery:
      1. Remove ``ppt/slideMasters/slideMaster1.xml`` + its rels.
      2. Empty ``<p:sldMasterIdLst>`` in ``ppt/presentation.xml``.
      3. Remove the master relationship from ``ppt/_rels/presentation.xml.rels``.
      4. Remove the master Override from ``[Content_Types].xml``.
      5. Optionally keep ``ppt/theme/theme1.xml`` (for Level-1 tests).

    The result opens in python-pptx with ``len(prs.slide_masters) == 0``.
    """
    with zipfile.ZipFile(src_pptx, "r") as zin:
        names = zin.namelist()
        # --- Strip master parts ---
        strip_prefixes = ("ppt/slideMasters/",)
        # Theme handling
        if not keep_theme:
            strip_prefixes += ("ppt/theme/",)

        with zipfile.ZipFile(dst_pptx, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                if any(name.startswith(p) for p in strip_prefixes):
                    continue  # skip master parts (and optionally theme)

                data = zin.read(name)

                # --- Empty <p:sldMasterIdLst> in presentation.xml ---
                if name == "ppt/presentation.xml":
                    root = etree.parse(BytesIO(data)).getroot()
                    for sml in root.iter(f"{{{_NS_P}}}sldMasterIdLst"):
                        for child in list(sml):
                            sml.remove(child)
                    data = etree.tostring(
                        root, xml_declaration=True, encoding="UTF-8", standalone=True
                    )

                # --- Remove master relationship from presentation.xml.rels ---
                elif name == "ppt/_rels/presentation.xml.rels":
                    root = etree.parse(BytesIO(data)).getroot()
                    for rel in list(root):
                        target = rel.get("Target", "")
                        if "slideMaster" in target:
                            root.remove(rel)
                    data = etree.tostring(
                        root, xml_declaration=True, encoding="UTF-8", standalone=True
                    )

                # --- Remove master Override from [Content_Types].xml ---
                elif name == "[Content_Types].xml":
                    root = etree.parse(BytesIO(data)).getroot()
                    for override in list(root):
                        part_name = override.get("PartName", "")
                        if "slideMaster" in part_name:
                            root.remove(override)
                    data = etree.tostring(
                        root, xml_declaration=True, encoding="UTF-8", standalone=True
                    )

                info = zin.getinfo(name)
                out_info = zipfile.ZipInfo(filename=name, date_time=info.date_time)
                out_info.compress_type = info.compress_type
                out_info.external_attr = info.external_attr
                zout.writestr(out_info, data)

    return dst_pptx


@pytest.fixture
def default_pptx(tmp_path):
    """BT-142 Phase 1.5: synthesize a minimal valid PPTX in tmp_path.

    Replaces the legacy ``template/default.pptx`` dependency. The synthesized
    PPTX has one slide_master + one blank layout + a default theme — the same
    fallback the production code uses when ``default_template_path=None``.
    """
    p = tmp_path / "default.pptx"
    p.write_bytes(_build_minimal_pptx_bytes(None))
    return str(p)


@pytest.fixture
def masterless_with_theme(tmp_path, default_pptx):
    """A-L1: no master, but ppt/theme/theme1.xml survives in the zip."""
    dst = str(tmp_path / "masterless_theme.pptx")
    return _make_masterless_fixture(default_pptx, dst, keep_theme=True)


@pytest.fixture
def masterless_no_theme(tmp_path, default_pptx):
    """A-L2/A-L3: no master, no theme part."""
    dst = str(tmp_path / "masterless_no_theme.pptx")
    return _make_masterless_fixture(default_pptx, dst, keep_theme=False)


# ---------------------------------------------------------------------------
# Fixture builder self-test (R-FIXTURE-1 mitigation)
# ---------------------------------------------------------------------------

class TestFixtureBuilder:
    def test_masterless_fixture_has_no_master(self, masterless_with_theme):
        prs = Presentation(masterless_with_theme)
        assert len(prs.slide_masters) == 0

    def test_masterless_with_theme_keeps_theme_part(self, masterless_with_theme):
        theme = _salvage_theme_part(masterless_with_theme)
        assert theme is not None

    def test_masterless_no_theme_has_no_theme_part(self, masterless_no_theme):
        theme = _salvage_theme_part(masterless_no_theme)
        assert theme is None


# ---------------------------------------------------------------------------
# Level 1: salvage theme1.xml from the zip
# ---------------------------------------------------------------------------

class TestLevel1Salvage:
    def test_salvage_returns_theme_bytes(self, masterless_with_theme):
        theme_bytes = _salvage_theme_part(masterless_with_theme)
        assert theme_bytes is not None
        assert b"themeElements" in theme_bytes or b"clrScheme" in theme_bytes

    def test_salvage_returns_none_when_no_theme(self, masterless_no_theme):
        assert _salvage_theme_part(masterless_no_theme) is None

    def test_salvaged_theme_has_colors_and_fonts(self, masterless_with_theme):
        theme_bytes = _salvage_theme_part(masterless_with_theme)
        colors, fonts = parse_theme_xml(theme_bytes)
        # default.pptx's theme has all 12 color roles + major/minor fonts.
        assert len(colors) >= 6  # at least dk1, lt1, dk2, lt2 + some accents
        assert "major_latin" in fonts or "minor_latin" in fonts


# ---------------------------------------------------------------------------
# Level 2: scavenge slide styles
# ---------------------------------------------------------------------------

class TestLevel2Scavenge:
    def test_scavenge_finds_fonts_in_slides(self, masterless_no_theme, default_pptx):
        """Level 2 should produce a synthesized theme from slide explicit styles."""
        result = _scavenge_slide_styles(masterless_no_theme, default_pptx)
        # The masterless_no_theme fixture is derived from default.pptx whose
        # slides carry explicit <a:rPr> elements → scavenge should succeed.
        if result is not None:
            # Verify the synthesized theme has a fontScheme.
            assert b"fontScheme" in result or b"majorFont" in result

    def test_scavenge_returns_none_when_no_styles(self, tmp_path, default_pptx):
        """Level 2 returns None when slides have no explicit styles at all."""
        # Build a fixture whose slides have no <a:rPr> / <p:spPr> styles.
        # We reuse _make_masterless_fixture but then strip all style attributes.
        import zipfile as zf
        from lxml import etree as ET

        base = str(tmp_path / "no_styles.pptx")
        _make_masterless_fixture(default_pptx, base, keep_theme=False)

        # Post-process: strip all <a:rPr> and <a:solidFill> from slide XML.
        stripped = str(tmp_path / "no_styles_stripped.pptx")
        with zf.ZipFile(base, "r") as zin, zf.ZipFile(stripped, "w", zf.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename.startswith("ppt/slides/") and info.filename.endswith(".xml"):
                    root = ET.parse(BytesIO(data)).getroot()
                    # Remove all <a:rPr> elements.
                    for rpr in root.iter(f"{{{_NS_A}}}rPr"):
                        parent = rpr.getparent()
                        if parent is not None:
                            parent.remove(rpr)
                    # Remove all <a:solidFill> inside <p:spPr>.
                    _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
                    for spPr in root.iter(f"{{{_NS_P}}}spPr"):
                        for sf in spPr.iter(f"{{{_NS_A}}}solidFill"):
                            spPr.remove(sf)
                    data = ET.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                out_info = zf.ZipInfo(filename=info.filename, date_time=info.date_time)
                out_info.compress_type = info.compress_type
                out_info.external_attr = info.external_attr
                zout.writestr(out_info, data)

        result = _scavenge_slide_styles(stripped, default_pptx)
        assert result is None  # no recoverable styles → None


# ---------------------------------------------------------------------------
# repair_if_needed — the cascade
# ---------------------------------------------------------------------------

class TestRepairIfNeeded:
    def test_no_repair_needed_when_master_present(self, default_pptx):
        prs = Presentation(default_pptx)
        result = repair_if_needed(prs, default_pptx, default_pptx)
        assert result.level == "none"
        assert result.mutated is False

    def test_level1_repair_when_theme_survives(self, masterless_with_theme, default_pptx):
        prs = Presentation(masterless_with_theme)
        result = repair_if_needed(prs, masterless_with_theme, default_pptx)
        assert result.level == "L1"
        assert result.mutated is True
        assert result.theme_source == "salvaged"
        assert result.repaired_path is not None
        assert Path(result.repaired_path).exists()

    def test_level3_repair_when_no_theme_no_styles(self, tmp_path, default_pptx):
        """Genuine Level 3: no theme part + no recoverable slide styles."""
        import zipfile as zf
        from lxml import etree as ET

        # Build a fixture with no master, no theme, and no slide styles.
        base = str(tmp_path / "bare.pptx")
        _make_masterless_fixture(default_pptx, base, keep_theme=False)
        stripped = str(tmp_path / "bare_stripped.pptx")
        with zf.ZipFile(base, "r") as zin, zf.ZipFile(stripped, "w", zf.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename.startswith("ppt/slides/") and info.filename.endswith(".xml"):
                    root = ET.parse(BytesIO(data)).getroot()
                    for rpr in root.iter(f"{{{_NS_A}}}rPr"):
                        parent = rpr.getparent()
                        if parent is not None:
                            parent.remove(rpr)
                    _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
                    for spPr in root.iter(f"{{{_NS_P}}}spPr"):
                        for sf in spPr.iter(f"{{{_NS_A}}}solidFill"):
                            spPr.remove(sf)
                    data = ET.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                out_info = zf.ZipInfo(filename=info.filename, date_time=info.date_time)
                out_info.compress_type = info.compress_type
                out_info.external_attr = info.external_attr
                zout.writestr(out_info, data)

        prs = Presentation(stripped)
        result = repair_if_needed(prs, stripped, default_pptx)
        assert result.mutated is True
        assert result.level == "L3"
        # BT-142 Phase 1.5: theme_source renamed from "default" to "minimal-in-code"
        # (or "default" if legacy path with explicit default_template_path — both accepted)
        assert result.theme_source in ("default", "minimal-in-code")
        assert result.repaired_path is not None

    def test_repaired_file_has_master(self, masterless_with_theme, default_pptx):
        prs = Presentation(masterless_with_theme)
        result = repair_if_needed(prs, masterless_with_theme, default_pptx)
        repaired_prs = Presentation(result.repaired_path)
        assert len(repaired_prs.slide_masters) > 0

    def test_repaired_file_has_layouts(self, masterless_with_theme, default_pptx):
        prs = Presentation(masterless_with_theme)
        result = repair_if_needed(prs, masterless_with_theme, default_pptx)
        repaired_prs = Presentation(result.repaired_path)
        assert len(repaired_prs.slide_layouts) > 0

    def test_original_file_untouched(self, masterless_with_theme, default_pptx):
        original_bytes = Path(masterless_with_theme).read_bytes()
        prs = Presentation(masterless_with_theme)
        repair_if_needed(prs, masterless_with_theme, default_pptx)
        after_bytes = Path(masterless_with_theme).read_bytes()
        assert original_bytes == after_bytes


# ---------------------------------------------------------------------------
# Cascade priority
# ---------------------------------------------------------------------------

class TestCascadePriority:
    def test_level1_wins_over_level2(self, masterless_with_theme, default_pptx):
        """When theme1.xml is present in the zip, Level 1 is used (not L2/L3)."""
        prs = Presentation(masterless_with_theme)
        result = repair_if_needed(prs, masterless_with_theme, default_pptx)
        assert result.level == "L1"


# ---------------------------------------------------------------------------
# parse_theme_xml — pure function test (CRIT-2)
# ---------------------------------------------------------------------------

class TestParseThemeXml:
    def test_parses_default_theme(self, default_pptx):
        with zipfile.ZipFile(default_pptx) as z:
            theme_bytes = z.read("ppt/theme/theme1.xml")
        colors, fonts = parse_theme_xml(theme_bytes)
        assert len(colors) > 0
        # Should have standard roles.
        assert "dk1" in colors or "lt1" in colors

    def test_empty_bytes_returns_empty(self):
        colors, fonts = parse_theme_xml(b"")
        assert colors == {}
        assert fonts == {}

    def test_garbage_bytes_returns_empty(self):
        colors, fonts = parse_theme_xml(b"not xml at all")
        assert colors == {}
        assert fonts == {}
