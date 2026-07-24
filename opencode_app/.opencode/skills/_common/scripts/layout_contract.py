"""layout_contract.py — shared template-contract / layout-matching layer (PLAN-GIT-72).

The pure, presentation-agnostic logic that derives a "render contract" from a
template (embedded-schema-first, sidecar fallback) and matches a slide_type to
the best layout by placeholder-composition fingerprint. Physically lives in the
shared ``_common/scripts`` package so all three skills (generate-slide,
generate-template, template-modifier) import it from one place instead of
parasitically reaching into ``generate-slide-skill/scripts/ppt_builder.py``.

Extracted from ``ppt_builder.py`` (PLAN-GIT-72 Phase 2 / architecture-review C1
closure). The 13 symbols here are pure — they operate on the contract dict /
read embedded JSON and do NOT bind to a live fill-side ``Presentation`` object
(except :func:`_live_layout_count`, which opens a template read-only to count
layouts — a contract probe, not fill rendering).

Invariant: this module MUST NOT import back into ``ppt_builder`` (enforced by
the PLAN-GIT-72 test matrix grep).
"""
from __future__ import annotations

import logging
import re
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from contract_adapter import embedded_schema_to_contract
from schema_extractor import TemplateExtractionError, read_embedded_schema
from template_introspector import get_contract

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_LAYOUT_NAME_MAP: Dict[str, List[str]] = {
    "title_slide": ["Title Slide"],
    "closing_slide": ["End"],
    "section_header_slide": ["Section Header"],
    "content_slide": ["Title and Content"],
    "two_content_slide": ["7_Two Content"],
    "comparison_slide": ["Comparison"],
    "content_image_slide": ["Picture with Caption"],
    "chart_slide": ["Blank"],
}

# Issue #44 (P1): ideal placeholder-composition fingerprint per slide_type.
# Used by _resolve_layout_by_fingerprint() so the engine fills ANY template by
# placeholder composition — layout NAMES become a tie-breaker / fallback, not
# the primary key (DESIGN §6 A2/A3).
_SLIDE_TYPE_FINGERPRINT: Dict[str, List[str]] = {
    "title_slide": ["TITLE", "SUBTITLE"],
    "closing_slide": ["TITLE", "SUBTITLE"],
    "section_header_slide": ["TITLE", "SUBTITLE"],
    "content_slide": ["TITLE", "OBJECT"],
    "two_content_slide": ["TITLE", "OBJECT", "OBJECT"],
    "comparison_slide": ["TITLE", "OBJECT", "OBJECT"],
    "content_image_slide": ["TITLE", "PICTURE"],
    "chart_slide": ["TITLE"],
}

# Type-satisfaction relation: which layout placeholder types can SERVE a given
# ideal fingerprint type. A content/body (OBJECT) placeholder is versatile — it
# can host text, pictures, tables or charts — so it satisfies several ideal
# types. This lets e.g. a [TITLE, OBJECT] layout serve a [TITLE, PICTURE] ideal.
_SERVES_LAYOUT: Dict[str, Tuple[str, ...]] = {
    "TITLE": ("TITLE",),
    "SUBTITLE": ("SUBTITLE",),
    "PICTURE": ("PICTURE", "OBJECT"),
    "CHART": ("CHART", "OBJECT"),
    "TABLE": ("TABLE", "OBJECT"),
    "MEDIA": ("MEDIA", "OBJECT"),
    "OBJECT": ("OBJECT",),
}

_LAYOUT_TYPES_NEEDING_SIDEBYSIDE = {"two_content_slide", "comparison_slide"}


# ---------------------------------------------------------------------------
# Pure helpers
# ---------------------------------------------------------------------------
def _normalize_layout_name(name: str) -> str:
    return re.sub(r"^\d+_", "", name).strip().lower()


def _composition_diff(ideal: List[str], layout_fp: List[str]) -> Tuple[int, int]:
    """Return ``(missing, extra)`` between an ideal fingerprint and a layout's.

    ``missing`` = ideal types that no layout placeholder can serve (each ideal
    type consumes one distinct *serving* placeholder). ``extra`` = layout
    placeholders left unconsumed after satisfying the ideal. A content/body
    (OBJECT) placeholder is versatile — it can also serve PICTURE/TABLE/CHART.
    """
    avail: Dict[str, int] = {}
    for t in layout_fp:
        avail[t] = avail.get(t, 0) + 1
    matched = 0
    # Satisfy non-OBJECT ideal types first so OBJECT placeholders are reserved.
    for it in ideal:
        if it == "OBJECT":
            continue
        for lt in _SERVES_LAYOUT.get(it, (it,)):
            if avail.get(lt, 0) > 0:
                avail[lt] -= 1
                matched += 1
                break
    object_need = sum(1 for it in ideal if it == "OBJECT")
    served = min(object_need, avail.get("OBJECT", 0))
    matched += served
    missing = len(ideal) - matched
    extra = len(layout_fp) - matched
    return missing, extra


def _name_affinity(layout_name: str, candidate_names: List[str]) -> int:
    """``2`` = exact (case-insensitive) name match, ``1`` = normalized, ``0`` = none."""
    cname = layout_name.lower()
    if any(cname == c.lower() for c in candidate_names):
        return 2
    norm = _normalize_layout_name(layout_name)
    if any(norm == _normalize_layout_name(c) for c in candidate_names):
        return 1
    return 0


def _content_placeholders_stacked(layout_contract: Dict[str, Any]) -> int:
    """Return 1 if this layout's content placeholders are vertically stacked,
    0 if horizontally separated (side-by-side).

    Used as a geometric tie-breaker so two_content/comparison slides prefer
    side-by-side layouts over vertically stacked ones.
    """
    phs = [p for p in layout_contract.get("placeholders", [])
           if p.get("type") == "OBJECT"]
    if len(phs) < 2:
        return 0
    lefts = sorted(p.get("left_in", 0) for p in phs[:2])
    return 0 if (lefts[1] - lefts[0]) > 1.0 else 1


# ---------------------------------------------------------------------------
# Contract entrypoints
# ---------------------------------------------------------------------------
def _resolve_layout_by_fingerprint(
    slide_type: str,
    contract: Dict[str, Any],
) -> Tuple[Optional[int], Optional[str]]:
    """Match ``slide_type`` to the best contract layout by composition.

    Returns ``(layout_index, None)`` on success, or ``(None, reason)`` on
    degradation (no layout can satisfy the ideal composition). Among
    composition-compatible layouts, ranking is: name affinity (highest) → fewest
    surplus placeholders → largest content area → lowest index. Names are a
    tie-breaker, not the primary key (DESIGN §6 A2, issue #44).
    """
    ideal = _SLIDE_TYPE_FINGERPRINT.get(slide_type)
    layouts = (contract or {}).get("layouts", [])
    if ideal is None:
        return None, f"no fingerprint defined for slide_type '{slide_type}'"
    if not layouts:
        return None, "contract has no layouts"

    candidate_names = _LAYOUT_NAME_MAP.get(slide_type, [])
    need_side_by_side = slide_type in _LAYOUT_TYPES_NEEDING_SIDEBYSIDE
    scored: List[Tuple[int, int, int, int, float, int, str]] = []
    for L in layouts:
        missing, extra = _composition_diff(ideal, L.get("fingerprint", []))
        affinity = _name_affinity(L.get("name", ""), candidate_names)
        stacked = _content_placeholders_stacked(L) if need_side_by_side else 0
        # sort key: (missing, -affinity, extra, stacked, -area, index)
        # — min() picks the lowest missing, then highest affinity, fewest
        # extras, side-by-side over stacked, largest area, lowest index.
        scored.append((
            missing, -affinity, extra, stacked,
            -float(L.get("content_area_in2", 0)), L["index"], L.get("name", ""),
        ))

    full = [s for s in scored if s[0] == 0]
    if not full:
        best = min(scored)
        return None, (
            f"no layout satisfies fingerprint {ideal} "
            f"(closest: '{best[6]}' missing {best[0]})"
        )
    best = min(full)
    return best[5], None


def servable_slide_types(contract: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """Report which engine slide_types a template's contract can serve (#45).

    Used by the content/outline stage to constrain itself to layouts the
    template actually provides (never emit a ``slide_type`` that would degrade).
    For each of the 8 slide types, returns ``{"available": bool, ...}`` — when
    available, the selected layout name/index; when not, the degradation reason.
    """
    layouts = (contract or {}).get("layouts", [])
    report: Dict[str, Dict[str, Any]] = {}
    for slide_type in _SLIDE_TYPE_FINGERPRINT:
        idx, reason = _resolve_layout_by_fingerprint(slide_type, contract)
        if idx is not None and 0 <= idx < len(layouts):
            report[slide_type] = {
                "available": True,
                "layout": layouts[idx].get("name", ""),
                "index": idx,
                "content_area_in2": layouts[idx].get("content_area_in2", 0),
            }
        else:
            report[slide_type] = {"available": False, "reason": reason}
    return report


def available_layouts(contract: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Report ALL layouts a template's contract exposes (GIT-93 Phase 1).

    Unlike :func:`servable_slide_types` (which reports only the 8 standard
    engine slide_types and whether the contract can serve each), this returns
    **every** layout in the contract so the agent can target any of them via a
    ``layout_name`` field — including layouts whose placeholder composition
    does not fingerprint-match any of the 8 standard types (e.g. Agenda,
    Timeline, Team, Quote layouts in a designer deck).

    Each entry: ``{"name", "index", "fingerprint", "content_area_in2"}``.
    """
    layouts = (contract or {}).get("layouts", [])
    return [
        {
            "name": L.get("name", ""),
            "index": L.get("index"),
            "fingerprint": list(L.get("fingerprint", [])),
            "content_area_in2": L.get("content_area_in2", 0),
        }
        for L in layouts
    ]


def classify_layout_fingerprint(fp: List[str]) -> Optional[str]:
    """Map an arbitrary placeholder-composition fingerprint to the nearest
    standard engine slide_type (GIT-93 Phase 1).

    Reuses :func:`_composition_diff` against the 8 ideal fingerprints and
    returns the slide_type with the fewest missing placeholders (ties broken
    by fewest extras, then dict order). Returns ``None`` when ``fp`` is empty.
    Useful as a semantic hint for non-standard layouts so the agent can pick a
    sensible field set (e.g. a ``[TITLE, OBJECT]`` Agenda layout →
    ``"content_slide"`` → use ``body``).
    """
    if not fp:
        return None
    best: Optional[Tuple[int, int, str]] = None  # (missing, extra, slide_type)
    for slide_type, ideal in _SLIDE_TYPE_FINGERPRINT.items():
        missing, extra = _composition_diff(ideal, fp)
        candidate = (missing, extra, slide_type)
        if best is None or candidate[:2] < best[:2]:
            best = candidate
    return best[2] if best else None


# ---------------------------------------------------------------------------
# Embedded-schema probe helpers + render-contract entrypoint (US-4.1)
# ---------------------------------------------------------------------------
def _live_layout_count(template_path: str) -> Optional[int]:
    """The live template's layout count, or ``None`` if it can't be read."""
    try:
        return len(Presentation(template_path).slide_layouts)
    except Exception:  # pragma: no cover - defensive
        return None


def _warn_if_embedded_stale(template_path: str, contract: Dict[str, Any]) -> None:
    """M5 staleness guard: warn if the embedded schema's layout count diverges
    from the live template (catches edit-without-re-embed — the embedded JSON has
    no mtime-invalidation, unlike the sidecar cache). Cheap structural check;
    non-fatal (never blocks the render, never falls back).

    Opens the template **once** (via :func:`_live_layout_count`) and reuses the
    count for both the check and the message (code-review MINOR-1 — the prior
    refactor opened it twice).
    """
    n_embedded = len(contract.get("layouts", []))
    live = _live_layout_count(template_path)
    if live is not None and live != n_embedded:
        logger.warning(
            "Stale embedded schema in %s: describes %d layouts, live template has %d "
            "(template may have been edited after embed); re-run generate-template-skill",
            template_path, n_embedded, live,
        )


def get_render_contract(template_path: str) -> Dict[str, Any]:
    """Return the render contract for ``template_path`` (US-4.1).

    Prefers the embedded ``ppt/template_schema.json`` (via the
    :mod:`contract_adapter` bridge) and falls back to the mtime-cached sidecar
    introspection contract (:func:`get_contract`). Provenance is tagged on the
    returned dict as ``_source ∈ {"embedded", "sidecar"}`` and logged.

    Failure handling (architecture review M4):
      - embedded JSON absent (legacy template) -> silent sidecar fallback.
      - embedded JSON malformed -> ``read_embedded_schema`` already warns;
        sidecar fallback.
      - corrupt/non-zip input -> ``TemplateExtractionError`` caught here, warned,
        sidecar fallback.

    May raise if the sidecar contract itself fails (the caller's try/except then
    degrades to name-based layout matching — backward compatible).
    """
    try:
        schema = read_embedded_schema(template_path)
    except TemplateExtractionError as exc:
        logger.warning(
            "Embedded schema unreadable for %s (%s); falling back to sidecar contract",
            template_path, exc,
        )
        schema = None
    except Exception as exc:  # defensive — never block the render
        logger.warning(
            "Embedded schema read failed for %s (%s); sidecar fallback",
            template_path, exc,
        )
        schema = None

    if schema is not None:
        try:
            contract = embedded_schema_to_contract(schema)
            contract["_source"] = "embedded"
            _warn_if_embedded_stale(template_path, contract)
            logger.info(
                "Render contract (embedded): %d layouts, ratio %s",
                len(contract.get("layouts", [])),
                contract.get("slide_size", {}).get("ratio", "?"),
            )
            return contract
        except Exception as exc:
            logger.warning(
                "Embedded schema -> contract failed for %s (%s); sidecar fallback",
                template_path, exc,
            )

    contract = get_contract(str(template_path))
    contract["_source"] = "sidecar"
    logger.info(
        "Render contract (sidecar): %d layouts, ratio %s",
        len(contract.get("layouts", [])),
        contract.get("slide_size", {}).get("ratio", "?"),
    )
    return contract
