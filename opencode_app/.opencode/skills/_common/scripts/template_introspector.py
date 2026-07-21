"""
template_introspector.py
========================
Template introspection engine — the shared foundation of Capability A
(template-agnostic filling) and Capability B (``pptx-template-modifier-skill``).

Accepts **any** ``.pptx``, introspects its structure, and emits a JSON
``contract`` (per ``DESIGN-template-agnostic.md`` §4) describing the slide
size, theme, and every layout's placeholders, placeholder-composition
*fingerprint*, and content area.

The contract is cached alongside the template as ``<stem>.contract.json`` and
re-used as long as the template file's ``mtime`` is unchanged (decision §1).

Public API:
    introspect(template_path) -> dict   # always-fresh introspection
    get_contract(template_path) -> dict # cache-aware (mtime invalidation)

Usage:
    from template_introspector import get_contract
    contract = get_contract("template/default.pptx")
"""

from __future__ import annotations

import json
import logging
import math
import os
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger(__name__)

_EMU_PER_INCH = 914400

# Placeholder types that are template "chrome" — never part of a content
# fingerprint (they appear on nearly every layout and carry no slide content).
_CHROME_TYPES = {
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.HEADER,
}

# Canonical type names for the contract + fingerprint. BODY and OBJECT are both
# generic content/body placeholders → unified to "OBJECT" so fingerprint
# matching is robust (DESIGN §4 + issue #44 fingerprint table both use OBJECT).
_TYPE_CANONICAL: Dict[Any, str] = {
    PP_PLACEHOLDER.TITLE: "TITLE",
    PP_PLACEHOLDER.CENTER_TITLE: "TITLE",
    PP_PLACEHOLDER.VERTICAL_TITLE: "TITLE",
    PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
    PP_PLACEHOLDER.BODY: "OBJECT",
    PP_PLACEHOLDER.OBJECT: "OBJECT",
    PP_PLACEHOLDER.VERTICAL_BODY: "OBJECT",
    PP_PLACEHOLDER.VERTICAL_OBJECT: "OBJECT",
    PP_PLACEHOLDER.PICTURE: "PICTURE",
    PP_PLACEHOLDER.BITMAP: "PICTURE",
    PP_PLACEHOLDER.CHART: "CHART",
    PP_PLACEHOLDER.TABLE: "TABLE",
    PP_PLACEHOLDER.ORG_CHART: "OBJECT",
    PP_PLACEHOLDER.MEDIA_CLIP: "MEDIA",
}

# Placeholder canonical types that count toward the text/content body area used
# for density calibration (Capability A) and over-limit checks (Capability B).
_CONTENT_AREA_TYPES = {"OBJECT"}

_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------
def _canonical_type(ph_type: Any) -> str:
    """Map a python-pptx placeholder type to a canonical contract name.

    ``None`` (an untyped ``<p:ph>``) is treated as a generic content placeholder
    → ``"OBJECT"``.
    """
    if ph_type is None:
        return "OBJECT"
    return _TYPE_CANONICAL.get(ph_type, getattr(ph_type, "name", "OTHER"))


def _is_chrome(ph_type: Any) -> bool:
    return ph_type in _CHROME_TYPES


def _inches(emu: Optional[int]) -> float:
    if emu is None:
        return 0.0
    return round(emu / _EMU_PER_INCH, 4)


def _compute_ratio(width_emu: int, height_emu: int) -> str:
    """Return the simplest ``W:H`` integer ratio for the given EMU dimensions."""
    if height_emu <= 0:
        return f"{width_emu}:0"
    g = math.gcd(width_emu, height_emu)
    w = width_emu // g
    h = height_emu // g
    # Keep the numbers readable (e.g. 16:9 rather than 7600:4275).
    return f"{w}:{h}"


def _build_slide_size(prs: Presentation) -> Dict[str, Any]:
    w_emu = int(prs.slide_width)
    h_emu = int(prs.slide_height)
    return {
        "width_emu": w_emu,
        "height_emu": h_emu,
        "width_in": round(w_emu / _EMU_PER_INCH, 2),
        "height_in": round(h_emu / _EMU_PER_INCH, 2),
        "ratio": _compute_ratio(w_emu, h_emu),
    }


def _build_theme(prs: Presentation) -> Dict[str, Any]:
    """Extract theme colors (hex per role) and major/minor Latin fonts."""
    colors: Dict[str, str] = {}
    fonts: Dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(_THEME_REL)
        theme_xml = etree.parse(BytesIO(theme_part.blob)).getroot()
        elements = theme_xml.find(f"{{{_NS_A}}}themeElements")
        if elements is not None:
            clr_scheme = elements.find(f"{{{_NS_A}}}clrScheme")
            if clr_scheme is not None:
                for child in clr_scheme:
                    role = etree.QName(child).localname
                    for color_elem in child:
                        color_tag = etree.QName(color_elem).localname
                        if color_tag == "srgbClr":
                            colors[role] = "#" + color_elem.get("val", "")
                        elif color_tag == "sysClr":
                            val = color_elem.get("lastClr") or color_elem.get("val", "")
                            if val:
                                colors[role] = "#" + val
            font_scheme = elements.find(f"{{{_NS_A}}}fontScheme")
            if font_scheme is not None:
                for label, key in (("majorFont", "major_latin"), ("minorFont", "minor_latin")):
                    font_elem = font_scheme.find(f"{{{_NS_A}}}{label}")
                    if font_elem is not None:
                        latin = font_elem.find(f"{{{_NS_A}}}latin")
                        if latin is not None and latin.get("typeface"):
                            fonts[key] = latin.get("typeface")
    except Exception as exc:  # pragma: no cover - defensive; theme is best-effort
        logger.warning("Theme extraction failed (%s); emitting empty theme", exc)
    return {"colors": colors, "fonts": fonts}


def placeholder_record(ph: Any) -> Optional[Dict[str, Any]]:
    """Extract a canonical placeholder record, or ``None`` for chrome placeholders.

    Used internally by :func:`_build_layout` to build each layout's placeholder
    list (issue #46).
    """
    ph_type = ph.placeholder_format.type
    if _is_chrome(ph_type):
        return None
    canonical = _canonical_type(ph_type)
    return {
        "idx": ph.placeholder_format.idx,
        "name": ph.name,
        "type": canonical,
        "left_in": _inches(ph.left),
        "top_in": _inches(ph.top),
        "width_in": _inches(ph.width),
        "height_in": _inches(ph.height),
    }


def _build_layout(entry: Any, index: int) -> Dict[str, Any]:
    """Build the contract entry for a single ``SlideLayout``."""
    placeholders: List[Dict[str, Any]] = []
    fingerprint: List[str] = []
    content_area_in2 = 0.0

    for ph in entry.placeholders:
        record = placeholder_record(ph)
        if record is None:
            continue
        canonical = record["type"]
        placeholders.append(record)
        fingerprint.append(canonical)
        if canonical in _CONTENT_AREA_TYPES:
            w = record["width_in"]
            h = record["height_in"]
            if w and h:
                content_area_in2 += w * h

    return {
        "index": index,
        "name": entry.name,
        "placeholders": placeholders,
        "fingerprint": fingerprint,
        "content_area_in2": round(content_area_in2, 2),
    }


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------
def introspect(template_path: str) -> Dict[str, Any]:
    """Introspect ``template_path`` and return a fresh contract dict.

    Always re-reads the template (no cache). See :func:`get_contract` for the
    cache-aware entry point used by the engine.
    """
    path = Path(template_path)
    prs = Presentation(str(path))

    layouts = [
        _build_layout(layout, index)
        for index, layout in enumerate(prs.slide_layouts)
    ]

    contract: Dict[str, Any] = {
        "source_file": path.name,
        "source_mtime": os.path.getmtime(str(path)),
        "slide_size": _build_slide_size(prs),
        "theme": _build_theme(prs),
        "layouts": layouts,
    }
    logger.info(
        "Introspected %s: %d layouts, %s",
        path.name, len(layouts), contract["slide_size"].get("ratio"),
    )
    return contract


def contract_path_for(template_path: str) -> Path:
    """Return the cache file path derived from the template path."""
    p = Path(template_path)
    return p.with_suffix(p.suffix + ".contract.json") if p.suffix else p.with_suffix(".contract.json")


def get_contract(template_path: str) -> Dict[str, Any]:
    """Return the contract for ``template_path``, using the mtime cache.

    Re-introspects only when the template's ``mtime`` differs from the cached
    ``source_mtime`` (decision §1). The cached contract is written next to the
    template as ``<stem>.pptx.contract.json``.
    """
    cache_path = contract_path_for(template_path)
    current_mtime = os.path.getmtime(template_path)

    if cache_path.exists():
        try:
            with open(cache_path, "r", encoding="utf-8") as fh:
                cached = json.load(fh)
            if cached.get("source_mtime") == current_mtime:
                logger.info(
                    "Contract cache hit (mtime unchanged): %s", cache_path.name
                )
                return cached
            logger.info("Contract cache stale (mtime changed); re-introspecting")
        except (json.JSONDecodeError, OSError) as exc:
            logger.warning("Contract cache unreadable (%s); re-introspecting", exc)

    contract = introspect(template_path)
    try:
        with open(cache_path, "w", encoding="utf-8") as fh:
            json.dump(contract, fh, indent=2, ensure_ascii=False)
    except OSError as exc:  # pragma: no cover - defensive
        logger.warning("Could not write contract cache (%s): %s", cache_path, exc)
    return contract
