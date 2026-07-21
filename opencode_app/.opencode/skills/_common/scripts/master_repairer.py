"""US-4.8 Phase 1 — Master repair cascade for masterless PPTX files.

When a user-supplied ``.pptx`` has no slide master (Scenario A), this module
repairs it via a **three-level cascade** (Chain of Responsibility):

- **Level 1** — salvage ``ppt/theme/theme1.xml`` from the zip (exact fidelity).
- **Level 2** — scavenge explicit styles from slide XML (best-effort fidelity).
- **Level 3** — fallback to an **in-code minimal theme** (no user styling).

BT-142 Phase 1.5: the original Level 3 strategy copied ``default.pptx`` for its
master+layouts skeleton. Per the user's "no bundled default" invariant, this
module now synthesizes a minimal valid PPTX in-memory using ``python-pptx``
(``Presentation()`` creates one master + one layout automatically), and uses
an in-code minimal theme XML as the L2 scavenge base. The ``default_template_path``
parameter is now **optional** — when ``None``, the in-code fallbacks are used.
"""
from __future__ import annotations

import logging
import os
import shutil
import tempfile
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

from lxml import etree
from io import BytesIO

logger = logging.getLogger(__name__)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_THEME_PATH = "ppt/theme/theme1.xml"
_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
_CONTENT_TYPES_XML = "[Content_Types].xml"
_SLIDES_DIR = "ppt/slides/"

# Font-size tiers for Level-2 scavenge (in centipoints; OOXML sz attribute).
# ≥28pt = title tier; <28pt = body tier. 1pt = 100 centipoints.
_TITLE_SIZE_THRESHOLD = 2800


# ---------------------------------------------------------------------------
# BT-142 Phase 1.5 — in-code minimal theme XML (replaces default.pptx dependency)
# ---------------------------------------------------------------------------
# A complete OOXML <a:theme> element with all 12 OPC color roles (dk1/lt1/dk2/lt2/
# accent1-6/hlink/folHlink) + majorFont/minorFont (Calibri Light / Calibri).
# Neutral gray palette — the L3 "no user styling recoverable" fallback.
# Used as: (a) the L2 scavenge structural base; (b) the L3 fallback theme.
_MINIMAL_THEME_XML = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Minimal Theme">
  <a:themeElements>
    <a:clrScheme name="Minimal">
      <a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>
      <a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>
      <a:dk2><a:srgbClr val="44546A"/></a:dk2>
      <a:lt2><a:srgbClr val="E7E6E6"/></a:lt2>
      <a:accent1><a:srgbClr val="4472C4"/></a:accent1>
      <a:accent2><a:srgbClr val="ED7D31"/></a:accent2>
      <a:accent3><a:srgbClr val="A5A5A5"/></a:accent3>
      <a:accent4><a:srgbClr val="FFC000"/></a:accent4>
      <a:accent5><a:srgbClr val="5B9BD5"/></a:accent5>
      <a:accent6><a:srgbClr val="70AD47"/></a:accent6>
      <a:hlink><a:srgbClr val="0563C1"/></a:hlink>
      <a:folHlink><a:srgbClr val="954F72"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="Minimal">
      <a:majorFont>
        <a:latin typeface="Calibri Light"/><a:ea typeface=""/><a:cs typeface=""/>
      </a:majorFont>
      <a:minorFont>
        <a:latin typeface="Calibri"/><a:ea typeface=""/><a:cs typeface=""/>
      </a:minorFont>
    </a:fontScheme>
    <a:fmtScheme name="Minimal">
      <a:fillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
      </a:fillStyleLst>
      <a:lnStyleLst>
        <a:ln w="6350" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
        <a:ln w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
        <a:ln w="19050" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
      </a:lnStyleLst>
      <a:effectStyleLst>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst/></a:effectStyle>
      </a:effectStyleLst>
      <a:bgFillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
      </a:bgFillStyleLst>
    </a:fmtScheme>
  </a:themeElements>
  <a:objectDefaults/>
  <a:extraClrSchemeLst/>
</a:theme>"""


def _build_minimal_pptx_bytes(
    theme_bytes: Optional[bytes],
    slide_width_emu: int = 12192000,
    slide_height_emu: int = 6858000,
) -> bytes:
    """BT-142 Phase 1.5 — synthesize a minimal valid PPTX in-memory.

    Uses ``python-pptx``'s ``Presentation()`` constructor which creates a
    Presentation with one slide_master + one blank slide_layout automatically
    (no external template file needed). If ``theme_bytes`` is provided
    (Level 1/2), the synthesized PPTX's ``ppt/theme/theme1.xml`` is replaced
    with the salvaged/scavenged theme content.

    Slide dimensions default to 16:9 (12192000 × 6858000 EMU = 13.33in × 7.5in),
    matching modern PowerPoint defaults. Override with ``slide_width_emu`` /
    ``slide_height_emu`` for other aspect ratios.

    Returns the PPTX as bytes (in-memory buffer; no disk I/O).
    """
    from pptx import Presentation  # local import — keeps module-import side-effect-free
    from pptx.util import Emu

    prs = Presentation()  # creates default master + blank layout
    # Default python-pptx is 4:3 (9144000 × 6858000); normalize to 16:9.
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)
    buf = BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    # BT-142 Phase 1.5: always replace python-pptx's default theme (which uses
    # the legacy Office 2007-2010 color palette — accent1=4F81BD) with our
    # _MINIMAL_THEME_XML (modern Office 2013+ palette — accent1=4472C4). This
    # ensures consistent theme colors whether theme_bytes is provided (L1/L2)
    # or not (L3). When theme_bytes IS provided, it wins over _MINIMAL_THEME_XML.
    effective_theme = theme_bytes if theme_bytes is not None else _MINIMAL_THEME_XML

    # Replace ppt/theme/theme1.xml with the effective theme.
    out_buf = BytesIO()
    with zipfile.ZipFile(BytesIO(pptx_bytes), "r") as zin, \
            zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename == _THEME_PATH:
                data = effective_theme
            out_info = zipfile.ZipInfo(filename=info.filename, date_time=info.date_time)
            out_info.compress_type = info.compress_type
            out_info.external_attr = info.external_attr
            zout.writestr(out_info, data)
    return out_buf.getvalue()


@dataclass
class RepairResult:
    """Outcome of ``repair_if_needed``."""

    level: Literal["none", "L1", "L2", "L3"]
    mutated: bool
    theme_source: str  # "salvaged", "scavenged", "default", "n/a"
    repaired_path: Optional[str] = None  # the derived file path if mutated


def repair_if_needed(
    prs: Any,
    template_path: str,
    default_template_path: Optional[str] = None,
) -> RepairResult:
    """Check if ``prs`` needs master repair (Scenario A); if so, repair it.

    **CRIT-1**: Must be called BEFORE ``get_render_contract`` in
    ``generate_ppt_from_data``, so the contract is fetched from the repaired prs.

    **BT-142 Phase 1.5**: ``default_template_path`` is now **optional**. When
    ``None`` (the new default), the L2 scavenge base and L3 fallback use the
    in-code ``_MINIMAL_THEME_XML`` constant + ``_build_minimal_pptx_bytes()``
    synthesized PPTX — no external ``default.pptx`` required. Callers that
    still pass a path (legacy) get the original behavior.

    Returns a :class:`RepairResult`. If ``mutated`` is True, the caller should
    reload ``prs`` from ``repaired_path``.
    """
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    if masters:
        return RepairResult(level="none", mutated=False, theme_source="n/a")

    logger.info(
        "Template has no slide master (Scenario A) — starting repair cascade for %s",
        template_path,
    )

    # --- Level 1: salvage theme1.xml from the zip ---
    theme_bytes = _salvage_theme_part(template_path)
    level: str
    theme_source: str

    if theme_bytes is not None:
        level = "L1"
        theme_source = "salvaged"
        logger.info("Level 1: salvaged ppt/theme/theme1.xml from %s", template_path)
    else:
        # --- Level 2: scavenge styles from slide XML ---
        theme_bytes = _scavenge_slide_styles(template_path, default_template_path)
        if theme_bytes is not None:
            level = "L2"
            theme_source = "scavenged"
            logger.info("Level 2: scavenged theme from slide styles in %s", template_path)
        else:
            # --- Level 3: in-code minimal fallback (was: default.pptx) ---
            level = "L3"
            theme_source = "minimal-in-code"
            theme_bytes = None
            logger.info("Level 3: using in-code minimal theme (no user styling recoverable)")

    # Build the derived file: synthesize minimal PPTX in-memory, optionally replace theme.
    repaired_path = _build_repaired_file(
        template_path, default_template_path, theme_bytes, level
    )

    return RepairResult(
        level=level,
        mutated=True,
        theme_source=theme_source,
        repaired_path=repaired_path,
    )


# ---------------------------------------------------------------------------
# Level 1: Salvage ppt/theme/theme1.xml from the zip
# ---------------------------------------------------------------------------

def _salvage_theme_part(pptx_path: str) -> Optional[bytes]:
    """Level 1: read ``ppt/theme/theme1.xml`` directly from the PPTX zip.

    The theme part often survives when the master part is stripped, because it
    is an independent zip entry (``ppt/theme/theme1.xml``), not embedded inside
    the master.

    Returns the raw theme XML bytes, or ``None`` if the part is absent.
    """
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            if _THEME_PATH in z.namelist():
                return z.read(_THEME_PATH)
    except (zipfile.BadZipFile, OSError) as exc:
        logger.warning("Level 1 salvage failed (%s)", exc)
    return None


# ---------------------------------------------------------------------------
# Level 2: Scavenge explicit styles from slide XML
# ---------------------------------------------------------------------------

@dataclass
class _ScavengedTheme:
    """Aggregated style data from slide XML."""

    major_font: Optional[str] = None  # title-tier typeface
    minor_font: Optional[str] = None  # body-tier typeface
    colors: List[str] = field(default_factory=list)  # all srgbClr hex values


def _scavenge_slide_styles(
    pptx_path: str, default_template_path: Optional[str] = None
) -> Optional[bytes]:
    """Level 2: aggregate explicit styles from surviving slide XML.

    Walks ``ppt/slides/slideN.xml`` for ``<a:rPr>`` (font/size/color) and
    ``<p:spPr>`` fill colors. Aggregates into a dominant major/minor font +
    top accent colors, then synthesizes a ``<a:theme>`` by overriding the
    base theme's clrScheme + fontScheme.

    BT-142 Phase 1.5: when ``default_template_path`` is ``None`` (the new
    default), the in-code ``_MINIMAL_THEME_XML`` is used as the structural
    base — no external ``default.pptx`` required.

    Returns the synthesized theme XML bytes, or ``None`` if no slide styles
    are recoverable.
    """
    try:
        scavenged = _aggregate_slide_styles(pptx_path)
    except Exception as exc:
        logger.warning("Level 2 scavenge failed (%s)", exc)
        return None

    if not scavenged.major_font and not scavenged.minor_font and not scavenged.colors:
        logger.info("Level 2: no explicit styles found in slides")
        return None

    # Load base theme: in-code minimal XML preferred (Phase 1.5); legacy path
    # reads default.pptx if explicitly provided.
    if default_template_path is not None:
        try:
            with zipfile.ZipFile(default_template_path, "r") as z:
                base_theme_bytes = z.read(_THEME_PATH)
        except Exception as exc:
            logger.warning(
                "Level 2: could not read default theme from %s (%s); using in-code minimal",
                default_template_path, exc,
            )
            base_theme_bytes = _MINIMAL_THEME_XML
    else:
        base_theme_bytes = _MINIMAL_THEME_XML

    return _build_synthetic_theme(base_theme_bytes, scavenged)


def _aggregate_slide_styles(pptx_path: str) -> _ScavengedTheme:
    """Parse all slide XML files and aggregate font/color frequencies."""
    typeface_by_tier: Dict[str, Counter] = {"title": Counter(), "body": Counter()}
    all_colors: List[str] = []

    with zipfile.ZipFile(pptx_path, "r") as z:
        slide_files = sorted(
            n for n in z.namelist() if n.startswith(_SLIDES_DIR) and n.endswith(".xml")
        )
        for slide_path in slide_files:
            try:
                root = etree.parse(BytesIO(z.read(slide_path))).getroot()
            except Exception:
                continue
            # Walk all <a:rPr> elements for typeface + size + color.
            for rpr in root.iter(f"{{{_NS_A}}}rPr"):
                typeface = rpr.get("typeface")
                sz = rpr.get("sz")
                # Determine tier by font size.
                tier = "title"
                if sz:
                    try:
                        tier = "title" if int(sz) >= _TITLE_SIZE_THRESHOLD else "body"
                    except ValueError:
                        pass
                if typeface:
                    typeface_by_tier[tier][typeface] += 1
                # Collect srgbClr children (text-run colors).
                for clr in rpr.iter(f"{{{_NS_A}}}srgbClr"):
                    val = clr.get("val", "").upper()
                    if val and len(val) == 6:
                        all_colors.append(val)
            # Walk shape-level fill colors ONLY (avoid double-counting rPr colors).
            # Use p:spPr namespace to target shape properties, not text runs.
            _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
            for spPr in root.iter(f"{{{_NS_P}}}spPr"):
                for clr in spPr.iter(f"{{{_NS_A}}}srgbClr"):
                    val = clr.get("val", "").upper()
                    if val and len(val) == 6:
                        all_colors.append(val)

    major_font = typeface_by_tier["title"].most_common(1)[0][0] if typeface_by_tier["title"] else None
    minor_font = typeface_by_tier["body"].most_common(1)[0][0] if typeface_by_tier["body"] else None

    return _ScavengedTheme(
        major_font=major_font,
        minor_font=minor_font,
        colors=all_colors,
    )


def _build_synthetic_theme(
    base_theme_bytes: bytes, scavenged: _ScavengedTheme
) -> bytes:
    """Override default theme's clrScheme + fontScheme with scavenged values.

    Keeps the base theme's fmtScheme and other elements intact (structural
    validity), only replacing the color palette and font names.

    **Fidelity ceiling (m-7):** only accent1–accent6 are overridden from
    scavenged colors. The background/text roles (dk1/lt1/dk2/lt2) keep
    default.pptx's values — a distinctive user background is NOT recovered.
    This is the accepted best-effort limit of Level 2.
    """
    root = etree.parse(BytesIO(base_theme_bytes)).getroot()
    elements = root.find(f"{{{_NS_A}}}themeElements")
    if elements is None:
        return base_theme_bytes  # can't override; return default as-is

    # --- Override fontScheme ---
    if scavenged.major_font or scavenged.minor_font:
        font_scheme = elements.find(f"{{{_NS_A}}}fontScheme")
        if font_scheme is not None:
            if scavenged.major_font:
                major = font_scheme.find(f"{{{_NS_A}}}majorFont")
                if major is not None:
                    latin = major.find(f"{{{_NS_A}}}latin")
                    if latin is not None:
                        latin.set("typeface", scavenged.major_font)
            if scavenged.minor_font:
                minor = font_scheme.find(f"{{{_NS_A}}}minorFont")
                if minor is not None:
                    latin = minor.find(f"{{{_NS_A}}}latin")
                    if latin is not None:
                        latin.set("typeface", scavenged.minor_font)

    # --- Override clrScheme (map top colors to accent roles) ---
    if scavenged.colors:
        color_counts = Counter(scavenged.colors)
        # Exclude pure black/white (they're dk1/lt1 defaults).
        non_mono = [c for c, _ in color_counts.most_common(12) if c not in ("000000", "FFFFFF")]
        clr_scheme = elements.find(f"{{{_NS_A}}}clrScheme")
        if clr_scheme is not None and non_mono:
            accent_roles = ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
            for i, role in enumerate(accent_roles):
                if i < len(non_mono):
                    elem = clr_scheme.find(f"{{{_NS_A}}}{role}")
                    if elem is not None:
                        srgb = elem.find(f"{{{_NS_A}}}srgbClr")
                        if srgb is not None:
                            srgb.set("val", non_mono[i])

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ---------------------------------------------------------------------------
# Derived file builder
# ---------------------------------------------------------------------------

def _build_repaired_file(
    template_path: str,
    default_template_path: Optional[str],
    theme_bytes: Optional[bytes],
    level: str = "L3",
) -> str:
    """Build the derived repaired PPTX; optionally replace the theme.

    BT-142 Phase 1.5: when ``default_template_path`` is ``None`` (the new
    default), the derived file is synthesized in-memory via
    ``_build_minimal_pptx_bytes()`` — a minimal valid PPTX with one master +
    one layout. If ``theme_bytes`` is provided (Level 1/2), the synthesized
    PPTX's ``ppt/theme/theme1.xml`` is replaced with the salvaged/scavenged
    content. Level 3 (``theme_bytes is None``) uses the synthesized PPTX
    verbatim (with its default in-code minimal theme).

    Legacy path: when ``default_template_path`` is explicitly provided, the
    original behavior is preserved (copy default.pptx, optionally replace theme).

    The derived file is named ``<stem>_repaired.pptx`` beside the original.
    """
    src = Path(template_path)
    out_dir = src.parent
    repaired_name = f"{src.stem}_repaired{src.suffix}"
    repaired_path = out_dir / repaired_name

    fd, tmp_name = tempfile.mkstemp(suffix=".pptx", dir=str(out_dir))
    os.close(fd)
    try:
        if default_template_path is not None:
            # Legacy path: copy default.pptx, optionally replacing the theme.
            if theme_bytes is not None:
                with zipfile.ZipFile(default_template_path, "r") as zin, \
                        zipfile.ZipFile(tmp_name, "w", zipfile.ZIP_DEFLATED) as zout:
                    for info in zin.infolist():
                        data = zin.read(info.filename)
                        if info.filename == _THEME_PATH:
                            data = theme_bytes  # replace theme
                        out_info = zipfile.ZipInfo(filename=info.filename, date_time=info.date_time)
                        out_info.compress_type = info.compress_type
                        out_info.external_attr = info.external_attr
                        zout.writestr(out_info, data)
            else:
                # Level 3 legacy: just copy default.pptx verbatim.
                shutil.copy2(default_template_path, tmp_name)
        else:
            # BT-142 Phase 1.5 new path: synthesize minimal PPTX in-memory.
            pptx_bytes = _build_minimal_pptx_bytes(theme_bytes)
            with open(tmp_name, "wb") as f:
                f.write(pptx_bytes)
        os.replace(tmp_name, str(repaired_path))
    except Exception:
        if os.path.exists(tmp_name):
            os.remove(tmp_name)
        raise

    logger.info("Repaired template: %s (level=%s)", repaired_path, level)
    return str(repaired_path)
