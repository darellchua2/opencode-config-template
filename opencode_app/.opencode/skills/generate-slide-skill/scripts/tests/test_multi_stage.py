"""Tests for the multi-stage generation building blocks (#21 / #24).

The outline → critique → detail stages themselves are LLM-driven at the agent
level. These tests cover the deterministic building blocks that support them:
the outline artifact store and schema-gating between stages.
"""
import pytest
from pptx import Presentation

from outline_store import (
    cleanup_all,
    latest_outline,
    load_outline,
    load_outline_mode,
    save_outline,
)
from schema_validator import validate_slide_data_list

import outline_store  # imported for monkeypatching the temp-dir module global


# ============================================================
# Outline artifact store (foundation for interactive checkpoint #24)
# ============================================================
class TestOutlineStore:
    def test_save_and_load_roundtrip(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("1. [title_slide] Hello", deck_id="abc")
        assert p.exists()
        assert load_outline(p).startswith("1. [title_slide]")

    def test_save_creates_md_file(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("outline text", deck_id="deck42")
        assert p.suffix == ".md"
        assert "deck42" in p.name

    def test_load_missing_file_raises(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            load_outline(tmp_path / "nope.md")

    def test_latest_outline_finds_most_recent(self, tmp_path, monkeypatch):
        import time
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        save_outline("first", deck_id="a")
        time.sleep(0.01)
        save_outline("second", deck_id="b")
        latest = latest_outline()
        assert latest is not None
        assert load_outline(latest) == "second"

    def test_latest_outline_none_when_empty(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        assert latest_outline() is None

    def test_save_empty_outline(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("", deck_id="empty")
        assert load_outline(p) == ""

    def test_cleanup_all_clears_entire_dir(self, tmp_path, monkeypatch):
        # The whole temp dir is cleared, including agent-written temp JSON
        # (e.g. slide_data.json), not just outline artifacts.
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        save_outline("first", deck_id="a")
        save_outline("second", deck_id="b")
        (tmp_path / "slide_data.json").write_text("[]", encoding="utf-8")
        assert cleanup_all() == 3
        assert latest_outline() is None
        assert not (tmp_path / "slide_data.json").exists()

    def test_cleanup_all_safe_when_empty_or_missing(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        assert cleanup_all() == 0  # empty dir is a no-op
        tmp_path.rmdir()
        assert cleanup_all() == 0  # missing dir is a no-op, never raises

    # ----- density-mode header persistence (Stage 2 → Stage 3 traceability) --
    def test_save_with_mode_records_header(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("1. [title_slide] Hi", deck_id="m1", mode="standard")
        text = load_outline(p)
        assert text.startswith("<!-- mode: standard -->")
        assert "1. [title_slide] Hi" in text

    def test_load_outline_mode_returns_recorded_mode(self, tmp_path, monkeypatch):
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("outline", deck_id="m2", mode="text-heavy")
        assert load_outline_mode(p) == "text-heavy"

    def test_load_outline_mode_none_when_no_mode(self, tmp_path, monkeypatch):
        # Backward compat: artifacts saved without a mode return None.
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("outline", deck_id="m3")
        assert load_outline_mode(p) is None

    def test_load_outline_mode_none_on_missing_file(self, tmp_path):
        assert load_outline_mode(tmp_path / "never_saved.md") is None

    def test_save_with_mode_none_writes_no_header(self, tmp_path, monkeypatch):
        # Explicit mode=None must behave identically to the old signature.
        monkeypatch.setattr(outline_store, "_TEMP_DIR", tmp_path)
        p = save_outline("plain outline", deck_id="m4", mode=None)
        assert load_outline(p) == "plain outline"


# ============================================================
# Schema gating between stages (#21)
# ============================================================
class TestSchemaGating:
    """Every stage's JSON output must be schema-validated before continuing."""

    def test_valid_stage_output_passes_gate(self):
        stage_output = [
            {"slide_type": "title_slide", "title": "T", "notes": "n"},
            {"slide_type": "content_slide", "title": "C", "body": "**x** - y", "notes": "n"},
        ]
        result = validate_slide_data_list(stage_output, strict=True)
        assert result.is_valid, result.error_messages()

    def test_invalid_stage_output_blocks_gate(self):
        # Missing required chart fields would block progression to the next stage.
        stage_output = [
            {"slide_type": "title_slide", "title": "T", "notes": "n"},
            {"slide_type": "chart_slide", "title": "C", "chart_type": "bar",
             "categories": ["a"], "notes": "n"},  # missing series
        ]
        result = validate_slide_data_list(stage_output, strict=True)
        assert not result.is_valid

    def test_gate_returns_actionable_errors(self):
        stage_output = [{"slide_type": "chart_slide", "title": "C", "notes": "n",
                         "chart_type": "banana", "categories": ["a"],
                         "series": [{"name": "S", "values": [1]}]}]
        result = validate_slide_data_list(stage_output, strict=True)
        msg = result.errors[0].format()
        # The error names the slide index, field, and valid options.
        assert "slide[0]" in msg
        assert "chart_type" in msg


# ============================================================
# Full pipeline order: resolve -> validate -> render
# ============================================================
class TestPipelineOrder:
    def test_resolve_then_validate_then_render(self, template_path, output_path, tmp_path):
        from ppt_builder import generate_ppt_from_data
        from resolvers import resolve_slide_data_list

        def search_fn(query, config):
            return {"categories": ["2022", "2023", "2024"],
                    "series": [{"name": "Revenue", "values": [10, 12, 15]}],
                    "source": "https://example.com"}

        deck = [
            {"slide_type": "title_slide", "title": "Report", "notes": "n"},
            {"slide_type": "chart_slide", "title": "Growth", "chart_type": "bar",
             "data_query": "revenue", "notes": "n"},
        ]
        # Stage: resolve placeholders
        resolved = resolve_slide_data_list(deck, {"chart_data": {"search_fn": search_fn}})
        # Stage: schema gate (strict)
        gate = validate_slide_data_list(resolved, strict=True)
        assert gate.is_valid, gate.error_messages()
        # Stage: render (cleanup_temp=False isolates this test from the global
        # temp dir so it only exercises resolve -> validate -> render).
        generate_ppt_from_data(resolved, template_path=template_path, output_path=output_path, cleanup_temp=False)
        prs = Presentation(output_path)
        assert len(prs.slides) == 2
        # Chart slide got real data via the resolver.
        chart_slide = prs.slides[1]
        charts = [s for s in chart_slide.shapes if s.has_chart]
        assert len(charts) == 1
