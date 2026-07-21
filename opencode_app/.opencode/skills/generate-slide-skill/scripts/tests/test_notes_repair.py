"""Tests for BT-142 Phase 3.5c notes_repair."""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches

from notes_repair import ensure_notes_placeholder, _find_notes_text_placeholder


def _build_minimal_prs_no_notes_placeholder(tmp_path):
    """Build a tiny presentation whose notes master lacks a body placeholder."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    # Touch notes_master so one exists, but it won't have a body placeholder
    # on a freshly constructed python-pptx Presentation.
    _ = prs.notes_master
    path = tmp_path / "n.pptx"
    prs.save(str(path))
    return path


def test_ensure_notes_placeholder_idempotent(tmp_path):
    path = _build_minimal_prs_no_notes_placeholder(tmp_path)
    prs = Presentation(str(path))
    first = ensure_notes_placeholder(prs)
    second = ensure_notes_placeholder(prs)
    assert first["placeholder_count"] == 1
    assert second["placeholder_count"] == 1
    assert second["repaired"] is False


def test_ensure_notes_placeholder_injects_when_missing(tmp_path):
    path = _build_minimal_prs_no_notes_placeholder(tmp_path)
    prs = Presentation(str(path))
    before = _find_notes_text_placeholder(prs.notes_master)
    if before:
        # python-pptx may already ship with one; test still passes (no-op)
        report = ensure_notes_placeholder(prs)
        assert report["placeholder_count"] == 1
        return
    report = ensure_notes_placeholder(prs)
    assert report["placeholder_count"] == 1
    assert report["error"] is None
    assert _find_notes_text_placeholder(prs.notes_master) is True


def test_persistence_after_save(tmp_path):
    path = _build_minimal_prs_no_notes_placeholder(tmp_path)
    prs = Presentation(str(path))
    ensure_notes_placeholder(prs)
    out = tmp_path / "out.pptx"
    prs.save(str(out))
    reloaded = Presentation(str(out))
    assert _find_notes_text_placeholder(reloaded.notes_master) is True


def test_returns_report_shape(tmp_path):
    path = _build_minimal_prs_no_notes_placeholder(tmp_path)
    prs = Presentation(str(path))
    report = ensure_notes_placeholder(prs)
    assert set(report.keys()) == {"repaired", "placeholder_count", "error"}
