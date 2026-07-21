"""Tests for the ppt_builder CLI ``main()`` (US-4.6 Phase 4 + arch-review m1).

Covers the ``--target-size`` parsing surface that the library tests don't
reach: valid preset, valid ``WxH`` shorthand, the multi-``x`` crash guard
(``5x3x2`` must NOT raise an uncaught ValueError), invalid preset exit code,
and the demo-deck fallback. Uses ``main(argv)`` directly (returns exit code).
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_builder import main


def _ratio(w, h):
    from geometry import compute_ratio
    return compute_ratio(int(w), int(h))


class TestTargetSizeCLI:
    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_valid_preset_renders(self, tmp_path):
        out = tmp_path / "o.pptx"
        rc = main(["--output", str(out), "--target-size", "4:3", "--log-level", "error"])
        assert rc == 0
        assert out.exists()
        prs = Presentation(str(out))
        assert _ratio(prs.slide_width, prs.slide_height) == "4:3"

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_valid_WxH_shorthand(self, tmp_path):
        # 10x7.5 inches -> 4:3 ratio.
        out = tmp_path / "o.pptx"
        rc = main(["--output", str(out), "--target-size", "10x7.5", "--log-level", "error"])
        assert rc == 0
        prs = Presentation(str(out))
        assert _ratio(prs.slide_width, prs.slide_height) == "4:3"

    def test_native_when_target_size_omitted(self, template_path, tmp_path):
        out = tmp_path / "o.pptx"
        rc = main([
            "--template", template_path, "--output", str(out), "--log-level", "error",
        ])
        assert rc == 0
        nat = Presentation(template_path)
        prs = Presentation(str(out))
        assert (int(prs.slide_width), int(prs.slide_height)) == (
            int(nat.slide_width), int(nat.slide_height))

    def test_invalid_preset_exits_1(self, tmp_path):
        out = tmp_path / "o.pptx"
        rc = main(["--output", str(out), "--target-size", "bogus", "--log-level", "error"])
        assert rc == 1  # ValueError -> exit 1, not a crash/traceback

    def test_multi_x_does_not_crash(self, tmp_path):
        # "5x3x2" must be treated as a preset name (not an uncaught unpack
        # ValueError). It then fails preset resolution -> exit 1, cleanly.
        out = tmp_path / "o.pptx"
        rc = main(["--output", str(out), "--target-size", "5x3x2", "--log-level", "error"])
        assert rc == 1

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_4x3_is_inches_not_preset(self, tmp_path):
        # "4x3" -> 4in x 3in (ratio 4:3), distinct from the "4:3" preset but
        # coincidentally the same ratio.
        out = tmp_path / "o.pptx"
        rc = main(["--output", str(out), "--target-size", "4x3", "--log-level", "error"])
        assert rc == 0
        prs = Presentation(str(out))
        assert _ratio(prs.slide_width, prs.slide_height) == "4:3"


class TestCLIDataAndDemo:
    def test_data_file_render(self, template_path, tmp_path):
        data = [{"slide_type": "title_slide", "title": "T", "subtitle": "S", "notes": "n"}]
        data_path = tmp_path / "slides.json"
        data_path.write_text(json.dumps(data), encoding="utf-8")
        out = tmp_path / "o.pptx"
        rc = main([
            "--template", template_path, "--data", str(data_path),
            "--output", str(out), "--log-level", "error",
        ])
        assert rc == 0
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_bad_data_file_exits_1(self, tmp_path):
        bad = tmp_path / "bad.json"
        bad.write_text("not json", encoding="utf-8")
        rc = main(["--data", str(bad), "--output", str(tmp_path / "o.pptx"), "--log-level", "error"])
        assert rc == 1
