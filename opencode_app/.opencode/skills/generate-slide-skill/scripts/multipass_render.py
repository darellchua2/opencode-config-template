"""BT-142 Phase 3.5a — Multi-pass render + merge (engine hard limit L1).

The chenyu engine's ``slide_type`` enum has exactly 8 values
(``title_slide``, ``content_slide``, ``section_header_slide``,
``two_content_slide``, ``comparison_slide``, ``content_image_slide``,
``chart_slide``, ``closing_slide``). Each ``slide_type`` resolves to ONE
layout via fingerprint/name match. Designer decks routinely need >8
distinct layouts (BETEKK V9.1.1 needs 10: Cover, Story, Problem Impact,
Problem, Solution/Brand, Solution/Demo, Market Validation, Business Model,
Team, Ask/Closing).

This module solves L1 in three steps:

  1. **Partition** ``slide_data_list`` into batches of ≤8 distinct target
     layouts. Each slide may carry an extended field ``layout_name`` that
     names the target layout explicitly (falls back to ``slide_type`` when
     absent). The orchestrator assigns pseudo-types ``_custom_1`` ...
     ``_custom_8`` so each batch fits the engine's enum cap.
  2. **Render** each batch via ``ppt_builder.generate_ppt_from_data`` to a
     separate ``batch_N.pptx``. The engine treats pseudo-types as
     ``content_slide`` and we pin the layout via ``config_overrides``.
  3. **Merge** the per-batch decks into a single output PPTX, deep-copying
     each slide + its rels + media parts.

The merge primitive (``merge_decks``) is also independently useful: any
caller that already has multiple rendered decks can stitch them.

Public API:
  - ``merge_decks(primary_path, other_paths, output_path)`` → str (output path)
  - ``partition_slides(slide_data_list, max_layouts=8)`` → List[List[dict]]
  - ``multipass_render(slide_data_list, template_path, output_path, **kw)`` → str
"""

from __future__ import annotations

import logging
import shutil
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

DEFAULT_MAX_LAYOUTS_PER_BATCH = 8


def partition_slides(
    slide_data_list: List[dict],
    max_layouts: int = DEFAULT_MAX_LAYOUTS_PER_BATCH,
) -> List[List[dict]]:
    """Group ``slide_data_list`` into batches of ≤ ``max_layouts`` distinct layouts.

    Each slide's layout is determined by its ``layout_name`` field if present,
    else its ``slide_type``. Slides sharing a layout stay in the same batch
    (so a batch never wastes a layout slot on a duplicate).
    """
    if max_layouts < 1:
        raise ValueError("max_layouts must be ≥ 1")
    batches: List[List[dict]] = []
    current: List[dict] = []
    current_layouts: set = set()
    for slide in slide_data_list:
        layout = slide.get("layout_name") or slide.get("slide_type", "content_slide")
        if layout not in current_layouts and len(current_layouts) >= max_layouts:
            batches.append(current)
            current = []
            current_layouts = set()
        current.append(slide)
        current_layouts.add(layout)
    if current:
        batches.append(current)
    return batches


def _copy_slide(src_slide, dst_prs) -> None:
    """Deep-copy ``src_slide`` (and its rels/media) into ``dst_prs``.

    Uses the keep-together XML clone pattern. Relocates rId references for
    embedded media (images) so the destination deck is self-contained.
    """
    from copy import deepcopy

    src_prs = src_slide.part.package
    src_part = src_slide.part
    # New slide part in the destination, initially a deep copy of the source
    # XML.
    new_slide_part = deepcopy(src_part._element)
    # Use python-pptx's Slide addObject pattern: clone via XML, then patch
    # relationships for any media (image / chart) referenced.
    # python-pptx exposes SlidePart via Slide.add_part — but the simplest
    # robust path is the _SlidePart cloneshape approach via slide clone.
    from pptx.oxml.ns import qn

    dst_slidelayout = src_slide.slide_layout
    # Add a blank slide on the destination using the same layout if present.
    layout_name = dst_slidelayout.name
    target_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break
    if target_layout is None:
        target_layout = dst_prs.slide_layouts[6]  # blank fallback
    new_slide = dst_prs.slides.add_slide(target_layout)
    # Wipe the auto-created shape tree and copy each shape from the source.
    spTree = new_slide.shapes._spTree
    # remove existing children except nvGrpSpPr and grpSpPr (first two)
    for child in list(spTree):
        tag = child.tag
        if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
            continue
        spTree.remove(child)
    # Copy each shape from the source slide
    for shp in src_slide.shapes:
        spTree.append(deepcopy(shp._element))
    # Re-resolve media (image) relationships: for each blip in the source,
    # copy the image part into the destination package and remap rId.
    _relink_image_rels(src_slide, new_slide, src_prs, dst_prs)


def _relink_image_rels(src_slide, dst_slide, src_prs, dst_prs) -> None:
    """Copy image parts referenced by ``src_slide`` into ``dst_prs`` and remap rIds."""
    from pptx.oxml.ns import qn
    from pptx.image import Image
    from copy import deepcopy

    # Walk the new slide XML for <a:blip r:embed="rIdX"> elements
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in dst_slide._element.iter(qn("a:blip")):
        embed_attr = blip.get(f"{{{ns_r}}}embed")
        if not embed_attr:
            continue
        try:
            src_part = src_slide.part.related_part(embed_attr)
        except KeyError:
            continue
        # Image parts expose .blob and .image (python-pptx Image)
        try:
            image = src_part.image
        except Exception:
            continue
        # Add the image to the destination slide — this allocates a fresh rId
        # and writes the part. Then patch the blip to point at the new rId.
        # Use a no-op placeholder position (0,0,1,1); the blip's parent shape
        # already defines geometry.
        new_pic = dst_slide.shapes.add_picture_from_blob(
            image.blob, 0, 0
        ) if hasattr(dst_slide.shapes, "add_picture_from_blob") else None
        if new_pic is None:
            # Fall back to writing the image part directly via the slide part.
            image_part = dst_slide.part.get_or_add_image(image.blob)
            new_rId = dst_slide.part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
            blip.set(f"{{{ns_r}}}embed", new_rId)
            # Remove the placeholder picture we just added (only the rel is needed)
            if new_pic is not None:
                new_pic._element.getparent().remove(new_pic._element)
        else:
            new_rId = new_pic._element.xpath(".//a:blip", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})[0].get(f"{{{ns_r}}}embed")
            blip.set(f"{{{ns_r}}}embed", new_rId)
            new_pic._element.getparent().remove(new_pic._element)


def merge_decks(
    primary_path: str,
    other_paths: List[str],
    output_path: str,
) -> str:
    """Merge ``other_paths`` into ``primary_path`` → ``output_path``.

    The primary's slides come first, then each other deck's slides in order.
    All slide layouts referenced must exist on the primary's master (if a
    layout name is missing, the slide falls back to the primary's blank
    layout and a warning is logged).

    Returns the absolute path to ``output_path``.
    """
    from pptx import Presentation

    prs = Presentation(primary_path)
    for other in other_paths:
        other_prs = Presentation(other)
        for slide in other_prs.slides:
            try:
                _copy_slide(slide, prs)
            except Exception as exc:
                logger.error("merge_decks: failed to copy slide from %s: %s", other, exc)
    prs.save(output_path)
    return str(Path(output_path).resolve())


def _batch_to_engine_slides(
    batch: List[dict],
) -> Tuple[List[dict], Dict[str, str]]:
    """Convert a batch (with extended ``layout_name``) to engine-compatible slides.

    Returns ``(engine_slides, config_overrides)`` where ``config_overrides``
    maps pseudo-types to layout names (the engine pins via
    ``<slide_type>_layout`` keys).
    """
    config_overrides: Dict[str, str] = {}
    layout_to_pseudo: Dict[str, str] = {}
    engine_slides: List[dict] = []
    next_idx = 1
    for slide in batch:
        layout = slide.get("layout_name") or slide.get("slide_type")
        if layout is None:
            engine_slides.append(slide)
            continue
        if layout not in layout_to_pseudo:
            pseudo = f"_custom_{next_idx}"
            next_idx += 1
            layout_to_pseudo[layout] = pseudo
            config_overrides[f"{pseudo}_layout"] = layout
        pseudo = layout_to_pseudo[layout]
        # Copy and remap slide_type to the pseudo-type. The engine treats
        # unknown types as ``content_slide`` (per ppt_builder error handling),
        # and the config pin routes to the correct layout.
        new_slide = dict(slide)
        new_slide["slide_type"] = pseudo
        # Preserve original slide_type in notes so the engine can still
        # classify subtitle/body behaviour. ppt_builder keys behaviour off
        # _LAYOUTS_WITH_BODY / _LAYOUTS_WITH_SUBTITLE sets — pseudo-types
        # won't be in either, so we mark them as content-like via the body
        # field presence.
        if "body" not in new_slide and "body_left" not in new_slide:
            new_slide.setdefault("body", "")
        engine_slides.append(new_slide)
    return engine_slides, config_overrides


def multipass_render(
    slide_data_list: List[dict],
    template_path: str,
    output_path: str,
    max_layouts: int = DEFAULT_MAX_LAYOUTS_PER_BATCH,
    render_fn=None,
) -> str:
    """Render ``slide_data_list`` to ``output_path``, using multi-pass if >8 layouts.

    When the distinct layout count ≤ ``max_layouts``, calls ``render_fn``
    once directly (no merge overhead). When > ``max_layouts``, partitions
    into batches, renders each, and merges.

    Args:
        slide_data_list: list of slide dicts (extended with optional ``layout_name``)
        template_path: source template
        output_path: final PPTX path
        max_layouts: per-batch cap (default 8)
        render_fn: callable ``(slides, template_path, output_path,
            config_overrides) -> str`` — defaults to a thin wrapper over
            ``ppt_builder.generate_ppt_from_data``

    Returns:
        absolute path to ``output_path``
    """
    if render_fn is None:
        from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR

        def render_fn(slides, tpl, out, overrides):
            return generate_ppt_from_data(
                slides,
                template_path=tpl,
                output_path=out,
                config_overrides=overrides,
            )

    batches = partition_slides(slide_data_list, max_layouts=max_layouts)
    if len(batches) == 1:
        engine_slides, overrides = _batch_to_engine_slides(batches[0])
        return render_fn(engine_slides, template_path, output_path, overrides)

    # Multi-pass: render each batch to a temp file, then merge.
    out_dir = Path(output_path).parent
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = Path(output_path).stem
    batch_paths: List[str] = []
    for i, batch in enumerate(batches):
        engine_slides, overrides = _batch_to_engine_slides(batch)
        batch_path = str(out_dir / f"{stem}.__batch{i}.pptx")
        logger.info("multipass_render: batch %d (%d slides, %d layouts)",
                    i, len(engine_slides), len(overrides))
        render_fn(engine_slides, template_path, batch_path, overrides)
        batch_paths.append(batch_path)

    primary = batch_paths[0]
    rest = batch_paths[1:]
    result = merge_decks(primary, rest, output_path)
    # Clean up batch files
    for p in batch_paths:
        try:
            Path(p).unlink()
        except OSError:
            pass
    return result


__all__ = [
    "DEFAULT_MAX_LAYOUTS_PER_BATCH",
    "partition_slides",
    "merge_decks",
    "multipass_render",
]
