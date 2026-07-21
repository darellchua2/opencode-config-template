"""Tests for schema_extractor.py (US-1.1).

Covers the Task-6 acceptance areas from PLAN-GIT-48:
  - bundled template extract -> validate passes
  - count invariant (63 layouts, every layout has a components array)
  - non-placeholder elements captured (layout 2: >=5 components vs 3 placeholders)
  - master components non-empty
  - polygon values in [0,1], exactly 4 points
  - font-cardinality rule (C1): non-text components have no font
  - group nesting recurses
  - empty/blank layout does not crash
  - synthetic edge cases (group, table, chart, zero-shape container)
  - negative test (non-PPTX input raises TemplateExtractionError)
  - second-template robustness (default Presentation())
"""
import os
from typing import Any, List, Optional

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

from schema_extractor import (
    COMPONENT_TYPE_ENUM,
    PLACEHOLDER_TYPE_ENUM,
    TemplateExtractionError,
    TITLE_SOURCES,
    TitleInference,
    _BUILTIN_FONTS,
    _classify_shape,
    _EMBEDDED_SCHEMA_PATH,
    _extract_components,
    _extract_text_fonts,
    _font_fallback,
    _IdCounter,
    _infer_title,
    _inject_json_default,
    _signed_area,
    build_extraction_summary,
    embed_schema,
    extract_schema,
    main,
    map_shape_type,
    normalize_polygon,
    read_embedded_schema,
    validate_template_schema,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------
@pytest.fixture
def schema(template_path):
    return extract_schema(template_path)


# ---------------------------------------------------------------------------
# (1) Bundled template extract -> validate passes
# ---------------------------------------------------------------------------
class TestExtractAndValidate:
    def test_validates_clean(self, schema):
        result = validate_template_schema(schema)
        assert result.is_valid, result.error_messages()

    def test_top_level_keys(self, schema):
        for key in (
            "template_metadata",
            "slide_master",
            "slide_layouts",
            "component_type_enum",
            "placeholder_type_enum",
        ):
            assert key in schema, f"missing top-level key {key}"

    def test_metadata_fields(self, schema):
        meta = schema["template_metadata"]
        assert meta["title"] and isinstance(meta["title"], str)
        assert meta["schema_version"]
        assert meta["generated_by"] == "opencode-pptx-subagent/schema_extractor"
        assert meta["generated_at"]
        dims = meta["slide_dimensions"]
        for k in ("width_emu", "height_emu", "width_inches", "height_inches", "aspect_ratio"):
            assert k in dims
        assert dims["width_emu"] > dims["height_emu"] > 0  # 16:9 landscape

    def test_enums_self_documenting(self, schema):
        assert schema["component_type_enum"] == COMPONENT_TYPE_ENUM
        assert schema["placeholder_type_enum"] == PLACEHOLDER_TYPE_ENUM

    def test_every_component_has_valid_type_confidence(self, schema):
        # US-1.3: type_confidence is always emitted and always high/low.
        containers = schema["slide_layouts"] + [schema["slide_master"]]
        seen = set()
        for container in containers:
            for c in container["components"]:
                assert c["type_confidence"] in ("high", "low"), c["id"]
                seen.add(c["type_confidence"])
        assert seen, "template produced no components"


# ---------------------------------------------------------------------------
# (2) Count invariant + every layout has a components array
# ---------------------------------------------------------------------------
class TestCountInvariant:
    def test_layout_count_matches_pptx(self, template_path):
        prs = Presentation(template_path)
        schema = extract_schema(template_path)
        assert len(schema["slide_layouts"]) == len(prs.slide_layouts)

    def test_every_layout_has_components(self, schema):
        for L in schema["slide_layouts"]:
            assert isinstance(L["components"], list), L.get("layout_name")
            for key in ("layout_id", "layout_name", "layout_index"):
                assert key in L

    def test_layout_indices_sequential(self, schema):
        idxs = [L["layout_index"] for L in schema["slide_layouts"]]
        assert idxs == list(range(len(idxs)))


# ---------------------------------------------------------------------------
# (3) Non-placeholder elements captured (the core US-1.1 gap)
# ---------------------------------------------------------------------------
class TestNonPlaceholderCapture:
    def test_layout2_captures_non_placeholders(self, schema):
        """Layout 2 has 5 total shapes but only 3 placeholders -> 2 dropped
        before US-1.1. Now all 5 must appear as components."""
        layout2 = schema["slide_layouts"][2]
        assert len(layout2["components"]) >= 5

    def test_master_components_non_empty(self, schema):
        assert len(schema["slide_master"]["components"]) >= 1

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_component_types_beyond_placeholder_present(self, schema):
        """Extraction must yield at least one non-placeholder type (image/shape)."""
        all_types = set()
        for L in schema["slide_layouts"]:
            for c in L["components"]:
                all_types.add(c["type"])
        for c in schema["slide_master"]["components"]:
            all_types.add(c["type"])
        assert "placeholder" in all_types
        assert all_types - {"placeholder"}, "no non-placeholder types captured"


# ---------------------------------------------------------------------------
# (4) Polygon correctness
# ---------------------------------------------------------------------------
class TestPolygon:
    def test_every_polygon_is_four_points_in_range(self, schema):
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                poly = c["polygon"]
                assert len(poly) == 4, c["id"]
                for pt in poly:
                    assert set(pt.keys()) == {"x", "y"}, c["id"]
                    assert 0.0 <= pt["x"] <= 1.0, (c["id"], pt)
                    assert 0.0 <= pt["y"] <= 1.0, (c["id"], pt)

    def test_global_id_uniqueness(self, schema):
        ids = [c["id"] for L in schema["slide_layouts"] for c in L["components"]]
        ids += [c["id"] for c in schema["slide_master"]["components"]]
        assert len(ids) == len(set(ids)), "component ids not globally unique"

    def test_z_order_nonneg_int(self, schema):
        for container in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in container["components"]:
                assert isinstance(c["z_order"], int) and c["z_order"] >= 0


# ---------------------------------------------------------------------------
# (5) Font-cardinality rule (architecture review C1)
# ---------------------------------------------------------------------------
TEXT_TYPES = {"textbox", "placeholder"}


class TestFontCardinality:
    def test_non_text_components_omit_font(self, schema):
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                if c["type"] not in TEXT_TYPES:
                    assert "font" not in c, f"{c['id']} ({c['type']}) must not carry font"

    def test_text_components_carry_populated_font(self, schema):
        # US-1.4: text components carry a populated font object (all keys present),
        # not the old empty {} stub. Values may be null (inherited) but the keys
        # must be there.
        found_text = False
        required_keys = {"family", "size_pt", "weight", "color",
                         "alignment", "is_available", "fallback"}
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                if c["type"] in TEXT_TYPES:
                    found_text = True
                    assert "font" in c, c["id"]
                    assert required_keys <= set(c["font"].keys()), (c["id"], c["font"])
                    assert isinstance(c["font"]["is_available"], bool), c["id"]
                    assert "runs" in c and isinstance(c["runs"], list), c["id"]
        assert found_text, "template produced no text components"


# ---------------------------------------------------------------------------
# (5b) Font detection (US-1.4)
# ---------------------------------------------------------------------------
def _deck_with_textbox(tmp_path, runs):
    """Build a 1-slide .pptx with a textbox whose paragraphs carry the given
    runs. `runs` is a list of (text, family, size_pt, bold, rgb_or_None).
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(0, 0, 9144000, 9144000)
    tf = tb.text_frame
    first = True
    for text, family, size_pt, bold, rgb in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        if family is not None:
            run.font.name = family
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if rgb is not None:
            run.font.color.rgb = RGBColor(*rgb)
    out = tmp_path / "font_deck.pptx"
    prs.save(str(out))
    return str(out)


class TestFontFallback:
    def test_builtin_or_null_returns_none(self):
        assert _font_fallback("Arial", "Arial") is None
        assert _font_fallback(None, "Arial") is None

    def test_mapped_family(self):
        assert _font_fallback("Helvetica", "Arial") == "Arial"
        assert _font_fallback("Roboto", "Calibri") == "Arial"

    def test_unmapped_uses_theme_default(self):
        assert _font_fallback("Acme Corp Font", "Calibri") == "Calibri"
        assert _font_fallback("Acme Corp Font", "NotBuiltin") == "Arial"

    def test_builtin_set_membership(self):
        assert "Calibri" in _BUILTIN_FONTS
        assert "Roboto" not in _BUILTIN_FONTS


class TestExtractTextFonts:
    def test_populated_first_run_summary(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(0, 0, 9144000, 9144000)
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Hi"
        run.font.name = "Roboto"
        run.font.size = Pt(18)
        run.font.bold = True
        summary, runs = _extract_text_fonts(tb, "Arial")
        assert summary["family"] == "Roboto"
        assert summary["size_pt"] == 18.0
        assert summary["weight"] == "bold"
        assert summary["is_available"] is False
        assert summary["fallback"] == "Arial"
        assert runs and runs[0]["text"] == "Hi"
        assert runs[0]["font"]["family"] == "Roboto"

    def test_inherited_fields_are_null(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(0, 0, 9144000, 9144000)
        tb.text_frame.text = "Inherited"  # no explicit font props
        summary, runs = _extract_text_fonts(tb, "Arial")
        assert summary["family"] is None
        assert summary["size_pt"] is None
        assert summary["weight"] is None
        assert summary["is_available"] is True  # null family -> available
        assert summary["fallback"] is None

    def test_rgb_color_hex_theme_color_null(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(0, 0, 9144000, 9144000)
        r1 = tb.text_frame.paragraphs[0].add_run()
        r1.text = "rgb"
        r1.font.color.rgb = RGBColor(0x1A, 0x2B, 0x3C)
        summary, runs = _extract_text_fonts(tb, "Arial")
        assert summary["color"] == "#1A2B3C"
        assert runs[0]["font"]["color"] == "#1A2B3C"
        # a no-color run -> null (no crash)
        tb2 = slide.shapes.add_textbox(0, 0, 9144000, 9144000)
        tb2.text_frame.text = "nocolor"
        s2, _ = _extract_text_fonts(tb2, "Arial")
        assert s2["color"] is None

    def test_empty_textbox(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(0, 0, 9144000, 9144000)  # no text
        summary, runs = _extract_text_fonts(tb, "Arial")
        assert set(summary.keys()) == {
            "family", "size_pt", "weight", "color", "alignment",
            "is_available", "fallback",
        }
        assert summary["family"] is None
        assert summary["is_available"] is True  # null family -> available
        assert summary["fallback"] is None
        assert runs == []


class TestFontExtractionIntegration:
    def test_custom_and_builtin_runs(self, tmp_path):
        # Per-component font extraction via _extract_components on real shapes.
        deck = _deck_with_textbox(tmp_path, [
            ("Custom", "Roboto", 20, True, (0x10, 0x20, 0x30)),
            ("Builtin", "Calibri", None, None, None),
        ])
        from pptx import Presentation as P
        slide = P(deck).slides[0]
        comps = _extract_components(slide.shapes, 9144000, 5143500, _IdCounter(), "Arial")
        text = [c for c in comps if c["type"] in TEXT_TYPES]
        assert text, "textbox not captured"
        t = text[0]
        assert t["font"]["family"] == "Roboto"
        assert t["font"]["is_available"] is False
        assert t["font"]["fallback"] == "Arial"
        assert t["font"]["color"] == "#102030"
        fams = {r["font"]["family"] for r in t["runs"]}
        assert "Roboto" in fams and "Calibri" in fams

    def test_missing_fonts_aggregated_and_warned(self):
        # Decoupled from the bundled template (the shipped template may use only
        # built-in fonts). A schema carrying a non-built-in font must yield a
        # non-fatal validate warning per font. The component-level detection
        # (font.family / is_available / fallback) is covered by
        # test_custom_and_builtin_runs; this pins the validate -> warning path.
        schema = _schema_with(_ok_component())
        schema["template_metadata"]["missing_fonts"] = [
            {"family": "Roboto", "is_available": False, "fallback": "Arial"},
            {"family": "Oswald", "is_available": False, "fallback": "Arial"},
        ]
        result = validate_template_schema(schema)
        assert result.is_valid  # non-fatal warning, not an error
        assert any("non-built-in font" in w.reason and "Roboto" in w.reason
                   for w in result.warnings)
        assert any("Oswald" in w.reason for w in result.warnings)

    def test_all_builtin_yields_no_missing(self, tmp_path):
        from pptx import Presentation
        deck = tmp_path / "default.pptx"
        Presentation().save(str(deck))  # Office default theme -> built-in fonts
        schema = extract_schema(str(deck))
        assert schema["template_metadata"]["missing_fonts"] == []
        result = validate_template_schema(schema)
        assert not any("non-built-in font" in w.reason for w in result.warnings)


class TestFontValidation:
    def test_ac4_bad_fallback_is_error(self):
        comp = _ok_component()
        comp["font"] = {
            "family": "Roboto", "is_available": False,
            "fallback": "NotABuiltinFont", "size_pt": None,
            "weight": None, "color": None, "alignment": None,
        }
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("AC4" in e.reason or "built-in font name" in e.reason for e in r.errors)

    def test_valid_fallback_passes(self):
        comp = _ok_component()
        comp["font"] = {
            "family": "Roboto", "is_available": False,
            "fallback": "Arial", "size_pt": None,
            "weight": None, "color": None, "alignment": None,
        }
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid, r.error_messages()

    def test_invariant_false_available_null_fallback(self):
        # M1: the is_available == (fallback is None) invariant must catch the
        # is_available=False / fallback=None quadrant (previously missed).
        comp = _ok_component()
        comp["font"] = {
            "family": "Roboto", "is_available": False,
            "fallback": None, "size_pt": None,
            "weight": None, "color": None, "alignment": None,
        }
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid  # invariant is a warning, not an error
        assert any("must equal (fallback is None)" in w.reason for w in r.warnings)

    def test_font_string_fields_type_checked(self):
        # m4: family/weight/color/alignment must be strings or null.
        comp = _ok_component()
        comp["font"] = {
            "family": 123, "is_available": True, "fallback": None,
            "size_pt": None, "weight": None, "color": None, "alignment": None,
        }
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("font.family" in e.reason for e in r.errors)


# ---------------------------------------------------------------------------
# (6) Validation catches violations
# ---------------------------------------------------------------------------
class TestValidationFinds:
    def test_rejects_font_on_non_text(self):
        bad = {
            "template_metadata": _ok_meta(),
            "slide_master": {"name": "M", "components": []},
            "slide_layouts": [{
                "layout_id": "x", "layout_name": "X", "layout_index": 0,
                "components": [_bad_component_with_font_on_image()],
            }],
            "component_type_enum": COMPONENT_TYPE_ENUM,
            "placeholder_type_enum": PLACEHOLDER_TYPE_ENUM,
        }
        r = validate_template_schema(bad)
        assert not r.is_valid
        assert any("must not carry a 'font'" in e.reason for e in r.errors)

    def test_rejects_polygon_out_of_range(self):
        comp = _ok_component()
        comp["polygon"][0]["x"] = 1.5
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("out of [0,1]" in e.reason for e in r.errors)

    def test_rejects_bad_type_enum(self):
        comp = _ok_component()
        comp["type"] = "bogus"
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("not in component_type_enum" in e.reason for e in r.errors)

    def test_rejects_bad_type_confidence(self):
        comp = _ok_component()
        comp["type_confidence"] = "maybe"
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("type_confidence" in e.reason for e in r.errors)

    def test_missing_type_confidence_is_valid(self):
        # Backward compatibility: type_confidence is optional.
        comp = _ok_component()
        comp.pop("type_confidence", None)
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid, r.error_messages()

    def test_shape_low_is_warning_not_error(self):
        # MINOR-2: shape/low must surface as a non-fatal warning so unrecognized
        # elements are "flagged for review" without blocking is_valid.
        comp = _ok_component()
        comp["type"] = "shape"
        comp["type_confidence"] = "low"
        # shape components must not carry a font (C1) -> drop it for a clean case.
        comp.pop("font", None)
        comp.pop("runs", None)
        comp.pop("content_template", None)
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid  # warning, not error
        assert any("flagged for review" in w.reason for w in r.warnings)

    def test_video_low_emits_no_warning(self):
        # MINOR-3: the "flagged for review" warning is shape-only (decision 6).
        # Indeterminate MEDIA (video/low) must pass completely silently so the
        # scoped boundary is locked against accidental widening.
        comp = _ok_component()
        comp["type"] = "video"
        comp["type_confidence"] = "low"
        # video is non-text -> must not carry font/runs/content_template (C1).
        comp.pop("font", None)
        comp.pop("runs", None)
        comp.pop("content_template", None)
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid
        assert r.warnings == []

    def test_rejects_non_object(self):
        r = validate_template_schema("not a dict")
        assert not r.is_valid


# ---------------------------------------------------------------------------
# Winding check (US-1.2): cross-product / shoelace signed-area verification
# ---------------------------------------------------------------------------
class TestWinding:
    def test_canonical_axis_aligned_passes(self):
        # TL->TR->BR->BL => positive signed area => valid.
        r = validate_template_schema(_schema_with(_ok_component()))
        assert r.is_valid, r.error_messages()

    def test_reversed_winding_is_error(self):
        comp = _ok_component()
        # BL -> BR -> TR -> TL (reversed) => negative signed area.
        comp["polygon"] = [
            {"x": 0.1, "y": 0.9}, {"x": 0.9, "y": 0.9},
            {"x": 0.9, "y": 0.1}, {"x": 0.1, "y": 0.1},
        ]
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("reversed winding" in e.reason for e in r.errors)

    def test_degenerate_is_warning_not_error(self):
        comp = _ok_component()
        comp["polygon"] = [{"x": 0.5, "y": 0.5}] * 4  # zero area
        r = validate_template_schema(_schema_with(comp))
        # degenerate => warning only => is_valid stays True
        assert r.is_valid
        assert any("degenerate" in w.reason for w in r.warnings)

    def test_non_axis_aligned_canonical_passes(self):
        # A tilted quadrilateral in canonical (CCW) winding — proves the check
        # is a real shoelace, not a rect-only shortcut.
        comp = _ok_component()
        comp["polygon"] = [
            {"x": 0.2, "y": 0.1}, {"x": 0.9, "y": 0.3},
            {"x": 0.7, "y": 0.9}, {"x": 0.1, "y": 0.7},
        ]
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid, r.error_messages()

    def test_collinear_is_warning(self):
        comp = _ok_component()
        # 4 collinear points => zero area, degenerate (warning).
        comp["polygon"] = [
            {"x": 0.1, "y": 0.5}, {"x": 0.4, "y": 0.5},
            {"x": 0.7, "y": 0.5}, {"x": 0.9, "y": 0.5},
        ]
        r = validate_template_schema(_schema_with(comp))
        assert r.is_valid  # warning, not error
        assert any("degenerate" in w.reason for w in r.warnings)

    def test_bundled_template_still_validates(self, schema):
        # Regression guard: real extraction must still pass (any real zero-area
        # shape would now be a non-fatal warning, keeping is_valid True).
        r = validate_template_schema(schema)
        assert r.is_valid, [e.format() for e in r.errors]

    def test_tilted_reversed_is_error(self):
        # A tilted quadrilateral in REVERSED winding => negative area => error.
        comp = _ok_component()
        comp["polygon"] = [
            {"x": 0.1, "y": 0.7}, {"x": 0.7, "y": 0.9},
            {"x": 0.9, "y": 0.3}, {"x": 0.2, "y": 0.1},
        ]
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("reversed winding" in e.reason for e in r.errors)


# ---------------------------------------------------------------------------
# _signed_area unit tests (forward-compat: n-point + guards)
# ---------------------------------------------------------------------------
class TestSignedArea:
    def test_canonical_rect_positive(self):
        poly = [{"x": 0.1, "y": 0.1}, {"x": 0.9, "y": 0.1},
                {"x": 0.9, "y": 0.9}, {"x": 0.1, "y": 0.9}]
        assert _signed_area(poly) > 0

    def test_reversed_negative(self):
        poly = [{"x": 0.1, "y": 0.9}, {"x": 0.9, "y": 0.9},
                {"x": 0.9, "y": 0.1}, {"x": 0.1, "y": 0.1}]
        assert _signed_area(poly) < 0

    def test_triangle_works(self):
        # n=3 (forward-compat for future non-rectangular vertices).
        poly = [{"x": 0.0, "y": 0.0}, {"x": 1.0, "y": 0.0}, {"x": 0.5, "y": 1.0}]
        assert _signed_area(poly) == pytest.approx(0.5)

    def test_fewer_than_three_zero(self):
        assert _signed_area([{"x": 0.1, "y": 0.1}, {"x": 0.9, "y": 0.9}]) == 0.0

    def test_degenerate_collinear_zero(self):
        poly = [{"x": 0.1, "y": 0.5}, {"x": 0.4, "y": 0.5},
                {"x": 0.7, "y": 0.5}, {"x": 0.9, "y": 0.5}]
        assert _signed_area(poly) == 0.0




def _ok_meta():
    return {
        "title": "T", "schema_version": "1.0.0",
        "generated_by": "x", "generated_at": "2026-01-01T00:00:00Z",
        "slide_dimensions": {
            "width_emu": 9144000, "height_emu": 5143500,
            "width_inches": 10.0, "height_inches": 5.625, "aspect_ratio": "16:9",
        },
    }


def _ok_component():
    return {
        "id": "comp_001", "type": "textbox", "name": "N",
        "type_confidence": "high",
        "polygon": [
            {"x": 0.1, "y": 0.1}, {"x": 0.9, "y": 0.1},
            {"x": 0.9, "y": 0.9}, {"x": 0.1, "y": 0.9},
        ],
        "z_order": 0,
        "placeholder_type": None, "font": {}, "runs": [],
        "content_template": "{{content}}",
    }


def _media_element(marker: str):
    """Build a minimal <p:pic> element carrying an <a:{marker}File> under <p:nvPr>.

    marker is 'audio' or 'video'. Used to attach a synthetic _element to a
    FakeShape so _classify_shape's MEDIA branch can be exercised end-to-end.
    """
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    xml = (
        f'<p:pic xmlns:p="{p}" xmlns:a="{a}" xmlns:r="{r}">'
        f'<p:nvPicPr><p:nvPr><a:{marker}File r:link="rId1"/></p:nvPr></p:nvPicPr>'
        f'</p:pic>'
    )
    return etree.fromstring(xml)


def _bad_component_with_font_on_image():
    c = _ok_component()
    c["id"] = "comp_002"
    c["type"] = "image"
    c["font"] = {"family": "Arial"}  # image must not carry font
    return c


def _schema_with(comp):
    return {
        "template_metadata": _ok_meta(),
        "slide_master": {"name": "M", "components": []},
        "slide_layouts": [{
            "layout_id": "x", "layout_name": "X", "layout_index": 0,
            "components": [comp],
        }],
        "component_type_enum": COMPONENT_TYPE_ENUM,
        "placeholder_type_enum": PLACEHOLDER_TYPE_ENUM,
    }


# ---------------------------------------------------------------------------
# (7) Synthetic edge cases — unit-test mapper + recursion with fake shapes
# ---------------------------------------------------------------------------
class FakePlaceholderFormat:
    def __init__(self, ptype: Any) -> None:
        self.type = ptype


class FakeShape:
    """Minimal stand-in for a python-pptx shape, for pure-logic unit tests."""

    def __init__(
        self,
        name: str,
        st: Any = None,
        *,
        is_placeholder: bool = False,
        has_table: bool = False,
        has_chart: bool = False,
        left: int = 0,
        top: int = 0,
        width: int = 0,
        height: int = 0,
        children: Optional[List["FakeShape"]] = None,
        ph_type: Any = None,
    ) -> None:
        self.name = name
        self.shape_type = st
        self.is_placeholder = is_placeholder
        self.has_table = has_table
        self.has_chart = has_chart
        self.left, self.top, self.width, self.height = left, top, width, height
        self._children = children or []
        self.placeholder_format = FakePlaceholderFormat(ph_type) if is_placeholder else None
        # _element is None by default; tests attach a synthetic lxml element
        # (e.g. _media_element) to exercise the MEDIA audio/video split.
        self._element = None

    @property
    def shapes(self):
        return self._children


SLIDE_W = 9144000  # EMU (10in)
SLIDE_H = 5143500  # EMU


class TestTypeMapper:
    def test_table_detected_before_shape_type(self):
        s = FakeShape("t", st=None, has_table=True)
        assert map_shape_type(s) == "table"

    def test_chart_detected(self):
        s = FakeShape("c", st=None, has_chart=True)
        assert map_shape_type(s) == "chart"

    def test_placeholder_takes_precedence(self):
        s = FakeShape("p", st=MSO_SHAPE_TYPE.PICTURE, is_placeholder=True, ph_type=PP_PLACEHOLDER.TITLE)
        assert map_shape_type(s) == "placeholder"

    def test_picture_image(self):
        s = FakeShape("pic", st=MSO_SHAPE_TYPE.PICTURE)
        assert map_shape_type(s) == "image"

    def test_group(self):
        s = FakeShape("g", st=MSO_SHAPE_TYPE.GROUP)
        assert map_shape_type(s) == "group"

    def test_textbox(self):
        s = FakeShape("tb", st=MSO_SHAPE_TYPE.TEXT_BOX)
        assert map_shape_type(s) == "textbox"

    def test_media_video(self):
        s = FakeShape("m", st=MSO_SHAPE_TYPE.MEDIA)
        assert map_shape_type(s) == "video"

    def test_smartart(self):
        s = FakeShape("sm", st=MSO_SHAPE_TYPE.IGX_GRAPHIC)
        assert map_shape_type(s) == "smartart"

    def test_unknown_degrades_to_shape(self):
        s = FakeShape("u", st=None)
        assert map_shape_type(s) == "shape"


# ---------------------------------------------------------------------------
# type_confidence + audio/video subtype (US-1.3)
# ---------------------------------------------------------------------------
class TestTypeConfidence:
    def test_recognized_types_are_high(self):
        cases = [
            (FakeShape("t", st=None, has_table=True), "table"),
            (FakeShape("c", st=None, has_chart=True), "chart"),
            (FakeShape("pic", st=MSO_SHAPE_TYPE.PICTURE), "image"),
            (FakeShape("g", st=MSO_SHAPE_TYPE.GROUP), "group"),
            (FakeShape("tb", st=MSO_SHAPE_TYPE.TEXT_BOX), "textbox"),
            (FakeShape("sm", st=MSO_SHAPE_TYPE.IGX_GRAPHIC), "smartart"),
        ]
        for shape, expected in cases:
            assert _classify_shape(shape) == (expected, "high")

    def test_placeholder_is_high(self):
        s = FakeShape("p", st=MSO_SHAPE_TYPE.PICTURE, is_placeholder=True, ph_type=PP_PLACEHOLDER.TITLE)
        assert _classify_shape(s) == ("placeholder", "high")

    def test_recognized_shapelike_types_are_shape_high(self):
        # AUTO_SHAPE / FREEFORM / LINE are recognized preset-geometry shapes.
        for st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.LINE):
            assert _classify_shape(FakeShape("x", st=st)) == ("shape", "high")

    def test_other_non_none_members_default_shape_high(self):
        # MAJOR-1 option (a): every non-None shape_type is recognized -> high.
        # CALLOUT / LINKED_PICTURE / TEXT_EFFECT / INK must NOT be mislabeled low.
        for st in (
            MSO_SHAPE_TYPE.CALLOUT,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
            MSO_SHAPE_TYPE.TEXT_EFFECT,
            MSO_SHAPE_TYPE.INK,
        ):
            assert _classify_shape(FakeShape("x", st=st)) == ("shape", "high"), st

    def test_web_video_maps_to_video_high(self):
        s = FakeShape("wv", st=MSO_SHAPE_TYPE.WEB_VIDEO)
        assert _classify_shape(s) == ("video", "high")

    def test_none_shape_type_is_shape_low(self):
        assert _classify_shape(FakeShape("u", st=None)) == ("shape", "low")

    def test_shape_type_exception_is_shape_low(self):
        class Boom:
            is_placeholder = False
            has_table = False
            has_chart = False

            @property
            def shape_type(self):
                raise RuntimeError("boom")

        assert _classify_shape(Boom()) == ("shape", "low")

    def test_media_without_marker_is_video_low(self):
        s = FakeShape("m", st=MSO_SHAPE_TYPE.MEDIA)  # _element is None
        assert _classify_shape(s) == ("video", "low")


class TestAudioVideo:
    def test_audio_marker_detected(self):
        s = FakeShape("a", st=MSO_SHAPE_TYPE.MEDIA)
        s._element = _media_element("audio")
        assert _classify_shape(s) == ("audio", "high")

    def test_video_marker_detected_high(self):
        s = FakeShape("v", st=MSO_SHAPE_TYPE.MEDIA)
        s._element = _media_element("video")
        assert _classify_shape(s) == ("video", "high")

    def test_audio_enum_reachable_via_wrapper(self):
        # The "audio" enum value is now reachable end-to-end (US-1.3 D2).
        s = FakeShape("a", st=MSO_SHAPE_TYPE.MEDIA)
        s._element = _media_element("audio")
        assert map_shape_type(s) == "audio"


class TestPolygonNormalizer:
    def test_four_points_tl_tr_br_bl(self):
        s = FakeShape("s", left=0, top=0, width=SLIDE_W, height=SLIDE_H)
        poly = normalize_polygon(s, SLIDE_W, SLIDE_H)
        assert poly[0] == {"x": 0.0, "y": 0.0}      # TL
        assert poly[1] == {"x": 1.0, "y": 0.0}      # TR
        assert poly[2] == {"x": 1.0, "y": 1.0}      # BR
        assert poly[3] == {"x": 0.0, "y": 1.0}      # BL

    def test_clamps_overflow(self):
        # shape extends beyond slide -> clamped to 1.0
        s = FakeShape("s", left=SLIDE_W // 2, top=SLIDE_H // 2,
                      width=SLIDE_W, height=SLIDE_H)
        poly = normalize_polygon(s, SLIDE_W, SLIDE_H)
        assert all(0.0 <= pt["x"] <= 1.0 for pt in poly)
        assert all(0.0 <= pt["y"] <= 1.0 for pt in poly)

    def test_zero_dims_safe(self):
        s = FakeShape("s")
        poly = normalize_polygon(s, 0, 0)
        assert poly == [{"x": 0.0, "y": 0.0}] * 4


class TestGroupRecursion:
    def test_group_children_flattened(self):
        child1 = FakeShape("c1", st=MSO_SHAPE_TYPE.TEXT_BOX,
                           left=0, top=0, width=1000, height=1000)
        child2 = FakeShape("c2", st=MSO_SHAPE_TYPE.PICTURE,
                           left=2000, top=2000, width=1000, height=1000)
        group = FakeShape("grp", st=MSO_SHAPE_TYPE.GROUP,
                          left=0, top=0, width=3000, height=3000,
                          children=[child1, child2])
        comps = _extract_components([group], SLIDE_W, SLIDE_H, _IdCounter())
        # group + 2 children = 3 components
        assert len(comps) == 3
        assert comps[0]["type"] == "group"
        assert {comps[1]["type"], comps[2]["type"]} == {"textbox", "image"}

    def test_z_order_unique_when_group_has_following_sibling(self):
        """Regression for code-review M1: a group followed by a sibling must
        not produce a z_order collision. Order: [shape0, group(c1,c2), shape3]."""
        shape0 = FakeShape("s0", st=MSO_SHAPE_TYPE.TEXT_BOX,
                           left=0, top=0, width=1000, height=1000)
        c1 = FakeShape("c1", st=MSO_SHAPE_TYPE.PICTURE, left=0, top=0, width=1000, height=1000)
        c2 = FakeShape("c2", st=MSO_SHAPE_TYPE.PICTURE, left=0, top=0, width=1000, height=1000)
        group = FakeShape("grp", st=MSO_SHAPE_TYPE.GROUP,
                          left=0, top=0, width=1000, height=1000, children=[c1, c2])
        shape3 = FakeShape("s3", st=MSO_SHAPE_TYPE.TEXT_BOX,
                           left=0, top=0, width=1000, height=1000)
        comps = _extract_components([shape0, group, shape3], SLIDE_W, SLIDE_H, _IdCounter())
        assert len(comps) == 5  # s0, group, c1, c2, s3
        z_orders = [c["z_order"] for c in comps]
        # Monotonic + unique (0..4) — the bug produced a duplicate '2'.
        assert z_orders == [0, 1, 2, 3, 4], z_orders
        assert len(z_orders) == len(set(z_orders))

    def test_empty_container_does_not_crash(self):
        comps = _extract_components([], SLIDE_W, SLIDE_H, _IdCounter())
        assert comps == []


# ---------------------------------------------------------------------------
# (8) Second-template robustness (default Presentation() != bundled template)
# ---------------------------------------------------------------------------
class TestSecondTemplate:
    def test_default_presentation_extracts(self, tmp_path):
        prs = Presentation()  # default blank deck (different master/layouts)
        out = tmp_path / "default.pptx"
        prs.save(str(out))
        schema = extract_schema(str(out))
        r = validate_template_schema(schema)
        assert r.is_valid, r.error_messages()
        assert len(schema["slide_layouts"]) == len(prs.slide_layouts)
        assert isinstance(schema["slide_master"]["components"], list)


# ---------------------------------------------------------------------------
# (9) Negative test — non-PPTX input raises domain error
# ---------------------------------------------------------------------------
class TestNegativePath:
    def test_non_pptx_raises_domain_error(self, tmp_path):
        bad = tmp_path / "not_a_deck.pptx"
        bad.write_bytes(b"this is not a zip file at all")
        with pytest.raises(TemplateExtractionError):
            extract_schema(str(bad))

    def test_missing_file_raises_domain_error(self):
        with pytest.raises(TemplateExtractionError):
            extract_schema("does_not_exist.pptx")


# ---------------------------------------------------------------------------
# (11) US-1.5 — embed / read the schema inside the PPTX zip
# ---------------------------------------------------------------------------
import hashlib  # noqa: E402
import zipfile  # noqa: E402


def _save_synthetic_deck(tmp_path, title=None):
    """A small real .pptx on disk (fast, deterministic fixture for zip tests)."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title is not None:
        slide.shapes.add_textbox(0, 0, 9144000, 9144000).text_frame.text = title
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return str(out)


class TestInjectJsonDefault:
    def test_injects_when_absent(self):
        xml = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>' \
              b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">' \
              b'<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>' \
              b'<Default Extension="xml" ContentType="application/xml"/>' \
              b'</Types>'
        out = _inject_json_default(xml)
        assert b'Extension="json"' in out and b'application/json' in out

    def test_idempotent_when_present(self):
        xml = b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">' \
              b'<Default Extension="json" ContentType="application/json"/></Types>'
        # already declared -> returned unchanged
        assert _inject_json_default(xml) == xml


class TestEmbedSchema:
    def test_round_trip_deep_equal(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        schema = extract_schema(deck)
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, schema, embedded)
        assert read_embedded_schema(embedded) == schema

    def test_content_types_first_valid_json_default_intact(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, extract_schema(deck), embedded)
        with zipfile.ZipFile(embedded) as z:
            names = z.namelist()
            assert names[0] == "[Content_Types].xml"
            ct = z.read("[Content_Types].xml").decode("utf-8")
        assert 'Extension="json"' in ct and 'application/json' in ct
        # the original Defaults (rels, xml) are preserved
        assert 'Extension="rels"' in ct and 'Extension="xml"' in ct

    def test_other_entries_decompressed_identical_only_schema_new(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, extract_schema(deck), embedded)
        with zipfile.ZipFile(deck) as zo, zipfile.ZipFile(embedded) as ze:
            onames = set(zo.namelist())
            enames = set(ze.namelist())
            assert onames.issubset(enames)                      # nothing dropped
            assert enames - onames == {_EMBEDDED_SCHEMA_PATH}   # only the schema is new
            for n in onames:
                if n == "[Content_Types].xml":
                    continue
                assert hashlib.md5(zo.read(n)).hexdigest() == hashlib.md5(ze.read(n)).hexdigest(), n

    def test_python_pptx_reopens_same_counts(self, tmp_path):
        from pptx import Presentation
        deck = _save_synthetic_deck(tmp_path)
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, extract_schema(deck), embedded)
        p_o, p_e = Presentation(deck), Presentation(embedded)
        assert len(p_o.slides) == len(p_e.slides)
        assert len(p_o.slide_layouts) == len(p_e.slide_layouts)

    def test_idempotent_re_embed(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        emb1 = str(tmp_path / "emb1.pptx")
        schema1 = extract_schema(deck)
        embed_schema(deck, schema1, emb1)
        # re-embed the already-embedded PPTX with a different schema
        schema2 = dict(schema1)
        schema2["__second__"] = True
        emb2 = str(tmp_path / "emb2.pptx")
        embed_schema(emb1, schema2, emb2)
        with zipfile.ZipFile(emb2) as z:
            assert z.namelist().count(_EMBEDDED_SCHEMA_PATH) == 1
        assert read_embedded_schema(emb2) == schema2

    def test_result_struct_and_minified(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        embedded = str(tmp_path / "deck.templated.pptx")
        result = embed_schema(deck, extract_schema(deck), embedded)
        assert result.output_path == embedded
        assert result.original_bytes > 0 and result.new_bytes > 0
        # embedded JSON is minified (no ", " / ": " insignificant whitespace)
        with zipfile.ZipFile(embedded) as z:
            payload = z.read(_EMBEDDED_SCHEMA_PATH).decode("utf-8")
        assert '", "' not in payload and '": "' not in payload

    def test_non_ascii_round_trip(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path, title="模板测试 — ✓")
        schema = extract_schema(deck)
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, schema, embedded)
        assert read_embedded_schema(embedded) == schema  # ensure_ascii=False round-trips

    def test_bundled_template_smoke(self, template_path, tmp_path):
        # realism: the real bundled template. NB: since US-4.1 the bundled
        # template ships PRE-TEMPLATED (embed_schema ran on it once), so re-
        # embedding is idempotent and delta_bytes may be 0 (a US-1.5 guarantee).
        # The load-bearing assertion is that the output carries a readable schema.
        embedded = str(tmp_path / "template.templated.pptx")
        embed_schema(template_path, extract_schema(template_path), embedded)
        assert read_embedded_schema(embedded) is not None


class TestReadEmbeddedSchema:
    def test_absent_returns_none(self, tmp_path):
        deck = _save_synthetic_deck(tmp_path)
        assert read_embedded_schema(deck) is None

    def test_malformed_json_returns_none(self, tmp_path, caplog):
        deck = _save_synthetic_deck(tmp_path)
        bad = str(tmp_path / "bad.pptx")
        # copy deck and write garbage at the schema path
        import shutil
        shutil.copy(deck, bad)
        with zipfile.ZipFile(bad, "a") as z:
            z.writestr(_EMBEDDED_SCHEMA_PATH, b"not-json{{")
        with caplog.at_level("WARNING"):
            assert read_embedded_schema(bad) is None
        # MINOR-2: the warning half of the contract must actually fire.
        assert any("malformed" in r.message for r in caplog.records), "expected a malformed-JSON warning"

    def test_non_dict_json_returns_none(self, tmp_path, caplog):
        # MINOR-2: valid JSON but not an object (e.g. an array) -> None + warning.
        import shutil
        deck = _save_synthetic_deck(tmp_path)
        bad = str(tmp_path / "nondict.pptx")
        shutil.copy(deck, bad)
        with zipfile.ZipFile(bad, "a") as z:
            z.writestr(_EMBEDDED_SCHEMA_PATH, b"[1, 2, 3]")
        with caplog.at_level("WARNING"):
            assert read_embedded_schema(bad) is None
        assert any("not a JSON object" in r.message for r in caplog.records)

    def test_missing_file_raises_domain_error(self):
        # MINOR-2: a missing file -> OSError -> TemplateExtractionError (mirrors extract_schema).
        with pytest.raises(TemplateExtractionError):
            read_embedded_schema("does_not_exist.pptx")

    def test_non_zip_raises_domain_error(self, tmp_path):
        bad = tmp_path / "not_a_zip.pptx"
        bad.write_bytes(b"definitely not a zip")
        with pytest.raises(TemplateExtractionError):
            read_embedded_schema(str(bad))


# ---------------------------------------------------------------------------
# (12) US-3.1 -- title_source provenance + extraction summary + CLI --summary
# ---------------------------------------------------------------------------
from pptx import Presentation as _Presentation  # noqa: E402


def _deck(tmp_path, name="deck.pptx", core_title=None, slide1_text=None):
    """Build a 1-slide deck controlling core.xml title and slide-1 text."""
    prs = _Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if slide1_text is not None:
        slide.shapes.add_textbox(0, 0, 9144000, 9144000).text_frame.text = slide1_text
    if core_title is not None:
        prs.core_properties.title = core_title
    out = tmp_path / name
    prs.save(str(out))
    return str(out)


class TestTitleInferenceSource:
    """US-3.1 MINOR-5: _infer_title returns a TitleInference with provenance."""

    def test_namedtuple_fields(self):
        inf = TitleInference("x", "filename")
        assert inf.title == "x"
        assert inf.source == "filename"

    def test_core_xml_wins(self, tmp_path):
        deck = _deck(tmp_path, core_title="Core Doc Title", slide1_text="Slide Title")
        prs = _Presentation(deck)
        inf = _infer_title(prs, deck)
        assert inf.title == "Core Doc Title"
        assert inf.source == "core_xml"

    def test_slide1_when_no_core_title(self, tmp_path):
        deck = _deck(tmp_path, slide1_text="First Slide Heading")
        prs = _Presentation(deck)
        inf = _infer_title(prs, deck)
        assert inf.title == "First Slide Heading"
        assert inf.source == "slide1"

    def test_filename_fallback(self, tmp_path):
        deck = _deck(tmp_path, name="my_template.pptx")
        prs = _Presentation(deck)
        inf = _infer_title(prs, deck)
        assert inf.title == "my_template"
        assert inf.source == "filename"

    def test_extract_emits_title_source(self, tmp_path):
        # extract_schema must carry title_source into template_metadata (Task 2).
        deck = _deck(tmp_path, core_title="From Core")
        schema = extract_schema(deck)
        meta = schema["template_metadata"]
        assert meta["title"] == "From Core"
        assert meta["title_source"] == "core_xml"

    def test_skill_write_back_to_user(self, tmp_path):
        # The skill layer sets title_source="user" after a user override; the
        # field round-trips through embed -> read unchanged.
        deck = _deck(tmp_path, core_title="Inferred")
        schema = extract_schema(deck)
        schema["template_metadata"]["title"] = "User Renamed"
        schema["template_metadata"]["title_source"] = "user"
        embedded = str(tmp_path / "deck.templated.pptx")
        embed_schema(deck, schema, embedded)
        back = read_embedded_schema(embedded)
        assert back["template_metadata"]["title_source"] == "user"


class TestTitleSourceValidation:
    """US-3.1 MAJOR-2: validate_template_schema enforces the title_source enum."""

    @staticmethod
    def _base_schema():
        # A minimal valid schema; title_source is the only varying field.
        return {
            "template_metadata": {
                "title": "T", "title_source": "core_xml", "schema_version": "1.0.0",
                "generated_by": "x", "generated_at": "now",
                "slide_dimensions": {"width_emu": 1, "height_emu": 1,
                                     "width_inches": 1.0, "height_inches": 1.0, "aspect_ratio": "1:1"},
            },
            "slide_master": {"name": "m", "components": []},
            "slide_layouts": [],
            "component_type_enum": list(COMPONENT_TYPE_ENUM),
            "placeholder_type_enum": list(PLACEHOLDER_TYPE_ENUM),
        }

    def test_valid_sources_pass(self):
        for src in TITLE_SOURCES:
            s = self._base_schema()
            s["template_metadata"]["title_source"] = src
            res = validate_template_schema(s)
            assert res.is_valid, (src, res.error_messages())

    def test_invalid_source_fails(self):
        s = self._base_schema()
        s["template_metadata"]["title_source"] = "garbage"
        res = validate_template_schema(s)
        assert not res.is_valid
        assert any("title_source" in m for m in res.error_messages())

    def test_absent_source_ok(self):
        # title_source is optional (legacy schemas without it still validate).
        s = self._base_schema()
        del s["template_metadata"]["title_source"]
        assert validate_template_schema(s).is_valid

    def test_bundled_template_validates_with_title_source(self, schema):
        # Regression: the real bundled template now emits title_source and still validates.
        assert schema["template_metadata"].get("title_source") in TITLE_SOURCES
        assert validate_template_schema(schema).is_valid


class TestExtractionSummary:
    """US-3.1 / US-3.3 AC2: build_extraction_summary content."""

    def test_contains_key_lines(self, schema):
        summary = build_extraction_summary(schema)
        assert isinstance(summary, str)
        assert "Title:" in summary
        assert "Layouts:" in summary
        assert "Slide master:" in summary
        # theme + fonts sections appear when the template carries them
        assert "Theme colors:" in summary
        assert "Font palette:" in summary

    def test_shows_title_source(self, tmp_path):
        deck = _deck(tmp_path, core_title="Deck Name")
        schema = extract_schema(deck)
        summary = build_extraction_summary(schema)
        assert "source: core_xml" in summary

    def test_pure_no_mutation(self, schema):
        before = dict(schema)
        build_extraction_summary(schema)
        assert schema == before  # the function must not mutate its input

    def test_handles_minimal_schema(self):
        # Defensive: missing optional keys do not crash.
        s = {"template_metadata": {"title": "X"}, "slide_master": {}, "slide_layouts": []}
        out = build_extraction_summary(s)
        assert "Title: X" in out
        assert "Layouts: 0" in out


class TestCliSummaryFlag:
    """US-3.1 Task 6: CLI --summary prints to stdout."""

    def test_summary_prints_and_exits_zero(self, template_path, tmp_path, capsys):
        out_json = str(tmp_path / "out.json")
        rc = main(["--input", template_path, "--output", out_json, "--summary"])
        assert rc == 0
        captured = capsys.readouterr()
        assert "Title:" in captured.out
        assert "Layouts:" in captured.out

    def test_no_summary_flag_is_silent(self, template_path, tmp_path, capsys):
        out_json = str(tmp_path / "out.json")
        rc = main(["--input", template_path, "--output", out_json])
        assert rc == 0
        captured = capsys.readouterr()
        assert "Title:" not in captured.out

    def test_embed_with_summary(self, template_path, tmp_path, capsys):
        # The full end-to-end CLI path (extract + embed + summary) exits 0.
        out_json = str(tmp_path / "out.json")
        out_pptx = str(tmp_path / "out.templated.pptx")
        rc = main([
            "--input", template_path, "--output", out_json,
            "--embed", "--output-pptx", out_pptx, "--summary",
        ])
        assert rc == 0
        captured = capsys.readouterr()
        assert "Layouts:" in captured.out
        assert read_embedded_schema(out_pptx) is not None
