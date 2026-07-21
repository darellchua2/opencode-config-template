"""Tests for BT-142 Phase 3.5b placeholder_backfill.

Uses lightweight shape fakes (same pattern as test_container_check) because
python-pptx's `placeholder_format` is read-only — adding a real placeholder
requires raw OOXML manipulation. The backfill logic is what we care about
here; integration tests with real placeholders are covered by Phase 6.12.
"""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from placeholder_backfill import backfill_slide, backfill_deck, BackfillReport


class _FakeRun:
    def __init__(self):
        self.text = ""


class _FakeParagraph:
    def __init__(self):
        self.runs = []

    def add_run(self):
        r = _FakeRun()
        self.runs.append(r)
        return r


class _FakeTextFrame:
    def __init__(self):
        self.text = ""
        self.paragraphs = [_FakeParagraph()]  # python-pptx always has ≥1

    def clear(self):
        self.text = ""
        self.paragraphs = [_FakeParagraph()]

    def add_paragraph(self):
        p = _FakeParagraph()
        self.paragraphs.append(p)
        return p


class _FakePhFmt:
    def __init__(self, idx, type_=None):
        self.idx = idx
        self.type = type_


class _FakeShape:
    def __init__(self, idx, text="", type_=None):
        self.name = f"PH {idx}"
        self.placeholder_format = _FakePhFmt(idx, type_)
        self.is_placeholder = True
        self._text = text
        self.has_text_frame = True
        self.text_frame = _FakeTextFrame()
        self.text_frame.text = text
        self.left = self.top = self.width = self.height = 0

    @property
    def text(self):
        return self._text


class _FakeShapes:
    def __init__(self, shapes):
        self._shapes = shapes

    def __iter__(self):
        return iter(self._shapes)

    def add_picture(self, *args, **kw):
        return _FakeShape(99, type_=18)


class _FakeSlide:
    def __init__(self, placeholders):
        self.shapes = _FakeShapes(placeholders)
        self.placeholders = placeholders


def test_backfill_body_slot_fills_empty_placeholder():
    slide = _FakeSlide([_FakeShape(idx=1, text=""), _FakeShape(idx=2, text="")])
    report = backfill_slide(slide, {
        "body_slots": [
            {"placeholder_idx": 1, "text": "Member A — CEO"},
            {"placeholder_idx": 2, "text": "Member B — CTO"},
        ]
    })
    assert "body[1]" in report.filled
    assert "body[2]" in report.filled
    assert len(report.errors) == 0


def test_backfill_skips_already_filled():
    slide = _FakeSlide([_FakeShape(idx=1, text="Existing content")])
    report = backfill_slide(slide, {
        "body_slots": [{"placeholder_idx": 1, "text": "Should not overwrite"}]
    })
    assert any("body[1]" in s for s in report.skipped)
    assert not any("body[1]" in f for f in report.filled)


def test_backfill_missing_placeholder_reports_error():
    slide = _FakeSlide([])  # no placeholders
    report = backfill_slide(slide, {
        "body_slots": [{"placeholder_idx": 99, "text": "x"}]
    })
    assert len(report.errors) >= 1
    assert any("99" in err for err in report.errors)


def test_backfill_handles_empty_slide_data():
    slide = _FakeSlide([_FakeShape(idx=1, text="")])
    report = backfill_slide(slide, {})
    assert report.filled == []
    assert report.errors == []


def test_backfill_deck_aligns_by_index():
    slide1 = _FakeSlide([_FakeShape(idx=1, text="")])
    slide2 = _FakeSlide([_FakeShape(idx=1, text="")])
    reports = backfill_deck(
        type("P", (), {"slides": [slide1, slide2]})(),
        [{"body_slots": [{"placeholder_idx": 1, "text": "x"}]},
         {"body_slots": [{"placeholder_idx": 1, "text": "y"}]}],
    )
    assert set(reports.keys()) == {0, 1}
    assert reports[0].filled == ["body[1]"]
    assert reports[1].filled == ["body[1]"]


def test_backfill_deck_handles_size_mismatch():
    slide1 = _FakeSlide([_FakeShape(idx=1, text="")])
    slide2 = _FakeSlide([_FakeShape(idx=1, text="")])
    reports = backfill_deck(
        type("P", (), {"slides": [slide1, slide2]})(),
        [{}],  # only 1 entry for 2 slides
    )
    assert len(reports[1].errors) >= 1


def test_backfill_report_to_dict_round_trip():
    r = BackfillReport(filled=["body[1]"], skipped=["body[2]"], errors=["err"])
    d = r.to_dict()
    assert d["filled"] == ["body[1]"]
    assert d["skipped"] == ["body[2]"]
    assert d["errors"] == ["err"]


# ---------------------------------------------------------------------------
# MED-1 (BT-142 Phase 8.5): _add_picture_into_placeholder coverage
# Uses real python-pptx slides + a real PNG image so the add_picture code
# path (fill AND fit sizing) is exercised end-to-end, not just the fake
# _FakeShapes.add_picture stub.
# ---------------------------------------------------------------------------
def _build_real_picture_slide():
    """Build a real slide with a PICTURE placeholder from python-pptx defaults."""
    from pptx import Presentation
    prs = Presentation()
    # Layout 8 = "Picture with Caption" — has TITLE(1) + PICTURE(18) + BODY(2)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    return prs, slide


def _make_test_png(tmp_path, w=200, h=150):
    """Create a small test PNG using PIL."""
    from PIL import Image
    p = tmp_path / "test_img.png"
    Image.new("RGB", (w, h), (100, 150, 200)).save(str(p))
    return str(p)


def test_add_picture_fill_sizing_adds_picture_to_placeholder(tmp_path):
    """MED-1: 'fill' sizing adds a picture filling the placeholder rect."""
    from placeholder_backfill import _add_picture_into_placeholder
    prs, slide = _build_real_picture_slide()
    # Find the PICTURE placeholder (type 18).
    pic_ph = None
    for ph in slide.placeholders:
        if "PICTURE" in str(ph.placeholder_format.type):
            pic_ph = ph
            break
    assert pic_ph is not None, "Layout 8 must have a PICTURE placeholder"
    img = _make_test_png(tmp_path)
    shapes_before = len(slide.shapes)
    _add_picture_into_placeholder(slide, pic_ph, img, sizing="fill")
    shapes_after = len(slide.shapes)
    assert shapes_after == shapes_before + 1, "add_picture should add exactly 1 shape"


def test_add_picture_fit_sizing_adds_picture_preserving_aspect(tmp_path):
    """MED-1: 'fit' sizing adds a picture letterboxed within the placeholder."""
    from placeholder_backfill import _add_picture_into_placeholder
    prs, slide = _build_real_picture_slide()
    pic_ph = None
    for ph in slide.placeholders:
        if "PICTURE" in str(ph.placeholder_format.type):
            pic_ph = ph
            break
    assert pic_ph is not None
    # Use a distinctly non-square image so fit math is exercised.
    img = _make_test_png(tmp_path, w=400, h=100)
    shapes_before = len(slide.shapes)
    _add_picture_into_placeholder(slide, pic_ph, img, sizing="fit")
    shapes_after = len(slide.shapes)
    assert shapes_after == shapes_before + 1


def test_add_picture_raises_on_missing_file(tmp_path):
    """MED-1: missing image file raises FileNotFoundError (not silent skip)."""
    from placeholder_backfill import _add_picture_into_placeholder
    import pytest
    prs, slide = _build_real_picture_slide()
    pic_ph = None
    for ph in slide.placeholders:
        if "PICTURE" in str(ph.placeholder_format.type):
            pic_ph = ph
            break
    assert pic_ph is not None
    bogus = str(tmp_path / "nonexistent.png")
    with pytest.raises(FileNotFoundError, match="image not found"):
        _add_picture_into_placeholder(slide, pic_ph, bogus, sizing="fill")
