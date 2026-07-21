"""Tests for chart styling: theme colors, fonts, gridlines, axis formatting."""
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

from ppt_builder import generate_ppt_from_data

EXPECTED_COLORS = [
    RGBColor(0x44, 0x72, 0xC4),  # accent1
    RGBColor(0xED, 0x7D, 0x31),  # accent2
    RGBColor(0xFF, 0xC0, 0x00),  # accent4
]
EXPECTED_FONT = "Calibri"
EXPECTED_GRIDLINE = RGBColor(0xE7, 0xE6, 0xE6)
EXPECTED_TEXT_COLOR = RGBColor(0x44, 0x54, 0x6A)


def _get_chart(output_path):
    prs = Presentation(output_path)
    for shape in prs.slides[0].shapes:
        if shape.has_chart:
            return shape.chart
    return None


class TestChartFonts:
    def test_chart_global_font_is_calibri(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.font.name == EXPECTED_FONT

    def test_legend_font_is_calibri(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.legend.font.name == EXPECTED_FONT

    def test_data_labels_font_is_calibri(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.plots[0].data_labels.font.name == EXPECTED_FONT

    def test_value_axis_tick_font_is_calibri(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.value_axis.tick_labels.font.name == EXPECTED_FONT

    def test_category_axis_tick_font_is_calibri(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.category_axis.tick_labels.font.name == EXPECTED_FONT


class TestChartColors:
    def test_bar_series_uses_theme_color(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        series = chart.plots[0].series[0]
        assert series.format.fill.fore_color.rgb == EXPECTED_COLORS[0]

    def test_line_multi_series_colors(self, line_chart_data, template_path, output_path):
        generate_ppt_from_data([line_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        for i, series in enumerate(chart.plots[0].series):
            if i < len(EXPECTED_COLORS):
                assert series.format.fill.fore_color.rgb == EXPECTED_COLORS[i]

    def test_data_label_text_color(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        labels = chart.plots[0].data_labels
        assert labels.font.color.rgb == EXPECTED_TEXT_COLOR

    def test_legend_text_color(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.legend.font.color.rgb == EXPECTED_TEXT_COLOR


class TestChartGridlines:
    def test_value_axis_has_major_gridlines(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        assert chart.value_axis.has_major_gridlines is True

    def test_gridline_color_matches_theme(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        gl = chart.value_axis.major_gridlines
        assert gl.format.line.color.rgb == EXPECTED_GRIDLINE


class TestChartNumberFormats:
    def test_pie_data_label_percentage_format(self, pie_chart_data, template_path, output_path):
        generate_ppt_from_data([pie_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        labels = chart.plots[0].data_labels
        assert labels.number_format == "0%"

    def test_bar_data_label_value_format(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        labels = chart.plots[0].data_labels
        assert labels.number_format == "#,##0.0"

    def test_custom_value_format(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
            "chart_options": {"value_format": "0"},
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        chart = _get_chart(output_path)
        labels = chart.plots[0].data_labels
        assert labels.number_format == "0"
