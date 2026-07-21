"""Backward compatibility tests: existing slide types must work unchanged."""
import pytest
from pptx import Presentation

from ppt_builder import generate_ppt_from_data


class TestBackwardCompat:
    def test_title_slide_still_works(self, template_path, output_path):
        data = [{"slide_type": "title_slide", "title": "Hello", "subtitle": "World"}]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Hello"

    def test_content_slide_still_works(self, template_path, output_path):
        data = [{
            "slide_type": "content_slide",
            "title": "Overview",
            "body": "**Point A** - desc A\n**Point B** - desc B",
        }]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Overview"

    def test_closing_slide_still_works(self, template_path, output_path):
        data = [{"slide_type": "closing_slide", "title": "Thanks", "subtitle": "Q&A"}]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Thanks"

    def test_section_header_still_works(self, template_path, output_path):
        data = [{"slide_type": "section_header_slide", "title": "Section 1"}]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 1

    def test_mixed_deck_without_charts(self, template_path, output_path):
        data = [
            {"slide_type": "title_slide", "title": "Deck", "subtitle": "2026"},
            {"slide_type": "content_slide", "title": "Page 1", "body": "**A** - x"},
            {"slide_type": "content_slide", "title": "Page 2", "body": "**B** - y"},
            {"slide_type": "closing_slide", "title": "End", "subtitle": "Bye"},
        ]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 4

    def test_mixed_deck_with_charts(self, mixed_deck, template_path, output_path):
        generate_ppt_from_data(mixed_deck, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 6

        chart_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_count += 1
        assert chart_count == 3

    def test_notes_preserved_on_all_types(self, template_path, output_path):
        data = [
            {"slide_type": "title_slide", "title": "T", "notes": "Title notes"},
            {"slide_type": "content_slide", "title": "C", "body": "**X** - y", "notes": "Content notes"},
            {
                "slide_type": "chart_slide",
                "title": "Chart",
                "chart_type": "bar",
                "categories": ["A"],
                "series": [{"name": "S", "values": [1]}],
                "notes": "Chart notes",
            },
        ]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert "Title notes" in prs.slides[0].notes_slide.notes_text_frame.text
        assert "Content notes" in prs.slides[1].notes_slide.notes_text_frame.text
        assert "Chart notes" in prs.slides[2].notes_slide.notes_text_frame.text

    def test_unknown_slide_type_skipped(self, template_path, output_path):
        data = [
            {"slide_type": "title_slide", "title": "Keep"},
            {"slide_type": "totally_unknown", "title": "Skip"},
            {"slide_type": "content_slide", "title": "Also Keep", "body": "**X** - y"},
        ]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path, default_closing=False)
        prs = Presentation(output_path)
        assert len(prs.slides) == 2
