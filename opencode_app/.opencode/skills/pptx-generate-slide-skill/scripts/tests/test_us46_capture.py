"""Tests for US-4.6 Phase 1 extraction additions (bullets + image rId + rules).

Layouts cannot receive shapes via the python-pptx high-level API, so the
capture helpers are unit-tested directly on slide-level shapes (the same
pattern the existing ``_extract_text_fonts`` tests use), plus the
``_build_component`` wiring is asserted. Covers:
  - text_properties.bullets capture: bulleted paragraph -> non-empty bullets
    with type/char/font; plain paragraph -> empty list (AC4 prerequisite).
  - image_properties capture: a picture -> partname (/ppt/media/...) + content
    type (C1); the partname resolves to real bytes in the package.
  - _build_component wiring: text components get text_properties.bullets;
    image components get image_properties.
  - validator cross-field rules: text_properties only on text components;
    image_properties only on image components; invalid bullet type -> warning.
"""
import zipfile

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from schema_extractor import (
    _IdCounter,
    _build_component,
    _extract_image_properties,
    _extract_paragraph_format,
    validate_template_schema,
)


def _slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _bullet_textbox(with_bullet=True, with_autonum=False):
    """A slide-level textbox with 2 paragraphs; the first optionally bulleted."""
    slide = _slide()
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.add_run().text = "First point"
    p2 = tf.add_paragraph()
    p2.add_run().text = "Second point"

    if with_bullet:
        pPr = p1._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        kind = "a:buAutoNum" if with_autonum else "a:buChar"
        attr = {"type": "arabicPeriod"} if with_autonum else {"char": "\u2022"}
        pPr.append(pPr.makeelement(qn(kind), attr))
    return tb, slide


def _picture_shape(tmp_path):
    from PIL import Image
    img_path = tmp_path / "pic.png"
    Image.new("RGB", (200, 150), (10, 20, 30)).save(str(img_path))

    slide = _slide()
    pic = slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(4), Inches(3))
    return pic, slide


# ---------------------------------------------------------------------------
# Bullet capture (AC4 prerequisite)
# ---------------------------------------------------------------------------
class TestBulletCapture:
    def test_bulleted_paragraph_captured(self):
        tb, _ = _bullet_textbox(with_bullet=True)
        bullets = _extract_paragraph_format(tb.text_frame)
        assert len(bullets) == 1
        b = bullets[0]
        assert b["type"] == "char"
        assert b["char"] == "\u2022"
        assert b["font"] == "Arial"

    def test_autonum_paragraph_captured(self):
        tb, _ = _bullet_textbox(with_bullet=True, with_autonum=True)
        bullets = _extract_paragraph_format(tb.text_frame)
        assert len(bullets) == 1
        assert bullets[0]["type"] == "autonum"

    def test_plain_paragraph_yields_empty(self):
        tb, _ = _bullet_textbox(with_bullet=False)
        assert _extract_paragraph_format(tb.text_frame) == []


# ---------------------------------------------------------------------------
# Image properties capture (C1)
# ---------------------------------------------------------------------------
class TestImageProperties:
    def test_picture_has_partname_and_content_type(self, tmp_path):
        pic, _ = _picture_shape(tmp_path)
        props = _extract_image_properties(pic)
        assert props["partname"].startswith("/ppt/media/")
        assert props["content_type"] == "image/png"
        assert props["width_px"] == 200 and props["height_px"] == 150

    def test_partname_resolves_to_real_bytes(self, tmp_path):
        pic, slide = _picture_shape(tmp_path)
        props = _extract_image_properties(pic)
        # The picture's part is reachable via the slide's package.
        partname = props["partname"]
        package = slide.part.package
        parts = [str(p.partname) for p in package.iter_parts()]
        assert partname in parts

    def test_non_picture_yields_empty(self):
        tb, _ = _bullet_textbox()
        # A textbox is not a picture -> no resolvable image -> empty props.
        assert _extract_image_properties(tb) == {}


# ---------------------------------------------------------------------------
# _build_component wiring
# ---------------------------------------------------------------------------
class TestBuildComponentWiring:
    def test_text_component_gets_text_properties(self):
        tb, slide = _bullet_textbox(with_bullet=True)
        comp = _build_component(tb, 9144000, 6858000, 0, _IdCounter())
        assert comp["type"] == "textbox"
        assert "text_properties" in comp
        assert len(comp["text_properties"]["bullets"]) == 1

    def test_plain_text_component_empty_bullets(self):
        tb, _ = _bullet_textbox(with_bullet=False)
        comp = _build_component(tb, 9144000, 6858000, 0, _IdCounter())
        assert comp["text_properties"]["bullets"] == []

    def test_image_component_gets_image_properties(self, tmp_path):
        pic, _ = _picture_shape(tmp_path)
        comp = _build_component(pic, 9144000, 6858000, 0, _IdCounter())
        assert comp["type"] == "image"
        assert comp["image_properties"]["partname"].startswith("/ppt/media/")


# ---------------------------------------------------------------------------
# Validator cross-field rules (m3 + C1 mirror)
# ---------------------------------------------------------------------------
def _minimal_text_comp(ctype="textbox", **extra):
    base = {
        "id": "comp_001", "type": ctype, "name": "n",
        "polygon": [{"x": 0.1, "y": 0.1}, {"x": 0.9, "y": 0.1},
                    {"x": 0.9, "y": 0.9}, {"x": 0.1, "y": 0.9}],
        "z_order": 0,
    }
    base.update(extra)
    return base


def _wrap(components):
    return {
        "template_metadata": {
            "title": "T", "schema_version": "1.1.0",
            "generated_by": "test", "generated_at": "2026-01-01T00:00:00Z",
            "slide_dimensions": {
                "width_emu": 9144000, "height_emu": 6858000,
                "width_inches": 10.0, "height_inches": 7.5, "aspect_ratio": "4:3",
            },
        },
        "slide_master": {"name": "m", "components": []},
        "slide_layouts": [{"layout_id": "l", "layout_name": "L", "layout_index": 0,
                           "components": components}],
        "component_type_enum": ["textbox"], "placeholder_type_enum": [None],
    }


class TestValidatorCrossFieldRules:
    def test_non_text_with_text_properties_is_error(self):
        comp = _minimal_text_comp(ctype="image", text_properties={"bullets": []})
        result = validate_template_schema(_wrap([comp]))
        assert any("text_properties" in e.reason for e in result.errors)

    def test_non_image_with_image_properties_is_error(self):
        comp = _minimal_text_comp(ctype="textbox", image_properties={"partname": "/x"})
        result = validate_template_schema(_wrap([comp]))
        assert any("image_properties" in e.reason for e in result.errors)

    def test_invalid_bullet_type_is_warning(self):
        comp = _minimal_text_comp(
            text_properties={"bullets": [{"type": "diamond"}]}
        )
        result = validate_template_schema(_wrap([comp]))
        assert any("bullets[0].type" in w.reason for w in result.warnings)

    def test_text_component_with_bullets_validates(self):
        comp = _minimal_text_comp(
            text_properties={"bullets": [{"type": "char", "char": "\u2022", "level": 1}]}
        )
        result = validate_template_schema(_wrap([comp]))
        assert result.is_valid
