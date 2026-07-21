"""Tests for US-4.3 — auto-chain extraction (templated output).

Verifies that every generated ``.pptx`` carries ``ppt/template_schema.json``
after save (AC2), that the schema is sourced from the **input template**
(arch-review M1 — not the rendered deck's cover), that a stale embedded input
schema is not laundered into the output (M2), and that the ``templating`` field
appears in the render report. Non-templated / templated / opt-out / failure
paths are covered.
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from ppt_builder import _ensure_output_templated, _embedded_schema_stale, generate_ppt_from_data
from schema_extractor import embed_schema, extract_schema, read_embedded_schema

_TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


def _strip_embedded(template_path, dst):
    """A non-templated copy: python-pptx save drops the embedded part."""
    Presentation(template_path).save(str(dst))
    return str(dst)


def _stale_fixture(template_path, dst):
    """A templated copy whose embedded schema has a WRONG layout count (M2)."""
    schema = extract_schema(template_path)
    schema["slide_layouts"] = schema["slide_layouts"][:5]  # truncate -> mismatch
    embed_schema(template_path, schema, str(dst))
    return str(dst)


def _identity_fixture(template_path, dst, identity_title):
    """A non-templated copy carrying a slide[0] with a distinct title (M1)."""
    prs = Presentation(template_path)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        if ph.placeholder_format and ph.placeholder_format.type in _TITLE_TYPES:
            ph.text_frame.paragraphs[0].text = identity_title
            break
    prs.save(str(dst))
    return str(dst)


def _render_report(output_path):
    rp = Path(output_path).with_name(Path(output_path).stem + ".render.json")
    return json.loads(rp.read_text(encoding="utf-8"))


# ============================================================
# _embedded_schema_stale (unit, M2)
# ============================================================
class TestEmbeddedStale:
    def test_not_stale_when_counts_match(self, template_path):
        n = len(extract_schema(template_path).get("slide_layouts", []))
        assert _embedded_schema_stale(template_path, n) is False

    def test_stale_when_counts_diverge(self, template_path):
        assert _embedded_schema_stale(template_path, 999) is True


# ============================================================
# Templated output (AC2) — end-to-end via generate_ppt_from_data
# ============================================================
class TestTemplatedOutput:
    def test_non_templated_input_yields_templated_output(self, template_path, tmp_path):
        tpl = _strip_embedded(template_path, tmp_path / "plain.pptx")
        assert read_embedded_schema(tpl) is None  # precondition: not templated
        out = str(tmp_path / "out.pptx")
        generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "T", "notes": "n"}],
            template_path=tpl, output_path=out,
        )
        assert read_embedded_schema(out) is not None  # AC2
        rep = _render_report(out)["templating"]
        assert rep["input_template_embedded"] is False
        assert rep["output_templated"] is True
        assert rep["schema_source"] == "extracted_input"

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_templated_input_copies_input_schema(self, template_path, tmp_path):
        out = str(tmp_path / "out.pptx")
        generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "T", "notes": "n"}],
            template_path=template_path, output_path=out,
        )
        assert read_embedded_schema(out) is not None
        rep = _render_report(out)["templating"]
        assert rep["input_template_embedded"] is True
        assert rep["output_templated"] is True
        assert rep["schema_source"] == "copied_input"

    def test_stale_input_is_not_laundered(self, template_path, tmp_path):
        tpl = _stale_fixture(template_path, tmp_path / "stale.pptx")
        out = str(tmp_path / "out.pptx")
        generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "T", "notes": "n"}],
            template_path=tpl, output_path=out,
        )
        rep = _render_report(out)["templating"]
        # M2: stale embedded schema -> fresh extraction, not a copied stale schema
        assert rep["input_template_embedded"] is True
        assert rep["schema_source"] == "extracted_input"
        # the output's embedded schema reflects the REAL layout count, not the
        # truncated (5) stale one
        schema = read_embedded_schema(out)
        assert len(schema["slide_layouts"]) > 5

    def test_schema_describes_template_not_deck(self, template_path, tmp_path):
        # M1: an input whose slide[0] is "TEMPLATE_ID"; the deck's cover is
        # "DECK_COVER". The output's embedded title must be the template's, not
        # the rendered deck's cover (which is what extracting from the OUTPUT
        # would leak).
        tpl = _identity_fixture(template_path, tmp_path / "identity.pptx", "TEMPLATE_ID")
        out = str(tmp_path / "out.pptx")
        generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "DECK_COVER", "notes": "n"}],
            template_path=tpl, output_path=out,
        )
        schema = read_embedded_schema(out)
        title = schema["template_metadata"]["title"]
        assert title != "DECK_COVER", "schema must describe the template, not the deck (M1)"

    def test_auto_template_disabled_leaves_output_plain(self, template_path, tmp_path):
        out = str(tmp_path / "out.pptx")
        generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "T", "notes": "n"}],
            template_path=template_path, output_path=out, auto_template=False,
        )
        assert read_embedded_schema(out) is None  # opt-out preserved
        rep = _render_report(out)["templating"]
        assert rep["output_templated"] is False
        assert rep["schema_source"] == "disabled"


# ============================================================
# Non-fatal failure path (unit)
# ============================================================
class TestEnsureTemplatedFailure:
    def test_bogus_template_does_not_raise(self, tmp_path):
        out = str(tmp_path / "out.pptx")
        Path(out).write_bytes(b"fake")  # output exists but template is bogus
        result = _ensure_output_templated(out, "/nonexistent/template.pptx", True)
        assert result["schema_source"] == "failed"
        assert result["output_templated"] is False
        assert "failed" in result["message"]
