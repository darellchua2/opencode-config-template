"""GIT-93 (slide_type decoupling) — layout_name targeting tests.

Covers the 7-phase plan (PLAN-GIT-93 Rev 2). Phase-by-phase accumulation.
The shared contract layer (``layout_contract``) is importable because
``conftest.py`` puts ``_common/scripts`` on ``sys.path``.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _layout(name, idx, fp, area=0.0):
    return {"name": name, "index": idx, "fingerprint": list(fp),
            "content_area_in2": area}


def _contract(layouts):
    return {"layouts": layouts, "slide_size": {"ratio": "16:9"}}


# ===========================================================================
# Phase 1 (#94) — available_layouts() + classify_layout_fingerprint
# ===========================================================================
class TestAvailableLayouts:
    def test_returns_all_layouts(self):
        from layout_contract import available_layouts
        layouts = [_layout(f"L{i}", i, ["TITLE"]) for i in range(35)]
        result = available_layouts(_contract(layouts))
        assert len(result) == 35
        assert result[0]["name"] == "L0"
        assert result[34]["index"] == 34

    def test_entry_has_four_keys(self):
        from layout_contract import available_layouts
        result = available_layouts(_contract([_layout("Agenda", 5, ["TITLE", "OBJECT"], 42.3)]))
        assert set(result[0].keys()) == {"name", "index", "fingerprint", "content_area_in2"}
        assert result[0]["fingerprint"] == ["TITLE", "OBJECT"]
        assert result[0]["content_area_in2"] == 42.3

    def test_empty_or_none_contract(self):
        from layout_contract import available_layouts
        assert available_layouts({}) == []
        assert available_layouts(None) == []

    def test_does_not_filter_to_servable_types(self):
        # A layout whose fingerprint matches NONE of the 8 standard types is
        # still listed (the whole point of GIT-93).
        from layout_contract import available_layouts
        result = available_layouts(_contract([_layout("Weird", 0, ["CHART", "TABLE", "MEDIA"])]))
        assert len(result) == 1
        assert result[0]["name"] == "Weird"

    def test_servable_slide_types_still_reports_eight(self):
        # T1.3 backward-compat gate: servable_slide_types unchanged.
        from layout_contract import servable_slide_types, _SLIDE_TYPE_FINGERPRINT
        report = servable_slide_types(_contract([_layout("Content", 0, ["TITLE", "OBJECT"])]))
        assert set(report.keys()) == set(_SLIDE_TYPE_FINGERPRINT.keys())
        assert len(report) == 8


class TestClassifyLayoutFingerprint:
    def test_title_object_maps_to_content_slide(self):
        from layout_contract import classify_layout_fingerprint
        assert classify_layout_fingerprint(["TITLE", "OBJECT"]) == "content_slide"

    def test_title_only_maps_to_chart_slide(self):
        from layout_contract import classify_layout_fingerprint
        # chart_slide's ideal fingerprint is [TITLE] (the only zero-missing match)
        assert classify_layout_fingerprint(["TITLE"]) == "chart_slide"

    def test_title_subtitle_maps_to_a_subtitle_type(self):
        from layout_contract import classify_layout_fingerprint
        # title_slide / closing_slide / section_header_slide all ideal [TITLE, SUBTITLE]
        assert classify_layout_fingerprint(["TITLE", "SUBTITLE"]) in (
            "title_slide", "closing_slide", "section_header_slide")

    def test_empty_returns_none(self):
        from layout_contract import classify_layout_fingerprint
        assert classify_layout_fingerprint([]) is None

    def test_exotic_layout_picks_nearest(self):
        from layout_contract import classify_layout_fingerprint
        # [TITLE, OBJECT, OBJECT] → two_content_slide / comparison_slide (missing 0)
        result = classify_layout_fingerprint(["TITLE", "OBJECT", "OBJECT"])
        assert result in ("two_content_slide", "comparison_slide")


# ===========================================================================
# Phase 2 (#95) — layout_name first-class + gate relaxation + AC6 relax
# ===========================================================================
_REPO_ROOT = Path(__file__).resolve().parents[5]
_DEFAULT_PPTX = _REPO_ROOT / "template" / "default.pptx"
# Fallback: BT-142 intentionally removed the bundled default.pptx from this
# repo. When it's absent, materialise python-pptx's default template — it
# ships the standard layout names these tests target ("Title Slide",
# "Title and Content", "Two Content") — into a temp path so the suite stays
# self-contained and the @_needs_default tests actually run instead of skipping.
if not _DEFAULT_PPTX.exists():
    import tempfile
    from pptx import Presentation as _Presentation
    _DEFAULT_PPTX = Path(tempfile.gettempdir()) / "git93_test_default.pptx"
    if not _DEFAULT_PPTX.exists():
        _Presentation().save(str(_DEFAULT_PPTX))
_needs_default = pytest.mark.skipif(
    not _DEFAULT_PPTX.exists(), reason="template/default.pptx not present"
)


@_needs_default
class TestLayoutNameTargeting:
    def test_unknown_slide_type_with_layout_name_renders(self, tmp_path):
        from ppt_builder import generate_ppt_from_data
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "agenda", "layout_name": "Title Slide",
                 "title": "Agenda Title", "notes": "n"}]
        result = generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(result)
        assert len(prs.slides) == 1  # NOT skipped (layout_name rescues it)

    def test_unknown_slide_type_without_layout_name_skipped(self, tmp_path):
        from ppt_builder import generate_ppt_from_data
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "agenda", "title": "X", "notes": "n"}]
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(out)
        assert len(prs.slides) == 0  # skipped (unknown type, no layout_name)

    def test_layout_name_overrides_fingerprint(self, tmp_path):
        from ppt_builder import generate_ppt_from_data
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        # Same slide_type, different layout_name → each uses its named layout
        data = [
            {"slide_type": "content_slide", "layout_name": "Title Slide",
             "title": "A", "notes": "n"},
            {"slide_type": "content_slide", "layout_name": "Title and Content",
             "title": "B", "notes": "n"},
        ]
        result = generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(result)
        names = [s.slide_layout.name for s in prs.slides]
        assert names == ["Title Slide", "Title and Content"]

    def test_pseudo_type_bug_fixed(self, tmp_path):
        # Pre-GIT-93: pseudo-type + config-pin was silently dropped (gate
        # rejected it before the config-pin branch). Now it renders.
        from ppt_builder import generate_ppt_from_data
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "_custom_1", "title": "Pseudo", "notes": "n"}]
        overrides = {"_custom_1_layout": "Title Slide"}
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out,
            config_overrides=overrides)
        prs = Presentation(out)
        assert len(prs.slides) == 1
        assert prs.slides[0].slide_layout.name == "Title Slide"


@_needs_default
class TestAC6Relaxation:
    def test_zero_servable_without_layout_name_raises(self, tmp_path):
        # A contract serving none of the 8 types + no layout_name → fatal.
        from ppt_builder import _validate_template
        from errors import TemplateError
        from pptx import Presentation
        prs = Presentation(str(_DEFAULT_PPTX))
        contract = {"layouts": [{"name": "Weird", "index": 0,
                                 "fingerprint": ["MEDIA", "TABLE"]}]}
        with pytest.raises(TemplateError):
            _validate_template(prs, "fake.pptx", contract,
                               [{"slide_type": "content_slide"}])

    def test_zero_servable_with_layout_name_warns_not_raises(self, tmp_path):
        # Same contract + >=1 slide carrying layout_name → warning, not fatal.
        from ppt_builder import _validate_template
        from pptx import Presentation
        prs = Presentation(str(_DEFAULT_PPTX))
        contract = {"layouts": [{"name": "Weird", "index": 0,
                                 "fingerprint": ["MEDIA", "TABLE"]}]}
        # Must NOT raise.
        _validate_template(prs, "fake.pptx", contract,
                           [{"slide_type": "x", "layout_name": "Weird"}])


# ===========================================================================
# Phase 3 (#98) — field-driven fill dispatch (MAJOR-1 sweep preserved)
# ===========================================================================
@_needs_default
class TestFieldDrivenFill:
    def test_field_driven_fill_unknown_type_body(self, tmp_path):
        # Unknown slide_type + body field → fills the layout's OBJECT/BODY placeholder
        from ppt_builder import generate_ppt_from_data, _OBJECT_TYPE, _BODY_TYPE
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "agenda", "layout_name": "Title and Content",
                 "title": "Agenda", "body": "**1.** Intro\n**2.** Demo", "notes": "n"}]
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(out)
        slide = prs.slides[0]
        body_texts = [
            p.text_frame.text for p in slide.placeholders
            if p.placeholder_format.type in (_OBJECT_TYPE, _BODY_TYPE)
        ]
        assert any("Intro" in t for t in body_texts), (
            f"body not filled; placeholder texts: {body_texts}")

    def test_subtitle_removed_when_empty(self, tmp_path):
        # MAJOR-1: a title_slide with NO subtitle must REMOVE the subtitle
        # placeholder (no master-inherited sample-text bleed).
        from ppt_builder import generate_ppt_from_data, _SUBTITLE_TYPE
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "title_slide", "title": "T", "notes": "n"}]
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(out)
        slide = prs.slides[0]
        subtitle_phs = [p for p in slide.placeholders
                        if p.placeholder_format.type == _SUBTITLE_TYPE]
        assert len(subtitle_phs) == 0  # removed, not left as inherited sample

    def test_subtitle_filled_when_present(self, tmp_path):
        # Counter-check: a title_slide WITH subtitle keeps (fills) it.
        from ppt_builder import generate_ppt_from_data, _SUBTITLE_TYPE
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "title_slide", "title": "T",
                 "subtitle": "My Subtitle", "notes": "n"}]
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(out)
        slide = prs.slides[0]
        subtitle_phs = [p for p in slide.placeholders
                        if p.placeholder_format.type == _SUBTITLE_TYPE]
        assert len(subtitle_phs) == 1
        assert "My Subtitle" in subtitle_phs[0].text_frame.text


# ===========================================================================
# Phase 4 (#96) — generic per-field schema for layout_name path
# ===========================================================================
class TestGenericSchema:
    def test_layout_name_slide_clean_passes(self):
        from schema_validator import validate_slide_data_list
        data = [{"slide_type": "agenda", "layout_name": "Agenda",
                 "title": "T", "body": "b", "notes": "n"}]
        res = validate_slide_data_list(data)
        assert res.is_valid

    def test_chart_missing_series_reports_error(self):
        # chart_type present but no categories/series → hard error (T4.3)
        from schema_validator import validate_slide_data_list
        data = [{"slide_type": "custom_chart", "layout_name": "Chart",
                 "title": "T", "chart_type": "bar", "notes": "n"}]
        res = validate_slide_data_list(data)
        assert not res.is_valid
        msgs = res.error_messages()
        assert any("series" in m for m in msgs)
        assert any("categories" in m for m in msgs)

    def test_missing_title_warns(self):
        # MAJOR-3: title/notes recommended (warn-if-absent) for layout_name slides
        from schema_validator import validate_slide_data_list
        data = [{"slide_type": "team", "layout_name": "Team",
                 "body": "b"}]  # no title, no notes
        res = validate_slide_data_list(data)
        # Warnings don't block is_valid, but must be present
        warns = res.warning_messages()
        assert any("title" in w for w in warns)
        assert any("notes" in w for w in warns)

    def test_unknown_field_warns(self):
        from schema_validator import validate_slide_data_list
        data = [{"slide_type": "x", "layout_name": "L",
                 "title": "T", "notes": "n", "bogus_field": 123}]
        res = validate_slide_data_list(data)
        warns = res.warning_messages()
        assert any("bogus_field" in w for w in warns)

    def test_all_field_specs_covers_union(self):
        from schemas import ALL_FIELD_SPECS
        # title/body/subtitle/body_left/body_right/chart_type present
        for f in ("title", "body", "subtitle", "body_left", "body_right",
                  "chart_type", "categories", "series", "notes"):
            assert f in ALL_FIELD_SPECS
        # extended fields
        for f in ("layout_name", "body_slots", "image_paths", "title_slots"):
            assert f in ALL_FIELD_SPECS


# ===========================================================================
# Phase 6 (#99) — overflow geometry by layout_name + density regression
# ===========================================================================
class TestOverflowGeometryByLayoutName:
    def test_unknown_type_layout_name_resolves_geometry(self):
        # NF-1/NF-2: geometry resolved by layout_name → layouts[i], reads "type"
        from overflow_check import _available_height_for_field
        contract = {"layouts": [
            {"name": "Agenda", "index": 0, "fingerprint": ["TITLE", "OBJECT"],
             "placeholders": [{"type": "OBJECT", "height_in": 2.0}]}
        ]}
        slide = {"slide_type": "agenda", "layout_name": "Agenda", "body": "x"}
        avail = _available_height_for_field("body", slide, contract)
        assert avail == 2.0  # real geometry, not None

    def test_unknown_type_no_layout_name_returns_none(self):
        from overflow_check import _available_height_for_field
        contract = {"layouts": [
            {"name": "Content", "index": 0, "fingerprint": ["TITLE", "OBJECT"],
             "placeholders": [{"type": "OBJECT", "height_in": 3.0}]}
        ]}
        # Unknown type, no layout_name, and Content's fingerprint won't match
        # an unknown slide_type → None (defer to image oracle).
        slide = {"slide_type": "totally_unknown_xyz", "body": "x"}
        avail = _available_height_for_field("body", slide, contract)
        assert avail is None

    def test_known_type_fingerprint_fallback(self):
        # Known slide_type (no layout_name) → fingerprint match resolves geometry
        from overflow_check import _available_height_for_field
        contract = {"layouts": [
            {"name": "Content", "index": 0, "fingerprint": ["TITLE", "OBJECT"],
             "placeholders": [{"type": "OBJECT", "height_in": 4.5}]}
        ]}
        slide = {"slide_type": "content_slide", "body": "x"}
        avail = _available_height_for_field("body", slide, contract)
        assert avail == 4.5


class TestDensityCountsUnknownType:
    def test_density_counts_unknown_type_body(self):
        # MINOR-1: count_slide_words already counts by field unconditionally —
        # lock the existing correct behavior for unknown types (regression guard).
        from density_mode import count_slide_words
        count = count_slide_words({"slide_type": "agenda", "body": "one two three four"})
        assert count == 4


# ===========================================================================
# Code-review follow-up tests (MINOR/NIT closure)
# ===========================================================================
class TestCodeReviewFollowups:
    def test_body_and_two_body_do_not_double_fill(self, tmp_path):
        # MINOR-1: a slide carrying BOTH body and body_left/body_right must not
        # run both fill paths (would overwrite). two-body wins; body is skipped.
        from ppt_builder import generate_ppt_from_data, _OBJECT_TYPE, _BODY_TYPE
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "content_slide", "layout_name": "Two Content",
                 "title": "T",
                 "body": "SHOULD NOT APPEAR",
                 "body_left": "left content", "body_right": "right content",
                 "notes": "n"}]
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out)
        prs = Presentation(out)
        texts = [p.text_frame.text for p in prs.slides[0].placeholders]
        joined = " ".join(texts)
        assert "SHOULD NOT APPEAR" not in joined  # body skipped (two-body wins)
        assert "left content" in joined
        assert "right content" in joined

    def test_overflow_layout_name_case_insensitive(self):
        # MINOR-2: layout_name differing by case still resolves geometry.
        from overflow_check import _available_height_for_field
        contract = {"layouts": [
            {"name": "Title and Content", "index": 0,
             "fingerprint": ["TITLE", "OBJECT"],
             "placeholders": [{"type": "OBJECT", "height_in": 3.0}]}
        ]}
        slide = {"slide_type": "x", "layout_name": "title and content", "body": "b"}
        assert _available_height_for_field("body", slide, contract) == 3.0

    def test_overflow_layout_name_prefix_normalized(self):
        # MINOR-2: layout_name with a numeric prefix (e.g. "1_two content")
        # matches contract name "Two Content" via normalization.
        from overflow_check import _available_height_for_field
        contract = {"layouts": [
            {"name": "1_Two Content", "index": 0,
             "fingerprint": ["TITLE", "OBJECT", "OBJECT"],
             "placeholders": [{"type": "OBJECT", "height_in": 2.5}]}
        ]}
        slide = {"slide_type": "x", "layout_name": "two content", "body": "b"}
        assert _available_height_for_field("body", slide, contract) == 2.5

    def test_layout_name_beats_config_pin(self, tmp_path):
        # NIT-1 policy: per-slide layout_name wins over deck-wide config-pin.
        from ppt_builder import generate_ppt_from_data
        from pptx import Presentation
        out = str(tmp_path / "out.pptx")
        data = [{"slide_type": "content_slide",
                 "layout_name": "Title Slide",   # should win
                 "title": "T", "notes": "n"}]
        # config-pin points the OTHER way
        overrides = {"content_slide_layout": "Title and Content"}
        generate_ppt_from_data(
            data, template_path=str(_DEFAULT_PPTX), output_path=out,
            config_overrides=overrides)
        prs = Presentation(out)
        assert prs.slides[0].slide_layout.name == "Title Slide"  # layout_name won

    def test_classify_fingerprint_tiebreak_pinned(self):
        # NIT-3: lock the deterministic tie-break winner for [TITLE, SUBTITLE].
        from layout_contract import classify_layout_fingerprint
        # Three ideals match with missing=0; dict order → title_slide wins.
        assert classify_layout_fingerprint(["TITLE", "SUBTITLE"]) == "title_slide"

    def test_strict_mode_promotes_layout_name_missing_notes(self):
        # MINOR-4: strict mode promotes notes-warning to error for layout_name slides too.
        from schema_validator import validate_slide_data_list
        data = [{"slide_type": "team", "layout_name": "Team", "title": "T"}]
        res = validate_slide_data_list(data, strict=True)
        assert not res.is_valid  # missing notes promoted to error
        assert any("notes" in m for m in res.error_messages())

