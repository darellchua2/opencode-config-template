"""Integration tests for US-4.2 text-fitting wired into the render path.

Exercises the fit-aware ``_set_text`` / ``_set_body_text`` (M1 base-size chain,
M3 inheritance-when-not-shrunk, auto-grow guard) and the ``<output>.render.json``
render report (AC3).

Template notes:
* The bundled template's title/subtitle placeholders are **auto-grow** (their
  reported base height < one line at the base size), so they never shrink — the
  box expands instead. The M3 inheritance test asserts this path.
* Shrinking is therefore demonstrated on **synthetic fixed-height textboxes**
  (controlled geometry), which is exactly the fixed-box case US-4.2 targets.
* The render-report path is tested end-to-end via ``title_slide`` (resolves) +
  the report-entry mapper unit-tested for the adjusted case.
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from ppt_builder import (
    _font_fit_report_entry,
    _set_body_text,
    _set_text,
    generate_ppt_from_data,
)
from text_fit import ROLE_BASE_PT, FontFit


def _blank_slide(prs):
    """A slide off the BLANK layout (or the emptiest layout), for textboxes."""
    blank = None
    for layout in prs.slide_layouts:
        if not list(layout.placeholders):
            blank = layout
            break
    return prs.slides.add_slide(blank or prs.slide_layouts[0])


def _textbox(slide, width_in, height_in):
    return slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(width_in), Inches(height_in)
    )


# ============================================================
# _set_body_text — template-derived sizing + shrink on a FIXED box
# ============================================================
class TestBodyTextFitting:
    def test_long_body_shrinks_on_fixed_box(self, template_path):
        prs = Presentation(template_path)
        slide = _blank_slide(prs)
        box = _textbox(slide, 6.0, 0.5)  # fixed short height → forces shrink
        long_body = "".join(
            ("**Point %d** - " + "word " * 20 + "\n") % i for i in range(5)
        )
        fit = _set_body_text(box, long_body, role="body", schema_font_map={}, layout=None)
        assert fit is not None
        assert fit.base_source == "role_ceiling"  # textbox → no schema/layout tier
        assert fit.adjusted is True  # long text in a short fixed box → shrunk
        assert fit.applied_size_pt < ROLE_BASE_PT["body"]
        # body always writes explicit template-derived run sizes
        sizes = [
            r.font.size.pt for p in box.text_frame.paragraphs
            for r in p.runs if r.font.size is not None
        ]
        assert sizes and max(sizes) <= ROLE_BASE_PT["body"] + 1e-6  # AC2

    def test_short_body_fits_at_base(self, template_path):
        prs = Presentation(template_path)
        slide = _blank_slide(prs)
        box = _textbox(slide, 8.0, 4.0)  # generous fixed box
        fit = _set_body_text(
            box, "**Short** - one concise line", role="body", schema_font_map={}, layout=None
        )
        assert fit is not None
        assert fit.adjusted is False
        assert fit.fits is True
        assert fit.applied_size_pt == pytest.approx(fit.base_size_pt)

    def test_returns_none_without_text_frame(self):
        assert _set_body_text(None, "x", role="body") is None


# ============================================================
# _set_text — M3: explicit size only on actual shrink
# ============================================================
class TestSingleTextFitting:
    def test_auto_grow_title_inherits_without_shrink(self, template_path):
        # The bundled template's title placeholder auto-grows (reported height
        # < one line at base) → estimator treats it as unbounded → no shrink,
        # no explicit size (M3 inheritance).
        prs = Presentation(template_path)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_ph = next(
            ph for ph in slide.placeholders
            if ph.placeholder_format and ph.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
            )
        )
        fit = _set_text(
            title_ph, "Short Title", role="title", schema_font_map={}, layout=prs.slide_layouts[0]
        )
        assert fit is not None
        assert fit.adjusted is False
        runs = [r for p in title_ph.text_frame.paragraphs for r in p.runs]
        assert runs
        assert all(r.font.size is None for r in runs), (
            "auto-grow title must inherit (no explicit size) per M3"
        )

    def test_fixed_box_long_text_shrinks_and_writes_explicit_size(self, template_path):
        prs = Presentation(template_path)
        slide = _blank_slide(prs)
        # 1.0in tall is a *fixed* box for the title role (one line @28pt ≈
        # 0.47in < 0.8in usable) → not treated as auto-grow → long text shrinks.
        box = _textbox(slide, 4.0, 1.0)
        long_text = "word " * 60
        fit = _set_text(box, long_text, role="title", schema_font_map={}, layout=None)
        assert fit is not None
        assert fit.adjusted or not fit.fits
        sizes = [
            r.font.size.pt for p in box.text_frame.paragraphs
            for r in p.runs if r.font.size is not None
        ]
        assert sizes, "shrunk text must write an explicit size"


# ============================================================
# Render report sidecar (AC3) — end-to-end + entry mapper
# ============================================================
class TestRenderReport:
    def test_report_written_with_valid_schema(self, template_path, output_path):
        data = [{
            "slide_type": "title_slide",
            "title": "Title",
            "subtitle": "Subtitle text",
            "notes": "n",
        }]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)

        report_path = Path(output_path).with_name(
            Path(output_path).stem + ".render.json"
        )
        assert report_path.exists(), "render report sidecar must be written (AC3)"
        report = json.loads(report_path.read_text(encoding="utf-8"))

        assert isinstance(report.get("slides"), list)
        assert len(report["slides"]) >= 1
        entry = report["slides"][0]
        assert {"index", "slide_type", "placeholders"} <= set(entry.keys())
        assert entry["placeholders"], "title fill should be recorded"
        ph = entry["placeholders"][0]
        for key in ("role", "field", "template_size_pt", "applied_size_pt",
                    "base_source", "font_size_adjusted", "fits", "lines_estimated"):
            assert key in ph, f"render report missing key: {key}"
        assert isinstance(ph["font_size_adjusted"], bool)

    def test_report_reflects_auto_grow_no_adjustment(self, template_path, output_path):
        # On this template the title auto-grows → the report honestly records
        # font_size_adjusted=False, fits=True.
        data = [{"slide_type": "title_slide", "title": "Some Title", "notes": "n"}]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        report = json.loads(
            Path(output_path).with_name(
                Path(output_path).stem + ".render.json"
            ).read_text(encoding="utf-8")
        )
        title_entries = [
            p for s in report["slides"] for p in s["placeholders"] if p["role"] == "title"
        ]
        assert title_entries
        assert title_entries[0]["font_size_adjusted"] is False

    def test_font_fit_report_entry_marks_adjustment(self):
        fit = FontFit(
            applied_size_pt=10.0, base_size_pt=14.0, base_source="role_ceiling",
            adjusted=True, fits=True, lines_estimated=8,
        )
        entry = _font_fit_report_entry("body", "body", fit)
        assert entry is not None
        assert entry["font_size_adjusted"] is True
        assert entry["applied_size_pt"] == 10.0
        assert entry["template_size_pt"] == 14.0
        assert entry["base_source"] == "role_ceiling"
        assert _font_fit_report_entry("body", "body", None) is None

    def test_report_never_blocks_render(self, template_path, output_path):
        out = generate_ppt_from_data(
            [{"slide_type": "title_slide", "title": "Hi", "subtitle": "Sub", "notes": "n"}],
            template_path=template_path, output_path=output_path,
        )
        assert Path(out).exists()
        assert Path(output_path).with_name(
            Path(output_path).stem + ".render.json"
        ).exists()


# ============================================================
# _build_schema_font_map (M1 tier 1) — pure unit
# ============================================================
class TestSchemaFontMap:
    def test_builds_role_to_size_first_wins(self):
        from ppt_builder import _build_schema_font_map
        schema = {"slide_layouts": [{
            "layout_name": "Content",
            "components": [
                {"type": "placeholder", "placeholder_type": "title", "font": {"size_pt": 32}},
                {"type": "placeholder", "placeholder_type": "body", "font": {"size_pt": 16}},
                {"type": "placeholder", "placeholder_type": "picture", "font": {"size_pt": 99}},
                {"type": "placeholder", "placeholder_type": "body", "font": {"size_pt": 20}},
                {"type": "shape", "font": {"size_pt": 5}},
                {"type": "placeholder", "placeholder_type": "subtitle"},
            ],
        }]}
        m = _build_schema_font_map(schema)
        # only text roles (title/subtitle/body); picture excluded; first body wins
        assert m == {"Content": {"title": 32.0, "body": 16.0}}

    def test_empty_on_missing_schema(self):
        from ppt_builder import _build_schema_font_map
        assert _build_schema_font_map(None) == {}
        assert _build_schema_font_map({}) == {}


# ============================================================
# _layout_line_spacing — Major regression: exact-point spacing
# ============================================================
class TestLayoutLineSpacing:
    def _slide_and_layout_ph(self, template_path):
        from ppt_builder import _matching_layout_placeholder
        prs = Presentation(template_path)
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        ph = next(p for p in slide.placeholders if getattr(p, "is_placeholder", False))
        layout_ph = _matching_layout_placeholder(ph, layout)
        assert layout_ph is not None, "layout[0] must have a matchable placeholder"
        return ph, layout, layout_ph

    def test_exact_point_spacing_falls_back_to_default(self, template_path):
        # Regression for the Major bug: an exact-point (Length/int-subclass)
        # line_spacing must NOT be misused as a 228600x multiplier.
        from pptx.util import Pt
        from ppt_builder import _layout_line_spacing
        from text_fit import LINE_SPACING_DEFAULT
        ph, layout, layout_ph = self._slide_and_layout_ph(template_path)
        layout_ph.text_frame.paragraphs[0].line_spacing = Pt(18)
        ls = _layout_line_spacing(ph, layout)
        assert ls == LINE_SPACING_DEFAULT, (
            f"exact-point spacing must fall back to {LINE_SPACING_DEFAULT}, got {ls}"
        )

    def test_multiple_spacing_is_honoured(self, template_path):
        from ppt_builder import _layout_line_spacing
        ph, layout, layout_ph = self._slide_and_layout_ph(template_path)
        layout_ph.text_frame.paragraphs[0].line_spacing = 1.5
        assert _layout_line_spacing(ph, layout) == 1.5

    def test_inherited_none_falls_back_to_default(self, template_path):
        from ppt_builder import _layout_line_spacing
        from text_fit import LINE_SPACING_DEFAULT
        ph, layout, layout_ph = self._slide_and_layout_ph(template_path)
        layout_ph.text_frame.paragraphs[0].line_spacing = None
        assert _layout_line_spacing(ph, layout) == LINE_SPACING_DEFAULT


# ============================================================
# Two-content body fit — both bodies recorded in the report
# ============================================================
class TestTwoContentFit:
    def test_both_bodies_recorded(self, template_path, output_path):
        data = [{
            "slide_type": "two_content_slide",
            "title": "Compare",
            "body_left": "**Left** - " + "word " * 40,
            "body_right": "**Right** - " + "word " * 40,
            "notes": "n",
        }]
        generate_ppt_from_data(
            data, template_path=template_path, output_path=output_path
        )
        report = json.loads(
            Path(output_path).with_name(
                Path(output_path).stem + ".render.json"
            ).read_text(encoding="utf-8")
        )
        body_entries = [
            p for s in report["slides"] for p in s["placeholders"] if p["role"] == "body"
        ]
        fields = {p["field"] for p in body_entries}
        assert {"body_left", "body_right"} <= fields, (
            f"both two-content bodies must be recorded, got {fields}"
        )
