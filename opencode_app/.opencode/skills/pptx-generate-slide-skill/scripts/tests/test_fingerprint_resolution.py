"""Tests for fingerprint-based layout resolution (issue #44, Phase 1).

Covers the #44 acceptance areas:
  * fingerprint match — exact composition + closest (surplus placeholders).
  * override precedence — config pin > fingerprint > name fallback.
  * ambiguity tie-break — name affinity, then content_area_in2.
  * degradation — unsatisfiable composition ⇒ (None, reason).
  * integration — a renamed template (no matching names) still resolves every
    slide_type via composition alone.
"""
import pytest

# PLAN-GIT-72: the pure contract symbols now live in layout_contract (_common);
# _select_layout stays in ppt_builder.
from layout_contract import (
    _SLIDE_TYPE_FINGERPRINT,
    _composition_diff,
    _name_affinity,
    _resolve_layout_by_fingerprint,
)
from ppt_builder import _select_layout
from pptx import Presentation
from template_introspector import introspect


# ============================================================
# _composition_diff
# ============================================================
class TestCompositionDiff:
    def test_exact_match(self):
        assert _composition_diff(["TITLE", "OBJECT"], ["TITLE", "OBJECT"]) == (0, 0)

    def test_layout_has_extra(self):
        # surplus placeholders cost "extra" but not "missing".
        assert _composition_diff(["TITLE", "OBJECT"], ["TITLE", "OBJECT", "OBJECT"]) == (0, 1)

    def test_missing_type(self):
        assert _composition_diff(["TITLE", "SUBTITLE"], ["TITLE", "OBJECT"]) == (1, 1)

    def test_object_serves_picture(self):
        # an OBJECT placeholder can serve a PICTURE ideal.
        assert _composition_diff(["TITLE", "PICTURE"], ["TITLE", "OBJECT"]) == (0, 0)

    def test_object_serves_chart(self):
        assert _composition_diff(["TITLE", "CHART"], ["TITLE", "OBJECT"]) == (0, 0)

    def test_picture_cannot_serve_object(self):
        # a PICTURE placeholder cannot serve a body (OBJECT) ideal.
        assert _composition_diff(["TITLE", "OBJECT"], ["TITLE", "PICTURE"]) == (1, 1)

    def test_repeated_object_need(self):
        # two_content needs two OBJECTs.
        assert _composition_diff(["TITLE", "OBJECT", "OBJECT"], ["TITLE", "OBJECT"]) == (1, 0)
        assert _composition_diff(["TITLE", "OBJECT", "OBJECT"], ["TITLE", "OBJECT", "OBJECT", "OBJECT"]) == (0, 1)


# ============================================================
# _name_affinity
# ============================================================
class TestNameAffinity:
    def test_exact_match(self):
        assert _name_affinity("Title Slide", ["Title Slide"]) == 2

    def test_case_insensitive(self):
        assert _name_affinity("title slide", ["Title Slide"]) == 2

    def test_normalized_match(self):
        # leading digits + underscore are stripped by the normalizer.
        assert _name_affinity("7_Two Content", ["7_Two Content"]) == 2
        assert _name_affinity("1_Title and Content", ["Title and Content"]) == 1

    def test_no_match(self):
        assert _name_affinity("End", ["Title Slide"]) == 0


# ============================================================
# _resolve_layout_by_fingerprint
# ============================================================
class TestResolveByFingerprint:
    def test_exact_composition_match(self):
        contract = {"layouts": [
            {"index": 0, "name": "Title and Content",
             "fingerprint": ["TITLE", "OBJECT"], "content_area_in2": 40},
        ]}
        idx, reason = _resolve_layout_by_fingerprint("content_slide", contract)
        assert idx == 0 and reason is None

    def test_closest_match_with_surplus(self):
        # ideal [TITLE, OBJECT]; layout has an extra body — still satisfiable.
        contract = {"layouts": [
            {"index": 3, "name": "Content Rich",
             "fingerprint": ["TITLE", "OBJECT", "OBJECT"], "content_area_in2": 40},
        ]}
        idx, _ = _resolve_layout_by_fingerprint("content_slide", contract)
        assert idx == 3

    def test_object_placeholder_serves_picture_ideal(self):
        contract = {"layouts": [
            {"index": 2, "name": "Pic", "fingerprint": ["TITLE", "OBJECT"], "content_area_in2": 10},
        ]}
        idx, _ = _resolve_layout_by_fingerprint("content_image_slide", contract)
        assert idx == 2

    def test_name_affinity_disambiguates_identical_composition(self):
        # Title Slide vs End share [TITLE, SUBTITLE]; name affinity decides.
        contract = {"layouts": [
            {"index": 0, "name": "Title Slide", "fingerprint": ["TITLE", "SUBTITLE"], "content_area_in2": 0},
            {"index": 5, "name": "End", "fingerprint": ["TITLE", "SUBTITLE"], "content_area_in2": 0},
        ]}
        assert _resolve_layout_by_fingerprint("title_slide", contract)[0] == 0
        assert _resolve_layout_by_fingerprint("closing_slide", contract)[0] == 5

    def test_content_area_tiebreak_when_affinity_equal(self):
        # two layouts, same composition, neither name matches → larger area wins.
        contract = {"layouts": [
            {"index": 0, "name": "L-A", "fingerprint": ["TITLE", "OBJECT"], "content_area_in2": 30},
            {"index": 1, "name": "L-B", "fingerprint": ["TITLE", "OBJECT"], "content_area_in2": 50},
        ]}
        assert _resolve_layout_by_fingerprint("content_slide", contract)[0] == 1

    def test_degradation_when_unsatisfiable(self):
        # no layout has an OBJECT/body for a content_slide.
        contract = {"layouts": [
            {"index": 0, "name": "Title Only", "fingerprint": ["TITLE"], "content_area_in2": 0},
        ]}
        idx, reason = _resolve_layout_by_fingerprint("content_slide", contract)
        assert idx is None
        assert "no layout satisfies" in reason

    def test_degradation_unknown_slide_type(self):
        idx, reason = _resolve_layout_by_fingerprint("no_such_type", {"layouts": []})
        assert idx is None
        assert "no fingerprint" in reason

    def test_empty_contract(self):
        idx, reason = _resolve_layout_by_fingerprint("content_slide", {"layouts": []})
        assert idx is None
        assert "no layouts" in reason


# ============================================================
# _select_layout — override precedence + integration
# ============================================================
class TestSelectLayoutPrecedence:
    def _setup(self, template_path):
        prs = Presentation(template_path)
        contract = introspect(template_path)
        exact, norm = {}, {}
        for layout in prs.slide_layouts:
            exact[layout.name.lower()] = layout
        return prs, contract, exact, norm

    def test_config_pin_overrides_fingerprint(self, template_path):
        prs, contract, exact, norm = self._setup(template_path)
        # Pin content_slide to a *different* content-bearing layout than the
        # fingerprint-preferred one. Chosen dynamically so the test does not
        # couple to any particular template's layout names.
        default_idx, _ = _resolve_layout_by_fingerprint("content_slide", contract)
        default_name = prs.slide_layouts[default_idx].name
        pinned_name = next(
            L["name"] for L in contract["layouts"]
            if "OBJECT" in (L.get("fingerprint") or []) and L["name"] != default_name
        )
        config = {"content_slide_layout": pinned_name}
        layout = _select_layout("content_slide", contract, config, prs, exact, norm, 1)
        assert layout is not None
        assert layout.name == pinned_name
        assert pinned_name != default_name  # the pin overrode the fingerprint default

    def test_unknown_slide_type_skipped(self, template_path):
        prs, contract, exact, norm = self._setup(template_path)
        layout = _select_layout("totally_unknown", contract, {}, prs, exact, norm, 1)
        assert layout is None

    def test_all_eight_types_resolve_on_real_template(self, template_path):
        """The current template: every slide_type resolves to a concrete layout."""
        prs, contract, exact, norm = self._setup(template_path)
        for slide_type in _SLIDE_TYPE_FINGERPRINT:
            layout = _select_layout(slide_type, contract, {}, prs, exact, norm, 1)
            assert layout is not None, f"{slide_type} failed to resolve"


class TestRenamedTemplateIntegration:
    def test_renamed_layouts_still_resolve(self, template_path):
        """Strip every layout NAME — fingerprint must still resolve each type
        by composition alone (the core Capability A guarantee)."""
        prs = Presentation(template_path)
        contract = introspect(template_path)
        # Scrub all names so name-affinity can never help.
        scrubbed = {
            "layouts": [
                {**L, "name": f"RENAMED-{L['index']}"}
                for L in contract["layouts"]
            ]
        }
        exact = {f"renamed-{L['index']}": layout
                 for L, layout in zip(contract["layouts"], prs.slide_layouts)}
        norm = exact
        resolved = 0
        for slide_type in _SLIDE_TYPE_FINGERPRINT:
            layout = _select_layout(slide_type, scrubbed, {}, prs, exact, norm, 1)
            if layout is not None:
                resolved += 1
        # Every type that the template can compositionally serve resolves.
        assert resolved == len(_SLIDE_TYPE_FINGERPRINT)

    def test_no_contract_falls_back_to_name_matching(self, template_path):
        """Without a contract the path is the original name-based matching."""
        prs = Presentation(template_path)
        exact, norm = {}, {}
        for layout in prs.slide_layouts:
            exact[layout.name.lower()] = layout
            from layout_contract import _normalize_layout_name
            norm.setdefault(_normalize_layout_name(layout.name), layout)
        # chart_slide maps to "Blank", which the bundled template spells "BLANK"
        # (matched case-insensitively by the name resolver).
        layout = _select_layout("chart_slide", None, {}, prs, exact, norm, 1)
        assert layout is not None
        assert layout.name.upper() == "BLANK"


# ============================================================
# servable_slide_types — the content-layer capability report (#45)
# ============================================================
class TestServableSlideTypes:
    def test_real_template_serves_all_eight(self, template_path):
        from layout_contract import servable_slide_types
        contract = introspect(template_path)
        report = servable_slide_types(contract)
        assert set(report) == set(_SLIDE_TYPE_FINGERPRINT)
        # The bundled template provides a layout for every engine slide_type.
        for slide_type, info in report.items():
            assert info["available"], f"{slide_type} not servable: {info.get('reason')}"
            assert info["layout"]

    def test_missing_type_reported_unavailable(self):
        from layout_contract import servable_slide_types
        # A title-only contract cannot serve content_slide (no OBJECT).
        contract = {"layouts": [
            {"index": 0, "name": "Title", "fingerprint": ["TITLE", "SUBTITLE"], "content_area_in2": 0},
        ]}
        report = servable_slide_types(contract)
        assert report["title_slide"]["available"] is True
        assert report["content_slide"]["available"] is False
        assert "no layout satisfies" in report["content_slide"]["reason"]

