"""Tests for ppt_builder.get_render_contract (US-4.1) — source swap + M4 provenance.

Covers the three resolution paths:
  - non-templated PPTX (no embedded JSON) -> silent sidecar fallback.
  - templated PPTX (valid embedded JSON)  -> embedded (adapter).
  - corrupt embedded JSON                 -> sidecar + warning.
Plus end-to-end render from a templated PPTX (AC3: uses layouts, embedded source).
"""
import shutil
import zipfile

import pytest
from pptx import Presentation

# PLAN-GIT-72: get_render_contract now lives in layout_contract (_common);
# generate_ppt_from_data stays in ppt_builder.
from layout_contract import get_render_contract
from ppt_builder import generate_ppt_from_data
from schema_extractor import (
    extract_schema,
    embed_schema,
    _EMBEDDED_SCHEMA_PATH,
)


# ---------------------------------------------------------------------------
# Fixtures: templated / corrupt / fresh (non-templated) PPTXs
# ---------------------------------------------------------------------------
def _templated(tmp_path, template_path, name="templated.pptx"):
    """A copy of the template WITH the schema embedded (a templated PPTX)."""
    out = str(tmp_path / name)
    embed_schema(template_path, extract_schema(template_path), out)
    return out


def _corrupt_embedded(tmp_path, template_path, name="corrupt.pptx"):
    """A fresh PPTX with GARBAGE at the embedded-schema path (no prior embed, so
    no duplicate-entry warning vs the now-pre-templated bundled template)."""
    out = _fresh_deck(tmp_path, name)
    with zipfile.ZipFile(out, "a") as z:
        z.writestr(_EMBEDDED_SCHEMA_PATH, b"not-json{{")
    return out


def _fresh_deck(tmp_path, name="fresh.pptx"):
    """A minimal PPTX with NO embedded JSON (guaranteed non-templated, regardless
    of whether the bundled template itself gets templated in Phase 4)."""
    out = str(tmp_path / name)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# Provenance / source paths (architecture review M4)
# ---------------------------------------------------------------------------
class TestRenderContractSource:
    def test_non_templated_silent_sidecar(self, tmp_path):
        t = _fresh_deck(tmp_path)
        c = get_render_contract(t)
        assert c["_source"] == "sidecar"
        assert c.get("layouts") is not None  # sidecar introspection succeeded

    def test_templated_uses_embedded(self, tmp_path, template_path):
        t = _templated(tmp_path, template_path)
        c = get_render_contract(t)
        assert c["_source"] == "embedded"
        assert c["layouts"], "embedded contract should have layouts"

    def test_corrupt_embedded_warns_and_falls_back(self, tmp_path, template_path, caplog):
        t = _corrupt_embedded(tmp_path, template_path)
        with caplog.at_level("WARNING"):
            c = get_render_contract(t)
        assert c["_source"] == "sidecar"
        # read_embedded_schema warns on the malformed JSON before returning None.
        assert any(
            "malformed" in r.message or "unreadable" in r.message or "sidecar" in r.message
            for r in caplog.records
        ), "expected a warning naming the fallback / malformed payload"

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_stale_embedded_warns_on_layout_count_mismatch(self, tmp_path, template_path, caplog):
        """M5 staleness guard: warn when the embedded schema's layout count != the
        live template's (catches edit-without-re-embed). The guard is warn-only —
        the embedded contract is still used (no fallback)."""
        # A fresh deck (~11 default layouts) carrying the BUNDLED schema (63 layouts)
        # -> the live layout count != the embedded count -> staleness warning.
        out = _fresh_deck(tmp_path, "stale.pptx")
        embed_schema(out, extract_schema(template_path), out)
        with caplog.at_level("WARNING"):
            c = get_render_contract(out)
        assert c["_source"] == "embedded"  # guard is warn-only, not a fallback
        assert any("Stale embedded schema" in r.message for r in caplog.records), (
            "expected a 'Stale embedded schema' warning on layout-count mismatch"
        )


# ---------------------------------------------------------------------------
# End-to-end: render from a templated PPTX (AC1/AC3 — embedded consumed, layouts used)
# ---------------------------------------------------------------------------
class TestRenderFromTemplated:
    def test_render_from_templated_produces_valid_deck(self, tmp_path, template_path):
        templated = _templated(tmp_path, template_path)
        # Uses slide types the bundled template serves (title + closing). NB:
        # content_slide is intentionally avoided — the bundled template lacks a
        # ['TITLE','OBJECT'] layout (a pre-existing template-coverage limitation,
        # unrelated to US-4.1; fails identically via the sidecar path).
        slide_data = [
            {"slide_type": "title_slide", "title": "Templated Render",
             "subtitle": "US-4.1", "notes": "KEY MESSAGE: embedded contract drives this."},
            {"slide_type": "closing_slide", "title": "Thank You",
             "notes": "KEY MESSAGE: sign-off."},
        ]
        out = str(tmp_path / "rendered.pptx")
        result = generate_ppt_from_data(slide_data, template_path=templated, output_path=out)
        # Output exists and is a valid PPTX with the expected slide count.
        prs = Presentation(result)
        assert len(prs.slides) == 2

    def test_render_from_non_templated_still_works(self, tmp_path, template_path):
        # Backward-compat: rendering against the bundled template still produces a
        # valid deck. NB: since US-4.1 the bundled template ships pre-templated, so
        # this exercises the EMBEDDED path; the sidecar fallback is covered by
        # test_non_templated_silent_sidecar (fresh, genuinely non-templated deck).
        slide_data = [
            {"slide_type": "title_slide", "title": "Bundled Render", "notes": "KEY MESSAGE: embedded path."},
            {"slide_type": "closing_slide", "title": "Thank You", "notes": "KEY MESSAGE: done."},
        ]
        out = str(tmp_path / "rendered_bundled.pptx")
        result = generate_ppt_from_data(slide_data, template_path=template_path, output_path=out)
        prs = Presentation(result)
        assert len(prs.slides) == 2
