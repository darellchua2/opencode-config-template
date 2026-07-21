"""Tests for the US-4.6 AC5 target-size gate + coordinate-path render (Phases 0+3).

The bundled template is 16:9 (read dynamically). The gate is RATIO-based (m1):
same ratio -> native path (no-op); different ratio -> coordinate path (resize +
geometry scaling). These tests assert the REAL behaviour now that Phase 3 lands
the coordinate executor (no longer a NotImplementedError stub).
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_builder import generate_ppt_from_data
from geometry import compute_ratio, resolve_target_size

_TITLE_DECK = [
    {"slide_type": "title_slide", "title": "T", "subtitle": "S", "notes": "n"},
    {"slide_type": "content_slide", "title": "C", "body": "- a\n- b", "notes": "n"},
]


def _native_size(template_path):
    prs = Presentation(template_path)
    return int(prs.slide_width), int(prs.slide_height)


def _ratio(w, h):
    return compute_ratio(w, h)


class TestTargetSizeGate:
    def test_none_is_native_noop(self, template_path, output_path):
        out = generate_ppt_from_data(_TITLE_DECK, template_path=template_path, output_path=output_path)
        nat_w, nat_h = _native_size(template_path)
        prs = Presentation(out)
        assert (int(prs.slide_width), int(prs.slide_height)) == (nat_w, nat_h)

    def test_same_ratio_is_native_noop(self, template_path, output_path):
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="16:9",
        )
        nat_w, nat_h = _native_size(template_path)
        prs = Presentation(out)
        assert (int(prs.slide_width), int(prs.slide_height)) == (nat_w, nat_h)

    def test_different_ratio_renders_at_target_size(self, template_path, output_path):
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="4:3",
        )
        prs = Presentation(out)
        tw, th = resolve_target_size("4:3")
        assert (int(prs.slide_width), int(prs.slide_height)) == (tw, th)
        assert _ratio(int(prs.slide_width), int(prs.slide_height)) == "4:3"
        assert len(prs.slides) == 2  # content preserved

    def test_square_ratio_renders_at_target_size(self, template_path, output_path):
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="1:1",
        )
        prs = Presentation(out)
        assert _ratio(int(prs.slide_width), int(prs.slide_height)) == "1:1"

    def test_explicit_same_ratio_dict_is_native(self, template_path, output_path):
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path,
            target_size={"width_in": 13.333, "height_in": 7.5},
        )
        nat_w, nat_h = _native_size(template_path)
        prs = Presentation(out)
        assert (int(prs.slide_width), int(prs.slide_height)) == (nat_w, nat_h)

    def test_invalid_target_size_raises_valueerror(self, template_path, output_path):
        with pytest.raises(ValueError, match="Invalid target_size"):
            generate_ppt_from_data(
                _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="bogus",
            )


# ---------------------------------------------------------------------------
# Coordinate-path output: render.json provenance (aspect_ratio) + M4 schema
# ---------------------------------------------------------------------------
class TestCoordinateReport:
    def test_render_json_records_aspect_ratio(self, template_path, output_path):
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="4:3",
        )
        report = json.loads(
            Path(output_path).with_name(Path(output_path).stem + ".render.json")
            .read_text(encoding="utf-8")
        )
        assert "aspect_ratio" in report
        ar = report["aspect_ratio"]
        tw, th = resolve_target_size("4:3")
        assert ar["target"] == [tw, th]
        assert "scale" in ar and "shapes_scaled" in ar

    def test_native_path_has_no_aspect_ratio_field(self, template_path, output_path):
        generate_ppt_from_data(_TITLE_DECK, template_path=template_path, output_path=output_path)
        report = json.loads(
            Path(output_path).with_name(Path(output_path).stem + ".render.json")
            .read_text(encoding="utf-8")
        )
        assert "aspect_ratio" not in report


class TestOutputSchemaRewriteM4:
    def test_output_embedded_schema_has_target_dims(self, template_path, output_path):
        from schema_extractor import read_embedded_schema
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="4:3",
        )
        schema = read_embedded_schema(out)
        assert schema is not None
        dims = schema["template_metadata"]["slide_dimensions"]
        tw, th = resolve_target_size("4:3")
        assert (dims["width_emu"], dims["height_emu"]) == (tw, th)
        assert dims["aspect_ratio"] == "4:3"


# ---------------------------------------------------------------------------
# AC3 integration: a placeholder re-normalizes to the same relative position
# after a 16:9 -> 4:3 render (proportional scaling end-to-end).
# ---------------------------------------------------------------------------
class TestProportionalScalingAC3:
    def test_title_placeholder_normalizes_to_same_relative_position(self, template_path, output_path):
        from pptx.enum.shapes import PP_PLACEHOLDER

        _TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)

        def _find_title(placeholders):
            for ph in placeholders:
                if ph.placeholder_format.type in _TITLE_TYPES:
                    return ph
            return None

        # Native (template) title placeholder normalized position (layout 0).
        nat_w, nat_h = _native_size(template_path)
        nat_prs = Presentation(template_path)
        nat_title = _find_title(nat_prs.slide_layouts[0].placeholders)
        assert nat_title is not None, "template title layout must have a title placeholder"
        nat_norm_x = (nat_title.left or 0) / nat_w
        nat_norm_y = (nat_title.top or 0) / nat_h

        # Render at 4:3 and read the output's title placeholder position.
        out = generate_ppt_from_data(
            _TITLE_DECK, template_path=template_path, output_path=output_path, target_size="4:3",
        )
        out_prs = Presentation(out)
        tw, th = resolve_target_size("4:3")
        out_title = _find_title(out_prs.slides[0].placeholders)
        assert out_title is not None
        out_norm_x = (out_title.left or 0) / tw
        out_norm_y = (out_title.top or 0) / th

        # Within 1% (AC3).
        assert abs(nat_norm_x - out_norm_x) <= 0.01
        assert abs(nat_norm_y - out_norm_y) <= 0.01


# ---------------------------------------------------------------------------
# Arch-review C1 regression: group shapes must scale LINEARLY (sx), not sx^2.
# Scaling only the group container (off/ext) leaves OOXML's child->slide map
# (chOff/chExt fixed) to place children at exactly sx of their original
# slide-space position. Verified empirically against the bundled engine.
# ---------------------------------------------------------------------------
class TestGroupScalingC1:
    @staticmethod
    def _child_slide_pos(grp, child):
        _A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        xfrm = grp._element.grpSpPr.find(_A + "xfrm")
        ch_off = xfrm.find(_A + "chOff")
        ch_ext = xfrm.find(_A + "chExt")
        cox, coy = int(ch_off.get("x")), int(ch_off.get("y"))
        cex, cey = int(ch_ext.get("cx")), int(ch_ext.get("cy"))
        px = grp.left + (child.left - cox) * grp.width / cex if cex else grp.left
        py = grp.top + (child.top - coy) * grp.height / cey if cey else grp.top
        return px, py

    def test_group_child_scales_linearly_not_quadratically(self, tmp_path):
        from pptx import Presentation
        from ppt_builder import _scale_shapes_geometry

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        grp = slide.shapes.add_group_shape()
        child = grp.shapes.add_textbox(914400, 914400, 1828800, 914400)  # child-space (1,1,2,1)in
        child.text_frame.text = "x"
        grp.left, grp.top = 914400, 914400         # off (1,1)in
        grp.width, grp.height = 3657600, 1828800   # ext (4,2)in

        before_x, before_y = self._child_slide_pos(grp, child)

        _scale_shapes_geometry(slide.shapes, 2.0, 2.0)

        after_x, after_y = self._child_slide_pos(grp, child)
        # AC2 / C1: child slide-space position must scale by exactly sx=2 (linear),
        # NOT sx^2=4 (the old double-scale bug).
        assert abs(after_x - 2 * before_x) <= 9144   # within 0.01in
        assert abs(after_y - 2 * before_y) <= 9144
        # Explicitly reject the old quadratic behaviour.
        assert abs(after_x - 4 * before_x) > 9144


