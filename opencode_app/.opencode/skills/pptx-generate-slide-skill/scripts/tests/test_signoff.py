"""Tests for the closing/title subtitle sign-off + inherited-sample suppression.

Covers:
* closing_slide presenter sign-off (presenter_name / presenter_email →
  "Prepared by: {name}\\n{email}"); partial variants; explicit-subtitle fallback.
* the no-content case REMOVES the subtitle placeholder so the layout's inherited
  sample text ("Prepared by: Lecturer Name" on closing, "Click to edit Master
  subtitle style" on title) cannot bleed in (an empty <a:p/> still inherits —
  verified empirically; removal is the reliable suppression).
"""
import pytest
from pptx import Presentation

from ppt_builder import generate_ppt_from_data


def _render(data, template_path, output_path):
    return generate_ppt_from_data(
        data, template_path=template_path, output_path=output_path,
        default_closing=False, cleanup_temp=False,
    )


def _subtitle_text(slide):
    """Text of the idx=1 (subtitle) placeholder, or None if removed."""
    for ph in slide.placeholders:
        if ph.placeholder_format and ph.placeholder_format.idx == 1:
            return ph.text_frame.text
    return None


def _has_marker(slide, markers):
    low = [m.lower() for m in markers]
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.lower()
            if any(m in t for m in low):
                return True
    return False


# ============================================================
# closing_slide — presenter sign-off
# ============================================================
class TestClosingSignoff:
    def test_name_and_email(self, template_path, output_path):
        out = _render([{
            "slide_type": "closing_slide", "title": "Thank You",
            "presenter_name": "Jane Doe", "presenter_email": "jane@x.com",
            "notes": "n",
        }], template_path, output_path)
        sub = _subtitle_text(Presentation(out).slides[0])
        # python-pptx reads an <a:br/> line break back as "\v" (0x0b); normalize.
        assert sub.replace("\v", "\n") == "Prepared by: Jane Doe\njane@x.com"

    def test_name_only(self, template_path, output_path):
        out = _render([{
            "slide_type": "closing_slide", "title": "Thank You",
            "presenter_name": "Jane Doe", "notes": "n",
        }], template_path, output_path)
        assert _subtitle_text(Presentation(out).slides[0]) == "Prepared by: Jane Doe"

    def test_email_only(self, template_path, output_path):
        out = _render([{
            "slide_type": "closing_slide", "title": "Thank You",
            "presenter_email": "jane@x.com", "notes": "n",
        }], template_path, output_path)
        assert _subtitle_text(Presentation(out).slides[0]) == "jane@x.com"

    def test_explicit_subtitle_fallback(self, template_path, output_path):
        # no presenter fields -> falls back to an explicit subtitle
        out = _render([{
            "slide_type": "closing_slide", "title": "Thank You",
            "subtitle": "Custom Sign-off", "notes": "n",
        }], template_path, output_path)
        assert _subtitle_text(Presentation(out).slides[0]) == "Custom Sign-off"

    def test_no_presenter_removes_placeholder(self, template_path, output_path):
        # neither presenter fields nor subtitle -> placeholder REMOVED, no bleed
        out = _render([{
            "slide_type": "closing_slide", "title": "Thank You", "notes": "n",
        }], template_path, output_path)
        slide = Presentation(out).slides[0]
        assert _subtitle_text(slide) is None  # placeholder removed
        assert not _has_marker(slide, ["Prepared by", "Lecturer", "Email address"])


# ============================================================
# title_slide — clear inherited "Click to edit Master subtitle style"
# ============================================================
class TestTitleSubtitle:
    def test_with_subtitle_filled(self, template_path, output_path):
        out = _render([{
            "slide_type": "title_slide", "title": "T", "subtitle": "My Subtitle",
            "notes": "n",
        }], template_path, output_path)
        assert _subtitle_text(Presentation(out).slides[0]) == "My Subtitle"

    def test_without_subtitle_removes_placeholder(self, template_path, output_path):
        out = _render([{
            "slide_type": "title_slide", "title": "T", "notes": "n",
        }], template_path, output_path)
        slide = Presentation(out).slides[0]
        assert _subtitle_text(slide) is None  # removed
        assert not _has_marker(slide, ["Click to edit", "Master subtitle"])
