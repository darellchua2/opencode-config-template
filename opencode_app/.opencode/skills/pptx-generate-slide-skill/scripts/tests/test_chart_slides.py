"""Tests for chart_slide creation, schema validation, and edge cases."""
import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from ppt_builder import generate_ppt_from_data


# ============================================================
# Chart Type Creation Tests
# ============================================================

class TestBarChartCreation:
    def test_bar_chart_creates_column_clustered(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    def test_bar_chart_has_title_in_placeholder(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert prs.slides[0].shapes.title.text == "Market Growth"

    def test_bar_chart_has_notes(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        assert "Revenue growing steadily" in prs.slides[0].notes_slide.notes_text_frame.text

    def test_bar_chart_has_legend(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.has_legend is True

    def test_bar_chart_has_data_labels(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.plots[0].has_data_labels is True

    def test_bar_chart_y_axis_scale(self, bar_chart_data, template_path, output_path):
        generate_ppt_from_data([bar_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.value_axis.minimum_scale == 0
        assert chart.value_axis.maximum_scale == 25


class TestPieChartCreation:
    def test_pie_chart_creates_pie_type(self, pie_chart_data, template_path, output_path):
        generate_ppt_from_data([pie_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == XL_CHART_TYPE.PIE

    def test_pie_chart_percentage_labels(self, pie_chart_data, template_path, output_path):
        generate_ppt_from_data([pie_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        labels = chart.plots[0].data_labels
        assert labels.show_percentage is True
        assert labels.show_value is False

    def test_pie_chart_legend_on_right(self, pie_chart_data, template_path, output_path):
        from pptx.enum.chart import XL_LEGEND_POSITION
        generate_ppt_from_data([pie_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.legend.position == XL_LEGEND_POSITION.RIGHT


class TestLineChartCreation:
    def test_line_chart_creates_line_markers(self, line_chart_data, template_path, output_path):
        generate_ppt_from_data([line_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == XL_CHART_TYPE.LINE_MARKERS

    def test_line_chart_multi_series(self, line_chart_data, template_path, output_path):
        generate_ppt_from_data([line_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert len(chart.plots[0].series) == 2

    def test_line_chart_categories_preserved(self, line_chart_data, template_path, output_path):
        generate_ppt_from_data([line_chart_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        cats = list(chart.plots[0].categories)
        assert cats == ["Q1", "Q2", "Q3", "Q4"]


# ============================================================
# Chart Type Variants Tests
# ============================================================

class TestChartTypeVariants:
    @pytest.mark.parametrize("chart_type,expected_xl", [
        ("bar", XL_CHART_TYPE.COLUMN_CLUSTERED),
        ("bar_stacked", XL_CHART_TYPE.COLUMN_STACKED),
        ("bar_horizontal", XL_CHART_TYPE.BAR_CLUSTERED),
        ("bar_horizontal_stacked", XL_CHART_TYPE.BAR_STACKED),
        ("pie", XL_CHART_TYPE.PIE),
        ("pie_exploded", XL_CHART_TYPE.PIE_EXPLODED),
        ("doughnut", XL_CHART_TYPE.DOUGHNUT),
        ("line", XL_CHART_TYPE.LINE),
        ("line_markers", XL_CHART_TYPE.LINE_MARKERS),
    ])
    def test_chart_type_mapping(self, chart_type, expected_xl, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": chart_type,
            "categories": ["A", "B", "C"],
            "series": [{"name": "S1", "values": [1, 2, 3]}],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == expected_xl


# ============================================================
# Schema / Edge Case Tests
# ============================================================

class TestChartSchemaEdgeCases:
    def test_invalid_chart_type_defaults_to_bar(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "nonexistent_type",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    def test_missing_chart_type_defaults_to_bar(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    def test_missing_categories_skips_chart(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "series": [{"name": "S1", "values": [1, 2]}],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        has_chart = any(s.has_chart for s in prs.slides[0].shapes)
        assert has_chart is False

    def test_missing_series_skips_chart(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        has_chart = any(s.has_chart for s in prs.slides[0].shapes)
        assert has_chart is False

    def test_empty_series_list_skips_chart(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
            "series": [],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        has_chart = any(s.has_chart for s in prs.slides[0].shapes)
        assert has_chart is False

    def test_no_chart_options_uses_defaults(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.has_legend is True
        assert chart.plots[0].has_data_labels is True

    def test_legend_none_disables_legend(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
            "chart_options": {"legend_position": "none"},
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.has_legend is False

    def test_data_labels_disabled(self, template_path, output_path):
        data = {
            "slide_type": "chart_slide",
            "title": "Test",
            "chart_type": "bar",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2]}],
            "chart_options": {"show_data_labels": False},
        }
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        chart = prs.slides[0].shapes[1].chart
        assert chart.plots[0].has_data_labels is False
