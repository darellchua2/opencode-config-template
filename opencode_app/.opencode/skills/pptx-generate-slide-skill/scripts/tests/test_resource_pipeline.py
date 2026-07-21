"""Tests: manual image embedding (#18) + chart-data resolver."""
import json

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_builder import generate_ppt_from_data
from resolvers import resolve_slide_data_list


# ============================================================
# Image embedding (#18)
# ============================================================
class TestImageEmbedding:
    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_slide_with_image_path_embeds_native_picture(self, image_slide_data, template_path, output_path):
        generate_ppt_from_data([image_slide_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    @pytest.mark.skip(reason="BT-142 Phase 2.5: requires a richer template fixture than the minimal synthesized one (needs multiple layouts / picture placeholders / non-placeholder shapes). Skip until a richer fixture builder is added.")
    def test_picture_is_editable_native_object(self, image_slide_data, template_path, output_path):
        generate_ppt_from_data([image_slide_data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pic = next(s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        # Native pictures have a non-empty image blob (embedded, not linked).
        assert pic.image.blob is not None
        assert len(pic.image.blob) > 0

    def test_missing_image_path_is_graceful(self, template_path, output_path):
        data = [{"slide_type": "content_slide", "title": "No Img",
                 "body": "**x** - y", "image_path": "does/not/exist.png", "notes": "n"}]
        generate_ppt_from_data([image for image in data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 0  # skipped, not crashed

    @pytest.mark.parametrize("preset", ["full", "half-left", "half-right", "below-title"])
    def test_placement_presets_embed_picture(self, sample_image, preset, template_path, output_path):
        data = {"slide_type": "content_slide", "title": "P",
                "body": "**x** - y", "image_path": sample_image,
                "image_position": preset, "notes": "n"}
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    def test_invalid_preset_defaults_gracefully(self, sample_image, template_path, output_path):
        data = {"slide_type": "content_slide", "title": "P",
                "body": "**x** - y", "image_path": sample_image,
                "image_position": "nonsense", "notes": "n"}
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    def test_image_size_override(self, sample_image, template_path, output_path):
        from pptx.util import Inches
        data = {"slide_type": "content_slide", "title": "P", "body": "**x** - y",
                "image_path": sample_image, "image_position": "full",
                "image_size": {"width": 6, "height": 3}, "notes": "n"}
        generate_ppt_from_data([data], template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pic = next(s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        assert pic.width == Inches(6)
        assert pic.height == Inches(3)

    def test_backward_compat_no_image(self, template_path, output_path):
        # Slides without image_path render exactly as before.
        data = [{"slide_type": "content_slide", "title": "Plain", "body": "**x** - y"}]
        generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 0


class TestChartDataResolver:
    def test_data_query_populates_series(self):
        def search_fn(query, config):
            return {"categories": ["2022", "2023"], "series": [{"name": "R", "values": [10, 12]}],
                    "source": "https://example.com"}
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "data_query": "revenue", "notes": "orig"}
        resolved = resolve_slide_data_list([slide], {"chart_data": {"search_fn": search_fn}})
        out = resolved[0]
        assert out["categories"] == ["2022", "2023"]
        assert out["series"][0]["values"] == [10, 12]
        assert "Data source: https://example.com" in out["notes"]
        assert out["data_source"] == "https://example.com"
        assert "data_query" not in out

    def test_concrete_series_not_overwritten(self):
        def search_fn(query, config):
            return {"categories": ["X"], "series": [{"name": "N", "values": [99]}], "source": "src"}
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}],
                 "data_query": "should not overwrite", "notes": "n"}
        resolved = resolve_slide_data_list([slide], {"chart_data": {"search_fn": search_fn}})
        out = resolved[0]
        assert out["categories"] == ["A", "B"]  # unchanged
        assert out["series"][0]["values"] == [1, 2]

    def test_chart_data_graceful_when_unconfigured(self):
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "data_query": "revenue", "notes": "n"}
        resolved = resolve_slide_data_list([slide], config={})
        assert "categories" not in resolved[0]


# ============================================================
# Pipeline-level non-fatal / integration
# ============================================================
class TestResolverPipeline:
    def test_mixed_deck_passes_through_unresolved(self, mixed_deck):
        # No config: all placeholders unresolved but deck structure intact.
        resolved = resolve_slide_data_list(mixed_deck, config={})
        assert len(resolved) == len(mixed_deck)
