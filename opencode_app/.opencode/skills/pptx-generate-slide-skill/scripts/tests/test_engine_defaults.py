"""Tests for engine-level default behaviours.

Covers:
  * #40 — generate_ppt_from_data auto-appends a Thank-You closing slide.
  * #37 — generate_ppt_from_data resolves resource placeholders by default.

The resolver pipeline itself is exercised in test_resource_pipeline.py; here we
only verify the engine wiring (the call happens, never crashes, opt-out works).
"""
from pptx import Presentation

from ppt_builder import generate_ppt_from_data


def _content(title: str) -> dict:
    return {"slide_type": "content_slide", "title": title, "body": "**a** - b", "notes": "n"}


# --- #40: auto-appended default closing slide -------------------------------

def test_auto_closing_appended_when_missing(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "Cover", "notes": "n"},
        _content("P1"),
        _content("P2"),
    ]
    generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
    prs = Presentation(output_path)
    assert len(prs.slides) == 4  # 3 input + 1 auto-closing
    # The auto-appended closing carries the default "Thank You" title
    # (template-agnostic: the layout name varies by template source).
    assert prs.slides[-1].shapes.title.text == "Thank You"


def test_auto_closing_skipped_when_already_closing(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "Cover", "notes": "n"},
        _content("P1"),
        {"slide_type": "closing_slide", "title": "Thanks", "notes": "n"},
    ]
    generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
    prs = Presentation(output_path)
    assert len(prs.slides) == 3  # already closed — no duplicate appended


def test_auto_closing_skipped_for_short_deck(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "Cover", "notes": "n"},
        _content("P1"),
    ]
    generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
    prs = Presentation(output_path)
    assert len(prs.slides) == 2  # N < 3 threshold respected


def test_auto_closing_disabled(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "Cover", "notes": "n"},
        _content("P1"),
        _content("P2"),
    ]
    generate_ppt_from_data(data, template_path=template_path, output_path=output_path, default_closing=False)
    prs = Presentation(output_path)
    assert len(prs.slides) == 3  # opt-out: nothing appended


# --- #37: resolve_placeholders wiring ---------------------------------------

def test_resolve_placeholders_default_renders(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "X", "notes": "n"},
        _content("P1"),
        {
            "slide_type": "content_image_slide", "title": "Img",
            "body": "**x** - y", "notes": "n",
        },
    ]
    generate_ppt_from_data(data, template_path=template_path, output_path=output_path)
    prs = Presentation(output_path)
    # default resolve_placeholders=True no-ops without config; deck still renders
    # (3 input + 1 auto-closing).
    assert len(prs.slides) == 4


def test_resolve_placeholders_disabled_renders(template_path, output_path):
    data = [
        {"slide_type": "title_slide", "title": "X", "notes": "n"},
        _content("P1"),
        _content("P2"),
    ]
    generate_ppt_from_data(
        data, template_path=template_path, output_path=output_path, resolve_placeholders=False,
    )
    prs = Presentation(output_path)
    assert len(prs.slides) == 4  # opt-out: placeholders ignored, closing still appended
