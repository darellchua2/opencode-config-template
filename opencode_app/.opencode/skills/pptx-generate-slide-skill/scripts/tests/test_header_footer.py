"""Tests for US-2.1 — header/footer detection + prompt-decision + inject helpers.

Covers:
* Detection on the bundled template (has_header=False, has_footer=True).
* Detection on synthetic masters (header-present → True; chrome-less → both
  False; both-false → needs_header_footer_prompt True) — exercises all branches
  of _detect_header_footer including the HEADER branch (code-review Major #1).
* ``needs_header_footer_prompt`` pure helper (all four boolean combinations).
* ``inject_default_header_zone`` — 4-point polygon + English note + explicit
  polygon assertions (len==4, {x,y}, [0,1] range) + schema validates.
"""

from pptx.enum.shapes import PP_PLACEHOLDER

from schema_extractor import (
    _detect_header_footer,
    extract_schema,
    inject_default_header_zone,
    needs_header_footer_prompt,
    validate_template_schema,
)


class _FakePh:
    """Minimal placeholder stand-in for _detect_header_footer tests."""

    def __init__(self, ph_type):
        self._pf = type("_PF", (), {"type": ph_type})()
        self.placeholder_format = self._pf


def _fake_prs(placeholders):
    """Build a minimal prs-like object with the given master placeholders."""
    master = type("_M", (), {"placeholders": placeholders})()
    return type("_P", (), {"slide_masters": [master]})()


class TestDetect:
    def test_bundled_has_footer_no_header(self, template_path):
        schema = extract_schema(template_path)
        hf = schema["template_metadata"]["header_footer"]
        assert hf["has_header"] is False
        assert hf["has_footer"] is True

    def test_bundled_does_not_need_prompt(self, template_path):
        schema = extract_schema(template_path)
        # has_footer=True → prompt NOT needed
        assert needs_header_footer_prompt(schema) is False

    # -- synthetic-master tests (code-review Major #1: exercise all branches) --

    def test_synthetic_header_present(self):
        """A master with a HEADER placeholder → has_header=True (untested branch)."""
        prs = _fake_prs([_FakePh(PP_PLACEHOLDER.HEADER), _FakePh(PP_PLACEHOLDER.FOOTER)])
        result = _detect_header_footer(prs)
        assert result == {"has_header": True, "has_footer": True}

    def test_synthetic_chromeless_both_absent(self):
        """A master with no chrome placeholders → both False."""
        prs = _fake_prs([])
        result = _detect_header_footer(prs)
        assert result == {"has_header": False, "has_footer": False}

    def test_synthetic_both_absent_needs_prompt(self):
        """End-to-end: chrome-less detection → needs_header_footer_prompt True."""
        prs = _fake_prs([])
        hf = _detect_header_footer(prs)
        schema = {"template_metadata": {"header_footer": hf}}
        assert needs_header_footer_prompt(schema) is True


class TestNeedsPrompt:
    def test_both_false_prompts(self):
        schema = {"template_metadata": {"header_footer": {"has_header": False, "has_footer": False}}}
        assert needs_header_footer_prompt(schema) is True

    def test_footer_present_no_prompt(self):
        schema = {"template_metadata": {"header_footer": {"has_header": False, "has_footer": True}}}
        assert needs_header_footer_prompt(schema) is False

    def test_header_present_no_prompt(self):
        schema = {"template_metadata": {"header_footer": {"has_header": True, "has_footer": False}}}
        assert needs_header_footer_prompt(schema) is False

    def test_empty_header_footer_no_prompt_when_missing(self):
        # no header_footer key at all → treat as absent → prompt
        schema = {"template_metadata": {}}
        assert needs_header_footer_prompt(schema) is True


class TestInject:
    def test_inject_creates_default_zone(self):
        schema = {"template_metadata": {"header_footer": {"has_header": False, "has_footer": False}}}
        inject_default_header_zone(schema)
        header = schema["template_metadata"]["header_footer"]["header"]

        # shape
        assert header["source"] == "user_default"
        assert isinstance(header["note"], str) and len(header["note"]) > 0

        # polygon: exactly 4 points, each {x,y}, all in [0,1]
        poly = header["polygon"]
        assert isinstance(poly, list)
        assert len(poly) == 4
        for pt in poly:
            assert set(pt.keys()) == {"x", "y"}
            assert 0 <= pt["x"] <= 1
            assert 0 <= pt["y"] <= 1

        # winding: TL→TR→BR→BL (top strip)
        assert poly[0] == {"x": 0, "y": 0}
        assert poly[1] == {"x": 1, "y": 0}
        assert poly[2]["y"] > poly[0]["y"]  # bottom edge below top

    def test_inject_idempotent(self):
        schema = {"template_metadata": {"header_footer": {}}}
        inject_default_header_zone(schema)
        first = schema["template_metadata"]["header_footer"]["header"]
        inject_default_header_zone(schema)  # second call
        second = schema["template_metadata"]["header_footer"]["header"]
        assert first is second  # setdefault — no overwrite

    def test_inject_schema_still_validates(self, template_path):
        schema = extract_schema(template_path)
        inject_default_header_zone(schema)
        result = validate_template_schema(schema)
        assert result.is_valid, result.error_messages()
