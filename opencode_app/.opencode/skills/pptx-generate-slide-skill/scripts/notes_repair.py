"""BT-142 Phase 3.5c — Notes-master repair (engine hard limit L5).

Templates promoted via Capability C (or any user-built template) may lack a
notes-text placeholder on the notes master. When that happens,
``slide.notes_slide.notes_text_frame`` returns ``None`` and speaker notes are
silently dropped by ``ppt_builder._set_notes``.

This module probes a presentation's notes master, adds a notes-text placeholder
if one is missing, and exposes a single high-level entry point
``ensure_notes_placeholder(prs)`` that ``ppt_builder`` calls once before filling
notes on any slide.

Algorithm (per Phase 3.5c spec):
  1. Probe ``prs.notes_master`` — if absent, python-pptx auto-creates one on
     first access; nothing further to do.
  2. Look for an existing notes-text placeholder (``type==BODY``, ``idx==1``).
  3. If absent, inject ``<p:ph type="body" idx="1"/>`` under the notes
     master's ``<p:cSld>`` via raw OOXML (python-pptx has no public API for
     adding placeholders to a notes master).
  4. Return a report: ``{repaired: bool, notes_master_placeholder_count: int}``.

The repair is idempotent and minimal — only the body placeholder is added; no
other notes-master shape is touched.
"""

from __future__ import annotations

import logging
from typing import Any, Dict
from xml.etree import ElementTree as ET

from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

_NOTES_TEXT_TYPE = "body"
_NOTES_TEXT_IDX = "1"

_NOTES_PLACEHOLDER_TAG = qn("p:ph")
_SP_TAG = qn("p:sp")
_TXBODY_TAG = qn("a:txBody")
_CSLD_TAG = qn("p:cSld")
_SP_SHAPE_TAG = qn("p:sp")
_NVSPPR_TAG = qn("p:nvSpPr")
_NVPICPR_TAG = qn("p:cNvPr")


def _find_notes_text_placeholder(notes_master) -> bool:
    """Return True iff the notes master already has a body-text placeholder."""
    for ph in notes_master.placeholders:
        try:
            if ph.placeholder_format.type is not None and (
                str(ph.placeholder_format.type).endswith(_NOTES_TEXT_TYPE)
                or getattr(ph.placeholder_format, "idx", None) == int(_NOTES_TEXT_IDX)
            ):
                return True
        except Exception:
            continue
    return False


def _build_notes_text_sp_xml() -> str:
    """Return OOXML for a minimal notes body placeholder shape.

    Coordinates are EMU and match the default notes-master body geometry
    (matches what PowerPoint emits for a fresh notes master). Tuned so the
    placeholder is visible but does not overlap the slide image thumbnail.
    """
    return (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:nvSpPr>'
        '<p:cNvPr id="2" name="Notes Placeholder 2"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr>'
        '</p:nvSpPr>'
        '<p:spPr>'
        '<a:xfrm>'
        '<a:off x="685800" y="4343400"/>'
        '<a:ext cx="5486400" cy="3600450"/>'
        '</a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</p:spPr>'
        '<p:txBody>'
        '<a:bodyPr vert="horz" wrap="square" anchor="t"/>'
        '<a:lstStyle/>'
        '<a:p><a:endParaRPr lang="en-US"/></a:p>'
        '</p:txBody>'
        '</p:sp>'
    )


def _inject_notes_text_placeholder(notes_master) -> None:
    """Append a notes body placeholder shape to ``notes_master``."""
    csld = notes_master._element.find(_CSLD_TAG)
    if csld is None:
        raise RuntimeError("notes master <p:cSld> not found — malformed notes master")
    sp_tree = csld.find(qn("p:spTree"))
    if sp_tree is None:
        raise RuntimeError("notes master <p:spTree> not found — malformed notes master")

    new_sp = ET.fromstring(_build_notes_text_sp_xml())
    sp_tree.append(new_sp)


def ensure_notes_placeholder(prs) -> Dict[str, Any]:
    """Ensure ``prs`` has a notes-text placeholder; inject one if missing.

    Called by ``ppt_builder`` before filling speaker notes. Safe to call
    repeatedly (idempotent). Never raises on success; on failure logs and
    returns ``repaired=False`` so the caller can soft-fall.

    Returns:
        ``{"repaired": bool, "placeholder_count": int, "error": Optional[str]}``
    """
    report: Dict[str, Any] = {"repaired": False, "placeholder_count": 0, "error": None}
    try:
        notes_master = prs.notes_master
    except Exception as exc:
        report["error"] = f"could not access notes master: {exc}"
        logger.warning("notes_repair: %s", report["error"])
        return report

    if _find_notes_text_placeholder(notes_master):
        report["placeholder_count"] = 1
        logger.debug("notes_repair: notes-text placeholder already present")
        return report

    try:
        _inject_notes_text_placeholder(notes_master)
        report["repaired"] = True
        report["placeholder_count"] = 1
        logger.info("notes_repair: injected notes-text placeholder on notes master")
    except Exception as exc:
        report["error"] = f"failed to inject placeholder: {exc}"
        logger.error("notes_repair: %s", report["error"])
    return report


__all__ = ["ensure_notes_placeholder"]
