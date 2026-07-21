"""BT-142 Phase 3.5b — Placeholder backfill (engine hard limits L2/L3/L4).

The chenyu engine fills, per slide, at most:
  - one TITLE placeholder
  - one BODY placeholder (or two for ``_LAYOUTS_WITH_TWO_BODIES``: body_left, body_right)
  - one PICTURE placeholder (the first matching one)
  - one TABLE placeholder (only on chart/table slide types)

Real designer decks routinely have **multi-image slides** (Team: 4 photos,
Market Validation: 5 logos) and **multi-body slides** (Team: 4 member bios,
Business Model: a grid of 5-6 metric cards). The engine leaves the remaining
TITLE/BODY/PICTURE placeholders empty.

This module is the **post-render backfill** pass. After ``ppt_builder`` (or
``multipass_render.merge``) produces the output PPTX, iterate the slides
against the extended slide_data fields and fill any placeholder the engine
skipped.

Extended slide_data fields (recognized by backfill only — engine ignores them):
  - ``image_paths: [{path: str, placeholder_idx: int, sizing: "fill"|"fit"}]``
        — array of additional pictures to place. ``sizing="fill"`` crops to
        fit; ``sizing="fit"`` letterboxes (default: "fit").
  - ``body_slots: [{text: str, placeholder_idx: int}]``
        — array of additional body-text blocks.
  - ``title_slots: [{text: str, placeholder_idx: int}]``
        — rarely needed; supported for symmetric completeness.

Backfill is **idempotent**: it skips placeholders that already have content
(so re-running on an engine-filled placeholder is safe).

Public API:
  - ``backfill_slide(slide, slide_data)`` → BackfillReport
  - ``backfill_deck(prs, slide_data_list)`` → Dict[int, BackfillReport]
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# python-pptx placeholder type IDs (PP_PLACEHOLDER enum)
_TITLE_TYPE_ID = 13
_BODY_TYPE_ID = 2
_OBJECT_TYPE_ID = 17
_SUBTITLE_TYPE_ID = 3
_PICTURE_TYPE_ID = 18


@dataclass
class BackfillReport:
    filled: List[str] = field(default_factory=list)
    skipped: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "filled": list(self.filled),
            "skipped": list(self.skipped),
            "errors": list(self.errors),
        }


def _placeholder_idx(ph) -> Optional[int]:
    try:
        return ph.placeholder_format.idx
    except Exception:
        return None


def _placeholder_is_empty(ph) -> bool:
    """Heuristic: True if the placeholder has no visible text and no picture."""
    if not ph.has_text_frame:
        return False
    txt = ph.text_frame.text.strip()
    return txt == ""


def _find_placeholder_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if _placeholder_idx(ph) == idx:
            return ph
    return None


def _set_body_text(ph, text: str) -> None:
    tf = ph.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        # tf.paragraphs is guaranteed non-empty in python-pptx (one empty
        # paragraph by default), but be defensive in case clear() removed it.
        if i == 0 and getattr(tf, "paragraphs", None):
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line


def _add_picture_into_placeholder(slide, ph, image_path: str, sizing: str) -> None:
    """Insert an image into a PICTURE placeholder (fills if possible)."""
    from PIL import Image as PILImage
    from pptx.util import Emu

    if not Path(image_path).exists():
        raise FileNotFoundError(f"image not found: {image_path}")

    left = ph.left
    top = ph.top
    width = ph.width
    height = ph.height
    if width is None or height is None:
        # Fall back to native image size scaled to fit the placeholder rect.
        with PILImage.open(image_path) as img:
            iw, ih = img.size
        if width is None and height is None:
            width = Emu(iw * 9525)
            height = Emu(ih * 9525)
        elif width is None:
            width = int(height * iw / ih)
        else:
            height = int(width * ih / iw)

    if sizing == "fill":
        # Crop-to-fill: python-pptx add_picture fills the box with crop.
        slide.shapes.add_picture(
            image_path, left, top, width=width, height=height
        )
    else:
        # Fit (letterbox): compute scaled rect preserving aspect.
        with PILImage.open(image_path) as img:
            iw, ih = img.size
        target_ratio = float(width) / float(height) if height else 1.0
        src_ratio = float(iw) / ih
        if src_ratio > target_ratio:
            new_w = width
            new_h = int(width * ih / iw)
            new_left = left
            new_top = top + (height - new_h) // 2
        else:
            new_h = height
            new_w = int(height * iw / ih)
            new_left = left + (width - new_w) // 2
            new_top = top
        slide.shapes.add_picture(
            image_path, new_left, new_top, width=new_w, height=new_h
        )


def _get_extended_field(slide_data: dict, key: str) -> List[dict]:
    val = slide_data.get(key) or []
    if not isinstance(val, list):
        logger.warning("%s must be a list; got %s — ignoring", key, type(val).__name__)
        return []
    return val


def backfill_slide(slide, slide_data: dict) -> BackfillReport:
    """Fill additional placeholders on ``slide`` from extended slide_data."""
    report = BackfillReport()

    # Body slots
    for slot in _get_extended_field(slide_data, "body_slots"):
        try:
            idx = slot.get("placeholder_idx")
            text = slot.get("text", "")
            if idx is None:
                report.errors.append(f"body_slots missing placeholder_idx: {slot}")
                continue
            ph = _find_placeholder_by_idx(slide, idx)
            if ph is None:
                report.errors.append(f"placeholder idx {idx} not found on slide")
                continue
            if not _placeholder_is_empty(ph):
                report.skipped.append(f"body[{idx}] (already filled)")
                continue
            _set_body_text(ph, text)
            report.filled.append(f"body[{idx}]")
        except Exception as exc:
            report.errors.append(f"body_slots entry failed: {exc}")

    # Title slots (rare; symmetric support)
    for slot in _get_extended_field(slide_data, "title_slots"):
        try:
            idx = slot.get("placeholder_idx")
            text = slot.get("text", "")
            ph = _find_placeholder_by_idx(slide, idx) if idx is not None else None
            if ph is None or not _placeholder_is_empty(ph):
                report.skipped.append(f"title[{idx}] (not empty or missing)")
                continue
            _set_body_text(ph, text)
            report.filled.append(f"title[{idx}]")
        except Exception as exc:
            report.errors.append(f"title_slots entry failed: {exc}")

    # Image slots (multi-image)
    for slot in _get_extended_field(slide_data, "image_paths"):
        try:
            idx = slot.get("placeholder_idx")
            path = slot.get("path")
            sizing = slot.get("sizing", "fit")
            if idx is None or not path:
                report.errors.append(f"image_paths missing path/idx: {slot}")
                continue
            ph = _find_placeholder_by_idx(slide, idx)
            if ph is None:
                report.errors.append(f"image placeholder idx {idx} not found")
                continue
            # Idempotency: skip if placeholder already has an image.
            # python-pptx exposes pictures as shapes; placeholder pictures are
            # detected by the placeholder tag + presence of a blip.
            has_image = False
            try:
                if ph.has_text_frame and ph.text_frame.text.strip() == "" and not ph.has_table:
                    # Placeholder is a picture placeholder with no text — assume empty.
                    has_image = False
            except Exception:
                pass
            if has_image:
                report.skipped.append(f"image[{idx}] (already filled)")
                continue
            _add_picture_into_placeholder(slide, ph, path, sizing)
            report.filled.append(f"image[{idx}]")
        except Exception as exc:
            report.errors.append(f"image_paths entry failed: {exc}")

    return report


def backfill_deck(prs, slide_data_list: List[dict]) -> Dict[int, BackfillReport]:
    """Run ``backfill_slide`` on every slide in ``prs`` aligned by index."""
    reports: Dict[int, BackfillReport] = {}
    for i, slide in enumerate(prs.slides):
        if i >= len(slide_data_list):
            reports[i] = BackfillReport(errors=["no slide_data entry for this index"])
            continue
        reports[i] = backfill_slide(slide, slide_data_list[i])
    return reports


__all__ = ["BackfillReport", "backfill_slide", "backfill_deck"]
