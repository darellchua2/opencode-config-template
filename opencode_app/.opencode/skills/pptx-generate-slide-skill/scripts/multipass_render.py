"""GIT-93 Phase 5 — single-pass render + deck merge.

Originally (BT-142 Phase 3.5a) this module worked around the engine's 8-type
``slide_type`` ceiling by partitioning slides into ≤8-layout batches,
pseudo-typing each batch (``_custom_N``), and merging. GIT-93 made
``layout_name`` a first-class field (Phase 2), so a single
``generate_ppt_from_data`` pass renders N distinct layouts — the batching /
pseudo-typing machinery is obsolete and was deleted.

This module now exposes:
  - ``merge_decks(primary_path, other_paths, output_path)`` → str (output path)
    — independently useful for stitching pre-rendered decks.
  - ``multipass_render(slide_data_list, template_path, output_path, render_fn=None)``
    → str — retained as a thin single-pass wrapper for callers that used the
    old API; delegates directly to ``generate_ppt_from_data``.

The merge primitives (``_copy_slide`` / ``_relink_image_rels``) are unchanged.
"""
from __future__ import annotations

import logging
import shutil
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


def _copy_slide(src_slide, dst_prs) -> None:
    """Deep-copy ``src_slide`` (and its rels/media) into ``dst_prs``.

    Uses the keep-together XML clone pattern. Relocates rId references for
    embedded media (images) so the destination deck is self-contained.
    """
    from copy import deepcopy
    from pptx.oxml.ns import qn

    src_prs = src_slide.part.package

    # Resolve the target layout by name on the destination master.
    layout_name = src_slide.slide_layout.name
    target_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break
    if target_layout is None:
        target_layout = dst_prs.slide_layouts[6]  # blank fallback
        logger.warning(
            "_copy_slide: layout %r not on dst master; using blank layout",
            layout_name,
        )

    # Add a new slide on the destination using the resolved layout.
    new_slide = dst_prs.slides.add_slide(target_layout)

    # Wipe the auto-created shape tree (except nvGrpSpPr + grpSpPr wrappers)
    # and copy each shape from the source slide.
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag
        if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
            continue
        spTree.remove(child)
    for shp in src_slide.shapes:
        spTree.append(deepcopy(shp._element))

    # Re-resolve media (image) relationships: for each blip in the cloned
    # slide, look up the source's image part, copy its blob into the dst
    # package via get_or_add_image_part, and remap the rId.
    _relink_image_rels(src_slide, new_slide)


def _relink_image_rels(src_slide, dst_slide) -> None:
    """Copy image parts referenced by ``src_slide`` into ``dst_slide``'s package; remap rIds.

    python-pptx 1.0+ exposes ``SlidePart.get_or_add_image_part(file_like)`` which
    handles dedup (returns existing part if blob matches). We use it to add
    each source image to the destination, then update the cloned blip's
    ``r:embed`` attribute to point at the new rId.

    Robust against python-pptx's lazy target-part resolution: when a freshly-
    loaded PPTX has rels whose ``_target`` isn't yet a Part instance, we fall
    back to resolving via the rel's ``target_ref`` path + basename lookup in
    the source package's iter_parts().
    """
    from io import BytesIO
    from pptx.oxml.ns import qn

    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    embed_attr_name = f"{{{ns_r}}}embed"
    src_part = src_slide.part
    src_package = src_part.package

    # Cache source image parts by basename for fast lookup.
    src_image_parts_by_basename = {}
    for part in src_package.iter_parts():
        pname = str(part.partname)
        if "/image" in pname or "image" in pname.split("/")[-1].lower():
            src_image_parts_by_basename[pname.split("/")[-1]] = part

    for blip in dst_slide._element.iter(qn("a:blip")):
        old_rId = blip.get(embed_attr_name)
        if not old_rId:
            continue
        # Resolve the source image part. ``target_part`` can raise
        # AssertionError on freshly-loaded packages (python-pptx lazily
        # loads parts). Fall back to ``target_ref`` (always a string) +
        # basename lookup.
        src_image_part = None
        try:
            if old_rId in src_part.rels:
                src_rel = src_part.rels[old_rId]
                try:
                    if not src_rel.is_external:
                        src_image_part = src_rel.target_part
                except (AssertionError, Exception):
                    pass
                if src_image_part is None and not src_rel.is_external:
                    target_ref = src_rel.target_ref
                    basename = target_ref.split("/")[-1]
                    src_image_part = src_image_parts_by_basename.get(basename)
        except (KeyError, Exception):
            pass
        if src_image_part is None:
            logger.debug("_relink_image_rels: no source image for rId=%s; skipping", old_rId)
            continue
        # Copy the image bytes into the destination package
        try:
            blob = src_image_part.blob
            new_image_part = dst_slide.part.get_or_add_image_part(BytesIO(blob))
            new_rId = dst_slide.part.relate_to(
                new_image_part,
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            )
            blip.set(embed_attr_name, new_rId)
        except Exception as exc:
            logger.warning("_relink_image_rels: failed to copy image (rId=%s): %s",
                           old_rId, exc)


def merge_decks(
    primary_path: str,
    other_paths: List[str],
    output_path: str,
) -> str:
    """Merge ``other_paths`` into ``primary_path`` → ``output_path``.

    The primary's slides come first, then each other deck's slides in order.
    All slide layouts referenced must exist on the primary's master (if a
    layout name is missing, the slide falls back to the primary's blank
    layout and a warning is logged).

    Returns the absolute path to ``output_path``.
    """
    from pptx import Presentation

    prs = Presentation(primary_path)
    for other in other_paths:
        other_prs = Presentation(other)
        # Force-load all parts in the source package so their relationship
        # targets resolve to Part instances (python-pptx 1.0+ loads targets
        # lazily; without this, accessing ``target_part`` during save raises
        # AssertionError).
        _ = list(other_prs.part.package.iter_parts())
        for slide in other_prs.slides:
            try:
                _copy_slide(slide, prs)
            except Exception as exc:
                logger.error("merge_decks: failed to copy slide from %s: %s", other, exc)
    prs.save(output_path)
    return str(Path(output_path).resolve())


def multipass_render(
    slide_data_list: List[dict],
    template_path: str,
    output_path: str,
    render_fn: Optional[Any] = None,
) -> str:
    """Render ``slide_data_list`` to ``output_path`` in a single pass.

    GIT-93 Phase 5: ``layout_name`` is now native (Phase 2), so a single
    ``generate_ppt_from_data`` pass renders N distinct layouts — no batching
    or pseudo-typing is needed. This wrapper is retained for callers that
    used the old >8-layout multipass API; it delegates directly.

    Args:
        slide_data_list: list of slide dicts (may carry ``layout_name``).
        template_path: source template.
        output_path: final PPTX path.
        render_fn: callable ``(slides, template_path, output_path,
            config_overrides) -> str`` — defaults to a thin wrapper over
            ``ppt_builder.generate_ppt_from_data``.

    Returns:
        absolute path to ``output_path``.
    """
    if render_fn is None:
        from ppt_builder import generate_ppt_from_data

        def render_fn(slides, tpl, out, overrides):
            return generate_ppt_from_data(
                slides,
                template_path=tpl,
                output_path=out,
                config_overrides=overrides,
            )

    # Single pass — no batching. config_overrides is empty: layout_name is
    # resolved natively by _select_layout, and pseudo-typing is obsolete.
    return render_fn(slide_data_list, template_path, output_path, {})


__all__ = [
    "merge_decks",
    "multipass_render",
]
