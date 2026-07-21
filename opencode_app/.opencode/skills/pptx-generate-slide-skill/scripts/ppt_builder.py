"""
ppt_builder.py
==============
PPT engine using template.pptx Slide Master layouts with proper placeholders.

Loads the template, adds new slides from named layouts (resolved by name, not
index, so layout reordering does not break it), fills placeholders by type
(TITLE, SUBTITLE, OBJECT), and saves the result.

Layouts are matched by name via ``_LAYOUT_NAME_MAP``; ``default.config.json``
may override the layout name for ``title_slide`` / ``content_slide``.

Usage:
    from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR

    result = generate_ppt_from_data(
        slide_data_list,
        output_path=str(DEFAULT_OUTPUT_DIR / "report.pptx"),
    )
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt

# PLAN-GIT-72 (US-5.2): the shared extraction/contract/schema infra now lives
# in the sibling `_common/scripts` package. Put it on sys.path before importing
# those modules. Self-bootstrap so callers need not each replicate this (a DRY
# bootstrap helper is deferred to Phase B / m5).
_COMMON_SCRIPTS = str(Path(__file__).resolve().parent.parent.parent / "_common" / "scripts")
if _COMMON_SCRIPTS not in sys.path:
    sys.path.insert(0, _COMMON_SCRIPTS)

from schema_validator import ValidationError, validate_slide_data_list
from resolvers import resolve_slide_data_list
from schema_extractor import (
    embed_schema,
    extract_schema,
    parse_theme_xml,
    read_embedded_schema,
)
# PLAN-GIT-72 (Phase 2): the pure contract / layout-matching layer now lives in
# the shared layout_contract module. Import the symbols this fill-side engine
# still references (the rest moved there and are imported directly by consumers).
from layout_contract import (
    _LAYOUT_NAME_MAP,
    _SLIDE_TYPE_FINGERPRINT,
    _live_layout_count,
    _normalize_layout_name,
    _resolve_layout_by_fingerprint,
    get_render_contract,
    servable_slide_types,
)
# US-4.6 (m2): shared pure polygon/EMU primitives + target-size resolver + the
# AC5 ratio no-op gate.
from geometry import (
    _EMU_PER_INCH,
    aspect_ratios_match,
    compute_ratio,
    denormalize_polygon,
    normalize_polygon,
    resolve_target_size,
)
from errors import TemplateError  # noqa: E402 — US-4.8/MINOR-2: relocated to _common
from text_fit import (
    LINE_SPACING_DEFAULT,
    MIN_FONT_SIZE_PT,
    ROLE_BASE_PT,
    TEXT_PADDING_IN,
    FontFit,
    fit_font_size,
)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)

_SCRIPT_DIR = Path(__file__).resolve().parent
# BT-142 Phase 2.3: the bundled default template is REMOVED. Callers MUST
# supply a user template path — there is no fallback. This enforces the user's
# "no bundled default.pptx" invariant (Goal #1 of PLAN-BT-142).
# scripts → pptx-generate-slide-skill → skills → .opencode → repo root
_REPO_ROOT = _SCRIPT_DIR.parents[3]
DEFAULT_OUTPUT_DIR = Path.cwd() / "output"

# _TEMPLATE_FILE is intentionally None — there is no bundled default.
# Passing template_path=None to generate_ppt_from_data raises TemplateError.
_TEMPLATE_FILE = None

# TemplateError is imported from _common/scripts/errors.py (US-4.8/MINOR-2).
# It is re-exported here for back-compat: callers that do
# ``from ppt_builder import TemplateError`` continue to work.


_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_SUBTITLE_TYPE = PP_PLACEHOLDER.SUBTITLE
_BODY_TYPE = PP_PLACEHOLDER.BODY
_OBJECT_TYPE = PP_PLACEHOLDER.OBJECT
_PICTURE_TYPE = PP_PLACEHOLDER.PICTURE

# _EMU_PER_INCH is imported from geometry (single source of truth, US-4.6 m2);
# it converts live-placeholder geometry to the inches the text-fit estimator
# consumes (ppt_builder.py _box_inches).

# Maps the embedded-schema ``placeholder_type`` enum to a text-fit role
# (US-4.2). Only text-bearing placeholders participate in font-fitting.
_PT_TO_ROLE = {"title": "title", "subtitle": "subtitle", "body": "body"}

# Body desc-run size as a fraction of the title-run size — preserves the
# historical 12/14 visual hierarchy (0.857) now that sizes are template-derived.
_BODY_DESC_RATIO = 0.85

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# PLAN-GIT-72: _LAYOUT_NAME_MAP / _SLIDE_TYPE_FINGERPRINT / _SERVES_LAYOUT moved
# to layout_contract (shared). _LAYOUTS_WITH_* below are fill-specific (stay).

_LAYOUTS_WITH_SUBTITLE = {
    "title_slide", "closing_slide",
}
_LAYOUTS_WITH_BODY = {
    "content_slide", "content_image_slide",
}
_LAYOUTS_WITH_TWO_BODIES = {
    "two_content_slide", "comparison_slide",
}
_LAYOUTS_WITH_CHART = {
    "chart_slide",
}

# Slide types whose body/content area is the primary selection concern
# (used for the content_area_in2 tie-break).
_CONTENT_SLIDE_TYPES = {
    "content_slide", "two_content_slide", "comparison_slide", "content_image_slide",
}

_CHART_TYPE_MAP: Dict[str, Any] = {
    "bar":                    XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked":            XL_CHART_TYPE.COLUMN_STACKED,
    "bar_horizontal":         XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_horizontal_stacked": XL_CHART_TYPE.BAR_STACKED,
    "pie":                    XL_CHART_TYPE.PIE,
    "pie_exploded":           XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut":               XL_CHART_TYPE.DOUGHNUT,
    "line":                   XL_CHART_TYPE.LINE,
    "line_markers":           XL_CHART_TYPE.LINE_MARKERS,
}

_LEGEND_POSITION_MAP: Dict[str, Any] = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "right":  XL_LEGEND_POSITION.RIGHT,
    "top":    XL_LEGEND_POSITION.TOP,
    "left":   XL_LEGEND_POSITION.LEFT,
}

_CHART_COLORS: List[RGBColor] = [
    RGBColor(0x44, 0x72, 0xC4),
    RGBColor(0xED, 0x7D, 0x31),
    RGBColor(0xFF, 0xC0, 0x00),
    RGBColor(0x5B, 0x9B, 0xD5),
    RGBColor(0x70, 0xAD, 0x47),
    RGBColor(0x95, 0x4F, 0x72),
    RGBColor(0x44, 0x54, 0x6A),
    RGBColor(0xA5, 0xA5, 0xA5),
]

_CHART_FONT_NAME = "Calibri"
_CHART_GRIDLINE_COLOR = RGBColor(0xE7, 0xE6, 0xE6)
_CHART_AXIS_COLOR = RGBColor(0x44, 0x54, 0x6A)
_CHART_TEXT_COLOR = RGBColor(0x44, 0x54, 0x6A)

_CHART_DEFAULT_TYPE = "bar"

_PIE_CHART_TYPES = {"pie", "pie_exploded", "doughnut"}
_BAR_CHART_TYPES = {
    "bar", "bar_stacked", "bar_horizontal", "bar_horizontal_stacked",
}


def _slide_dims_emu(slide: Any) -> Tuple[int, int]:
    """Return ``(width_emu, height_emu)`` of the presentation's slide size."""
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def _chart_bbox(slide: Any) -> Tuple[int, int, int, int]:
    """Compute a chart bounding box responsive to the slide's actual size.

    The former hard-coded ``_CHART_X/Y/CX/CY`` constants were sized for
    13.33x7.5in widescreen but overflowed on smaller 16:9 templates (e.g.
    10x5.625in).  Margins are proportional so the chart fits any slide size.
    """
    sw, sh = _slide_dims_emu(slide)
    margin_x = max(int(sw * 0.07), int(Inches(0.5)))
    y = max(int(sh * 0.25), int(Inches(1.4)))
    bottom_margin = int(Inches(0.3))
    cx = sw - 2 * margin_x
    cy = sh - y - bottom_margin
    return margin_x, y, cx, cy


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a '#RRGGBB' or 'RRGGBB' string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _extract_chart_theme(slide: Any) -> Optional[Dict[str, Any]]:
    """Extract accent/dk/lt colors + minor font from the live prs theme.

    Returns a dict with keys ``series`` (list of RGBColor), ``gridline``,
    ``axis_text``, ``font`` — or ``None`` on any failure (caller falls back
    to the hardcoded ``_CHART_*`` constants).
    """
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        prs = slide.part.package.presentation_part.presentation
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(RT.THEME)
        colors, fonts = parse_theme_xml(theme_part.blob)

        accent_roles = ("accent1", "accent2", "accent3", "accent4", "accent5", "accent6")
        series = []
        for role in accent_roles:
            hex_val = colors.get(role, "")
            if hex_val:
                series.append(_hex_to_rgb(hex_val))
        # Pad with dk2 if fewer than 3 accents were found.
        if len(series) < 3:
            dk2 = colors.get("dk2", "44546A")
            series.append(_hex_to_rgb(dk2))

        return {
            "series": series,
            "gridline": _hex_to_rgb(colors.get("lt2", "E7E6E6")),
            "axis_text": _hex_to_rgb(colors.get("dk2", "44546A")),
            "font": fonts.get("minor_latin") or "Calibri",
        }
    except Exception:
        return None


# --- Image placement (#18) -------------------------------------------------
_IMAGE_DEFAULT_PRESET = "below-title"
_VALID_IMAGE_PRESETS = {"full", "below-title", "half-left", "half-right"}


def _image_bbox(slide: Any, preset_key: str) -> Dict[str, int]:
    """Compute an image bounding box responsive to the slide's actual size."""
    sw, sh = _slide_dims_emu(slide)
    y = max(int(sh * 0.25), int(Inches(1.4)))
    bottom_margin = int(Inches(0.3))
    cy = sh - y - bottom_margin
    if preset_key in ("half-left", "half-right"):
        margin = int(Inches(0.5))
        half_w = (sw - 3 * margin) // 2
        if preset_key == "half-left":
            return {"x": margin, "y": y, "cx": half_w, "cy": cy}
        return {"x": 2 * margin + half_w, "y": y, "cx": half_w, "cy": cy}
    margin_x = max(int(sw * 0.07), int(Inches(0.5)))
    return {"x": margin_x, "y": y, "cx": sw - 2 * margin_x, "cy": cy}


def _build_layout_index(prs: Presentation) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    exact: Dict[str, Any] = {}
    normalized: Dict[str, Any] = {}
    for layout in prs.slide_layouts:
        nm = layout.name
        exact[nm.lower()] = layout
        norm_key = _normalize_layout_name(nm)
        normalized.setdefault(norm_key, layout)
    return exact, normalized


def _resolve_layout(
    candidate_names: List[str],
    exact: Dict[str, Any],
    normalized: Dict[str, Any],
) -> Optional[Any]:
    for cand in candidate_names:
        if cand.lower() in exact:
            return exact[cand.lower()]
    for cand in candidate_names:
        key = _normalize_layout_name(cand)
        if key in normalized:
            return normalized[key]
    return None


# PLAN-GIT-72: _composition_diff / _name_affinity / _LAYOUT_TYPES_NEEDING_SIDEBYSIDE
# / _content_placeholders_stacked / _resolve_layout_by_fingerprint moved to
# layout_contract (shared). _select_layout below consumes them via that import.


def _select_layout(
    slide_type: str,
    contract: Optional[Dict[str, Any]],
    config: Dict[str, Any],
    prs: Presentation,
    exact_idx: Dict[str, Any],
    norm_idx: Dict[str, Any],
    page_num: int,
) -> Optional[Any]:
    """Resolve a slide_type to a concrete ``SlideLayout`` (issue #44).

    Precedence: config pin (``<slide_type>_layout``) → fingerprint match →
    name-based fallback → degradation (skip + warn). Without a contract the path
    is the original name-based matching, so behaviour is backward compatible.
    """
    if slide_type not in _SLIDE_TYPE_FINGERPRINT and slide_type not in _LAYOUT_NAME_MAP:
        logger.warning("Page %d: unknown slide_type '%s', skipped", page_num, slide_type)
        return None

    # 1. Config pin — explicit layout name (highest precedence; all 8 types).
    pinned = config.get(f"{slide_type}_layout")
    if pinned:
        layout = _resolve_layout([pinned], exact_idx, norm_idx)
        if layout is not None:
            return layout
        logger.warning(
            "Page %d: config pin '%s' not found; falling back", page_num, pinned)

    # 2. Fingerprint match (contract-aware; names are a tie-breaker).
    if contract:
        idx, reason = _resolve_layout_by_fingerprint(slide_type, contract)
        if idx is not None:
            return prs.slide_layouts[idx]
        logger.warning(
            "Page %d: fingerprint degradation for '%s': %s",
            page_num, slide_type, reason)

    # 3. Name-based fallback (backward-compatible safety net).
    candidates = _LAYOUT_NAME_MAP.get(slide_type)
    if candidates:
        layout = _resolve_layout(candidates, exact_idx, norm_idx)
        if layout is not None:
            return layout

    # 4. Degradation: nothing usable.
    logger.warning(
        "Page %d: no layout matched for slide_type '%s', skipped", page_num, slide_type)
    return None


# PLAN-GIT-72: servable_slide_types moved to layout_contract (shared).


def _load_config(template_path: Optional[str] = None) -> Dict[str, Any]:
    """Load optional layout-name pin file next to the user-supplied template.

    BT-142 Phase 2.3: previously read a fixed ``template/default.config.json``
    shipped with the bundled default. Since there is no longer a bundled
    default, the config file is now looked up next to whatever template the
    user supplied: ``<template_stem>.config.json`` in the same directory.

    Returns an empty dict if the user supplied no template or no config file
    exists alongside it (the common case — fingerprint matching handles
    layout resolution without pins).
    """
    if template_path is None:
        return {}
    p = Path(template_path)
    config_path = p.with_suffix(".config.json")  # e.g. my_deck.config.json
    if config_path.exists():
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (OSError, ValueError) as exc:
            logger.warning("Could not load config %s (%s)", config_path, exc)
    return {}


def _validate_template(
    prs: Presentation,
    template_path: str,
    contract: Optional[Dict[str, Any]] = None,
) -> None:
    """US-4.7 pre-flight: abort with :class:`TemplateError` on severe problems.

    Severe (fatal): no slide master / zero layouts / serves none of the 8 slide
    types. Openability is checked by the caller (``Presentation(...)`` is wrapped).
    Minor issues (missing fonts, no header/footer, small content area, no embedded
    schema) are NOT checked here — they stay non-fatal warnings elsewhere.
    """
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    if not masters:
        raise TemplateError(
            f"Template has no slide master — cannot generate slides. Auto-repair "
            f"was attempted but failed (see logged repair level). Provide a "
            f"valid .pptx that has a slide master, or omit --template to use "
            f"the default template. Path: {template_path}"
        )

    if len(prs.slide_layouts) == 0:
        raise TemplateError(
            f"Template has no slide layouts — cannot add any slide: {template_path}"
        )

    # Servability (AC6) can only be judged when the contract is available. When
    # contract build failed we cannot make this determination, so skip (warned
    # upstream) rather than risk a false fatal.
    if contract is not None:
        try:
            served = servable_slide_types(contract)
        except Exception:
            served = {}
        available = [t for t, info in served.items() if info and info.get("available")]
        if not available:
            raise TemplateError(
                f"Template serves none of the 8 slide types — cannot generate "
                f"any slide: {template_path}"
            )


def _remove_all_slides(prs: Presentation) -> int:
    count = len(prs.slides)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].attrib.get(f"{{{_REL_NS}}}id")
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return count


def _find_placeholder(slide: Any, ph_type: Any) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            return ph
    return None


def _find_placeholders(slide: Any, ph_type: Any) -> List[Any]:
    return [ph for ph in slide.placeholders if ph.placeholder_format.type == ph_type]


def _find_title_placeholder(slide: Any) -> Optional[Any]:
    for ph_type in _TITLE_TYPES:
        ph = _find_placeholder(slide, ph_type)
        if ph:
            return ph
    return None


def _find_body_placeholder(slide: Any) -> Optional[Any]:
    ph = _find_placeholder(slide, _BODY_TYPE)
    if ph:
        return ph
    objects = _find_placeholders(slide, _OBJECT_TYPE)
    return objects[0] if objects else None


def _matching_layout_placeholder(shape: Any, layout: Any) -> Optional[Any]:
    """The layout placeholder matching ``shape`` by ``placeholder_format.idx``.

    Returns ``None`` for non-placeholder shapes (textboxes — python-pptx raises
    on ``.placeholder_format``), an absent idx, no layout, or no match. Shared by
    the M1 base-size and line-spacing readers to avoid duplicating the idx walk
    (and the ``is_placeholder`` guard that avoids the ValueError).
    """
    if not getattr(shape, "is_placeholder", False) or layout is None:
        return None
    pf = getattr(shape, "placeholder_format", None)
    idx = pf.idx if pf is not None else None
    if idx is None:
        return None
    try:
        for ph in layout.placeholders:
            lpf = getattr(ph, "placeholder_format", None)
            if lpf is not None and lpf.idx == idx:
                return ph
    except Exception:  # best-effort; never block the render
        return None
    return None


def _layout_sample_font_pt(shape: Any, layout: Any) -> Optional[float]:
    """Probe the **layout's** matching placeholder for an explicit run size.

    Layout placeholders usually carry "Click to add…" sample text with an
    explicit font size — far more reliable than the post-``add_slide`` empty
    slide placeholder, which python-pptx does not resolve through ``lstStyle``
    inheritance. Returns ``None`` when no explicit size is found.
    (US-4.2 / arch-review M1 tier 2.)
    """
    ph = _matching_layout_placeholder(shape, layout)
    if ph is None or not getattr(ph, "has_text_frame", False):
        return None
    try:
        for p in ph.text_frame.paragraphs:
            for run in p.runs:
                size = run.font.size
                if size is not None:
                    return float(size.pt)
    except Exception:  # best-effort; never block the render
        return None
    return None


def _build_schema_font_map(schema: Optional[Dict[str, Any]]) -> Dict[str, Dict[str, float]]:
    """Build ``{layout_name: {role: size_pt}}`` from an embedded schema.

    Only the first explicit ``font.size_pt`` per role per layout is kept. Empty
    when the schema is absent (the caller then relies on layout-sample → role
    ceilings). (US-4.2 / arch-review M1 tier 1.)
    """
    result: Dict[str, Dict[str, float]] = {}
    if not schema:
        return result
    for layout in schema.get("slide_layouts") or []:
        name = layout.get("layout_name")
        role_map: Dict[str, float] = {}
        for comp in layout.get("components") or []:
            if comp.get("type") != "placeholder":
                continue
            role = _PT_TO_ROLE.get(comp.get("placeholder_type") or "")
            if not role or role in role_map:
                continue
            size = (comp.get("font") or {}).get("size_pt")
            if size:
                role_map[role] = float(size)
        if name:
            result[name] = role_map
    return result


def _resolve_base_font_pt(
    shape: Any,
    role: str,
    schema_font_map: Dict[str, Dict[str, float]],
    layout: Any,
) -> Tuple[float, str]:
    """Resolve the base font size (pt) for ``role`` via the M1 chain.

    Returns ``(base_pt, source)`` where ``source ∈ {"schema","layout_sample",
    "role_ceiling"}``. The conservative role ceilings (body 14 / title 28 /
    subtitle 18) ensure the resolved base never exceeds a typical template's
    real size when the upper tiers miss.
    """
    role_map = (schema_font_map or {}).get(getattr(layout, "name", "") or "") or {}
    schema_pt = role_map.get(role)
    if schema_pt and schema_pt > 0:
        return float(schema_pt), "schema"
    sample_pt = _layout_sample_font_pt(shape, layout)
    if sample_pt and sample_pt > 0:
        return float(sample_pt), "layout_sample"
    return ROLE_BASE_PT.get(role, ROLE_BASE_PT["body"]), "role_ceiling"


def _layout_line_spacing(shape: Any, layout: Any) -> float:
    """Best-effort line-spacing factor from the layout placeholder (else 1.2).

    Reads the matched layout placeholder's first explicit ``paragraph
    .line_spacing``. python-pptx returns ``None`` for ``lstStyle``-inherited
    spacing, so this almost always falls back to the default — but when a
    layout sets spacing explicitly it is honoured (US-4.2 Details). Returns a
    float multiplier (single spacing = 1.0).
    """
    ph = _matching_layout_placeholder(shape, layout)
    if ph is None or not getattr(ph, "has_text_frame", False):
        return LINE_SPACING_DEFAULT
    try:
        for p in ph.text_frame.paragraphs:
            ls = p.line_spacing
            if ls is not None:
                # python-pptx: a plain ``float`` is a line-height multiple
                # (e.g. 1.2); a ``Length`` (Pt/Centipoints) is an *exact
                # points* value and is an ``int`` subclass — NOT expressible
                # as a multiplier without the font size, so fall back to the
                # default rather than misusing the EMU/centipoint integer
                # (e.g. Pt(18) -> 228600) as a 228600x multiplier.
                return float(ls) if isinstance(ls, float) else LINE_SPACING_DEFAULT
    except Exception:
        pass
    return LINE_SPACING_DEFAULT


def _box_inches(shape: Any) -> Tuple[float, float]:
    """Return ``(width_in, height_in)`` of a shape; (0,0) if unavailable."""
    try:
        w = int(shape.width or 0) / _EMU_PER_INCH
        h = int(shape.height or 0) / _EMU_PER_INCH
        return w, h
    except Exception:
        return 0.0, 0.0


def _effective_box_height(h_in: float, base_pt: float, line_spacing: float) -> float:
    """Height used for vertical fit-checking, with an auto-grow guard.

    python-pptx reports ``None`` for an inherited ``text_frame.auto_size``
    (it does not resolve the ``lstStyle`` cascade), so auto-grow placeholders
    (titles/subtitles that expand to fit) can't be detected directly. Their
    reported base height is often smaller than a single line at the base size.
    When that happens we treat the box as **unbounded** (returns ``inf``) so the
    estimator does not false-shrink — the placeholder grows instead, and only
    ``word_wrap`` governs width. Tall fixed boxes (body content) keep their real
    height and shrink normally.
    """
    one_line_in = base_pt * line_spacing / 72.0
    if h_in <= 0 or (h_in - 2 * TEXT_PADDING_IN) < one_line_in:
        return float("inf")
    return h_in


def _set_text(
    shape: Any,
    text: str,
    *,
    role: str = "title",
    schema_font_map: Optional[Dict[str, Dict[str, float]]] = None,
    layout: Any = None,
) -> Optional[FontFit]:
    """Fill a single-string placeholder (title/subtitle) with text-fitting.

    Resolves the base size (M1 chain), shrinks in −2pt steps to the 8pt floor
    when the text overflows the box, and **writes an explicit ``run.font.size``
    only when the estimator actually shrank** (M3) — otherwise the placeholder
    keeps inheriting the layout/master size (today's behaviour), which preserves
    title/subtitle appearance on the common non-overflow path. Sets
    ``word_wrap=True``. Returns the :class:`FontFit` (or ``None`` on failure).
    """
    if not shape or not shape.has_text_frame:
        return None
    try:
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        if not (text or "").strip():
            return None
        base, source = _resolve_base_font_pt(shape, role, schema_font_map or {}, layout)
        w_in, h_in = _box_inches(shape)
        ls = _layout_line_spacing(shape, layout)
        fit = fit_font_size(
            text, base, w_in, _effective_box_height(h_in, base, ls),
            line_spacing=ls, base_source=source,
        )
        p = tf.paragraphs[0]
        p.text = text
        if fit.adjusted:  # M3: explicit size only on actual shrink
            for run in p.runs:
                run.font.size = Pt(fit.applied_size_pt)
        return fit
    except Exception as exc:
        logger.warning("Failed to set text: %s", exc)
        return None


def _parse_line(line: str) -> Tuple[str, str]:
    clean = re.sub(r"\*\*", "", line.strip())
    if not clean:
        return ("", "")
    for sep in [" \u2014 ", " - ", ": "]:
        if sep in clean:
            parts = clean.split(sep, 1)
            return (parts[0].strip(), parts[1].strip())
    return (clean, "")


def _set_notes(slide: Any, notes_text: str) -> bool:
    text = (notes_text or "").strip()
    if not text:
        return False
    try:
        slide.notes_slide.notes_text_frame.text = text
        return True
    except Exception as exc:
        logger.warning("Failed to set notes: %s", exc)
        return False


def _set_body_text(
    shape: Any,
    text: str,
    *,
    role: str = "body",
    schema_font_map: Optional[Dict[str, Dict[str, float]]] = None,
    layout: Any = None,
) -> Optional[FontFit]:
    """Fill a body placeholder with text-fitting (US-4.2).

    Body base size is **template-derived** (M1 chain), replacing the previous
    ``Pt(14)``/``Pt(12)`` hardcode (Q2). Each line keeps the bold-title/desc
    split: the title run is sized at the fitted ``applied_pt`` (bold), the desc
    run at ``applied_pt × 0.85`` (preserving the historical 12/14 hierarchy).
    Because body is the primary text-fitting target, an explicit size is always
    written here (template-derived); when the estimator does not shrink,
    ``applied == base``, so non-overflow bodies are appearance-stable whenever
    the resolved base matches the historical ceiling (body 14 → 12/14). Sets
    ``word_wrap=True``. Returns the :class:`FontFit` (or ``None`` on failure).
    """
    if not shape or not shape.has_text_frame:
        return None
    try:
        lines = [l for l in (text or "").split("\n") if l.strip()]
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        if not lines:
            return None
        base, source = _resolve_base_font_pt(shape, role, schema_font_map or {}, layout)
        w_in, h_in = _box_inches(shape)
        ls = _layout_line_spacing(shape, layout)
        fit = fit_font_size(
            text, base, w_in, _effective_box_height(h_in, base, ls),
            line_spacing=ls, base_source=source,
        )
        title_size = fit.applied_size_pt
        desc_size = max(title_size * _BODY_DESC_RATIO, MIN_FONT_SIZE_PT)

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            title_part, desc_part = _parse_line(line)

            if title_part:
                run = p.add_run()
                run.text = title_part
                run.font.bold = True
                run.font.size = Pt(title_size)
            if desc_part:
                run = p.add_run()
                run.text = f" \u2014 {desc_part}" if title_part else desc_part
                run.font.size = Pt(desc_size)

        return fit
    except Exception as exc:
        logger.warning("Failed to set body text: %s", exc)
        return None


def _apply_series_colors(chart: Any, chart_type_key: str, colors: Optional[List[RGBColor]] = None) -> None:
    palette = colors if colors else _CHART_COLORS
    is_pie = chart_type_key in _PIE_CHART_TYPES
    try:
        plot = chart.plots[0]
        if is_pie:
            for idx, point in enumerate(plot.series[0].points):
                color = palette[idx % len(palette)]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
        else:
            for idx, series in enumerate(plot.series):
                color = palette[idx % len(palette)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
                if chart_type_key.startswith("line"):
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2.5)
    except Exception as exc:
        logger.warning("Failed to apply series colors: %s", exc)


def _add_chart_to_slide(slide: Any, slide_data: Dict[str, Any]) -> bool:
    chart_type_key = slide_data.get("chart_type", _CHART_DEFAULT_TYPE)
    if chart_type_key not in _CHART_TYPE_MAP:
        logger.warning(
            "Unknown chart_type '%s', defaulting to '%s'",
            chart_type_key, _CHART_DEFAULT_TYPE,
        )
        chart_type_key = _CHART_DEFAULT_TYPE

    categories = slide_data.get("categories", [])
    series_list = slide_data.get("series", [])

    if not categories or not series_list:
        logger.warning(
            "Chart slide missing categories or series, skipping chart"
        )
        return False

    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for s in series_list:
        name = s.get("name", "")
        values = list(s.get("values", []))
        chart_data.add_series(name, values)

    xl_type = _CHART_TYPE_MAP[chart_type_key]
    try:
        cx_chart_x, cx_chart_y, cx_chart_cx, cx_chart_cy = _chart_bbox(slide)
        graphic_frame = slide.shapes.add_chart(
            xl_type,
            cx_chart_x, cx_chart_y, cx_chart_cx, cx_chart_cy,
            chart_data,
        )
    except Exception as exc:
        logger.error("Failed to create chart: %s", exc)
        return False

    chart = graphic_frame.chart
    options = slide_data.get("chart_options", {})

    # Extract theme-derived chart styling (falls back to constants on failure).
    theme = _extract_chart_theme(slide)
    font_name = theme["font"] if theme else _CHART_FONT_NAME
    text_color = theme["axis_text"] if theme else _CHART_TEXT_COLOR
    gridline_color = theme["gridline"] if theme else _CHART_GRIDLINE_COLOR
    axis_color = theme["axis_text"] if theme else _CHART_AXIS_COLOR
    series_palette = theme["series"] if theme else None

    chart.has_title = False
    chart.font.name = font_name
    chart.font.size = Pt(11)

    legend_pos_key = options.get("legend_position", "bottom")
    if legend_pos_key == "none":
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = _LEGEND_POSITION_MAP.get(
            legend_pos_key, XL_LEGEND_POSITION.BOTTOM,
        )
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = font_name
        chart.legend.font.color.rgb = text_color

    is_pie = chart_type_key in _PIE_CHART_TYPES
    is_bar = chart_type_key in _BAR_CHART_TYPES

    try:
        plot = chart.plots[0]
        show_labels = options.get("show_data_labels", True)
        plot.has_data_labels = show_labels
        if show_labels:
            labels = plot.data_labels
            labels.font.size = Pt(10)
            labels.font.name = font_name
            labels.font.color.rgb = text_color
            if is_pie:
                labels.show_percentage = True
                labels.show_value = False
                labels.show_category_name = False
                labels.number_format = "0%"
            else:
                labels.show_value = True
                labels.show_percentage = False
                labels.number_format = options.get("value_format", "#,##0.0")
                if is_bar:
                    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception as exc:
        logger.warning("Failed to set data labels: %s", exc)

    if not is_pie:
        try:
            val_axis = chart.value_axis
            val_axis.has_major_gridlines = True
            val_axis.major_gridlines.format.line.color.rgb = gridline_color
            val_axis.major_gridlines.format.line.width = Pt(0.75)
            val_axis.tick_labels.font.size = Pt(10)
            val_axis.tick_labels.font.name = font_name
            val_axis.tick_labels.font.color.rgb = axis_color
            val_axis.format.line.color.rgb = axis_color
            val_axis.tick_labels.number_format = options.get("y_axis_format", "#,##0.0")
            if options.get("y_axis_min") is not None:
                val_axis.minimum_scale = options["y_axis_min"]
            if options.get("y_axis_max") is not None:
                val_axis.maximum_scale = options["y_axis_max"]
            if options.get("y_axis_major_unit") is not None:
                val_axis.major_unit = options["y_axis_major_unit"]
            if options.get("y_axis_title"):
                val_axis.has_title = True
                val_axis.axis_title.text_frame.text = options["y_axis_title"]
                val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                val_axis.axis_title.text_frame.paragraphs[0].font.name = font_name
                val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = axis_color

            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(10)
            cat_axis.tick_labels.font.name = font_name
            cat_axis.tick_labels.font.color.rgb = axis_color
            cat_axis.format.line.color.rgb = axis_color
            if options.get("x_axis_title"):
                cat_axis.has_title = True
                cat_axis.axis_title.text_frame.text = options["x_axis_title"]
                cat_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                cat_axis.axis_title.text_frame.paragraphs[0].font.name = font_name
                cat_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = axis_color
        except Exception as exc:
            logger.warning("Failed to set axis options: %s", exc)

    _apply_series_colors(chart, chart_type_key, colors=series_palette)

    logger.info(
        "  Chart: type=%s, categories=%d, series=%d",
        chart_type_key, len(categories), len(series_list),
    )
    return True


def _find_picture_placeholder(slide: Any) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.type == _PICTURE_TYPE:
            return ph
    return None


def _add_image_to_slide(slide: Any, slide_data: Dict[str, Any]) -> bool:
    """Embed a native, editable PowerPoint picture from ``image_path`` (#18).

    Placement order:
      1. If the layout has a PICTURE placeholder, fill it natively.
      2. Otherwise place using a named preset (``image_position``) or an
         explicit ``image_size`` override, in the free space below the title.

    Images are embedded (not linked) so the PPTX stays self-contained and the
    picture remains fully editable in PowerPoint.
    """
    image_path = slide_data.get("image_path")
    if not image_path:
        return False

    p = Path(image_path)
    if not p.exists():
        logger.warning("image_path not found: %s, skipping image", image_path)
        return False

    preset_key = slide_data.get("image_position", _IMAGE_DEFAULT_PRESET)
    if preset_key not in _VALID_IMAGE_PRESETS:
        logger.warning(
            "Unknown image_position '%s', defaulting to '%s'",
            preset_key, _IMAGE_DEFAULT_PRESET,
        )
        preset_key = _IMAGE_DEFAULT_PRESET

    try:
        pic_ph = _find_picture_placeholder(slide)
        if pic_ph is not None:
            try:
                pic_ph.insert_picture(str(p))
                logger.info("  Image (placeholder): %s", p.name)
                return True
            except Exception:
                # Fall back to free placement at the placeholder's frame box.
                box = {
                    "x": pic_ph.left, "y": pic_ph.top,
                    "cx": pic_ph.width, "cy": pic_ph.height,
                }
                slide.shapes.add_picture(str(p), box["x"], box["y"], box["cx"], box["cy"])
                logger.info("  Image (placeholder frame): %s", p.name)
                return True

        box = _image_bbox(slide, preset_key)
        cx, cy = box["cx"], box["cy"]
        size = slide_data.get("image_size")
        if isinstance(size, dict):
            if size.get("width"):
                cx = Inches(size["width"])
            if size.get("height"):
                cy = Inches(size["height"])
        slide.shapes.add_picture(str(p), box["x"], box["y"], cx, cy)
        logger.info("  Image (%s): %s", preset_key, p.name)
        return True
    except Exception as exc:
        logger.error("Failed to embed image '%s': %s", image_path, exc)
        return False


_DEFAULT_CLOSING_NOTES = (
    'KEY MESSAGE: Thank you — close warmly and open the floor for questions.\n'
    '"Thank you all for your time today."\n'
    'Pause. Make eye contact across the room.\n'
    '"I hope this gave you a clear picture of where we are and where we are headed."\n'
    'TRANSITION: "I would love to take any questions you have."\n'
    'COACHING: Warm, unhurried close. Be ready for: "Can you share the deck?" '
    '— yes, I will send it after.'
)


def _ensure_default_closing(
    slide_data_list: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    # #40 / US-6: append a default Thank-You closing when the deck is large
    # enough and does not already end on a closing_slide. Non-destructive —
    # returns a new list, leaving the caller's list untouched.
    if not isinstance(slide_data_list, list) or len(slide_data_list) < 3:
        return slide_data_list
    last_type = (slide_data_list[-1] or {}).get("slide_type", "")
    if last_type == "closing_slide":
        return slide_data_list
    closing = {
        "slide_type": "closing_slide",
        "title": "Thank You",
        "notes": _DEFAULT_CLOSING_NOTES,
    }
    logger.info("Auto-appending default closing slide (default_closing=True)")
    return list(slide_data_list) + [closing]


def _compose_signoff(slide_data: Dict[str, Any]) -> str:
    """Compose the closing-slide sign-off text from optional presenter fields.

    ``presenter_name`` / ``presenter_email`` → ``"Prepared by: {name}\\n{email}"``
    (empty parts omitted). Returns ``""`` when neither is set, which the render
    loop uses to **clear** the subtitle placeholder — overriding the layout's
    inherited sample text (e.g. ``"Prepared by: Lecturer Name\\nEmail address"``)
    so it never bleeds into the slide (US-4.3 follow-up).
    """
    name = (slide_data.get("presenter_name") or "").strip()
    email = (slide_data.get("presenter_email") or "").strip()
    parts: List[str] = []
    if name:
        parts.append(f"Prepared by: {name}")
    if email:
        parts.append(email)
    return "\n".join(parts)


def _remove_placeholder(shape: Any) -> bool:
    """Detach a placeholder shape from its slide (element removal).

    Suppresses inherited layout sample text: an empty slide placeholder still
    inherits the layout's text (e.g. closing's ``"Prepared by: Lecturer Name"`` —
    verified: an empty ``<a:p/>`` shows the layout's runs in PowerPoint), so when
    there is no content we remove the placeholder element entirely — nothing
    remains to render the inherited text. Best-effort (never raises).
    """
    try:
        el = getattr(shape, "element", None)
        if el is not None and el.getparent() is not None:
            el.getparent().remove(el)
            return True
    except Exception as exc:  # never block the render
        logger.warning("Failed to remove placeholder: %s", exc)
    return False


# PLAN-GIT-72: _live_layout_count moved to layout_contract (imported above;
# _embedded_schema_stale below still consumes it).


def _embedded_schema_stale(template_path: str, embedded_layout_count: int) -> bool:
    """True if the live template's layout count diverges from the embedded
    schema's (edit-without-re-embed — the embedded JSON has no mtime-invalidation,
    unlike the sidecar cache). Cheap structural check; returns ``False`` if the
    live count cannot be determined (never forces a re-extract on uncertainty).
    Shared by the M5 staleness warning and US-4.3's auto-template (M2).
    """
    live = _live_layout_count(template_path)
    return live is not None and live != embedded_layout_count


# PLAN-GIT-72: _warn_if_embedded_stale + get_render_contract moved to
# layout_contract (shared). get_render_contract is imported above and consumed
# by generate_ppt_from_data.


def _font_fit_report_entry(
    role: str, field: str, fit: Optional[FontFit]
) -> Optional[Dict[str, Any]]:
    """Convert a :class:`FontFit` into a render-report placeholder record."""
    if fit is None:
        return None
    return {
        "role": role,
        "field": field,
        "template_size_pt": round(fit.base_size_pt, 2),
        "applied_size_pt": round(fit.applied_size_pt, 2),
        "base_source": fit.base_source,
        "font_size_adjusted": fit.adjusted,
        "fits": fit.fits,
        "lines_estimated": fit.lines_estimated,
    }


def _write_render_report(pptx_path: str, report: Dict[str, Any]) -> None:
    """Write the US-4.2 render report sidecar ``<output>.render.json``.

    Records per-slide per-placeholder font-fit decisions (the
    ``font_size_adjusted`` flag + original/applied sizes). Failure is debug-log
    only — it never blocks a successful render. Written atomically
    (temp + rename) so a stale report from a failed prior run never persists
    (arch-review m3).
    """
    import json
    import os
    import tempfile

    out = Path(pptx_path)
    report_path = out.with_name(out.stem + ".render.json")
    try:
        fd, tmp = tempfile.mkstemp(
            prefix=out.stem + ".render.", suffix=".json", dir=str(out.parent)
        )
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as fh:
                json.dump(report, fh, indent=2, ensure_ascii=False)
            os.replace(tmp, str(report_path))
        except Exception:
            try:
                os.unlink(tmp)
            except OSError:
                pass
            raise
        logger.info("Render report: %s", report_path)
    except Exception as exc:  # pragma: no cover - never block the render
        logger.debug("Render report not written (%s): %s", report_path, exc)


def _ensure_output_templated(
    output_path: str, template_path: str, auto_template: bool
) -> Dict[str, Any]:
    """US-4.3: embed ``ppt/template_schema.json`` into the OUTPUT pptx after save.

    python-pptx drops the unmodeled part on ``prs.save``, so the output never
    carries embedded JSON unless re-embedded here (PLAN-GIT-63). The schema is
    sourced from the **input template** (arch-review M1 — never the rendered
    output, so ``_infer_title`` reads the template's own identity, not the deck's
    cover slide): copy the input's embedded schema when valid **and** non-stale
    (arch-review M2 — a stale embedded schema is not laundered into a trusted
    output), else extract fresh from the input. ``embed_schema`` is already
    atomic (temp + ``os.replace``), so it is called in place (m3). Non-fatal.
    Returns a ``templating`` report fragment.
    """
    result: Dict[str, Any] = {
        "input_template_embedded": False,
        "output_templated": False,
        "schema_source": "failed",
        "message": "",
    }
    if not auto_template:
        result["schema_source"] = "disabled"
        result["message"] = "auto_template disabled"
        return result
    try:
        embedded = read_embedded_schema(template_path)
        result["input_template_embedded"] = embedded is not None
        if embedded is not None and not _embedded_schema_stale(
            template_path, len(embedded.get("slide_layouts") or [])
        ):
            schema = embedded
            source = "copied_input"
        else:
            # M1: extract from the INPUT template (its master+layouts are
            # byte-identical to the output's; _infer_title reads the template's
            # own slides[0], not the rendered deck's cover).
            schema = extract_schema(template_path)
            source = "extracted_input"
        embed_schema(output_path, schema, output_path)  # m3: atomic, in place
        result["output_templated"] = True
        result["schema_source"] = source
        result["message"] = f"output templated ({source})"
        logger.info("US-4.3 auto-template: %s (%s)", output_path, source)
    except Exception as exc:  # never block the render
        result["schema_source"] = "failed"
        result["message"] = f"templating failed: {exc}"
        logger.debug("Auto-template failed for %s: %s", output_path, exc)
    return result


def _scale_shapes_geometry(shapes: Any, sx: float, sy: float) -> int:
    """Proportionally scale every shape's left/top/width/height by (sx, sy).

    Ratio-scaling preserves off-slide bleeds (no clamp) and repositions both
    placeholders and background chrome. Shapes with ``None`` geometry (purely
    inherited) are skipped. Returns the count of shapes processed without error.

    **Group shapes (US-4.6 arch-review C1 fix):** do NOT recurse into a group's
    children — only scale the group's container (``off``/``ext``). The OOXML
    child→slide affine map ``slide_pos = off + (child − chOff) · (ext/chExt)``
    then scales every child's *slide-space* position by exactly ``sx``/``sy``
    (``chOff``/``chExt`` stay fixed), so children are placed correctly without a
    second scaling pass. (Recursing AND scaling the container previously produced
    an ``sx²`` quadratic distortion — verified fixed.) Nested groups scale
    transitively through each level's map.
    """
    n = 0
    for sh in shapes:
        try:
            if sh.left is not None:
                sh.left = int(round(sh.left * sx))
            if sh.top is not None:
                sh.top = int(round(sh.top * sy))
            if sh.width is not None:
                sh.width = int(round(sh.width * sx))
            if sh.height is not None:
                sh.height = int(round(sh.height * sy))
            n += 1
        except Exception:  # never let one bad shape abort the resize
            pass
    return n


def _apply_target_resize(
    prs: Presentation, target_w: int, target_h: int
) -> Dict[str, Any]:
    """US-4.6 coordinate-path prep (architecture-review M1 + M2): resize the
    canvas to the target aspect ratio and proportionally scale every master +
    layout shape so backgrounds, chrome, and placeholder positions all move to
    the new dimensions. The shared render loop then fills target-geometry
    placeholders exactly as on the native path (styling/bullets inherit via
    ``add_slide`` — AC4 satisfied without manual re-application).

    Mutates ``prs`` in memory (the template file is untouched — ``prs`` was
    loaded from it; the result is saved to the output). Returns an
    ``aspect_ratio`` report fragment for render.json.
    """
    native_w, native_h = int(prs.slide_width), int(prs.slide_height)
    sx = target_w / native_w if native_w else 1.0
    sy = target_h / native_h if native_h else 1.0

    prs.slide_width = int(target_w)
    prs.slide_height = int(target_h)

    scaled = 0
    for master in prs.slide_masters:
        scaled += _scale_shapes_geometry(master.shapes, sx, sy)
    for layout in prs.slide_layouts:
        scaled += _scale_shapes_geometry(layout.shapes, sx, sy)

    logger.info(
        "US-4.6 resized %dx%d -> %dx%d (sx=%.4f sy=%.4f); scaled %d master/layout shapes",
        native_w, native_h, target_w, target_h, sx, sy, scaled,
    )
    return {
        "native": [native_w, native_h],
        "target": [int(target_w), int(target_h)],
        "scale": [round(sx, 6), round(sy, 6)],
        "shapes_scaled": scaled,
    }


def _rewrite_output_schema_target(
    output_path: str, target_w: int, target_h: int
) -> None:
    """US-4.6 M4: rewrite the OUTPUT's embedded schema ``slide_dimensions`` to
    the target size so the target-sized output is self-describing / re-usable as
    a target-sized template. Non-fatal: a no-op when the output has no embedded
    schema (e.g. extraction failed). US-4.3 already embeds a schema (freshly
    extracted for a non-templated input — satisfying M3); this post-step only
    rewrites the size fields.
    """
    try:
        schema = read_embedded_schema(output_path)
        if not schema:
            return
        dims = schema.setdefault("template_metadata", {}).setdefault("slide_dimensions", {})
        dims["width_emu"] = int(target_w)
        dims["height_emu"] = int(target_h)
        dims["width_inches"] = round(int(target_w) / _EMU_PER_INCH, 4)
        dims["height_inches"] = round(int(target_h) / _EMU_PER_INCH, 4)
        dims["aspect_ratio"] = compute_ratio(int(target_w), int(target_h))
        embed_schema(output_path, schema, output_path)  # atomic temp + replace
        logger.info("US-4.6 rewrote output schema slide_dimensions -> %s", dims["aspect_ratio"])
    except Exception as exc:  # never block a successful render
        logger.debug("US-4.6 output-schema rewrite skipped (%s): %s", output_path, exc)


def generate_ppt_from_data(
    slide_data_list: List[Dict[str, Any]],
    template_path: Optional[str] = None,
    output_path: str = "output.pptx",
    prompt_text: str = "",
    validate: bool = True,
    strict: bool = False,
    cleanup_temp: bool = True,
    resolve_placeholders: bool = True,
    default_closing: bool = True,
    config_overrides: Optional[Dict[str, str]] = None,
    auto_template: bool = True,
    target_size: Optional[Any] = None,
) -> str:
    # US-4.6 (AC5): a target aspect ratio. ``None`` (or a ratio equal to the
    # template's native ratio) -> the default US-4.1 ``add_slide(layout)``
    # native path (no-op for this story). A different ratio -> the coordinate-
    # path PREP (_apply_target_resize). ``target_size`` accepts a preset
    # ("16:9"/"4:3"/"1:1") or an explicit {width_in,height_in}/{width_emu,
    # height_emu} dict (see geometry.resolve_target_size).
    # #37: resolve resource placeholders (data_query) into concrete assets
    # BEFORE validation, so the validator sees materialized data.
    # Graceful no-op when resolver.config.json is absent.
    if resolve_placeholders:
        slide_data_list = resolve_slide_data_list(slide_data_list)

    # #40 / US-6: guarantee a Thank-You closing slide on decks of N >= 3.
    if default_closing:
        slide_data_list = _ensure_default_closing(slide_data_list)

    # Phase 1 Track A: defensive validation. Catches malformed input with a
    # clear ValidationError instead of a cryptic crash in the render loop.
    if validate:
        result = validate_slide_data_list(slide_data_list, strict=strict)
        for msg in result.warning_messages():
            logger.warning("Validation: %s", msg)
        # Strict mode (agent pre-flight gate): any schema error is fatal.
        if strict and not result.is_valid:
            raise ValidationError(result.errors)
        # Non-strict mode: surface per-slide schema errors as warnings for
        # visibility, but keep the engine's existing graceful-degradation
        # behaviour (skip slide / default chart / skip chart). Only abort on
        # unrecoverable top-level structure, which would otherwise crash the
        # render loop cryptically.
        if not strict:
            for msg in result.error_messages():
                logger.warning("Validation (degraded): %s", msg)
        if not isinstance(slide_data_list, list):
            raise ValidationError(
                result.errors if result.errors
                else "slide_data_list must be a JSON array"
            )

    # BT-142 Phase 2.3: template_path is required (no bundled default).
    # Passing None or "auto" raises TemplateError with an actionable message.
    if template_path and template_path != "auto":
        template = Path(template_path)
    else:
        raise TemplateError(
            "template_path is required — no bundled default template is shipped. "
            "Pass a user-supplied .pptx path to use as the Slide Master template."
        )
    output = Path(output_path)

    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")

    # #46 (P3): state machine ① — discard any derived template_new.pptx left
    # from a prior run so the base template.pptx is re-evaluated fresh each
    # request. Inline (no cross-skill import); the full lifecycle lives in the
    # pptx-template-modifier-skill and is wired by P4 when cloning is implemented.
    _derived = template.with_name(template.stem + "_new" + template.suffix)
    if _derived.exists():
        try:
            _derived.unlink()
            logger.info("Discarded leftover derived template: %s", _derived.name)
        except OSError as exc:  # pragma: no cover - defensive
            logger.debug("Could not delete leftover %s: %s", _derived, exc)

    if not output.is_absolute():
        output = DEFAULT_OUTPUT_DIR / output
    output.parent.mkdir(parents=True, exist_ok=True)

    config = _load_config(str(template))
    # #47 (P4): merge caller-supplied layout overrides (e.g. cloned-layout pins
    # from pptx-template-modifier-skill's resolve_and_clone). Caller overrides win.
    if config_overrides:
        config = {**config, **config_overrides}

    logger.info("Loading template: %s", template.name)
    # US-4.7 (AC3): a corrupt / non-PPTX file must surface a clear TemplateError
    # instead of a raw python-pptx traceback.
    try:
        prs = Presentation(str(template))
    except Exception as exc:
        raise TemplateError(
            f"Could not open template as PPTX ({exc.__class__.__name__}: {exc}): "
            f"{template}"
        ) from exc
    logger.info("Template: %d slides", len(prs.slides))

    # US-4.8 (CRIT-1): if the template has no slide master (Scenario A), repair
    # it BEFORE get_render_contract so the contract describes the repaired prs.
    # The repair copies default.pptx's master skeleton + optionally replaces
    # the theme with salvaged/scavenged user styling (3-level cascade).
    repair_report: Optional[Dict[str, Any]] = None
    try:
        masters_check = list(prs.slide_masters)
    except Exception:
        masters_check = []
    if not masters_check:
        try:
            from master_repairer import repair_if_needed
            # BT-142 Phase 1.5 + 2.3: pass default_template_path=None to use the
            # in-code minimal theme fallback (no bundled default.pptx dependency).
            repair = repair_if_needed(prs, str(template), None)
            if repair.mutated and repair.repaired_path:
                template = Path(repair.repaired_path)
                prs = Presentation(str(template))
                repair_report = {
                    "level": repair.level,
                    "theme_source": repair.theme_source,
                }
                logger.info(
                    "Template repaired (level=%s, theme=%s): %s",
                    repair.level, repair.theme_source, template,
                )
        except Exception as exc:
            logger.warning(
                "Template repair failed (%s); continuing with original template "
                "(_validate_template will raise if the master is still absent)", exc
            )

    logger.info("Template loaded: %d layouts", len(prs.slide_layouts))

    # #43 (P0): auto-introspect the template into a JSON contract before render.
    # US-4.1: prefer the embedded JSON (via the adapter); fall back to the
    # mtime-cached sidecar contract. Non-fatal: on any failure the engine falls
    # back to name-based layout matching (backward compatible). Provenance is
    # tagged on the contract as ``_source`` and logged inside get_render_contract.
    contract = None
    try:
        contract = get_render_contract(str(template))
    except Exception as exc:  # pragma: no cover - defensive; never block render
        logger.warning("Template contract unavailable (%s); using name matching", exc)

    # US-4.7 (AC4/AC5/AC6): severe template problems abort before the render loop.
    # Structural checks (no master / zero layouts) always run; servability (AC6)
    # runs only when the contract is available. See ``_validate_template``.
    _validate_template(prs, str(template), contract)

    # US-4.2: build a {layout_name: {role: size_pt}} map from the embedded
    # schema (best-effort) for the M1 base-size resolution chain. Absent/corrupt
    # → empty → the chain falls back to layout-sample → role ceilings.
    schema_font_map: Dict[str, Dict[str, float]] = {}
    try:
        schema_font_map = _build_schema_font_map(read_embedded_schema(str(template)))
    except Exception as exc:  # pragma: no cover - best-effort
        logger.debug("Schema font map unavailable (%s); using role ceilings", exc)

    render_report: Dict[str, Any] = {"slides": []}

    removed = _remove_all_slides(prs)
    logger.info("Cleared %d example slides", removed)

    exact_idx, norm_idx = _build_layout_index(prs)

    # US-4.6 (AC5): decide native vs coordinate path by ASPECT RATIO (m1).
    # None or same-ratio target -> native path (no-op for this story). A
    # different ratio -> coordinate-path PREP (resize + geometry scaling, M2)
    # applied to ``prs`` in place, then the SAME shared render loop runs. The
    # coordinate path is a prep step, not a second render strategy (M1).
    target_dims: Optional[Tuple[int, int]] = None
    if target_size is not None:
        try:
            target_dims = resolve_target_size(target_size)
        except ValueError as exc:
            raise ValueError(f"Invalid target_size: {exc}") from exc
    coordinate_report: Optional[Dict[str, Any]] = None
    if target_dims is not None and not aspect_ratios_match(
        int(prs.slide_width), int(prs.slide_height), target_dims[0], target_dims[1]
    ):
        coordinate_report = _apply_target_resize(prs, target_dims[0], target_dims[1])

    for page_num, slide_data in enumerate(slide_data_list, start=1):
        slide_type = slide_data.get("slide_type", "")

        # Resolve layout: config pin → fingerprint match → name fallback (#44).
        layout = _select_layout(
            slide_type, contract, config, prs, exact_idx, norm_idx, page_num
        )
        if layout is None:
            continue  # degradation warning already logged

        layout_idx = prs.slide_layouts.index(layout)
        try:
            slide = prs.slides.add_slide(layout)
            logger.info("Page %d: added slide from layout[%d] '%s'", page_num, layout_idx, layout.name)

            slide_report: List[Dict[str, Any]] = []
            title_text = slide_data.get("title", "")

            # Always try to fill title placeholder
            if title_text:
                title_ph = _find_title_placeholder(slide)
                if title_ph:
                    fit = _set_text(
                        title_ph, title_text,
                        role="title", schema_font_map=schema_font_map, layout=layout,
                    )
                    if fit and fit.adjusted:
                        logger.info(
                            "  Title auto-shrunk %g→%gpt", fit.base_size_pt, fit.applied_size_pt
                        )
                    entry = _font_fit_report_entry("title", "title", fit)
                    if entry:
                        slide_report.append(entry)
                    logger.info("  Title: \"%s\"", title_text)

            # Fill subtitle (for title, agenda, closing, section-header-sub).
            # With content -> fill (overrides inheritance). Without content ->
            # REMOVE the placeholder so the layout's inherited sample text (e.g.
            # closing's "Prepared by: Lecturer Name", title's "Click to edit
            # Master subtitle style") cannot bleed in (an empty <a:p/> still
            # inherits — verified empirically).
            if slide_type in _LAYOUTS_WITH_SUBTITLE:
                sub_ph = _find_placeholder(slide, _SUBTITLE_TYPE)
                if sub_ph:
                    if slide_type == "closing_slide":
                        # closing sign-off from presenter fields; fall back to an
                        # explicit subtitle; both empty -> removed below.
                        subtitle_text = _compose_signoff(slide_data) or slide_data.get("subtitle", "")
                    else:
                        subtitle_text = slide_data.get("subtitle", "")
                    if subtitle_text:
                        fit = _set_text(
                            sub_ph, subtitle_text,
                            role="subtitle", schema_font_map=schema_font_map, layout=layout,
                        )
                        if fit and fit.adjusted:
                            logger.info(
                                "  Subtitle auto-shrunk %g→%gpt",
                                fit.base_size_pt, fit.applied_size_pt,
                            )
                        entry = _font_fit_report_entry("subtitle", "subtitle", fit)
                        if entry:
                            slide_report.append(entry)
                        logger.info("  Subtitle: \"%s\"", subtitle_text[:50])
                    else:
                        _remove_placeholder(sub_ph)
                        logger.info("  Subtitle: removed (no content; no inherited-sample bleed)")

            # Fill body text (for content slides)
            if slide_type in _LAYOUTS_WITH_BODY:
                body_text = slide_data.get("body", "")
                if body_text:
                    body_ph = _find_body_placeholder(slide)
                    if body_ph:
                        fit = _set_body_text(
                            body_ph, body_text,
                            role="body", schema_font_map=schema_font_map, layout=layout,
                        )
                        if fit:
                            entry = _font_fit_report_entry("body", "body", fit)
                            if entry:
                                slide_report.append(entry)
                            if fit.adjusted:
                                logger.info(
                                    "  Body auto-shrunk %g→%gpt",
                                    fit.base_size_pt, fit.applied_size_pt,
                                )
                            if not fit.fits:
                                logger.warning(
                                    "  Body still overflows at %gpt floor (AC1 best-effort)",
                                    fit.applied_size_pt,
                                )
                        logger.info("  Body: %d lines", len([l for l in body_text.split("\n") if l.strip()]))

            # Fill two body areas (for two-content slides)
            if slide_type in _LAYOUTS_WITH_TWO_BODIES:
                body_left = slide_data.get("body_left", "")
                body_right = slide_data.get("body_right", "")
                objects = _find_placeholders(slide, _OBJECT_TYPE)
                # BODY placeholders serve the same role as OBJECT (the
                # introspector normalizes BODY→OBJECT in the contract, but
                # _find_placeholders uses the raw python-pptx type). Fall
                # back so BODY-based templates fill correctly.
                if len(objects) < 2:
                    body_phs = _find_placeholders(slide, _BODY_TYPE)
                    if len(body_phs) > len(objects):
                        objects = body_phs
                if len(objects) >= 2:
                    if body_left:
                        fit = _set_body_text(
                            objects[0], body_left,
                            role="body", schema_font_map=schema_font_map, layout=layout,
                        )
                        entry = _font_fit_report_entry("body", "body_left", fit)
                        if entry:
                            slide_report.append(entry)
                        logger.info("  Body-left: %d lines", len([l for l in body_left.split("\n") if l.strip()]))
                    if body_right:
                        fit = _set_body_text(
                            objects[1], body_right,
                            role="body", schema_font_map=schema_font_map, layout=layout,
                        )
                        entry = _font_fit_report_entry("body", "body_right", fit)
                        if entry:
                            slide_report.append(entry)
                        logger.info("  Body-right: %d lines", len([l for l in body_right.split("\n") if l.strip()]))
                elif len(objects) == 1 and (body_left or body_right):
                    fit = _set_body_text(
                        objects[0], body_left or body_right,
                        role="body", schema_font_map=schema_font_map, layout=layout,
                    )
                    entry = _font_fit_report_entry("body", "body_left", fit)
                    if entry:
                        slide_report.append(entry)
                    logger.warning(
                        "  Two-content slide has only 1 content placeholder; "
                        "body_left/body_right merged into one")
                elif body_left or body_right:
                    logger.warning(
                        "  Two-content slide has no content placeholders; "
                        "body_left/body_right dropped")

            # Add chart (for chart slides)
            if slide_type in _LAYOUTS_WITH_CHART:
                _add_chart_to_slide(slide, slide_data)

            # Embed image (any slide carrying image_path) — #18
            if slide_data.get("image_path"):
                _add_image_to_slide(slide, slide_data)

            # Fill speaker notes (must be English; only visible in Presenter View)
            notes_text = slide_data.get("notes", "")
            if _set_notes(slide, notes_text):
                logger.info("  Notes: %d chars", len(notes_text))

            render_report["slides"].append({
                "index": page_num,
                "slide_type": slide_type,
                "placeholders": slide_report,
            })

        except Exception as exc:
            logger.error("Page %d failed: %s", page_num, exc)

    prs.save(str(output))
    logger.info("Saved: %s (%d slides)", output.resolve(), len(prs.slides))

    # US-4.3: ensure the output carries ppt/template_schema.json (python-pptx
    # strips the unmodeled part on save). Schema sourced from the INPUT template
    # (M1); staleness-aware (M2); atomic (m3). Non-fatal. AC2.
    render_report["templating"] = _ensure_output_templated(
        str(output), str(template), auto_template
    )

    # US-4.8 (MAJOR-7): record repair provenance in the render sidecar.
    if repair_report is not None:
        render_report["templating"]["repair"] = repair_report

    # US-4.6 (M4): a target-sized output is self-describing at the TARGET size.
    # Rewrite the output's embedded schema slide_dimensions to target so the
    # deck can be re-used as a target-sized template; record source->target
    # provenance in render.json. Non-fatal.
    if coordinate_report is not None:
        _rewrite_output_schema_target(str(output), target_dims[0], target_dims[1])
        render_report["aspect_ratio"] = coordinate_report

    # US-4.2: write the per-slide font-fit render report sidecar (AC3). Never
    # blocks a successful render.
    _write_render_report(str(output), render_report)

    # Auto-cleanup pipeline temp artifacts (outline checkpoints, agent-written
    # temp JSON) so they never accumulate on disk. Lazy import + try/except keeps
    # cleanup from ever affecting a successful render. Pass cleanup_temp=False to
    # retain them (e.g. while debugging a failed run).
    if cleanup_temp:
        try:
            from outline_store import cleanup_all
            removed = cleanup_all()
            if removed:
                logger.info("Cleaned up %d temp artifact(s)", removed)
        except Exception as exc:  # cleanup must never break a successful render
            logger.debug("Temp cleanup skipped: %s", exc)
    return str(output.resolve())


def _demo_deck() -> List[Dict[str, Any]]:
    return [
        {
            "slide_type": "title_slide",
            "title": "AI Empowering Finance",
            "subtitle": "2026 Q1",
            "notes": (
                "KEY MESSAGE: Open with energy — set the stakes in one line.\n"
                "\"Hold the slide for two seconds before you speak.\"\n"
                "\"Good [morning/afternoon], I'm [Name]. Today I want to show you how AI is already transforming finance — not in theory, but in the numbers.\"\n"
                "Pause. Let the tagline land.\n"
                "\"We'll walk through where it delivers the clearest ROI today.\"\n"
                "TRANSITION: \"Let me start with the core scenarios.\"\n"
                "COACHING: Eye contact, confident. Do not read the slide. Be ready for: \"Is this hype or real?\" — lead with the 80 percent figure."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Core AI Scenarios",
            "body": (
                "**Automated Reporting** \u2014 RPA tools auto-generate monthly reports, cutting manual effort by 80%\n"
                "**Smart Reconciliation** \u2014 AI matches bank transactions at 99.5% accuracy\n"
                "**Fraud Detection** \u2014 Real-time anomaly detection with automated alerts\n"
                "**Tax Optimization** \u2014 ML identifies savings opportunities across tax structures"
            ),
            "notes": (
                "KEY MESSAGE: Four high-impact scenarios where AI already delivers measurable ROI.\n"
                "\"Let's make this concrete. These aren't edge cases — this is everyday finance.\"\n"
                "\"Automated reporting alone removes eighty percent of the manual effort behind every monthly close.\"\n"
                "Pause. Let the number land.\n"
                "\"Smart reconciliation now matches transactions at ninety-nine-point-five percent accuracy, and fraud detection flags anomalies in real time.\"\n"
                "\"Ask your CFO: how much would one missed discrepancy cost?\"\n"
                "TRANSITION: \"Here is how we roll this out.\"\n"
                "COACHING: Matter-of-fact tone, don't over-sell. Be ready for: \"What about false positives?\" — answer: tuned thresholds, human-in-the-loop review."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Roadmap",
            "body": (
                "**Phase 1: Pilot** \u2014 Deploy in 2 business units by Q2\n"
                "**Phase 2: Scale** \u2014 Expand to all departments by Q4\n"
                "**Phase 3: Full Deployment** \u2014 Organization-wide adoption by 2027"
            ),
            "notes": (
                "KEY MESSAGE: A phased, low-risk rollout — pilot, scale, then full adoption.\n"
                "\"We don't boil the ocean. We pilot in two units first, prove the numbers, then scale.\"\n"
                "\"By Q4 every department is on board, and full organisation-wide adoption lands in 2027.\"\n"
                "Walk the three phases left to right.\n"
                "TRANSITION: Open for questions.\n"
                "COACHING: Keep it tight, end with confidence. Be ready for: \"What could delay Phase 2?\" — answer: only change-management, never the technology."
            ),
        },
    ]


def _is_number_literal(s: str) -> bool:
    """True if ``s`` parses as a float (used by the ``--target-size WxH`` guard)."""
    try:
        float(s)
        return True
    except (ValueError, TypeError):
        return False


def main(argv: Optional[List[str]] = None) -> int:
    """CLI entry. Reads a slide-data JSON array from ``--data`` (or runs the demo
    deck) and writes a PPTX. ``--target-size`` triggers the US-4.6 coordinate
    path (preset ``16:9``/``4:3``/``1:1`` or ``WxH`` inches). Exit codes:
    0 success, 1 validation error, 2 runtime error (US-5.1 convention).
    """
    parser = argparse.ArgumentParser(
        prog="ppt_builder",
        description="Generate a PPTX from a slide-data JSON array using a template.",
    )
    parser.add_argument("--template", "-t", help="path to the template .pptx (defaults to the bundled template).")
    parser.add_argument("--output", "-o", default="output.pptx", help="output .pptx path.")
    parser.add_argument("--data", "-d", help="path to a JSON file containing a slide-data array.")
    parser.add_argument(
        "--target-size",
        help="US-4.6 target aspect ratio — a preset (16:9/4:3/1:1) or 'WxH' inches "
             "(e.g. '10x7.5'). Omit to render at the template's native size.",
    )
    parser.add_argument("--log-level", default="info", help="log level (debug/info/warn/error).")
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=getattr(logging, args.log_level.upper(), logging.INFO),
        format="%(levelname)s %(name)s: %(message)s",
    )

    # Resolve target_size: preset string OR "WxH" inches shorthand -> dict.
    # The WxH shorthand requires EXACTLY two numeric segments (guard against
    # "5x3x2" which would otherwise raise an uncaught ValueError on unpack).
    target_size: Any = None
    if args.target_size:
        ts = args.target_size.strip()
        parts = ts.split("x")
        if (
            len(parts) == 2
            and all(_is_number_literal(p.strip()) for p in parts)
        ):
            target_size = {"width_in": float(parts[0].strip()),
                           "height_in": float(parts[1].strip())}
        else:
            target_size = ts  # preset name (validated downstream by resolve_target_size)

    if args.data:
        try:
            slide_data_list = json.loads(Path(args.data).read_text(encoding="utf-8"))
        except Exception as exc:
            print(f"error: cannot read --data {args.data}: {exc}", file=sys.stderr)
            return 1
    else:
        print("No --data given; running the demo deck.", file=sys.stderr)
        slide_data_list = _demo_deck()

    try:
        result = generate_ppt_from_data(
            slide_data_list,
            template_path=args.template,
            output_path=args.output,
            target_size=target_size,
        )
    except (ValueError, TemplateError) as exc:  # validation / template error (bad input)
        print(f"error: {exc}", file=sys.stderr)
        return 1
    except Exception as exc:  # runtime error
        print(f"error: {exc}", file=sys.stderr)
        return 2

    print(result)  # stdout: the output path only (US-5.3 convention)
    return 0


if __name__ == "__main__":
    sys.exit(main())
