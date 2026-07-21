"""BT-142 Phase 3.4.1 — Capability C: Promote Designer Slides to Empty Master.

When a user supplies a PPTX whose Slide Master is **empty** (zero or one
blank layout, no placeholders) and all branding is baked per-shape on each
hand-crafted slide, current routing fails:

  - ``pptx-generate-template-skill`` extracts the schema but reports "stock Office
    theme, single blank layout" — useless for fill
  - ``pptx-generate-slide-skill`` cannot fill against a blank layout
  - ``pptx-template-modifier-skill`` Capability B borrows a layout **from a donor**
    — but here the donor IS the user's deck, whose own layouts are also blank

End-user intent (mimic PowerPoint): open Slide Master view → Add Layout →
design the layout (or right-click an existing slide → "Add to Master"). The
slides become children of the master, inheriting theme/fonts/colors while
exposing editable placeholders.

Capability C reverse-engineers each designer slide's structure into a named
layout on the Slide Master, with proper text/picture/table placeholders
(not loose shapes). Decorative brand shapes (cards, accent bars, dividers)
are preserved as non-placeholder layout shapes so the visual identity
carries over.

Pipeline (per Phase 3.4.1 algorithm):
  1. Cluster shapes by role (headline → TITLE, body text → BODY, image →
     PICTURE, table → TABLE, everything else → decorative)
  2. Allocate placeholder indices (TITLE=0, BODY=1..n, PICTURE/TABLE by
     position)
  3. Promote to layout (clone the slide as a new SlideLayout under the
     master; convert each clustered shape to its placeholder type)
  4. Inherit theme from master (strip per-shape rPr overrides)
  5. Strip source slide (template = 0 slides, N layouts)

Plus: rewrite the Slide Master's theme XML from the source's per-shape
palette + dominant font (not stock Office), then run ``container_check`` on
each promoted layout.

Public API:
  - ``promote_designer_slides(source_pptx, output_path, layout_names=None,
     theme_override=None, **kw) -> PromotionReport``
"""

from __future__ import annotations

import logging
import re
import shutil
import statistics
from collections import Counter, defaultdict
from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlideLayoutPart

# Use lxml.etree throughout (python-pptx's oxml layer is lxml-based; stdlib
# xml.etree.ElementTree elements cannot be mixed with oxml CT_* elements).
from lxml import etree as _lxml
ET = _lxml  # alias so the rest of the module keeps its ET usage but uses lxml

# Local import (pptx-template-modifier-skill/scripts/ is on sys.path at call time)
try:
    from container_check import container_violations, Violation as ContainerViolation
except ImportError:
    container_violations = None  # type: ignore
    ContainerViolation = None  # type: ignore

try:
    from contrast_check import contrast_violations, ContrastViolation
except ImportError:
    contrast_violations = None  # type: ignore
    ContrastViolation = None  # type: ignore

try:
    from vision_extractor import (
        VisionSlideSchema,
        fallback_xml_background,
    )
except ImportError:
    VisionSlideSchema = None  # type: ignore
    def fallback_xml_background(slide, theme=None):  # type: ignore
        # Signature MUST match the real vision_extractor.fallback_xml_background
        # (which takes an optional theme dict). Without the theme kwarg, the
        # call site at _resolve_slide_background would TypeError on import
        # failure — BT-142 review caught this P0 bug.
        return None

logger = logging.getLogger(__name__)


class ContainerFitBlocked(Exception):
    """Raised when ``container_critical_blocks=True`` and a layout has critical overflow.

    Propagates out of ``promote_designer_slides`` (NOT caught by the per-slide
    handler) so the build halts and the user sees a clear failure instead of
    a silently-partial template. BT-142 review P0 fix.
    """


class ContrastBlocked(Exception):
    """Raised when ``contrast_critical_blocks=True`` and unfixed critical contrast remains.

    Same propagation rules as :class:`ContainerFitBlocked`.
    """


# python-pptx PP_PLACEHOLDER type enum integer values
PH_TITLE = 13
PH_BODY = 2
PH_SUBTITLE = 3
PH_OBJECT = 17
PH_PICTURE = 18
PH_TABLE = 19
PH_CHART = 20


@dataclass
class ShapeRole:
    """Classification of one source shape into a layout role."""

    shape_id: int
    name: str
    role: str  # "title" | "body" | "picture" | "table" | "decorative"
    placeholder_idx: Optional[int] = None
    text: str = ""
    font_size_pt: Optional[float] = None
    fill_rgb: Optional[str] = None


@dataclass
class LayoutPlan:
    """Per-source-slide promotion plan."""

    source_slide_index: int
    layout_name: str
    roles: List[ShapeRole] = field(default_factory=list)


@dataclass
class PromotionReport:
    """Result of a ``promote_designer_slides`` call."""

    output_path: str
    layouts_created: int
    placeholders_total: int
    theme_applied: Dict[str, str]  # {"major_font": "...", "accent1": "#...", ...}
    container_violations: Dict[str, List[dict]] = field(default_factory=dict)
    contrast_violations: Dict[str, List[dict]] = field(default_factory=dict)
    skipped_slides: List[int] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "output_path": self.output_path,
            "layouts_created": self.layouts_created,
            "placeholders_total": self.placeholders_total,
            "theme_applied": dict(self.theme_applied),
            "container_violations": {k: list(v) for k, v in self.container_violations.items()},
            "contrast_violations": {k: list(v) for k, v in self.contrast_violations.items()},
            "skipped_slides": list(self.skipped_slides),
            "errors": list(self.errors),
        }


# --------------------------------------------------------------------------
# Stage 1: shape clustering
# --------------------------------------------------------------------------

def _shape_font_size_pt(shape) -> Optional[float]:
    """Return the largest run's font size in pt, or None."""
    if not shape.has_text_frame:
        return None
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            fs = run.font.size
            if fs is not None:
                sizes.append(fs.pt)
    if not sizes:
        return None
    return max(sizes)


def _shape_dominant_fill_rgb(shape) -> Optional[str]:
    """Hex RGB string ("#RRGGBB") of the shape's fill, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return f"#{rgb}"
    except Exception:
        pass
    return None


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _shape_has_table(shape) -> bool:
    return hasattr(shape, "has_table") and shape.has_table


def _shape_has_image(shape) -> bool:
    """Heuristic: True if the shape is a picture or contains a blip."""
    tag = shape._element.tag
    if tag.endswith("}pic"):
        return True
    # AutoShapes with picture fills etc. — keep conservative, return False
    return False


def _classify_shape(shape) -> str:
    """Return one of: 'title', 'body', 'picture', 'table', 'decorative'."""
    if _shape_has_table(shape):
        return "table"
    if _shape_has_image(shape):
        return "picture"
    text = _shape_text(shape)
    if not text:
        return "decorative"
    font_pt = _shape_font_size_pt(shape)
    # Title heuristic: largest text shape, usually at the top
    if font_pt is not None and font_pt >= 28:
        return "title"
    return "body"


def cluster_shapes_by_role(slide) -> List[ShapeRole]:
    """Cluster all shapes on ``slide`` into layout roles."""
    raw: List[Tuple[int, ShapeRole]] = []
    for shp in slide.shapes:
        try:
            shape_id = shp.shape_id
            name = shp.name
        except Exception:
            continue
        role = _classify_shape(shp)
        sr = ShapeRole(
            shape_id=shape_id,
            name=name,
            role=role,
            text=_shape_text(shp),
            font_size_pt=_shape_font_size_pt(shp),
            fill_rgb=_shape_dominant_fill_rgb(shp),
        )
        raw.append((shape_id, sr))
    # Allocate placeholder indices: TITLE=0, BODY=1..n by reading order
    # (top-to-bottom, left-to-right), PICTURE/TABLE by position.
    title_allocated = False
    body_idx = 1
    picture_idx = 10  # leave room for bodies
    table_idx = 20
    for _, sr in raw:
        if sr.role == "title" and not title_allocated:
            sr.placeholder_idx = 0
            title_allocated = True
        elif sr.role == "title":
            # Demote extra "title-like" shapes to body
            sr.role = "body"
            sr.placeholder_idx = body_idx
            body_idx += 1
        elif sr.role == "body":
            sr.placeholder_idx = body_idx
            body_idx += 1
        elif sr.role == "picture":
            sr.placeholder_idx = picture_idx
            picture_idx += 1
        elif sr.role == "table":
            sr.placeholder_idx = table_idx
            table_idx += 1
    return [sr for _, sr in raw]


# --------------------------------------------------------------------------
# Stage: theme extraction (from per-shape branding)
# --------------------------------------------------------------------------

def _hex_to_ooxml_rgb(hex_str: str) -> str:
    """'#2DD4BF' → '2DD4BF'."""
    return hex_str.lstrip("#").upper()


def extract_theme_from_shapes(
    slides: List,
) -> Dict[str, str]:
    """Infer major/minor font + 8 OPC color roles from per-shape branding.

    Returns a dict suitable for ``apply_theme_xml``:
      ``{"major_font": "Century Gothic", "minor_font": "Century Gothic",
         "dk1": "#09090B", "lt1": "#FFFFFF", "dk2": "#18181B",
         "lt2": "#27272A", "accent1": "#2DD4BF", "accent2": "#FB923C",
         "accent3": "#FFFF00", "accent4": "#0D9488", ...}``
    """
    font_counter: Counter = Counter()
    fill_counter: Counter = Counter()
    for slide in slides:
        for shp in slide.shapes:
            fs = _shape_font_size_pt(shp)
            if fs and _shape_text(shp):
                # Use font name from runs
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font.name
                        if font:
                            # Weight by run text length so headings count more
                            font_counter[font] += max(1, len(run.text))
            rgb = _shape_dominant_fill_rgb(shp)
            if rgb:
                fill_counter[rgb] += 1

    major_font = (font_counter.most_common(1)[0][0] if font_counter else "Calibri")
    minor_font = major_font  # designer decks typically use one font for both

    top_fills = [c for c, _ in fill_counter.most_common(8)]
    # Pad with Office defaults if we didn't find 8 distinct fills
    defaults = ["#2DD4BF", "#FB923C", "#FFFF00", "#0D9488", "#09090B", "#FFFFFF", "#18181B", "#27272A"]
    while len(top_fills) < 8:
        top_fills.append(defaults[len(top_fills)])

    # Sort fills heuristically into OPC roles:
    #   darkest → dk1, next → dk2
    #   lightest → lt1, next → lt2
    #   remaining sorted by frequency → accent1..accent4
    def brightness(hex_str: str) -> int:
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return 0
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return r + g + b

    sorted_by_brightness = sorted(top_fills, key=brightness)
    dk1 = sorted_by_brightness[0]
    lt1 = sorted_by_brightness[-1]
    remaining = [c for c in top_fills if c not in (dk1, lt1)]
    # Pick dk2 (next darkest) and lt2 (next lightest)
    sorted_rem = sorted(remaining, key=brightness)
    dk2 = sorted_rem[0] if sorted_rem else dk1
    lt2 = sorted_rem[-1] if sorted_rem else lt1
    accents = [c for c in remaining if c not in (dk2, lt2)][:4]
    while len(accents) < 4:
        accents.append(defaults[4 + len(accents)])

    return {
        "major_font": major_font,
        "minor_font": minor_font,
        "dk1": dk1,
        "lt1": lt1,
        "dk2": dk2,
        "lt2": lt2,
        "accent1": accents[0],
        "accent2": accents[1],
        "accent3": accents[2],
        "accent4": accents[3],
        "accent5": defaults[0],
        "accent6": defaults[1],
        "hlink": defaults[0],
        "folHlink": defaults[1],
    }


# --------------------------------------------------------------------------
# Stage: theme XML rewrite
# --------------------------------------------------------------------------

def apply_theme_xml(prs, theme: Dict[str, str]) -> None:
    """Rewrite the Slide Master's theme XML with the extracted palette/fonts."""
    master = prs.slide_masters[0]
    master_part = master.part
    # The theme is related to the slideMaster via the "theme" relationship.
    theme_part = None
    for rel in master_part.rels.values():
        if "theme" in rel.reltype.lower() and not rel.is_external:
            theme_part = rel.target_part
            break
    if theme_part is None:
        # Fall back to scanning all parts for theme*.xml
        for part in master_part.package.iter_parts():
            if "theme" in str(part.partname) and str(part.partname).endswith(".xml"):
                theme_part = part
                break
    if theme_part is None:
        logger.warning("apply_theme_xml: theme*.xml not found in package")
        return
    _rewrite_theme_part(theme_part, theme)


def _rewrite_theme_part(theme_part, theme: Dict[str, str]) -> None:
    """Edit the theme part's <a:clrScheme> and <a:fontScheme> in place."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # Theme parts are XmlParts — access via _element (lxml tree); python-pptx
    # serializes it on save. If _element is unavailable (plain Part), fall
    # back to blob-parse + reload pattern.
    tree = getattr(theme_part, "_element", None)
    if tree is None:
        from lxml import etree as _lxml
        tree = _lxml.fromstring(theme_part.blob)
        _serialize_back = True
    else:
        _serialize_back = False
    # clrScheme
    clr_scheme = tree.find(".//{{{}}}clrScheme".format(ns_a))
    if clr_scheme is not None:
        role_to_child = {
            "dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2",
            "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
            "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
            "hlink": "hlink", "folHlink": "folHlink",
        }
        for role, tag in role_to_child.items():
            if role not in theme:
                continue
            elem = clr_scheme.find("{{{}}}{}".format(ns_a, tag))
            if elem is None:
                continue
            srgb = elem.find("{{{}}}srgbClr".format(ns_a))
            if srgb is None:
                for child in list(elem):
                    elem.remove(child)
                srgb = ET.SubElement(elem, "{{{}}}srgbClr".format(ns_a))
            srgb.set("val", _hex_to_ooxml_rgb(theme[role]))
    # fontScheme
    font_scheme = tree.find(".//{{{}}}fontScheme".format(ns_a))
    if font_scheme is not None:
        major = font_scheme.find("{{{}}}majorFont".format(ns_a))
        minor = font_scheme.find("{{{}}}minorFont".format(ns_a))
        for parent, font_name in ((major, theme.get("major_font")),
                                   (minor, theme.get("minor_font"))):
            if parent is None or not font_name:
                continue
            latin = parent.find("{{{}}}latin".format(ns_a))
            if latin is not None:
                latin.set("typeface", font_name)
    # If we parsed from blob (plain Part), write back via a fresh blob.
    if _serialize_back:
        from lxml import etree as _lxml
        new_xml = _lxml.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
        # Plain Part has no public blob setter; rely on the load_blob / private
        # _blob attribute if available. Skip silently if neither works.
        for attr in ("_blob", "blob"):
            if hasattr(theme_part, attr) and isinstance(getattr(theme_part, attr), (bytes, bytearray)):
                try:
                    setattr(theme_part, attr, bytes(new_xml))
                    return
                except (AttributeError, TypeError):
                    continue
        logger.warning("apply_theme_xml: cannot write theme part back (no blob setter)")


# --------------------------------------------------------------------------
# Stage: layout creation (clone slide → layout, convert shapes → placeholders)
# --------------------------------------------------------------------------

def _make_placeholder_xml(shape_element, ph_type: str, idx: int) -> None:
    """Mutate ``shape_element`` in place: convert to a placeholder.

    Replaces ``<p:nvSpPr>/<p:nvPr>`` content with ``<p:ph type="..." idx="..."/>``.
    The shape body (geometry, fill, text-frame) is preserved.
    """
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nv_sp_pr = shape_element.find(f"{{{ns_p}}}nvSpPr")
    if nv_sp_pr is None:
        return
    nv_pr = nv_sp_pr.find(f"{{{ns_p}}}nvPr")
    if nv_pr is None:
        return
    # Remove existing placeholder children
    for child in list(nv_pr):
        nv_pr.remove(child)
    ph = ET.SubElement(nv_pr, f"{{{ns_p}}}ph")
    if ph_type:
        ph.set("type", ph_type)
    ph.set("idx", str(idx))


def _inject_layout_background(layout_element, hex_color: str) -> bool:
    """Inject ``<p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val='HEX'/>``.

    Works on any ``<p:cSld>``-bearing element — Slide Master OR Slide Layout.
    Replaces any existing ``<p:bg>`` (including ``<p:bgRef>`` style references
    like ``idx="1001"`` which PowerPoint uses for the master default).
    Returns True on success.
    """
    if not hex_color:
        return False
    h = hex_color.lstrip("#").upper()
    if len(h) != 6:
        return False
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = layout_element.find(f"{{{ns_p}}}cSld")
    if cSld is None:
        return False
    # Drop existing <p:bg> (covers both <p:bgPr> and <p:bgRef> forms)
    existing_bg = cSld.find(f"{{{ns_p}}}bg")
    if existing_bg is not None:
        cSld.remove(existing_bg)
    bg_xml = (
        f'<p:bg xmlns:p="{ns_p}" xmlns:a="{ns_a}">'
        f'<p:bgPr><a:solidFill><a:srgbClr val="{h}"/></a:solidFill>'
        f'<a:effectLst/></p:bgPr></p:bg>'
    )
    try:
        bg_elem = parse_xml(bg_xml)
        # <p:bg> must be the first child of <p:cSld> (ECMA-376 schema order)
        cSld.insert(0, bg_elem)
        return True
    except Exception as exc:
        logger.warning("_inject_layout_background failed: %s", exc)
        return False


def _inject_master_background(prs, hex_color: str) -> bool:
    """Inject ``<p:bg>`` on the Slide Master element (the layouts' parent).

    Without this, the master's default ``<p:bgRef idx="1001"><a:schemeClr
    val="bg1"/></p:bgRef>`` resolves to the theme's ``lt1`` (white in a
    dark-mode deck) — so the master thumbnail in PowerPoint's Slide Master
    view shows white, and any layout that doesn't override ``<p:bg>``
    inherits white. Setting the master bg to the deck's dominant color
    makes the 'base' feel consistent with the brand.

    Returns True on success.
    """
    if not prs.slide_masters:
        return False
    master_element = prs.slide_masters[0]._element
    return _inject_layout_background(master_element, hex_color)


def _compute_dominant_master_bg(
    slides: List[Any],
    theme: Dict[str, str],
) -> str:
    """Pick the most common slide background color across ``slides``.

    Used as the Slide Master background. Falls back to theme dk1 (designer
    decks are typically dark-mode) when no slide yields a resolvable bg.
    """
    from collections import Counter
    counts: Counter = Counter()
    for slide in slides:
        bg_hex, _ = _resolve_slide_background(slide, None, theme)
        if bg_hex:
            counts[bg_hex] += 1
    if counts:
        return counts.most_common(1)[0][0]
    return theme.get("dk1", "#FFFFFF")


def _resolve_slide_background(
    slide,
    vision_schema: Optional[Any],
    theme: Dict[str, str],
) -> Tuple[str, str]:
    """Return ``(bg_hex, source)`` for a slide.

    Priority:
      1. Vision-derived dominant_bg_hex (if confidence ≥ 0.5)
      2. XML fallback (largest covering shape's fill)
      3. Theme dk1 (last resort — usually matches the designer intent)
    """
    if (
        vision_schema is not None
        and getattr(vision_schema, "confidence", 0) >= 0.5
        and getattr(vision_schema, "dominant_bg_hex", "")
    ):
        return vision_schema.dominant_bg_hex, "vision"
    xml_bg = fallback_xml_background(slide, theme=theme)
    if xml_bg:
        return xml_bg, "xml"
    # Last resort: theme dk1 (designer decks are typically dark-mode)
    return theme.get("dk1", "#FFFFFF"), "theme"


_EMPTY_LAYOUT_TEMPLATE = (
    '<p:sldLayout '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'type="title" preserve="1" name="{name}">'
    '<p:cSld>'
    '<p:spTree>'
    '<p:nvGrpSpPr>'
    '<p:cNvPr id="1" name=""/>'
    '<p:cNvGrpSpPr/>'
    '<p:nvPr/>'
    '</p:nvGrpSpPr>'
    '<p:grpSpPr>'
    '<a:xfrm>'
    '<a:off x="0" y="0"/>'
    '<a:ext cx="0" cy="0"/>'
    '<a:chOff x="0" y="0"/>'
    '<a:chExt cx="0" cy="0"/>'
    '</a:xfrm>'
    '</p:grpSpPr>'
    '</p:spTree>'
    '</p:cSld>'
    '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
    '</p:sldLayout>'
)


def _promote_slide_to_layout(
    slide,
    master,
    layout_name: str,
    plan: LayoutPlan,
    bg_hex: str = "",
) -> Tuple[Any, List[ShapeRole]]:
    """Promote ``slide`` to a named layout on ``master`` per ``plan``.

    Returns ``(layout, roles_promoted)``. Uses the SlideLayoutPart pattern
    from ``layout_creator._clone_layout_into`` (proven on real PPTX files).

    Args:
        bg_hex: optional background color ("#RRGGBB") to inject as the
            layout's ``<p:cSld><p:bg>`` — typically from vision extraction
            or XML fallback. Empty string = leave background unset.
    """
    from copy import deepcopy

    master_part = master.part
    master_element = master._element
    package = master_part.package

    # 1. Build a new empty layout XML element with the target name.
    #    Use python-pptx's parse_xml so the result is an oxml/lxml element
    #    compatible with shape._element (CT_Shape) clones.
    new_element = parse_xml(_EMPTY_LAYOUT_TEMPLATE.format(name=layout_name))

    # 1b. Inject layout background (vision-derived or XML fallback).
    if bg_hex:
        _inject_layout_background(new_element, bg_hex)

    # 2. Copy each shape from the source slide into the new layout's spTree,
    #    converting clustered shapes to placeholders in-place.
    spTree = new_element.find(qn("p:cSld")).find(qn("p:spTree"))
    promoted: List[ShapeRole] = []
    role_by_shape_id = {r.shape_id: r for r in plan.roles}
    for shp in slide.shapes:
        try:
            sid = shp.shape_id
        except Exception:
            continue
        sr = role_by_shape_id.get(sid)
        new_elem = deepcopy(shp._element)
        if (
            sr
            and sr.role in ("title", "body", "picture", "table")
            and sr.placeholder_idx is not None
        ):
            type_map = {"title": "title", "body": "body", "picture": "pic", "table": "tbl"}
            _make_placeholder_xml(new_elem, type_map[sr.role], sr.placeholder_idx)
            promoted.append(sr)
        spTree.append(new_elem)

    # 3. Wrap in a new SlideLayoutPart with a fresh partname.
    new_partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
    new_part = SlideLayoutPart(new_partname, CT.PML_SLIDE_LAYOUT, package, new_element)

    # 4. Master → layout relationship.
    rId = master_part.relate_to(new_part, RT.SLIDE_LAYOUT)

    # 5. Register <p:sldLayoutId> in the master's list (so the layout shows
    #    up in the Slide Master's layout gallery).
    sld_layout_id_lst = master_element.get_or_add_sldLayoutIdLst()
    entry = sld_layout_id_lst._add_sldLayoutId(rId=rId)
    entry.set("id", str(_max_layout_id(sld_layout_id_lst)))

    # 6. Layout → master relationship (the layout's only structural rel).
    new_part.relate_to(master_part, RT.SLIDE_MASTER)

    # 7. Reload via python-pptx to expose the new layout as a SlideLayout
    #    object for container_check / inspection. We can't construct one
    #    directly from the part — but master.slide_layouts will now include
    #    it after the part registration.
    layouts = master.slide_layouts
    new_layout = None
    for layout in layouts:
        # Match by the part we just created (the SlideLayoutList is cached;
        # the last layout after registration is the one we just added).
        if layout.part is new_part:
            new_layout = layout
            break
    if new_layout is None:
        new_layout = list(layouts)[-1]
    # Explicitly pin the name (the parse_xml path sometimes loses it on save).
    try:
        new_layout.name = layout_name
    except Exception:
        pass
    return new_layout, promoted


def _max_layout_id(sld_layout_id_lst) -> int:
    """Next unique <p:sldLayoutId id> (ECMA-376 min 2147483648)."""
    max_id = 2147483647
    for entry in sld_layout_id_lst.sldLayoutId_lst:
        try:
            id_val = int(entry.get("id", "0"))
            if id_val > max_id:
                max_id = id_val
        except (TypeError, ValueError):
            continue
    return max_id + 1


# --------------------------------------------------------------------------
# Stage: strip source slides
# --------------------------------------------------------------------------

def _strip_slides(prs) -> int:
    """Remove every slide from ``prs`` (template ≠ deck). Returns count removed."""
    # python-pptx has no public drop_slides; manipulate the sldIdLst.
    sldIdLst = prs.slides._sldIdLst
    count = len(list(sldIdLst))
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    return count


# --------------------------------------------------------------------------
# Public API
# --------------------------------------------------------------------------

def promote_designer_slides(
    source_pptx: str,
    output_path: str,
    layout_names: Optional[Dict[int, str]] = None,
    theme_override: Optional[Dict[str, str]] = None,
    vision_results: Optional[List[Any]] = None,
    run_container_check: bool = True,
    run_contrast_check: bool = True,
    auto_fix_contrast: bool = True,
    container_critical_blocks: bool = True,
    contrast_critical_blocks: bool = False,
) -> PromotionReport:
    """Promote each slide in ``source_pptx`` into a named layout on its master.

    Args:
        source_pptx: input deck with empty master + N designed slides
        output_path: where to write the templated PPTX
        layout_names: optional ``{slide_index: "Layout Name"}`` map; defaults
            to ``"Slide <i+1>"``
        theme_override: optional theme dict (skips per-shape inference);
            see ``extract_theme_from_shapes`` for the shape
        vision_results: optional list of ``VisionSlideSchema`` (one per
            source slide, ordered by index) from ``vision_extractor``.
            When present and confidence ≥ 0.5, the schema's
            ``dominant_bg_hex`` overrides the XML-derived background for
            each promoted layout. Falls back to XML extraction otherwise.
        run_container_check: when True, runs ``container_check`` on each
            promoted layout and includes violations in the report
        run_contrast_check: when True, runs ``contrast_check`` on each
            promoted layout (WCAG 2.1 contrast verification)
        auto_fix_contrast: when True (and run_contrast_check is True),
            override each low-contrast placeholder's default text color to
            a high-contrast alternative; the violation is still reported
            with ``auto_fixed=True`` so the orchestrator knows what changed
        container_critical_blocks: when True (and run_container_check is
            True), raise on any critical violation; otherwise log + report
        contrast_critical_blocks: when True (and run_contrast_check is True),
            raise on any critical contrast violation; otherwise log + report
            (auto-fix is preferred over blocking)

    Returns:
        :class:`PromotionReport`
    """
    layout_names = layout_names or {}
    prs = Presentation(source_pptx)
    if not prs.slides:
        raise ValueError("source deck has no slides to promote")
    master = prs.slide_masters[0]

    # Stage A: extract + apply theme
    theme = theme_override or extract_theme_from_shapes(list(prs.slides))
    try:
        apply_theme_xml(prs, theme)
    except Exception as exc:
        logger.error("theme rewrite failed: %s", exc)

    # Stage A.5: inject Slide Master background
    # The master's default <p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef>
    # resolves to theme lt1 (white in a dark-mode deck) — making the master
    # thumbnail in PowerPoint look off-brand. Override with the dominant
    # slide bg so the 'base' matches the deck. Vision takes precedence here
    # too: if vision_results were provided, we already have authoritative bg
    # hex per slide; we just tally them.
    try:
        master_bg = _compute_dominant_master_bg(list(prs.slides), theme)
        _inject_master_background(prs, master_bg)
        logger.info("master background injected: %s", master_bg)
    except Exception as exc:
        logger.warning("master background injection failed: %s", exc)

    # Stage A.6: inject the same bg on PRE-EXISTING layouts + notes master.
    # Source decks typically ship with a "DEFAULT" blank layout whose
    # <p:bgRef idx="1001"> resolves to white even after the master is darkened
    # (the layout's own bg wins over inheritance). Without this pass, the
    # first thumbnail in PowerPoint's layouts gallery shows white. The notes
    # master has the same default and shows up in Presenter View / Notes
    # Page view. We only inject on layouts that DON'T already carry a solid
    # <p:bgPr> (promoted layouts already have one from Stage B below).
    try:
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for layout in prs.slide_layouts:
            bg = layout._element.find(f"{{{ns_p}}}cSld").find(f"{{{ns_p}}}bg")
            if bg is None:
                # No <p:bg> at all → will inherit master, OK
                continue
            bg_pr = bg.find(f"{{{ns_p}}}bgPr")
            if bg_pr is not None:
                # Already has an explicit solid fill — skip
                continue
            # Has <p:bgRef> (default style) → replace with solid master bg
            _inject_layout_background(layout._element, master_bg)
            logger.info("replaced <p:bgRef> on layout %r with solid %s",
                        getattr(layout, "name", "?"), master_bg)
        # Notes master: same treatment
        try:
            nm = prs.notes_master
            nm_bg = nm._element.find(f"{{{ns_p}}}cSld").find(f"{{{ns_p}}}bg")
            if nm_bg is not None:
                nm_bg_pr = nm_bg.find(f"{{{ns_p}}}bgPr")
                if nm_bg_pr is None:
                    _inject_layout_background(nm._element, master_bg)
                    logger.info("replaced <p:bgRef> on notes master with solid %s", master_bg)
        except Exception as exc:
            logger.warning("notes master bg injection skipped: %s", exc)
    except Exception as exc:
        logger.warning("pre-existing layout bg injection failed: %s", exc)

    # Stage B: cluster + promote each slide (with vision-derived bg if available)
    placeholders_total = 0
    skipped: List[int] = []
    errors: List[str] = []
    layout_container_violations: Dict[str, List[dict]] = {}
    layout_contrast_violations: Dict[str, List[dict]] = {}
    bg_sources: Dict[str, str] = {}  # for diagnostics: layout_name -> "vision"|"xml"|"theme"
    for i, slide in enumerate(prs.slides):
        layout_name = layout_names.get(i, f"Slide {i + 1}")
        roles = cluster_shapes_by_role(slide)
        if not any(r.role in ("title", "body", "picture", "table") for r in roles):
            skipped.append(i)
            continue
        plan = LayoutPlan(source_slide_index=i, layout_name=layout_name, roles=roles)
        # Resolve background (vision > XML fallback > theme dk1)
        vision_schema = (
            vision_results[i]
            if vision_results is not None and i < len(vision_results)
            else None
        )
        bg_hex, bg_source = _resolve_slide_background(slide, vision_schema, theme)
        bg_sources[layout_name] = bg_source
        try:
            layout, promoted = _promote_slide_to_layout(
                slide, master, layout_name, plan, bg_hex=bg_hex,
            )
            placeholders_total += sum(1 for r in promoted if r.placeholder_idx is not None)
            # Container-fit check (geometry)
            if run_container_check and container_violations is not None:
                vs = container_violations(layout)
                if vs:
                    layout_container_violations[layout_name] = [v.to_dict() for v in vs]
                    critical = [v for v in vs if v.severity == "critical"]
                    if critical and container_critical_blocks:
                        # Raise OUTSIDE the per-slide try/except below so it
                        # actually halts the build. Using a dedicated exception
                        # type so callers can distinguish "fatal policy
                        # violation" from "ordinary slide failure". The
                        # per-slide handler catches generic slide-construction
                        # errors (XML clones, shape iteration); policy-gate
                        # violations must propagate. BT-142 review P0 fix.
                        raise ContainerFitBlocked(
                            f"critical container-fit violations on layout '{layout_name}': "
                            f"{[v.detail for v in critical]}"
                        )
            # Contrast check (WCAG 2.1) — runs AFTER container check so the
            # resolved container fill is current. Auto-fix by default.
            if run_contrast_check and contrast_violations is not None:
                cvs = contrast_violations(layout, theme=theme, auto_fix=auto_fix_contrast)
                if cvs:
                    layout_contrast_violations[layout_name] = [v.to_dict() for v in cvs]
                    critical = [v for v in cvs if v.severity == "critical" and not v.auto_fixed]
                    if critical and contrast_critical_blocks:
                        raise ContrastBlocked(
                            f"critical contrast violations on layout '{layout_name}': "
                            f"{[v.detail for v in critical]}"
                        )
        except (ContainerFitBlocked, ContrastBlocked):
            # Policy-gate violations propagate up — do NOT swallow them in
            # the generic per-slide handler. The build must halt so the user
            # sees the failure instead of getting a partial template.
            raise
        except Exception as exc:
            errors.append(f"slide {i} ('{layout_name}'): {exc}")
            skipped.append(i)

    # Stage C: strip source slides
    try:
        _strip_slides(prs)
    except Exception as exc:
        errors.append(f"failed to strip source slides: {exc}")

    # Stage D: save
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    report = PromotionReport(
        output_path=str(Path(output_path).resolve()),
        layouts_created=len(prs.slide_layouts),
        placeholders_total=placeholders_total,
        theme_applied=theme,
        container_violations=layout_container_violations,
        contrast_violations=layout_contrast_violations,
        skipped_slides=skipped,
        errors=errors,
    )
    # Stash the per-layout bg source as a diagnostic (attached via to_dict's
    # extra fields if callers want it). Kept off the dataclass to avoid
    # breaking the typed contract.
    report._bg_sources = bg_sources  # type: ignore[attr-defined]
    return report


__all__ = [
    "ShapeRole",
    "LayoutPlan",
    "PromotionReport",
    "cluster_shapes_by_role",
    "extract_theme_from_shapes",
    "apply_theme_xml",
    "promote_designer_slides",
]
