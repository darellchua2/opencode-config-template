"""Regression tests for BT-142 architecture-review P0 bugs."""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from container_check import _is_text_placeholder
from designer_promoter import (
    ContainerFitBlocked,
    ContrastBlocked,
    fallback_xml_background,
)


# ---------------------------------------------------------------------------
# BUG-1: container_check._is_text_placeholder had `or True` defeating the
#         type filter — all placeholders were treated as text candidates.
# ---------------------------------------------------------------------------

class _FakePhFmt:
    def __init__(self, type_id):
        # Mimic PP_PLACEHOLDER enum: str() yields "TITLE (13)" etc.
        class _E:
            def __init__(self, name, val):
                self.name = name
                self.value = val
            def __str__(self):
                return f"PP_PLACEHOLDER.{self.name} ({self.value})"
        names = {13: "TITLE", 2: "BODY", 17: "OBJECT", 3: "SUBTITLE",
                 18: "PICTURE", 19: "TABLE", 20: "CHART"}
        self.type = _E(names.get(type_id, "OTHER"), type_id) if type_id else None


class _FakeShape:
    def __init__(self, type_id, is_placeholder=True):
        self.is_placeholder = is_placeholder
        self.placeholder_format = _FakePhFmt(type_id) if is_placeholder else None


def test_is_text_placeholder_returns_true_for_title():
    assert _is_text_placeholder(_FakeShape(13)) is True


def test_is_text_placeholder_returns_true_for_body():
    assert _is_text_placeholder(_FakeShape(2)) is True


def test_is_text_placeholder_returns_true_for_subtitle():
    assert _is_text_placeholder(_FakeShape(3)) is True


def test_is_text_placeholder_returns_true_for_object():
    assert _is_text_placeholder(_FakeShape(17)) is True


def test_is_text_placeholder_returns_false_for_picture():
    """BUG-1 regression: previously returned True due to trailing `or True`."""
    assert _is_text_placeholder(_FakeShape(18)) is False


def test_is_text_placeholder_returns_false_for_table():
    assert _is_text_placeholder(_FakeShape(19)) is False


def test_is_text_placeholder_returns_false_for_chart():
    assert _is_text_placeholder(_FakeShape(20)) is False


def test_is_text_placeholder_returns_false_for_non_placeholder():
    assert _is_text_placeholder(_FakeShape(13, is_placeholder=False)) is False


def test_is_text_placeholder_handles_none_type_as_text():
    """Generic 'OBJECT' placeholders (type=None) carry text in designer templates."""
    assert _is_text_placeholder(_FakeShape(None)) is True


# ---------------------------------------------------------------------------
# BUG-3: designer_promoter's ImportError fallback for fallback_xml_background
#         had signature (slide) but caller passed (slide, theme=...).
#         This would TypeError whenever the import failed.
# ---------------------------------------------------------------------------

def test_fallback_xml_background_importerror_signature_accepts_theme():
    """The fallback signature MUST accept the optional theme kwarg.

    Calls _resolve_slide_background's expectation: fallback_xml_background(slide, theme=theme).
    """
    import inspect
    sig = inspect.signature(fallback_xml_background)
    params = list(sig.parameters.keys())
    assert "slide" in params
    # 'theme' may be positional-or-keyword; just verify it accepts the kwarg
    # without TypeError. Calling with a None slide is safe (returns None).
    result = fallback_xml_background(None, theme={"dk1": "#000000"})
    assert result is None


def test_fallback_xml_background_no_kwarg_still_works():
    """Backward-compat: bare call without theme still succeeds."""
    result = fallback_xml_background(None)
    assert result is None


# ---------------------------------------------------------------------------
# BUG-2: container_critical_blocks / contrast_critical_blocks were caught
#         by the per-slide try/except, so they never halted the build.
# ---------------------------------------------------------------------------

def test_container_fit_blocked_exception_class_exists():
    """ContainerFitBlocked is importable and is an Exception subclass."""
    assert issubclass(ContainerFitBlocked, Exception)


def test_contrast_blocked_exception_class_exists():
    assert issubclass(ContrastBlocked, Exception)


def test_container_fit_blocked_propagates_through_promote(tmp_path):
    """End-to-end: with container_critical_blocks=True, critical overflow halts.

    Build a deck with a deliberately overflowing placeholder; verify the
    build aborts with ContainerFitBlocked (not swallowed).
    """
    from designer_promoter import promote_designer_slides
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    # Full-coverage dark shape so the slide passes the "has designed content" gate
    bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor.from_string("09090B")
    # Small card 2x1in at (1,1) -> bbox (1,1)-(3,2)
    card = s.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string("2DD4BF")
    # Body text wider than card AND centered on it (center 2.5,1.5 inside card)
    # but extending 0.5in past left and 1.5in past right -> >20px critical overflow
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(4), Inches(0.9))
    tb.text_frame.text = "Body text that overflows the small card on both sides"
    # Explicit body-sized font so cluster_shapes_by_role classifies as body (not title)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(14)
    src_path = tmp_path / "src.pptx"
    src.save(str(src_path))
    out_path = tmp_path / "out.pptx"
    try:
        promote_designer_slides(
            source_pptx=str(src_path),
            output_path=str(out_path),
            layout_names={0: "L1"},
            run_container_check=True,
            run_contrast_check=False,  # isolate container path
            container_critical_blocks=True,
        )
        # Should have raised; if we get here, the bug is back.
        assert False, "ContainerFitBlocked did NOT propagate — BUG-2 regressed"
    except ContainerFitBlocked:
        pass  # expected
    except Exception as e:
        assert False, f"Wrong exception type propagated: {type(e).__name__}: {e}"


def test_container_fit_non_blocking_reports_but_succeeds(tmp_path):
    """With container_critical_blocks=False, violations are reported but build completes."""
    from designer_promoter import promote_designer_slides
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor.from_string("09090B")
    card = s.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string("2DD4BF")
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(4), Inches(0.9))
    tb.text_frame.text = "Overflowing body"
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(14)
    src_path = tmp_path / "src.pptx"
    src.save(str(src_path))
    out_path = tmp_path / "out.pptx"
    # With blocking OFF, build completes and report includes the violation
    report = promote_designer_slides(
        source_pptx=str(src_path),
        output_path=str(out_path),
        layout_names={0: "L1"},
        run_container_check=True,
        run_contrast_check=False,
        container_critical_blocks=False,
    )
    assert report.errors == []
    # Container violation should be reported
    assert "L1" in report.container_violations
