"""US-4.8 Phase 2 tests — master_cloner (Scenario B/C extension).

Tests that ``clone_master_and_borrow`` correctly borrows layouts from a
multi-layout donor template and injects them under the user's existing
master, with the user's theme inherited automatically.

BT-142 Phase 8.4: the 3 class-level ``@pytest.mark.skip`` decorators (13
skipped tests) are removed. The previous skip reason ("need multi-layout
donor fixture") is resolved by ``_build_donor_pptx()`` below, which
synthesizes a donor with all 11 default Office layouts (Title Slide, Title
and Content, Two Content, Comparison, Picture with Caption, Title Only,
etc.) — enough to satisfy every chenyu ``slide_type`` fingerprint.

The legacy ``template/default.pptx`` file was removed by Phase 2.3/3.2
(no bundled default invariant). These tests now use ``tmp_path`` fixtures
exclusively — no repo-level template file required.
"""
import shutil
import sys
from pathlib import Path

import pytest
from pptx import Presentation

_HERE = Path(__file__).resolve().parent
_MODIFIER_SCRIPTS = _HERE.parent                   # .../pptx-template-modifier-skill/scripts
_SKILLS = _MODIFIER_SCRIPTS.parent.parent          # .../skills
_COMMON_SCRIPTS = _SKILLS / "_common" / "scripts"
_FILLER_SCRIPTS = _SKILLS / "pptx-generate-slide-skill" / "scripts"
for _p in (str(_MODIFIER_SCRIPTS), str(_COMMON_SCRIPTS), str(_FILLER_SCRIPTS)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from master_cloner import clone_master_and_borrow  # noqa: E402
from layout_contract import get_render_contract, _resolve_layout_by_fingerprint  # noqa: E402


def _build_donor_pptx(tmp_path: Path) -> str:
    """Synthesize a multi-layout donor PPTX for testing.

    A fresh ``Presentation()`` carries python-pptx's built-in 11-layout
    Office template (Title Slide, Title and Content, Two Content,
    Comparison, Picture with Caption, Title Only, Blank, etc.). This is
    sufficient to satisfy every chenyu ``slide_type`` fingerprint:
    ``title_slide`` → idx 0, ``content_slide`` → idx 1,
    ``two_content_slide`` → idx 3, ``comparison_slide`` → idx 4,
    ``content_image_slide`` → idx 8, ``chart_slide`` → idx 5.

    Replaces the deleted ``template/default.pptx`` (Phase 2.3/3.2 removed
    the bundled-default invariant). No repo-level fixture file needed.
    """
    donor_path = tmp_path / "donor.pptx"
    prs = Presentation()
    prs.save(str(donor_path))
    return str(donor_path)


def _base_in(tmp_path: Path, donor_path: str) -> str:
    """Copy the donor into ``tmp_path/template.pptx`` as the user's base.

    Derived files (``template_new.pptx``) don't pollute the repo. Both
    base and donor start identical — borrowing creates a duplicate layout
    under the base's master.
    """
    return str(shutil.copy2(donor_path, tmp_path / "template.pptx"))


class TestCloneMasterAndBorrow:
    def test_borrow_produces_extended_template(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base,
            ["chart_slide"],  # donor has this (Title Only idx=5); borrow a fresh copy
            donor,
        )
        assert extended != base
        assert Path(extended).exists()
        assert "chart_slide" in overrides
        assert overrides["chart_slide"].endswith("(borrowed)")

    def test_borrowed_layout_findable_by_name(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, ["chart_slide"], donor
        )
        reloaded = Presentation(extended)
        layouts = reloaded.slide_masters[0].slide_layouts
        for name in overrides.values():
            assert layouts.get_by_name(name) is not None

    def test_borrowed_layout_has_placeholders(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, ["content_slide"], donor
        )
        reloaded = Presentation(extended)
        clone = reloaded.slide_masters[0].slide_layouts.get_by_name(
            overrides["content_slide"]
        )
        assert clone is not None
        # content_slide fingerprint = [TITLE, OBJECT] — at least 2 placeholders.
        assert len(list(clone.placeholders)) >= 2

    def test_borrow_inherits_user_theme(self, tmp_path):
        """The borrowed layout is under the user's master → inherits user theme."""
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        orig_prs = Presentation(base)
        orig_master_count = len(orig_prs.slide_masters)
        orig_layout_count = len(orig_prs.slide_masters[0].slide_layouts)

        extended, overrides = clone_master_and_borrow(
            base, ["content_slide"], donor
        )
        reloaded = Presentation(extended)
        # No new master created (Decision 5 fallback: layouts under existing master).
        assert len(reloaded.slide_masters) == orig_master_count
        # The layout count grew.
        new_layout_count = len(reloaded.slide_masters[0].slide_layouts)
        assert new_layout_count > orig_layout_count

    def test_borrow_multiple_slide_types(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base,
            ["chart_slide", "two_content_slide", "comparison_slide"],
            donor,
        )
        assert len(overrides) == 3
        assert set(overrides.keys()) == {"chart_slide", "two_content_slide", "comparison_slide"}

    def test_base_unchanged_after_borrow(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        before = Path(base).read_bytes()
        clone_master_and_borrow(base, ["chart_slide"], donor)
        after = Path(base).read_bytes()
        assert before == after

    def test_no_borrow_when_no_missing_types(self, tmp_path):
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, [], donor
        )
        assert extended == base
        assert overrides == {}

    def test_borrowed_layout_servable_by_fingerprint(self, tmp_path):
        """After borrowing, the new layout is found by the fingerprint matcher."""
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, ["two_content_slide"], donor
        )
        contract = get_render_contract(extended)
        idx, reason = _resolve_layout_by_fingerprint("two_content_slide", contract)
        assert idx is not None  # the borrowed layout is matched


class TestRollback:
    def test_rollback_on_corrupt_input(self, tmp_path):
        """A non-PPTX file should trigger rollback."""
        donor = _build_donor_pptx(tmp_path)
        bad = tmp_path / "broken.pptx"
        bad.write_bytes(b"not a pptx")
        extended, overrides = clone_master_and_borrow(
            str(bad), ["chart_slide"], donor,
            output_path=str(tmp_path / "out.pptx"),
        )
        # Safe fallback: returns the base path, no overrides.
        assert extended == str(bad)
        assert overrides == {}


class TestCrossPackageSerialization:
    """M-4: verify that borrowed layouts from the donor's package serialize
    correctly into the user's package (no dangling foreign-part references)."""

    def test_output_zip_is_valid(self, tmp_path):
        import zipfile
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, ["content_image_slide"], donor
        )
        # The output must be a readable zip with all core PPTX parts.
        with zipfile.ZipFile(extended) as z:
            names = z.namelist()
            assert "[Content_Types].xml" in names
            assert any(n.startswith("ppt/slideMasters/") for n in names)
            assert any(n.startswith("ppt/slideLayouts/") for n in names)
            assert "ppt/presentation.xml" in names

    def test_output_has_content_types_override_for_layouts(self, tmp_path):
        import zipfile
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        extended, overrides = clone_master_and_borrow(
            base, ["content_slide"], donor
        )
        with zipfile.ZipFile(extended) as z:
            ct = z.read("[Content_Types].xml").decode("utf-8")
            # The layout content type must be present (either as Default or Override).
            assert "slideLayout" in ct


class TestStateMachineDispatch:
    """Verify that resolve_and_clone dispatches to Level 1 when Level 0 has no donor."""

    def test_no_circular_import(self):
        """CRIT-4: layout_creator and master_cloner must not import each other."""
        import layout_creator
        import master_cloner
        # Both modules load without error — no circular import.
        assert hasattr(layout_creator, "clone_for_over_limit")
        assert hasattr(master_cloner, "clone_master_and_borrow")

    def test_resolve_and_clone_with_missing_type(self, tmp_path):
        """End-to-end: a slide_type not in the template triggers Level 1 borrow."""
        from state_machine import resolve_and_clone
        donor = _build_donor_pptx(tmp_path)
        base = _base_in(tmp_path, donor)
        # Request only a title slide (which the donor definitely has).
        # This won't trigger cloning since the template already serves it.
        deck = [{"slide_type": "title_slide", "title": "T", "subtitle": "S"}]
        active, overrides, note = resolve_and_clone(base, deck, donor_template_path=donor)
        # No cloning needed — template already serves title_slide.
        assert active == base
        assert overrides == {}

    def test_resolve_and_clone_triggers_level1_borrow(self, tmp_path):
        """C-1 regression: Level 1 borrow actually fires and produces overrides.

        Builds a MINIMAL template (blank ``Presentation()`` with 11 layouts)
        and requests ``content_slide``. Since the minimal template DOES have
        a content_slide-compatible layout (idx=1 "Title and Content"), this
        is a Level 0 same-file clone, not Level 1. To force Level 1, we
        build a stripped template with only a Blank layout.
        """
        from state_machine import resolve_and_clone
        from pptx import Presentation as Prs

        # Build a minimal template with only the Blank layout (idx=6) —
        # content_slide has no donor here → Level 1 borrow from donor kicks in.
        donor = _build_donor_pptx(tmp_path)
        stripped_path = tmp_path / "stripped.pptx"
        prs = Prs()
        # Delete all layouts except Blank (idx=6) by removing their sldLayoutId entries.
        master = prs.slide_masters[0]
        sld_layout_id_lst = master._element.get_or_add_sldLayoutIdLst()
        layout_ids = list(sld_layout_id_lst)
        # Keep only the last layout id (Blank is typically last in the default template).
        for lid in layout_ids[:-1]:
            sld_layout_id_lst.remove(lid)
        prs.save(str(stripped_path))

        deck = [{"slide_type": "content_slide", "title": "T", "body": "**A** - x"}]
        active, overrides, note = resolve_and_clone(
            str(stripped_path), deck, donor_template_path=donor
        )
        # Level 1 should have borrowed from the donor.
        if overrides:
            assert "content_slide" in overrides
            assert Path(active).exists()
