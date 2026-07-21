"""Tests for dynamic layout cloning (issue #47 / DESIGN §7) — highest-risk phase.

Covers the #47 acceptance criteria:
  * reload-verification — a cloned layout is findable by name after save.
  * rollback — a failure during clone/verify deletes template_new.pptx and
    leaves the base untouched.
  * end-to-end — clone → render against template_new.pptx → output opens.
  * resize — content placeholders are enlarged.
  * existing tests stay green (verified separately by the full suite).
"""
import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from layout_creator import clone_for_over_limit
from state_machine import plan_resolution, resolve_and_clone, template_new_path_for
from template_introspector import introspect


def _base_in(tmp_path, template_path):
    """Copy the real template into tmp_path so derived files don't pollute the repo."""
    base = shutil.copy2(template_path, tmp_path / "template.pptx")
    return str(base)


def _text_heavy_deck():
    """A content slide whose body far exceeds the standard content area."""
    long_body = "\n".join(f"**Point {i}** - lorem ipsum dolor sit amet" for i in range(40))
    return [{"slide_type": "content_slide", "title": "Wall of text", "body": long_body}]


# ============================================================
# clone_for_over_limit
# ============================================================
class TestCloneForOverLimit:
    def test_produces_template_new_and_override(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        plan = plan_resolution(base, _text_heavy_deck(), clone_on="any")
        assert plan.needs_cloning  # precondition

        contract = introspect(base)
        out, overrides = clone_for_over_limit(base, plan, contract)

        assert Path(out).exists()
        assert out != base
        assert "content_slide" in overrides
        clone_name = overrides["content_slide"]
        assert clone_name.endswith("(extended)")

    def test_cloned_layout_findable_by_name(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        contract = introspect(base)
        plan = plan_resolution(base, _text_heavy_deck(), clone_on="any")
        out, overrides = clone_for_over_limit(base, plan, contract)

        reloaded = Presentation(out)
        layouts = reloaded.slide_masters[0].slide_layouts
        for name in overrides.values():
            assert layouts.get_by_name(name) is not None
        # The clone is an ADDITIONAL layout (count grew).
        assert len(layouts) > len(contract["layouts"])

    def test_content_placeholder_resized_larger(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        contract = introspect(base)
        plan = plan_resolution(base, _text_heavy_deck(), clone_on="any")
        out, overrides = clone_for_over_limit(base, plan, contract)

        # Donor content area (the layout the fingerprint selected for content_slide).
        donor_area_in2 = 0.0
        for L in contract["layouts"]:
            if L["name"] == overrides["content_slide"].replace(" (extended)", ""):
                donor_area_in2 = L["content_area_in2"]
        # Guard: if a fresh template changes the donor, this still computes the
        # right baseline rather than a hardcoded in² threshold.
        donor_area_emu2 = donor_area_in2 * 914400 * 914400

        reloaded = Presentation(out)
        clone = reloaded.slide_masters[0].slide_layouts.get_by_name(overrides["content_slide"])
        max_clone_area = 0
        for ph in clone.placeholders:
            if ph.width and ph.height and ph.placeholder_format.idx != 0:
                max_clone_area = max(max_clone_area, ph.width * ph.height)
        # Resize must enlarge the content placeholder beyond the donor's area
        # (template-agnostic — works for 10" and 13.33" slides alike).
        assert max_clone_area > donor_area_emu2

    def test_base_unchanged_after_clone(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        before = Path(base).read_bytes()
        contract = introspect(base)
        plan = plan_resolution(base, _text_heavy_deck(), clone_on="any")
        clone_for_over_limit(base, plan, contract)
        after = Path(base).read_bytes()
        assert before == after  # base is never written


# ============================================================
# Rollback
# ============================================================
class TestRollback:
    def test_verify_failure_deletes_template_new(self, template_path, tmp_path, monkeypatch):
        base = _base_in(tmp_path, template_path)
        contract = introspect(base)
        plan = plan_resolution(base, _text_heavy_deck(), clone_on="any")

        # Sabotage reload-verification so the clone is considered failed.
        def _boom(output_path, names):
            raise RuntimeError("injected verify failure")
        monkeypatch.setattr("layout_creator._verify_layouts", _boom)

        with pytest.raises(RuntimeError):
            clone_for_over_limit(base, plan, contract)

        # The derived file must be rolled back (deleted)…
        assert not template_new_path_for(base).exists()
        # …and the base untouched.
        reloaded = Presentation(base)
        assert len(reloaded.slide_masters[0].slide_layouts) == len(contract["layouts"])


# ============================================================
# resolve_and_clone — the full Capability B loop
# ============================================================
class TestResolveAndClone:
    def test_normal_deck_no_clone(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        deck = [
            {"slide_type": "title_slide", "title": "Deck", "subtitle": "2026"},
            {"slide_type": "content_slide", "title": "OK", "body": "**A** - x"},
            {"slide_type": "closing_slide", "title": "Thank You"},
        ]
        active, overrides, note = resolve_and_clone(base, deck)
        assert active == base
        assert overrides == {}
        assert note is None
        assert not template_new_path_for(base).exists()

    def test_over_limit_produces_clone_and_notification(self, template_path, tmp_path):
        base = _base_in(tmp_path, template_path)
        active, overrides, note = resolve_and_clone(base, _text_heavy_deck(), clone_on="any")
        assert active != base
        assert Path(active).exists()
        assert "content_slide" in overrides
        assert note is not None
        assert "template_new.pptx" in note

    def test_clone_failure_falls_back_to_base(self, template_path, tmp_path, monkeypatch):
        base = _base_in(tmp_path, template_path)
        monkeypatch.setattr(
            "layout_creator._verify_layouts",
            lambda *a, **k: (_ for _ in ()).throw(RuntimeError("boom")),
        )
        active, overrides, note = resolve_and_clone(base, _text_heavy_deck(), clone_on="any")
        # Safe fallback: render against the base, no overrides, no notification.
        assert active == base
        assert overrides == {}
        assert not template_new_path_for(base).exists()


# ============================================================
# End-to-end: clone → render → reopen
# ============================================================
class TestEndToEnd:
    def test_render_against_template_new_opens_cleanly(self, template_path, tmp_path):
        import sys
        base = _base_in(tmp_path, template_path)
        active, overrides, _note = resolve_and_clone(base, _text_heavy_deck(), clone_on="any")

        # Capability A renders against the derived template, pinned to the clone.
        # (pptx-generate-slide-skill/scripts is on sys.path via conftest's _FILLER_SCRIPTS.)
        from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR

        out = generate_ppt_from_data(
            _text_heavy_deck(),
            template_path=active,
            output_path=str(tmp_path / "rendered.pptx"),
            config_overrides=overrides,
            default_closing=False,
        )
        # The output opens cleanly in python-pptx (proxy for PowerPoint-open).
        result = Presentation(out)
        assert len(result.slides) == 1
        assert result.slides[0].shapes.title is not None
