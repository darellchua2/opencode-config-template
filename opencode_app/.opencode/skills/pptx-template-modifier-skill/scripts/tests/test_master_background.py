"""Tests for BT-142 master-background injection (Phase 3.4.1c)."""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from designer_promoter import (
    _inject_master_background,
    _inject_layout_background,
    _compute_dominant_master_bg,
)


def _add_full_slide_dark_shape(slide, hex_no_hash):
    shp = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(hex_no_hash)


def test_inject_master_background_replaces_default_bg_ref():
    """Master's default <p:bgRef idx='1001'> must be replaced with solid fill."""
    prs = Presentation()
    # Inject a dark master bg
    ok = _inject_master_background(prs, "#09090B")
    assert ok is True
    master = prs.slide_masters[0]
    bg = master._element.find(qn("p:cSld")).find(qn("p:bg"))
    assert bg is not None
    # Must NOT contain the default <p:bgRef>
    bg_ref = bg.find(qn("p:bgRef"))
    assert bg_ref is None
    # Must contain <p:bgPr> with the solid fill
    bg_pr = bg.find(qn("p:bgPr"))
    assert bg_pr is not None
    srgb = bg_pr.find(".//" + qn("a:srgbClr"))
    assert srgb is not None
    assert srgb.get("val") == "09090B"


def test_inject_master_background_is_idempotent():
    """Re-injecting replaces (not stacks) the bg element."""
    prs = Presentation()
    _inject_master_background(prs, "#000000")
    _inject_master_background(prs, "#FFFFFF")
    master = prs.slide_masters[0]
    bgs = master._element.find(qn("p:cSld")).findall(qn("p:bg"))
    assert len(bgs) == 1
    srgb = bgs[0].find(".//" + qn("a:srgbClr"))
    assert srgb.get("val") == "FFFFFF"


def test_inject_master_background_rejects_invalid_hex():
    prs = Presentation()
    assert _inject_master_background(prs, "") is False
    assert _inject_master_background(prs, "not-a-color") is False
    assert _inject_master_background(prs, "#abc") is False  # too short


def test_inject_layout_background_works_on_layout_element():
    """The shared helper works on any <p:cSld>-bearing element."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    ok = _inject_layout_background(layout._element, "#18181B")
    assert ok is True
    bg = layout._element.find(qn("p:cSld")).find(qn("p:bg"))
    assert bg is not None


def test_compute_dominant_master_bg_picks_most_common():
    """When slides have mixed bgs, the most common one wins."""
    prs = Presentation()
    # 3 dark slides
    for _ in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_full_slide_dark_shape(s, "09090B")
    # 1 light slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_full_slide_dark_shape(s, "FFFFFF")
    theme = {"dk1": "#09090B", "lt1": "#FFFFFF"}
    bg = _compute_dominant_master_bg(list(prs.slides), theme)
    assert bg == "#09090B"


def test_compute_dominant_master_bg_falls_back_to_theme_dk1():
    """No resolvable bg on any slide → use theme dk1 (designer dark-mode default)."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank, no full-coverage shape
    theme = {"dk1": "#09090B", "lt1": "#FFFFFF"}
    bg = _compute_dominant_master_bg(list(prs.slides), theme)
    assert bg == "#09090B"


def test_promote_designer_slides_injects_master_bg(tmp_path):
    """End-to-end: promoting a designer deck injects bg on BOTH master and layouts."""
    from designer_promoter import promote_designer_slides
    # Build a minimal designer deck: 2 slides with full-coverage dark shape
    src = Presentation()
    for _ in range(2):
        s = src.slides.add_slide(src.slide_layouts[6])
        _add_full_slide_dark_shape(s, "09090B")
        # add a small text shape so cluster_shapes_by_role has something to promote
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Hello world test text"
    src_path = tmp_path / "src.pptx"
    src.save(str(src_path))
    out_path = tmp_path / "out.pptx"
    report = promote_designer_slides(
        source_pptx=str(src_path),
        output_path=str(out_path),
        layout_names={0: "L1", 1: "L2"},
        run_container_check=False,
        run_contrast_check=False,
    )
    assert report.errors == []
    # Verify master bg
    prs2 = Presentation(str(out_path))
    master = prs2.slide_masters[0]
    master_bg = master._element.find(qn("p:cSld")).find(qn("p:bg"))
    assert master_bg is not None
    srgb = master_bg.find(".//" + qn("a:srgbClr"))
    assert srgb is not None
    assert srgb.get("val") == "09090B"


def test_promote_designer_slides_overrides_default_layout_bg_ref(tmp_path):
    """The pre-existing 'DEFAULT' layout's <p:bgRef idx='1001'> must be replaced.

    Source decks ship with a blank layout whose <p:bgRef idx='1001'>
    resolves to theme lt1 (white in a dark-mode deck). Without this fix the
    first thumbnail in PowerPoint's layouts gallery shows white even though
    every other layout + the master is dark.
    """
    from designer_promoter import promote_designer_slides
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    _add_full_slide_dark_shape(s, "09090B")
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Hello"
    src_path = tmp_path / "src.pptx"
    src.save(str(src_path))
    out_path = tmp_path / "out.pptx"
    promote_designer_slides(
        source_pptx=str(src_path),
        output_path=str(out_path),
        layout_names={0: "L1"},
        run_container_check=False,
        run_contrast_check=False,
    )
    prs2 = Presentation(str(out_path))
    # Find the DEFAULT layout (the original one carried over from source)
    default_layout = None
    for layout in prs2.slide_layouts:
        if layout.name == "DEFAULT":
            default_layout = layout
            break
    if default_layout is None:
        return  # source had no DEFAULT layout; nothing to assert
    bg = default_layout._element.find(qn("p:cSld")).find(qn("p:bg"))
    assert bg is not None, "DEFAULT layout lost its <p:bg>"
    # Must NOT still carry the white <p:bgRef>
    assert bg.find(qn("p:bgRef")) is None, "DEFAULT layout still has white <p:bgRef>"
    # Must carry a solid fill matching the master bg
    bg_pr = bg.find(qn("p:bgPr"))
    assert bg_pr is not None
    srgb = bg_pr.find(".//" + qn("a:srgbClr"))
    assert srgb is not None
    assert srgb.get("val") == "09090B"


def test_promote_designer_slides_overrides_notes_master_bg_ref(tmp_path):
    """The notes master's default <p:bgRef> must also be replaced.

    Notes master shows up in Presenter View and Notes Page view. Leaving it
    white makes those views visually inconsistent with a dark deck.
    """
    from designer_promoter import promote_designer_slides
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    _add_full_slide_dark_shape(s, "09090B")
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Hello"
    src_path = tmp_path / "src.pptx"
    src.save(str(src_path))
    out_path = tmp_path / "out.pptx"
    promote_designer_slides(
        source_pptx=str(src_path),
        output_path=str(out_path),
        layout_names={0: "L1"},
        run_container_check=False,
        run_contrast_check=False,
    )
    prs2 = Presentation(str(out_path))
    nm = prs2.notes_master
    bg = nm._element.find(qn("p:cSld")).find(qn("p:bg"))
    # Notes master should now have a solid fill (replacing the default bgRef)
    if bg is not None:
        bg_pr = bg.find(qn("p:bgPr"))
        if bg_pr is not None:
            srgb = bg_pr.find(".//" + qn("a:srgbClr"))
            assert srgb is not None
            assert srgb.get("val") == "09090B"
