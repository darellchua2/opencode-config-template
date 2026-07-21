"""Tests for BT-142 Phase 3.4.1b vision_extractor (XML fallback + aggregation).

Vision dispatch (soffice render + image-analyzer-subagent) is mocked — the
unit-testable surface is prompt construction, response aggregation, and the
XML fallback path. Live end-to-end runs are validated by Phase 6.11.
"""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches

from vision_extractor import (
    VisionSlideSchema,
    TextBlock,
    Region,
    build_image_analyzer_prompt,
    aggregate_vision_results,
    fallback_xml_background,
    _coerce_hex,
    _coerce_bbox,
)


# ---------------------------------------------------------------------------
# Coercion helpers
# ---------------------------------------------------------------------------

def test_coerce_hex_normalizes_lower_case():
    assert _coerce_hex("2dd4bf") == "#2DD4BF"
    assert _coerce_hex("#2dd4bf") == "#2DD4BF"
    assert _coerce_hex("#2DD4BF") == "#2DD4BF"


def test_coerce_hex_rejects_invalid():
    assert _coerce_hex("") == ""
    assert _coerce_hex(None) == ""
    assert _coerce_hex("#abc") == ""  # too short
    assert _coerce_hex("not-a-color") == ""


def test_coerce_bbox_filters_invalid():
    assert _coerce_bbox([0.1, 0.2, 0.3, 0.4]) == [0.1, 0.2, 0.3, 0.4]
    assert _coerce_bbox([0.1, 0.2, 0.3]) == [0.1, 0.2, 0.3]  # missing 4th OK
    assert _coerce_bbox([1.5, 0.2, 0.3, 0.4]) == [0.2, 0.3, 0.4]  # >1 filtered
    assert _coerce_bbox(None) == []
    assert _coerce_bbox("not-a-list") == []


# ---------------------------------------------------------------------------
# Prompt construction
# ---------------------------------------------------------------------------

def test_build_image_analyzer_prompt_contains_required_fields():
    p = build_image_analyzer_prompt(
        slide_index=2,
        slide_count=10,
        png_path="/tmp/slide-3.png",
        pptx_name="BETEKK.pptx",
    )
    assert "slide 2" in p.lower() or "slide index: 2" in p.lower()
    assert "10" in p  # slide count
    assert "/tmp/slide-3.png" in p
    assert "BETEKK.pptx" in p
    assert "dominant_bg_hex" in p
    assert "JSON" in p


# ---------------------------------------------------------------------------
# Response aggregation
# ---------------------------------------------------------------------------

def test_aggregate_clean_dict_responses():
    raw = [
        {
            "layout_archetype": "cover",
            "dominant_bg_hex": "#09090B",
            "text_blocks": [{"role": "headline", "color_hex": "#FFFFFF", "bbox_pct": [0.05, 0.1, 0.9, 0.15]}],
            "regions": [{"role": "card", "fill_hex": "#2DD4BF", "bbox_pct": [0.1, 0.5, 0.4, 0.4]}],
            "confidence": 0.9,
        },
        {
            "layout_archetype": "two_column",
            "dominant_bg_hex": "#18181B",
            "text_blocks": [],
            "regions": [],
            "confidence": 0.85,
        },
    ]
    out = aggregate_vision_results(raw)
    assert len(out) == 2
    assert out[0].dominant_bg_hex == "#09090B"
    assert out[0].layout_archetype == "cover"
    assert len(out[0].text_blocks) == 1
    assert out[0].text_blocks[0].color_hex == "#FFFFFF"
    assert out[0].confidence == 0.9
    assert out[1].dominant_bg_hex == "#18181B"


def test_aggregate_json_string_responses():
    import json
    raw = [
        json.dumps({"dominant_bg_hex": "#09090B", "confidence": 0.8}),
        "not valid json",
    ]
    out = aggregate_vision_results(raw)
    assert out[0].dominant_bg_hex == "#09090B"
    assert out[0].confidence == 0.8
    # Invalid entry falls back to empty schema
    assert out[1].dominant_bg_hex == ""
    assert out[1].confidence == 0.0


def test_aggregate_extracts_json_from_text():
    """Image-analyzer subagents sometimes wrap JSON in prose."""
    raw = [
        'Here is my analysis: {"dominant_bg_hex": "#18181B", "confidence": 0.75} — hope it helps.',
    ]
    out = aggregate_vision_results(raw)
    assert out[0].dominant_bg_hex == "#18181B"
    assert out[0].confidence == 0.75


def test_aggregate_handles_none_entries():
    """None entries (subagent dispatch failed) produce empty schemas."""
    out = aggregate_vision_results([None, None])
    assert len(out) == 2
    assert all(s.dominant_bg_hex == "" for s in out)
    assert all(s.confidence == 0.0 for s in out)


def test_aggregate_coerces_malformed_colors():
    raw = [{"dominant_bg_hex": "nothex", "confidence": "high"}]
    out = aggregate_vision_results(raw)
    assert out[0].dominant_bg_hex == ""
    # Invalid confidence falls back to 0.0
    assert out[0].confidence == 0.0


def test_aggregate_pads_short_input():
    """If fewer responses than slides, missing entries become empty schemas."""
    # No way to know target length — aggregate just returns what it got.
    # Caller's responsibility to align indices; we just don't crash.
    out = aggregate_vision_results([{"dominant_bg_hex": "#000000"}])
    assert len(out) == 1


# ---------------------------------------------------------------------------
# XML fallback (when vision unavailable)
# ---------------------------------------------------------------------------

def test_fallback_xml_background_detects_large_dark_shape(tmp_path):
    """When ≥50% of slide is covered by a dark shape, return its fill."""
    from pptx.dml.color import RGBColor
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-slide dark rectangle
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string("09090B")
    result = fallback_xml_background(slide)
    assert result is not None
    assert "09090B" in result.upper()


def test_fallback_xml_background_returns_none_when_no_large_shape():
    """Small shapes don't count as the slide background."""
    from pptx.dml.color import RGBColor
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Small accent shape
    card = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string("2DD4BF")
    result = fallback_xml_background(slide)
    assert result is None  # too small to be the background


def test_fallback_xml_background_handles_empty_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert fallback_xml_background(slide) is None


def test_fallback_xml_background_resolves_schemeClr_tx1_alias(tmp_path):
    """MED-2: <a:schemeClr val='tx1'> resolves to theme['dk1'] (ECMA-376 alias)."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string("000000")
    # Overwrite the srgbClr with schemeClr val="tx1" (alias for dk1)
    spPr = bg._element.find(qn("p:spPr"))
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid = spPr.find(f"{{{ns_a}}}solidFill")
    for child in list(solid):
        solid.remove(child)
    scheme = solid.makeelement(f"{{{ns_a}}}schemeClr", {"val": "tx1"})
    solid.append(scheme)
    theme = {"dk1": "1A1A1A", "lt1": "FFFFFF"}
    result = fallback_xml_background(slide, theme=theme)
    assert result is not None
    assert "1A1A1A" in result.upper()


def test_fallback_xml_background_resolves_schemeClr_bg1_alias(tmp_path):
    """MED-2: <a:schemeClr val='bg1'> resolves to theme['lt1'] (ECMA-376 alias)."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string("000000")
    spPr = bg._element.find(qn("p:spPr"))
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid = spPr.find(f"{{{ns_a}}}solidFill")
    for child in list(solid):
        solid.remove(child)
    scheme = solid.makeelement(f"{{{ns_a}}}schemeClr", {"val": "bg1"})
    solid.append(scheme)
    theme = {"dk1": "000000", "lt1": "FAFAFA"}
    result = fallback_xml_background(slide, theme=theme)
    assert result is not None
    assert "FAFAFA" in result.upper()


def test_fallback_xml_background_schemeClr_without_theme_returns_none(tmp_path):
    """MED-2: schemeClr without theme dict → can't resolve → None."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string("000000")
    spPr = bg._element.find(qn("p:spPr"))
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid = spPr.find(f"{{{ns_a}}}solidFill")
    for child in list(solid):
        solid.remove(child)
    scheme = solid.makeelement(f"{{{ns_a}}}schemeClr", {"val": "dk1"})
    solid.append(scheme)
    result = fallback_xml_background(slide, theme=None)
    assert result is None


# ---------------------------------------------------------------------------
# VisionSlideSchema serialization
# ---------------------------------------------------------------------------

def test_vision_slide_schema_to_dict():
    schema = VisionSlideSchema(
        slide_index=3,
        layout_archetype="team",
        dominant_bg_hex="#09090B",
        text_blocks=[TextBlock(role="body", color_hex="#FFFFFF", bbox_pct=[0.1, 0.2, 0.3, 0.4])],
        regions=[Region(role="card", fill_hex="#2DD4BF", bbox_pct=[0.5, 0.5, 0.4, 0.4])],
        confidence=0.92,
    )
    d = schema.to_dict()
    assert d["slide_index"] == 3
    assert d["dominant_bg_hex"] == "#09090B"
    assert len(d["text_blocks"]) == 1
    assert d["text_blocks"][0]["role"] == "body"
    assert d["confidence"] == 0.92
