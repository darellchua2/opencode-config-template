"""Dynamic layout cloning — Approach 2 (pure XML/part clone), issue #47 / DESIGN §7.

When P3's over-limit checker decides a layout must be added, this module
**deep-copies a donor layout** into a derived ``template_new.pptx``, resizes its
content placeholders, and verifies the result — rolling back on any failure so
the base ``template.pptx`` is **never corrupted**.

python-pptx 1.0.2 has no public ``SlideLayouts.add`` (DESIGN §1), so a layout is
created by: deep-copying the donor ``<p:sldLayout>`` element into a new
``SlideLayoutPart``, relating it to the master, registering a
``<p:sldLayoutId>`` entry, and re-attaching the donor's non-master rels.

Safety model: the clone works on an in-memory copy of the template and **saves
to a different file** (``template_new.pptx``). The base is therefore never
written; rollback = delete the derived file if save/reload-verify fails.
"""
from __future__ import annotations

import copy
import logging
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

# PLAN-GIT-72: shared contract layer now in _common (no more sibling-skill hack).
_COMMON_SCRIPTS = Path(__file__).resolve().parents[2] / "_common" / "scripts"
if str(_COMMON_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_COMMON_SCRIPTS))

from pptx.parts.slide import SlideLayoutPart  # noqa: E402
from errors import TemplateError  # noqa: E402 — US-4.8/MINOR-2: shared error from _common
from layout_contract import _SLIDE_TYPE_FINGERPRINT, _resolve_layout_by_fingerprint  # noqa: E402
from state_machine import (  # noqa: E402 (same package dir)
    ResolutionPlan,
    template_new_path_for,
)

logger = logging.getLogger(__name__)

_EMU_PER_INCH = 914400
_MARGIN_EMU = 457200  # 0.5"
# Placeholder ph/@type values that represent content/body (to be resized).
_CONTENT_PH_TYPES = {"", "body", "obj", "txt"}


def _max_layout_id(sld_layout_id_lst: Any) -> int:
    """Next unique ``<p:sldLayoutId id>`` (ECMA-376 min 2147483648)."""
    existing = []
    for entry in sld_layout_id_lst.sldLayoutId_lst:
        raw = entry.get("id")
        if raw is not None:
            try:
                existing.append(int(raw))
            except ValueError:
                pass
    return (max(existing) + 1) if existing else 2147483649


def _set_layout_name(element: Any, new_name: str) -> None:
    cSld = element.find(qn("p:cSld"))
    if cSld is not None:
        cSld.set("name", new_name)


def _resize_content_placeholders(element: Any, slide_w_emu: int, slide_h_emu: int) -> int:
    """Expand content/body placeholders toward the slide's bottom-right edge.

    Never shrinks. Returns the number of placeholders resized.
    """
    resized = 0
    cSld = element.find(qn("p:cSld"))
    if cSld is None:
        return 0
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return 0
    for sp in spTree.findall(qn("p:sp")):
        ph = sp.find(qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph"))
        if ph is None:
            continue
        if ph.get("type", "") not in _CONTENT_PH_TYPES:
            continue
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            continue
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            continue
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is None or ext is None:
            continue
        x = int(off.get("x", "0"))
        y = int(off.get("y", "0"))
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
        new_cx = max(cx, slide_w_emu - x - _MARGIN_EMU)
        new_cy = max(cy, slide_h_emu - y - _MARGIN_EMU)
        if (new_cx, new_cy) != (cx, cy):
            ext.set("cx", str(new_cx))
            ext.set("cy", str(new_cy))
            resized += 1
    return resized


def _clone_layout_into(
    prs: Any,
    donor_index: int,
    new_name: str,
    resize: bool = True,
) -> str:
    """Mutate ``prs``: add a cloned layout (7-step, DESIGN §7). Returns new_name.

    Raises on any failure (caller rolls back by discarding the unsaved prs /
    deleting the output file). The base file is never written.
    """
    masters = list(prs.slide_masters)
    if not masters:
        raise TemplateError(
            "Cannot clone layout: presentation has no slide master. "
            "A slide master is required to host cloned layouts."
        )
    master = masters[0]
    master_part = master.part
    master_element = master._element
    package = prs.part.package

    # Step 1: donor
    donor = prs.slide_layouts[donor_index]
    donor_part = donor.part

    # Step 2: deep-copy the donor <p:sldLayout> element
    new_element = copy.deepcopy(donor_part._element)
    _set_layout_name(new_element, new_name)

    # Step 6 (resize) — applied to the element before wrapping in a part
    if resize:
        n = _resize_content_placeholders(new_element, int(prs.slide_width), int(prs.slide_height))
        logger.info("Cloned layout '%s': resized %d content placeholder(s)", new_name, n)

    # Step 2b: construct a new SlideLayoutPart
    new_partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
    new_part = SlideLayoutPart(new_partname, CT.PML_SLIDE_LAYOUT, package, new_element)

    # Step 4: master → layout relationship
    rId = master_part.relate_to(new_part, RT.SLIDE_LAYOUT)

    # Step 5: register <p:sldLayoutId> in the master's list
    sld_layout_id_lst = master_element.get_or_add_sldLayoutIdLst()
    entry = sld_layout_id_lst._add_sldLayoutId(rId=rId)
    entry.set("id", str(_max_layout_id(sld_layout_id_lst)))

    # The clone's own master rel (its only structural rel)
    new_part.relate_to(master_part, RT.SLIDE_MASTER)

    # Step 3: copy the donor's non-master rels (background images, etc.)
    for _, rel in donor_part.rels.items():
        if rel.reltype == RT.SLIDE_MASTER:
            continue
        if rel.is_external:
            new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_part.relate_to(rel.target_part, rel.reltype)

    return new_name


def _verify_layouts(output_path: str, expected_names: List[str]) -> None:
    """Reload the saved file and confirm every cloned layout is findable.

    Raises ``RuntimeError`` if any expected layout is missing (triggers rollback).
    """
    reloaded = Presentation(output_path)
    masters = list(reloaded.slide_masters)
    if not masters:
        raise TemplateError(
            f"reload-verify failed: {output_path} has no slide master after save"
        )
    layouts = masters[0].slide_layouts
    for name in expected_names:
        found = layouts.get_by_name(name)
        if found is None:
            raise RuntimeError(
                f"reload-verify failed: cloned layout '{name}' not found in {output_path}"
            )


def clone_for_over_limit(
    template_path: str,
    plan: ResolutionPlan,
    contract: Dict[str, Any],
    output_path: Optional[str] = None,
) -> Tuple[str, Dict[str, str]]:
    """Clone an extended layout for each over-limit slide_type (issue #47).

    Returns ``(template_new_path, config_overrides)`` where ``config_overrides``
    maps each cloned ``slide_type`` → clone layout name (so Capability A renders
    against the extended layout). Rolls back (deletes the output) on any failure.
    """
    out = Path(output_path) if output_path else template_new_path_for(template_path)

    prs = Presentation(template_path)
    overrides: Dict[str, str] = {}
    seen: set = set()
    for req in plan.over_limit_slides:
        slide_type = req.slide_type
        if slide_type in seen or slide_type not in _SLIDE_TYPE_FINGERPRINT:
            continue
        donor_idx, _ = _resolve_layout_by_fingerprint(slide_type, contract)
        if donor_idx is None:
            logger.warning("No donor for over-limit '%s'; skipping clone", slide_type)
            continue
        new_name = f"{slide_type} (extended)"
        _clone_layout_into(prs, donor_idx, new_name, resize=True)
        overrides[slide_type] = new_name
        seen.add(slide_type)

    if not overrides:
        # Nothing clonable (e.g. only unknown types); don't write a derived file.
        return template_path, {}

    try:
        prs.save(str(out))
        _verify_layouts(str(out), list(overrides.values()))
        logger.info("Produced %s with %d extended layout(s)", out.name, len(overrides))
    except Exception:
        # Rollback: discard the derived file so the base stays authoritative.
        if out.exists():
            try:
                out.unlink()
                logger.warning("Rolled back: deleted failed %s", out.name)
            except OSError:
                pass
        raise
    return str(out), overrides
