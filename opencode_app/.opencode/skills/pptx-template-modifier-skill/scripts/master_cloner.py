"""US-4.8 Phase 2 — Extend a template with borrowed-geometry layouts.

When the user's template has a master but lacks layouts for some slide types
(Scenario B), this module borrows layouts from ``default.pptx`` and injects
them under the user's **existing** master.

**Spike result (T2.0):** Master cloning was found unnecessary — cross-file
LAYOUT cloning (deep-copy ``<p:sldLayout>`` from default.pptx → create a new
``SlideLayoutPart`` in the user's package → register under the user's master)
works correctly and the new layout inherits the user's theme automatically via
the master. This is the Decision 5 fallback path: simpler, safer, and the
shared-theme invariant holds trivially.

The module is imported **lazily** by ``state_machine.resolve_and_clone`` (CRIT-4
— no circular import with ``layout_creator``).
"""
from __future__ import annotations

import copy
import logging
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT

# PLAN-GIT-72: shared contract layer in _common.
_COMMON_SCRIPTS = Path(__file__).resolve().parents[2] / "_common" / "scripts"
if str(_COMMON_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_COMMON_SCRIPTS))

from pptx.parts.slide import SlideLayoutPart  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from layout_contract import (  # noqa: E402
    _resolve_layout_by_fingerprint,
    get_render_contract,
)
from state_machine import template_new_path_for  # noqa: E402 (same package dir)
from layout_creator import _max_layout_id, _resize_content_placeholders  # noqa: E402

logger = logging.getLogger(__name__)


def clone_master_and_borrow(
    template_path: str,
    missing_slide_types: List[str],
    default_template_path: str,
    output_path: Optional[str] = None,
) -> Tuple[str, Dict[str, str]]:
    """Borrow layouts from ``default.pptx`` for each missing slide_type.

    .. deprecated-in-name::
        The function name retains "clone_master" for historical continuity with
        the plan, but NO master is cloned (the T2.0 spike found it unnecessary).
        Only layouts are borrowed and injected under the user's existing master.

    For each slide_type that has no donor in the user's template, deep-copy the
    matching layout from ``default_template_path`` into the user's existing
    master. The new layout inherits the user's theme automatically (Decision 5
    fallback — no master cloning needed).

    Returns ``(extended_path, {slide_type: new_layout_name})``. Rolls back
    (deletes the derived file) on any failure.
    """
    if output_path is None:
        output_path = str(template_new_path_for(template_path))

    try:
        prs = Presentation(template_path)
    except Exception as exc:
        logger.error("clone_master_and_borrow: cannot open %s (%s)", template_path, exc)
        return template_path, {}

    masters = list(prs.slide_masters)
    if not masters:
        logger.error("clone_master_and_borrow: template has no master: %s", template_path)
        return template_path, {}

    master = masters[0]
    master_part = master.part
    master_element = master._element
    package = prs.part.package

    # BT-142 Phase 3.2: ``default_template_path`` is now REQUIRED and must be
    # a user-supplied donor PPTX (or a synthesized minimal PPTX via
    # ``_build_minimal_pptx_bytes(None)`` written to disk by the caller).
    # Per the user's "no bundled default.pptx" invariant, there is no implicit
    # fallback — fail loudly with an actionable error when no donor is supplied.
    if not default_template_path:
        raise TemplateError(
            "No donor slide master available — provide a template with a slide "
            "master, or generate one first via pptx-generate-template-skill. The "
            "pptx-template-modifier-skill cannot extend a masterless template "
            "without a donor to borrow layouts from."
        )

    # Open the donor PPTX to source layouts.
    default_prs = Presentation(default_template_path)
    default_contract = get_render_contract(default_template_path)

    # User slide dimensions for placeholder resize.
    user_w_emu = int(prs.slide_width)
    user_h_emu = int(prs.slide_height)

    overrides: Dict[str, str] = {}

    for slide_type in missing_slide_types:
        # Find the donor layout in default.pptx via fingerprint.
        donor_idx, reason = _resolve_layout_by_fingerprint(slide_type, default_contract)
        if donor_idx is None:
            logger.warning(
                "No donor layout for '%s' in default.pptx: %s — skipping",
                slide_type, reason,
            )
            continue

        donor_layout = default_prs.slide_layouts[donor_idx]
        donor_part = donor_layout.part
        donor_name = donor_layout.name or f"layout_{donor_idx}"
        new_name = f"{donor_name} (borrowed)"

        logger.info(
            "Borrowing layout '%s' (idx=%d) from default.pptx for slide_type '%s'",
            donor_name, donor_idx, slide_type,
        )

        try:
            new_element = copy.deepcopy(donor_part._element)

            # Rename the cloned layout.
            cSld = new_element.find(qn("p:cSld"))
            if cSld is not None:
                cSld.set("name", new_name)

            # Resize content placeholders to fill the user's slide dimensions.
            n = _resize_content_placeholders(new_element, user_w_emu, user_h_emu)
            logger.info("Resized %d placeholder(s) for '%s'", n, new_name)

            # Create the new SlideLayoutPart in the user's package.
            new_partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
            new_part = SlideLayoutPart(new_partname, CT.PML_SLIDE_LAYOUT, package, new_element)

            # Relate user master → new layout.
            rId = master_part.relate_to(new_part, RT.SLIDE_LAYOUT)

            # Register <p:sldLayoutId> in the master's list.
            sld_layout_id_lst = master_element.get_or_add_sldLayoutIdLst()
            entry = sld_layout_id_lst._add_sldLayoutId(rId=rId)
            entry.set("id", str(_max_layout_id(sld_layout_id_lst)))

            # New layout → master back-reference.
            new_part.relate_to(master_part, RT.SLIDE_MASTER)

            # Copy donor's non-master rels (images, etc.) — share parts, don't duplicate.
            for _, rel in donor_part.rels.items():
                if rel.reltype == RT.SLIDE_MASTER:
                    continue
                if rel.is_external:
                    new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_part.relate_to(rel.target_part, rel.reltype)

            overrides[slide_type] = new_name
            logger.info("Borrowed layout '%s' registered for '%s'", new_name, slide_type)

        except Exception as exc:
            logger.error("Failed to borrow layout for '%s': %s", slide_type, exc)
            continue

    if not overrides:
        logger.warning("No layouts could be borrowed — returning base template")
        return template_path, {}

    # Save and reload-verify.
    try:
        prs.save(output_path)
        reloaded = Presentation(output_path)
        reloaded_master = list(reloaded.slide_masters)
        if not reloaded_master:
            raise RuntimeError("reload-verify: no master after save")
        layouts = reloaded_master[0].slide_layouts
        for name in overrides.values():
            if layouts.get_by_name(name) is None:
                raise RuntimeError(f"reload-verify: borrowed layout '{name}' not found")
        logger.info(
            "Produced %s with %d borrowed layout(s)",
            Path(output_path).name, len(overrides),
        )
    except Exception as exc:
        logger.error("Reload-verify failed (%s) — rolling back", exc)
        out = Path(output_path)
        if out.exists():
            try:
                out.unlink()
            except OSError:
                pass
        return template_path, {}

    return output_path, overrides
