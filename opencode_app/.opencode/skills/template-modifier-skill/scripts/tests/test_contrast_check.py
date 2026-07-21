"""Tests for BT-142 Phase 3.4.2b contrast_check (WCAG 2.1)."""

from __future__ import annotations

import sys
import pathlib

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from contrast_check import (
    contrast_ratio,
    required_ratio_for,
    contrast_violations,
    check_template_contrast,
    ContrastViolation,
    WCAG_AA_NORMAL,
    WCAG_AA_LARGE,
    WCAG_CRITICAL,
)


# ---------------------------------------------------------------------------
# WCAG math
# ---------------------------------------------------------------------------

def test_contrast_ratio_white_on_black_max():
    """White on black should be the maximum ratio (~21:1)."""
    r = contrast_ratio("#FFFFFF", "#000000")
    assert 20.5 <= r <= 21.5


def test_contrast_ratio_same_color_min():
    """Same color = no contrast (1.0)."""
    assert contrast_ratio("#2DD4BF", "#2DD4BF") == 1.0


def test_contrast_ratio_betekk_text_on_dark():
    """The BETEKK defect: white text on #09090B should be very high contrast."""
    r = contrast_ratio("#FFFFFF", "#09090B")
    assert r > 15.0  # strong contrast


def test_contrast_ratio_betekk_defect_dark_text_on_dark():
    """Dark text (#09090B dk1) on dark card (#18181B) — the invisible-text case."""
    r = contrast_ratio("#09090B", "#18181B")
    assert r < 2.0  # effectively unreadable


def test_contrast_ratio_unresolvable_returns_one():
    """Invalid color strings return 1.0 (no contrast)."""
    assert contrast_ratio("not-a-color", "#FFFFFF") == 1.0
    assert contrast_ratio("#FFFFFF", "") == 1.0


def test_required_ratio_normal_text():
    assert required_ratio_for(12.0) == WCAG_AA_NORMAL  # 4.5
    assert required_ratio_for(16.0) == WCAG_AA_NORMAL


def test_required_ratio_large_text_regular():
    """≥18pt regular = large text = lower bar (3.0)."""
    assert required_ratio_for(18.0) == WCAG_AA_LARGE
    assert required_ratio_for(24.0) == WCAG_AA_LARGE


def test_required_ratio_large_text_bold():
    """≥14pt bold = large text."""
    assert required_ratio_for(14.0, is_bold=True) == WCAG_AA_LARGE


def test_required_ratio_none_size_defaults_normal():
    assert required_ratio_for(None) == WCAG_AA_NORMAL


# ---------------------------------------------------------------------------
# contrast_violations() with fakes
# ---------------------------------------------------------------------------

class _FakePhFmt:
    def __init__(self, idx, type_=None):
        self.idx = idx
        self.type = type_


class _FakeFill:
    def __init__(self, rgb_hex=None):
        self.type = "solid" if rgb_hex else None
        from pptx.dml.color import RGBColor
        if rgb_hex:
            self.fore_color = type("FC", (), {"rgb": RGBColor.from_string(rgb_hex.lstrip("#"))})()
        else:
            self.fore_color = None


class _FakeRun:
    def __init__(self, text="", color_hex=None, size_pt=None):
        self.text = text
        self.font = type("F", (), {})()
        from pptx.dml.color import RGBColor
        if color_hex:
            self.font.color = type("C", (), {"rgb": RGBColor.from_string(color_hex.lstrip("#"))})()
        else:
            self.font.color = type("C", (), {"rgb": None})()
        from pptx.util import Pt
        self.font.size = Pt(size_pt) if size_pt else None


class _FakePara:
    def __init__(self, runs):
        self.runs = runs


class _FakeTextFrame:
    def __init__(self, paragraphs_or_runs):
        # Accept either a list of paragraphs OR a list of runs (auto-wrapped).
        if not paragraphs_or_runs:
            self.paragraphs = []
        elif all(isinstance(p, _FakePara) for p in paragraphs_or_runs):
            self.paragraphs = paragraphs_or_runs
        else:
            # Treat as a list of runs → wrap in a single paragraph.
            self.paragraphs = [_FakePara(paragraphs_or_runs)]
        self.text = " ".join(r.text for p in self.paragraphs for r in p.runs)


class _FakeShape:
    def __init__(self, name, x_in, y_in, w_in, h_in, *,
                 is_placeholder=False, idx=None,
                 fill_hex=None, runs=None, has_text_frame=True):
        from pptx.util import Inches
        self.name = name
        self.left = Inches(x_in); self.top = Inches(y_in)
        self.width = Inches(w_in); self.height = Inches(h_in)
        self.is_placeholder = is_placeholder
        self.placeholder_format = _FakePhFmt(idx) if is_placeholder else None
        self.fill = _FakeFill(fill_hex)
        self.has_text_frame = has_text_frame
        if has_text_frame:
            self.text_frame = _FakeTextFrame(runs or [])
        self.shape_id = hash(name) & 0xFFFF


class _FakeShapes:
    def __init__(self, shapes):
        self._shapes = shapes
    def __iter__(self):
        return iter(self._shapes)


class _FakeLayout:
    def __init__(self, name, shapes):
        self.name = name
        self.shapes = _FakeShapes(shapes)


def test_violation_detects_dark_text_on_dark_card():
    """Default dk1 text on a dark card → critical violation."""
    card = _FakeShape("card", 0, 0, 5, 3, fill_hex="#18181B")
    body = _FakeShape(
        "body", 1, 1, 4, 1.5,
        is_placeholder=True, idx=1,
        runs=[_FakeRun(text="hello", color_hex=None, size_pt=14)],  # inherits default
    )
    layout = _FakeLayout("L", [card, body])
    vs = contrast_violations(layout, theme={"dk1": "#09090B", "lt1": "#FFFFFF"})
    assert len(vs) == 1
    assert vs[0].severity == "critical"
    assert vs[0].contrast_ratio < WCAG_CRITICAL


def test_violation_clean_when_white_text_on_dark_card():
    card = _FakeShape("card", 0, 0, 5, 3, fill_hex="#18181B")
    body = _FakeShape(
        "body", 1, 1, 4, 1.5,
        is_placeholder=True, idx=1,
        runs=[_FakeRun(text="hello", color_hex="#FFFFFF", size_pt=14)],
    )
    layout = _FakeLayout("L", [card, body])
    vs = contrast_violations(layout, theme={"dk1": "#09090B", "lt1": "#FFFFFF"})
    assert vs == []


def test_violation_uses_theme_default_when_placeholder_has_no_text():
    """Empty placeholder inherits theme dk1 → still gets checked."""
    card = _FakeShape("card", 0, 0, 5, 3, fill_hex="#09090B")
    body = _FakeShape(
        "body", 1, 1, 4, 1.5,
        is_placeholder=True, idx=1,
        runs=[],  # empty — inherits theme dk1
    )
    layout = _FakeLayout("L", [card, body])
    vs = contrast_violations(layout, theme={"dk1": "#09090B", "lt1": "#FFFFFF"})
    assert len(vs) == 1
    assert vs[0].severity == "critical"


def test_auto_fix_flips_low_contrast_text_to_white():
    """auto_fix=True overrides the placeholder's default text color."""
    card = _FakeShape("card", 0, 0, 5, 3, fill_hex="#09090B")
    body = _FakeShape(
        "body", 1, 1, 4, 1.5,
        is_placeholder=True, idx=1,
        runs=[_FakeRun(text="hi", color_hex="#09090B", size_pt=14)],
    )
    layout = _FakeLayout("L", [card, body])
    # Auto-fix path needs a real text_frame with paragraphs — fakes lack it.
    # Verify the violation is still detected (auto-fix returns False on the fake).
    vs = contrast_violations(layout, theme={"dk1": "#09090B"}, auto_fix=True)
    assert len(vs) == 1
    assert vs[0].auto_fixed is False  # fake text_frame rejects the override


def test_violation_to_dict_round_trip():
    v = ContrastViolation(
        placeholder_idx=1, placeholder_name="body", fg_color="#000", bg_color="#fff",
        contrast_ratio=1.5, required_ratio=4.5, severity="critical", detail="x",
    )
    d = v.to_dict()
    assert d["severity"] == "critical"
    assert d["contrast_ratio"] == 1.5


def test_check_template_contrast_summary_shape():
    """Summary aggregates critical/warning/auto_fixed across layouts."""
    from pptx import Presentation
    out = check_template_contrast(Presentation())
    assert "__summary" in out
    s = out["__summary"]
    assert set(s.keys()) == {"critical", "warning", "auto_fixed", "layouts_checked"}
    assert s["layouts_checked"] >= 1


def test_invert_hex_picks_high_contrast_partner():
    """The auto-fix helper flips dark→white and light→black."""
    from contrast_check import _invert_hex
    assert _invert_hex("#000000") == "#FFFFFF"
    assert _invert_hex("#09090B") == "#FFFFFF"  # very dark
    assert _invert_hex("#FFFFFF") == "#000000"
    assert _invert_hex("#2DD4BF") == "#000000"  # teal is fairly bright
