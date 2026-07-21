"""BT-142 Phase 3.4.2 — Static container-fit check for promoted layouts.

When ``designer_promoter`` (Phase 3.4.1) reverse-engineers per-shape branding
into text placeholders on a layout, a placeholder's bounding box may extend
beyond its visual container shape (orange card, teal panel, decorative band).
BETEKK V9.1.1 slide 4 is the canonical example: the body text spilled outside
its orange box.

This module performs a **pure-geometry** check (no rendering, no LLM) at
template-build time. For every TITLE/BODY placeholder on a layout, find the
nearest decorative shape whose bounding box geometrically contains the
placeholder's center point. If found, verify the placeholder's rect fits
inside the container's rect (with optional padding tolerance). Report
violations grouped by severity (critical > 20px overflow, warning 4-20px).

Algorithm (per Phase 3.4.2 spec):
  1. For each text placeholder on the layout, compute its bounding box in EMU
     and its center point.
  2. For each non-placeholder shape on the same layout (the candidate
     "container"), check whether the placeholder's center lies inside the
     candidate's rect (point-in-rect).
  3. Among matching candidates, pick the tightest (smallest area) as the
     assigned container.
  4. Compute the geometric overflow of the placeholder beyond the container
     rect (positive = placeholder extends outside container on that edge).
  5. Severity:
       - max_edge_overflow > 20 px  → critical
       - 4..20 px                  → warning
       - < 4 px                    → ok (within tolerance)
  6. Return a list of violations and a summary (critical_count, warning_count).

Public API:
  - ``container_violations(layout, tolerance_px=4, critical_px=20)`` → List[Violation]
  - ``check_template(prs, **kw)`` → Dict[layout_name, List[Violation]]
  - ``ContainerFitError`` → raised by ``designer_promoter`` on critical violations
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from typing import Iterable, List, Optional, Tuple
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from pptx.slide import SlideLayout
    from pptx.shapes.shapetree import Shape
    from pptx import Presentation

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400
EMU_PER_PX = 9525  # 96 DPI assumption used for px<->EMU conversion


class ContainerFitError(Exception):
    """Raised when a critical container-fit violation blocks a template build."""


@dataclass
class Rect:
    """Axis-aligned rectangle in EMU."""

    x: int
    y: int
    cx: int
    cy: int

    @property
    def left(self) -> int:
        return self.x

    @property
    def top(self) -> int:
        return self.y

    @property
    def right(self) -> int:
        return self.x + self.cx

    @property
    def bottom(self) -> int:
        return self.y + self.cy

    @property
    def center(self) -> Tuple[float, float]:
        return (self.x + self.cx / 2.0, self.y + self.cy / 2.0)

    @property
    def area(self) -> int:
        return self.cx * self.cy

    def contains_point(self, px: float, py: float) -> bool:
        return self.left <= px <= self.right and self.top <= py <= self.bottom

    def edge_overflow(self, other: "Rect") -> int:
        """Max EMU by which ``self`` extends outside ``other`` on any edge.

        Returns 0 (or negative) when ``self`` is fully inside ``other``.
        """
        left_o = other.left - self.left
        top_o = other.top - self.top
        right_o = self.right - other.right
        bottom_o = self.bottom - other.bottom
        return max(0, left_o, top_o, right_o, bottom_o)


@dataclass
class Violation:
    """A single container-fit violation on one placeholder."""

    placeholder_idx: int
    placeholder_name: str
    container_shape_id: int
    container_name: str
    overflow_emu: int
    overflow_px: int
    severity: str  # "critical" | "warning"
    detail: str = ""

    def to_dict(self) -> dict:
        return {
            "placeholder_idx": self.placeholder_idx,
            "placeholder_name": self.placeholder_name,
            "container_shape_id": self.container_shape_id,
            "container_name": self.container_name,
            "overflow_emu": self.overflow_emu,
            "overflow_px": self.overflow_px,
            "severity": self.severity,
            "detail": self.detail,
        }


def _shape_to_rect(shape) -> Optional[Rect]:
    try:
        return Rect(int(shape.left), int(shape.top), int(shape.width), int(shape.height))
    except (TypeError, AttributeError):
        return None


def _is_text_placeholder(shape) -> bool:
    """True iff shape is a placeholder carrying text (TITLE / BODY / OBJECT / SUBTITLE).

    Other placeholder types (PICTURE=18, TABLE=19, CHART=20, etc.) return
    False — they are NOT text candidates for container-fit purposes. The
    previous `or True` here defeated the type filter (every placeholder was
    treated as text); BT-142 review caught this P0 bug.
    """
    is_ph = getattr(shape, "is_placeholder", False) and shape.is_placeholder
    if not is_ph:
        return False
    try:
        type_ = shape.placeholder_format.type
        if type_ is None:
            # python-pptx returns None for generic "OBJECT" placeholders that
            # commonly carry text in designer templates — treat as text.
            return True
        # PP_PLACEHOLDER enum int values: TITLE=13, BODY=2, OBJECT=17, SUBTITLE=3
        # Use .value first (canonical), fall back to int() conversion, then
        # to string-parse of "NAME (NN)" format (older python-pptx).
        type_id = getattr(type_, "value", None)
        if type_id is None:
            try:
                type_id = int(type_)
            except (TypeError, ValueError):
                # str(enum) looks like "PICTURE (18)" — extract the int.
                s = str(type_).rsplit("(", 1)[-1].rstrip(") ").strip()
                type_id = int(s)
        return int(type_id) in (2, 3, 13, 17)
    except Exception:
        # Conservative: when we can't determine the type, treat as a text
        # candidate (the geometry check will simply find no container for
        # non-text placeholders, so no false positives ensue).
        return True


def _classify_severity(overflow_px: int, critical_px: int) -> str:
    if overflow_px > critical_px:
        return "critical"
    return "warning"


def container_violations(
    layout: "SlideLayout",
    tolerance_px: int = 4,
    critical_px: int = 20,
) -> List[Violation]:
    """Return all container-fit violations on ``layout``.

    Args:
        layout: a ``pptx.slide.SlideLayout``
        tolerance_px: edge slack (px) absorbed before flagging (default 4px)
        critical_px: threshold (px) above which severity becomes critical
            (default 20px)

    Returns:
        List of :class:`Violation` instances, possibly empty.
    """
    violations: List[Violation] = []
    shapes = list(layout.shapes)

    text_placeholders = [s for s in shapes if _is_text_placeholder(s)]
    # Candidate containers = non-placeholder shapes that have an explicit fill
    # (cards, panels, accent bands). Decorative-only outlines are also kept —
    # we err on the side of checking.
    candidate_containers = [
        s for s in shapes if not _is_text_placeholder(s) and _shape_to_rect(s) is not None
    ]

    for ph in text_placeholders:
        ph_rect = _shape_to_rect(ph)
        if ph_rect is None:
            continue
        cx, cy = ph_rect.center
        # Tightest container whose rect contains the placeholder's center.
        best_container = None
        best_rect: Optional[Rect] = None
        best_area = None
        for cand in candidate_containers:
            cand_rect = _shape_to_rect(cand)
            if cand_rect is None or not cand_rect.contains_point(cx, cy):
                continue
            if best_area is None or cand_rect.area < best_area:
                best_container = cand
                best_rect = cand_rect
                best_area = cand_rect.area
        if best_container is None or best_rect is None:
            continue
        # Subtract tolerance from each edge before measuring overflow.
        tol_emu = int(tolerance_px * EMU_PER_PX)
        shrunk = Rect(
            best_rect.left + tol_emu,
            best_rect.top + tol_emu,
            max(0, best_rect.cx - 2 * tol_emu),
            max(0, best_rect.cy - 2 * tol_emu),
        )
        overflow_emu = ph_rect.edge_overflow(shrunk)
        if overflow_emu <= 0:
            continue
        overflow_px = int(overflow_emu / EMU_PER_PX)
        severity = _classify_severity(overflow_px, critical_px)
        try:
            ph_idx = ph.placeholder_format.idx
            ph_name = ph.name
        except Exception:
            ph_idx, ph_name = -1, "<unnamed-placeholder>"
        violations.append(
            Violation(
                placeholder_idx=ph_idx,
                placeholder_name=ph_name,
                container_shape_id=getattr(best_container, "shape_id", -1),
                container_name=getattr(best_container, "name", "<unnamed-shape>"),
                overflow_emu=overflow_emu,
                overflow_px=overflow_px,
                severity=severity,
                detail=f"{ph_name} extends {overflow_px}px beyond {best_container.name}",
            )
        )
    return violations


def check_template(
    prs: "Presentation",
    tolerance_px: int = 4,
    critical_px: int = 20,
) -> dict:
    """Run ``container_violations`` on every layout in ``prs``.

    Returns:
        ``{layout_name: [violation_dict, ...], "__summary": {...}}``
    """
    summary = {"critical": 0, "warning": 0, "layouts_checked": 0}
    out = {}
    for layout in prs.slide_layouts:
        try:
            layout_name = layout.name
        except Exception:
            layout_name = f"<layout-{id(layout)}>"
        try:
            vs = container_violations(layout, tolerance_px=tolerance_px, critical_px=critical_px)
        except Exception as exc:
            logger.warning("container_check failed on layout %s: %s", layout_name, exc)
            vs = []
        out[layout_name] = [v.to_dict() for v in vs]
        summary["layouts_checked"] += 1
        for v in vs:
            if v.severity == "critical":
                summary["critical"] += 1
            else:
                summary["warning"] += 1
    out["__summary"] = summary
    return out


__all__ = [
    "Rect",
    "Violation",
    "ContainerFitError",
    "container_violations",
    "check_template",
]
