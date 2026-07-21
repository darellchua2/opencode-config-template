"""BT-142 Phase 3.4.2b — Static color-contrast check for promoted layouts.

Companion to ``container_check.py``. When ``designer_promoter`` (Phase 3.4.1)
reverse-engineers per-shape branding into text placeholders, two contrast
defects commonly appear — even when the original designer slide "looked fine":

  1. **Inherited-text defect.** The promoted placeholder has no explicit
     ``<a:rPr>` color, so it inherits the theme's ``dk1`` (typically near-
     black). If the placeholder sits on a dark card (``#09090B``,
     ``#18181B``), the text is invisible. The original designer shape
     probably had an explicit white run color that did not survive the
     shape→placeholder promotion.

  2. **Container-bg defect.** The placeholder was promoted from a shape that
     used the slide background as its effective background. After promotion,
     the placeholder sits on a different container shape (or none), and the
     old text color no longer contrasts.

This module performs a WCAG 2.1 contrast check at template-build time. For
every TITLE/BODY placeholder on a layout, resolve its effective foreground
color and the effective background it sits on, then compute the contrast
ratio. Flag ratios below 4.5:1 (normal text) or 3:1 (large text ≥18pt).

Public API:
  - ``contrast_ratio(fg_hex, bg_hex) -> float`` — WCAG 2.1 ratio
  - ``contrast_violations(layout, theme=None, auto_fix=False) -> List[Violation]``
  - ``check_template_contrast(prs, auto_fix=False) -> dict``
  - ``ContrastViolation`` dataclass with severity
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# WCAG 2.1 thresholds
WCAG_AA_NORMAL = 4.5      # normal text (< 18pt regular, < 14pt bold)
WCAG_AA_LARGE = 3.0       # large text (≥ 18pt regular, ≥ 14pt bold)
WCAG_CRITICAL = 3.0       # below this = effectively unreadable

LARGE_TEXT_REGULAR_PT = 18.0
LARGE_TEXT_BOLD_PT = 14.0


@dataclass
class ContrastViolation:
    """A single contrast defect on one placeholder."""

    placeholder_idx: int
    placeholder_name: str
    fg_color: str          # resolved foreground hex ("#RRGGBB")
    bg_color: str          # resolved background hex ("#RRGGBB")
    contrast_ratio: float
    required_ratio: float  # 4.5 (normal) or 3.0 (large)
    severity: str          # "critical" (< 3.0) | "warning" (< 4.5)
    detail: str = ""
    auto_fixed: bool = False  # True if auto_fix flipped the colors

    def to_dict(self) -> dict:
        return {
            "placeholder_idx": self.placeholder_idx,
            "placeholder_name": self.placeholder_name,
            "fg_color": self.fg_color,
            "bg_color": self.bg_color,
            "contrast_ratio": round(self.contrast_ratio, 2),
            "required_ratio": self.required_ratio,
            "severity": self.severity,
            "detail": self.detail,
            "auto_fixed": self.auto_fixed,
        }


# --------------------------------------------------------------------------
# WCAG 2.1 contrast math
# --------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> Optional[Tuple[int, int, int]]:
    """'#2DD4BF' or '2DD4BF' → (45, 212, 191). None on parse failure."""
    if not hex_str:
        return None
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return None
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def _linearize_channel(cs: float) -> float:
    """WCAG 2.1 sRGB → linear light conversion for one channel."""
    if cs <= 0.03928:
        return cs / 12.92
    return ((cs + 0.055) / 1.055) ** 2.4


def _relative_luminance(hex_str: str) -> Optional[float]:
    """WCAG 2.1 relative luminance (0..1). None if color unresolvable."""
    rgb = _hex_to_rgb(hex_str)
    if rgb is None:
        return None
    r, g, b = (c / 255.0 for c in rgb)
    R = _linearize_channel(r)
    G = _linearize_channel(g)
    B = _linearize_channel(b)
    return 0.2126 * R + 0.7152 * G + 0.0722 * B


def contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """WCAG 2.1 contrast ratio between two hex colors.

    Returns 1.0 (no contrast) if either color is unresolvable — callers
    should treat very low ratios as suspicious.
    """
    fg_l = _relative_luminance(fg_hex)
    bg_l = _relative_luminance(bg_hex)
    if fg_l is None or bg_l is None:
        return 1.0
    lighter = max(fg_l, bg_l)
    darker = min(fg_l, bg_l)
    return (lighter + 0.05) / (darker + 0.05)


def required_ratio_for(font_pt: Optional[float], is_bold: bool = False) -> float:
    """WCAG required ratio: 4.5 for normal text, 3.0 for large text."""
    if font_pt is None:
        return WCAG_AA_NORMAL
    threshold = LARGE_TEXT_BOLD_PT if is_bold else LARGE_TEXT_REGULAR_PT
    if font_pt >= threshold:
        return WCAG_AA_LARGE
    return WCAG_AA_NORMAL


# --------------------------------------------------------------------------
# Color resolution (effective fg/bg for a placeholder)
# --------------------------------------------------------------------------

def _shape_fill_hex(shape) -> Optional[str]:
    """Solid-fill hex of ``shape``, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return f"#{rgb}"
    except Exception:
        pass
    return None


def _shape_dominant_text_hex(shape) -> Optional[str]:
    """Largest-run text color hex of ``shape``, or None (inheritance)."""
    if not getattr(shape, "has_text_frame", False):
        return None
    candidates = []  # (font_size_pt, color_hex)
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                fs = run.font.size.pt if run.font.size else None
                try:
                    rgb = run.font.color.rgb
                    color_hex = f"#{rgb}" if rgb is not None else None
                except Exception:
                    color_hex = None
                if color_hex:
                    candidates.append((fs or 12.0, color_hex))
    except Exception:
        return None
    if not candidates:
        return None
    # Pick the color of the largest-font run (title-like runs dominate)
    candidates.sort(key=lambda c: c[0], reverse=True)
    return candidates[0][1]


def _shape_font_size_pt(shape) -> Optional[float]:
    """Largest run font size in pt, or None."""
    if not getattr(shape, "has_text_frame", False):
        return None
    sizes = []
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                fs = run.font.size
                if fs is not None:
                    sizes.append(fs.pt)
    except Exception:
        return None
    return max(sizes) if sizes else None


def _shape_is_text_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False) or not shape.is_placeholder:
        return False
    try:
        type_ = shape.placeholder_format.type
        # PP_PLACEHOLDER: TITLE=13, BODY=2, OBJECT=17, SUBTITLE=3
        type_id = int(str(type_).rsplit(".", 1)[-1]) if type_ is not None else None
        return type_id in (2, 3, 13, 17) or type_id is None
    except Exception:
        return True


def _shape_bounding_box(shape):
    try:
        return (
            int(shape.left), int(shape.top),
            int(shape.left) + int(shape.width),
            int(shape.top) + int(shape.height),
        )
    except (TypeError, AttributeError):
        return None


def _rect_contains_center(rect_a, rect_b) -> bool:
    """True iff rect_b's center is inside rect_a."""
    if rect_a is None or rect_b is None:
        return False
    a_l, a_t, a_r, a_b = rect_a
    b_l, b_t, b_r, b_b = rect_b
    cx = (b_l + b_r) / 2
    cy = (b_t + b_b) / 2
    return a_l <= cx <= a_r and a_t <= cy <= a_b


def _rect_area(rect_a):
    if rect_a is None:
        return 0
    return (rect_a[2] - rect_a[0]) * (rect_a[3] - rect_a[1])


def _find_container(shape, candidate_shapes) -> Optional[Any]:
    """Find the tightest non-placeholder shape whose bbox contains ``shape``'s center."""
    target_rect = _shape_bounding_box(shape)
    if target_rect is None:
        return None
    best = None
    best_area = None
    for cand in candidate_shapes:
        if _is_text_placeholder(cand):
            continue
        cand_rect = _shape_bounding_box(cand)
        if cand_rect is None:
            continue
        if not _rect_contains_center(cand_rect, target_rect):
            continue
        area = _rect_area(cand_rect)
        if best_area is None or area < best_area:
            best = cand
            best_area = area
    return best


def _is_text_placeholder(shape) -> bool:
    return _shape_is_text_placeholder(shape)


def _default_text_color_from_theme(theme: Optional[Dict[str, str]]) -> str:
    """Theme's dk1 (default dark text) — fallback when placeholder has no explicit color."""
    if theme and theme.get("dk1"):
        return theme["dk1"]
    return "#000000"  # Office default


def _default_bg_color_from_theme(theme: Optional[Dict[str, str]]) -> str:
    """Theme's lt1 (default light background) — fallback when no container fill."""
    if theme and theme.get("lt1"):
        return theme["lt1"]
    return "#FFFFFF"


def _invert_hex(hex_str: str) -> str:
    """'#09090B' → '#FFFFFF' (perceptual inversion, not just bitwise)."""
    rgb = _hex_to_rgb(hex_str)
    if rgb is None:
        return "#FFFFFF"
    # Use luminance to pick a high-contrast partner. Note: luminance 0 is
    # a valid value (pure black) — do NOT use `or` to default it.
    lum = _relative_luminance(hex_str)
    if lum is None:
        return "#FFFFFF"
    return "#FFFFFF" if lum < 0.5 else "#000000"


def _override_placeholder_text_color(ph, new_hex: str) -> bool:
    """Set ``ph``'s default run color to ``new_hex`` so newly-added content inherits it.

    Returns True on success. Implementation: ensure the placeholder's
    <p:txBody>/<a:lstStyle> endParaRPr + the first paragraph's defRPr carry
    the new color. This is the safest path — content added later via
    `text_frame.text = "..."` will pick it up.
    """
    try:
        from pptx.oxml.ns import qn
        from pptx.dml.color import RGBColor
        # Set on the text_frame's default paragraph properties
        tf = ph.text_frame
        # End-paragraph rPr (defines color for empty placeholders)
        for para in tf.paragraphs:
            pPr = para._pPr
            if pPr is None:
                pPr = para._p.get_or_add_pPr()
            defRPr = pPr.find(qn("a:defRPr"))
            if defRPr is None:
                defRPr = pPr.makeelement(qn("a:defRPr"), {})
                pPr.append(defRPr)
            solidFill = defRPr.find(qn("a:solidFill"))
            if solidFill is None:
                solidFill = defRPr.makeelement(qn("a:solidFill"), {})
                defRPr.append(solidFill)
            # Replace any existing srgbClr
            for child in list(solidFill):
                solidFill.remove(child)
            srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": new_hex.lstrip("#").upper()})
            solidFill.append(srgb)
        return True
    except Exception as exc:
        logger.warning("auto-fix failed on placeholder %s: %s", getattr(ph, "name", "?"), exc)
        return False


# --------------------------------------------------------------------------
# Main check API
# --------------------------------------------------------------------------

def contrast_violations(
    layout,
    theme: Optional[Dict[str, str]] = None,
    auto_fix: bool = False,
) -> List[ContrastViolation]:
    """Return all contrast violations on ``layout``.

    Args:
        layout: a ``pptx.slide.SlideLayout``
        theme: optional dict from ``designer_promoter.extract_theme_from_shapes``
            — used to resolve default text/bg colors when placeholders have
            no explicit color
        auto_fix: when True, override each low-contrast placeholder's default
            text color to a high-contrast alternative; the returned
            :class:`ContrastViolation` carries ``auto_fixed=True``

    Returns:
        List of violations (post-auto-fix; auto-fixed violations are still
        reported so the orchestrator knows what changed).
    """
    violations: List[ContrastViolation] = []
    shapes = list(layout.shapes)
    text_phs = [s for s in shapes if _is_text_placeholder(s)]
    candidate_containers = [s for s in shapes if not _is_text_placeholder(s)]

    default_fg = _default_text_color_from_theme(theme)
    default_bg = _default_bg_color_from_theme(theme)

    for ph in text_phs:
        # Resolve effective foreground
        fg = _shape_dominant_text_hex(ph) or default_fg
        # Resolve effective background
        container = _find_container(ph, candidate_containers)
        if container is not None:
            bg = _shape_fill_hex(container) or default_bg
        else:
            bg = default_bg
        # Compute ratio
        ratio = contrast_ratio(fg, bg)
        font_pt = _shape_font_size_pt(ph)
        required = required_ratio_for(font_pt)
        if ratio >= required:
            continue
        severity = "critical" if ratio < WCAG_CRITICAL else "warning"
        try:
            ph_idx = ph.placeholder_format.idx
            ph_name = ph.name
        except Exception:
            ph_idx, ph_name = -1, "<unnamed>"
        fixed = False
        if auto_fix:
            new_fg = _invert_hex(bg)
            # Sanity-check: don't apply a fix that's no better
            if contrast_ratio(new_fg, bg) >= required:
                fixed = _override_placeholder_text_color(ph, new_fg)
        violations.append(ContrastViolation(
            placeholder_idx=ph_idx,
            placeholder_name=ph_name,
            fg_color=fg,
            bg_color=bg,
            contrast_ratio=ratio,
            required_ratio=required,
            severity=severity,
            detail=f"{ph_name} {fg} on {bg} = {ratio:.2f}:1 (need {required:.1f}:1)",
            auto_fixed=fixed,
        ))
    return violations


def check_template_contrast(
    prs,
    theme: Optional[Dict[str, str]] = None,
    auto_fix: bool = False,
) -> dict:
    """Run ``contrast_violations`` on every layout in ``prs``.

    Returns:
        ``{layout_name: [violation_dict, ...], "__summary": {...}}``
    """
    summary = {"critical": 0, "warning": 0, "auto_fixed": 0, "layouts_checked": 0}
    out = {}
    for layout in prs.slide_layouts:
        try:
            layout_name = layout.name
        except Exception:
            layout_name = f"<layout-{id(layout)}>"
        try:
            vs = contrast_violations(layout, theme=theme, auto_fix=auto_fix)
        except Exception as exc:
            logger.warning("contrast_check failed on layout %s: %s", layout_name, exc)
            vs = []
        out[layout_name] = [v.to_dict() for v in vs]
        summary["layouts_checked"] += 1
        for v in vs:
            if v.severity == "critical":
                summary["critical"] += 1
            else:
                summary["warning"] += 1
            if v.auto_fixed:
                summary["auto_fixed"] += 1
    out["__summary"] = summary
    return out


__all__ = [
    "ContrastViolation",
    "contrast_ratio",
    "required_ratio_for",
    "contrast_violations",
    "check_template_contrast",
    "WCAG_AA_NORMAL",
    "WCAG_AA_LARGE",
    "WCAG_CRITICAL",
]
